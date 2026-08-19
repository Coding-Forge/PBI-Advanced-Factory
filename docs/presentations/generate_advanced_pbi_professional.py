"""Advanced Power BI -- Professional Program Deck (v2)

A full redesign of the curriculum overview deck: fewer, denser, more
purposeful slides; a refined enterprise-consulting visual language (deep
navy + muted gold, serif/sans pairing); a real diagram on every module
overview slide; condensed "at a glance" module deep-dive slides (full
topic-level detail preserved in speaker notes for anyone presenting live);
consistent wayfinding (arc progress rail, module rail, breadcrumbs); and
click-through navigation from the curriculum map / arc dividers into each
module.

Source content is read from content_model.json (extracted from the prior
deck) so the underlying curriculum copy is preserved verbatim; only the
structure, layout, and visual system are new.
"""
from pptx import Presentation
from pptx.util import Inches, Pt, Emu
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR
from pptx.enum.shapes import MSO_SHAPE, MSO_CONNECTOR
from pptx.oxml.ns import qn
from lxml import etree
import json
import os

HERE = os.path.dirname(__file__)
ICON_DIR = os.path.join(HERE, "assets", "icons")
IMG_DIR = os.path.join(HERE, "assets", "images")
DATA_PATH = os.path.join(HERE, "content_model.json")
OUT_PATH = os.path.join(HERE, "AdvancedPowerBI.pptx")

with open(DATA_PATH) as f:
    MODEL = json.load(f)

# ---------------------------------------------------------------- palette
NAVY = RGBColor(0x0A, 0x18, 0x2A)          # dark bg (title / dividers)
NAVY_PANEL = RGBColor(0x12, 0x24, 0x3A)    # slightly lighter dark panel
NAVY_LINE = RGBColor(0x2B, 0x3F, 0x57)     # hairline on dark bg
INK = RGBColor(0x1B, 0x24, 0x30)           # body text on light bg
MUTED = RGBColor(0x60, 0x71, 0x7F)         # secondary text
GOLD = RGBColor(0xC9, 0x96, 0x2F)          # muted brass/gold accent
GOLD_DEEP = RGBColor(0x8F, 0x66, 0x1E)     # deeper gold for light-bg text
LIGHT_BG = RGBColor(0xF7, 0xF8, 0xFA)      # content slide background
CARD_BG = RGBColor(0xFF, 0xFF, 0xFF)
WHITE = RGBColor(0xFF, 0xFF, 0xFF)
ICE = RGBColor(0xC7, 0xD2, 0xDD)           # muted light text on dark bg
LINE_GREY = RGBColor(0xDD, 0xE3, 0xEA)
FAINT = RGBColor(0xD7, 0xC9, 0xA9)         # faint gold for connectors

HEAD_FONT = "Cambria"
BODY_FONT = "Calibri"

prs = Presentation()
prs.slide_width = Inches(13.333)
prs.slide_height = Inches(7.5)
BLANK = prs.slide_layouts[6]
SW, SH = prs.slide_width, prs.slide_height
MARGIN = Inches(0.6)

ARC_NAMES = ["BUILD", "PRESENT", "PROTECT & OPERATE", "SCALE"]
ARC_FOR_MODULE = {1: 0, 2: 0, 3: 0, 4: 1, 5: 1, 6: 1, 7: 2, 8: 2, 9: 2, 10: 3, 11: 3}
MODULE_ICON = {
    1: "layers", 2: "formula", 3: "funnel", 4: "canvas_layout", 5: "gauge",
    6: "spark", 7: "shield", 8: "rocket", 9: "gauge", 10: "layers", 11: "gear",
}

# slide-number bookkeeping: filled in as we add slides so we can hyperlink
SLIDE_INDEX = {}   # key -> slide object added

# ============================================================ primitives

def add_slide():
    return prs.slides.add_slide(BLANK)


def set_bg(slide, color):
    slide.background.fill.solid()
    slide.background.fill.fore_color.rgb = color


def no_line(shape):
    shape.line.fill.background()


def add_rect(slide, x, y, w, h, color, line=False):
    shp = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, int(x), int(y), int(w), int(h))
    shp.fill.solid()
    shp.fill.fore_color.rgb = color
    if not line:
        no_line(shp)
    shp.shadow.inherit = False
    return shp


def add_rrect(slide, x, y, w, h, color, radius=0.06, line=False, line_color=None):
    shp = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, int(x), int(y), int(w), int(h))
    try:
        shp.adjustments[0] = radius
    except Exception:
        pass
    shp.fill.solid()
    shp.fill.fore_color.rgb = color
    if line_color:
        shp.line.color.rgb = line_color
        shp.line.width = Pt(1)
    elif not line:
        no_line(shp)
    shp.shadow.inherit = False
    return shp


def add_oval(slide, x, y, w, h, color, line_color=None):
    shp = slide.shapes.add_shape(MSO_SHAPE.OVAL, int(x), int(y), int(w), int(h))
    shp.fill.solid()
    shp.fill.fore_color.rgb = color
    if line_color:
        shp.line.color.rgb = line_color
        shp.line.width = Pt(1)
    else:
        no_line(shp)
    shp.shadow.inherit = False
    return shp


def add_connector(slide, x1, y1, x2, y2, color=MUTED, width=Pt(1.25), dash=None):
    conn = slide.shapes.add_connector(MSO_CONNECTOR.STRAIGHT, int(x1), int(y1), int(x2), int(y2))
    conn.line.color.rgb = color
    conn.line.width = width
    if dash:
        ln = conn.line._get_or_add_ln()
        d = etree.SubElement(ln, qn('a:prstDash'))
        d.set('val', dash)
    return conn


def add_text(slide, x, y, w, h, text, size=14, color=INK, bold=False, italic=False,
             align=PP_ALIGN.LEFT, font=BODY_FONT, anchor=MSO_ANCHOR.TOP, line_spacing=1.0,
             margin=0, letter_spacing=None):
    tb = slide.shapes.add_textbox(int(x), int(y), int(w), int(h))
    tf = tb.text_frame
    tf.word_wrap = True
    tf.vertical_anchor = anchor
    tf.margin_left = Emu(margin); tf.margin_right = Emu(margin)
    tf.margin_top = 0; tf.margin_bottom = 0
    lines = text.split("\n") if isinstance(text, str) else text
    for i, line in enumerate(lines):
        p = tf.paragraphs[0] if i == 0 else tf.add_paragraph()
        p.alignment = align
        p.line_spacing = line_spacing
        r = p.add_run()
        r.text = line
        r.font.size = Pt(size)
        r.font.bold = bold
        r.font.italic = italic
        r.font.color.rgb = color
        r.font.name = font
        if letter_spacing is not None:
            rPr = r._r.get_or_add_rPr()
            rPr.set('spc', str(letter_spacing))
    return tb


def add_icon(slide, icon_name, x, y, size, badge_color=NAVY, ring=None):
    """Circular badge with a centered icon glyph image."""
    badge = add_oval(slide, x, y, size, size, badge_color)
    if ring:
        badge.line.color.rgb = ring
        badge.line.width = Pt(1.25)
    pad = int(size * 0.26)
    path = os.path.join(ICON_DIR, f"{icon_name}.png")
    if os.path.exists(path):
        slide.shapes.add_picture(path, x + pad, y + pad, size - 2 * pad, size - 2 * pad)
    return badge


def add_kicker(slide, text, x, y, w, color=GOLD, size=11, font=BODY_FONT):
    return add_text(slide, x, y, w, Inches(0.3), text, size=size, color=color, bold=True,
                     font=font, letter_spacing=140)


def add_footer(slide, breadcrumb, page_num, dark=False, backlink_target=None):
    color = ICE if dark else MUTED
    add_text(slide, Inches(0.5), Inches(7.13), Inches(8.0), Inches(0.3), breadcrumb,
              size=9.5, color=color, font=BODY_FONT)
    add_text(slide, SW - Inches(1.0), Inches(7.13), Inches(0.5), Inches(0.3), str(page_num),
              size=9.5, color=color, align=PP_ALIGN.RIGHT, font=BODY_FONT)
    if backlink_target is not None:
        tb = add_text(slide, SW - Inches(2.55), Inches(7.13), Inches(1.9), Inches(0.3),
                       "\u2191 Curriculum Map", size=9.5, color=GOLD if not dark else GOLD,
                       bold=True, align=PP_ALIGN.RIGHT, font=BODY_FONT)
        tb.click_action.target_slide = backlink_target


def add_progress_rail(slide, current_idx, labels, x, y, w, dark=False):
    """Small horizontal wayfinding rail: a row of labeled segments, current one highlighted."""
    n = len(labels)
    gap = Inches(0.12)
    seg_w = int((w - gap * (n - 1)) / n)
    seg_h = Inches(0.06)
    label_color_on = GOLD
    label_color_off = ICE if dark else MUTED
    for i, label in enumerate(labels):
        sx = x + i * (seg_w + gap)
        active = (i == current_idx)
        add_rect(slide, sx, y, seg_w, seg_h, GOLD if active else (NAVY_LINE if dark else LINE_GREY))
        add_text(slide, sx, y + Inches(0.12), seg_w, Inches(0.24), label,
                  size=8, bold=active, color=label_color_on if active else label_color_off,
                  font=BODY_FONT, letter_spacing=60)


def add_header(slide, kicker, breadcrumb_right=None):
    add_rect(slide, 0, 0, SW, Inches(0.06), GOLD)
    add_kicker(slide, kicker, Inches(0.6), Inches(0.28), Inches(9.0))
    if breadcrumb_right:
        add_text(slide, SW - Inches(3.5), Inches(0.24), Inches(2.9), Inches(0.3), breadcrumb_right,
                  size=10, color=MUTED, align=PP_ALIGN.RIGHT, font=BODY_FONT)


# ============================================================ diagrams
# Each function draws inside a panel (x, y, w, h) already framed by caller.

def _panel_frame(slide, x, y, w, h, title):
    add_rrect(slide, x, y, w, h, CARD_BG, radius=0.05, line_color=LINE_GREY)
    add_text(slide, x + Inches(0.18), y + Inches(0.14), w - Inches(0.36), Inches(0.26), title,
              size=9.5, bold=True, color=MUTED, letter_spacing=70)


def diagram_star_schema(slide, x, y, w, h):
    _panel_frame(slide, x, y, w, h, "STAR SCHEMA")
    cx, cy = x + w // 2, y + int(h * 0.58)
    fact_w, fact_h = Inches(1.7), Inches(0.55)
    dim_w, dim_h = Inches(1.15), Inches(0.46)
    offs = [(-1.65, -0.78), (1.65, -0.78), (-1.65, 0.78), (1.65, 0.78)]
    labels = ["Dim Date", "Dim Product", "Dim Customer", "Dim Store"]
    for (ox, oy), label in zip(offs, labels):
        dx, dy = cx + Inches(ox) - dim_w // 2, cy + Inches(oy) - dim_h // 2
        add_connector(slide, cx, cy, dx + dim_w // 2, dy + dim_h // 2, color=FAINT, width=Pt(1.25))
    for (ox, oy), label in zip(offs, labels):
        dx, dy = cx + Inches(ox) - dim_w // 2, cy + Inches(oy) - dim_h // 2
        add_rrect(slide, dx, dy, dim_w, dim_h, LIGHT_BG, radius=0.2, line_color=RGBColor(0xB9, 0xC5, 0xD1))
        add_text(slide, dx, dy, dim_w, dim_h, label, size=9, color=INK, align=PP_ALIGN.CENTER, anchor=MSO_ANCHOR.MIDDLE)
    add_rrect(slide, cx - fact_w // 2, cy - fact_h // 2, fact_w, fact_h, GOLD_DEEP, radius=0.2)
    add_text(slide, cx - fact_w // 2, cy - fact_h // 2, fact_w, fact_h, "Fact Sales", size=10.5, bold=True,
              color=WHITE, align=PP_ALIGN.CENTER, anchor=MSO_ANCHOR.MIDDLE)
    add_text(slide, x + Inches(0.2), y + h - Inches(0.42), w - Inches(0.4), Inches(0.36),
              "One grain, one fact table -- descriptive context lives in the surrounding dimensions.",
              size=8.5, color=MUTED)


def _vertical_flow(slide, x, y, w, h, steps, footnote=None):
    n = len(steps)
    top = Inches(0.55)
    bottom = Inches(0.42) if footnote else Inches(0.16)
    gap = Inches(0.12)
    avail = h - top - bottom
    step_h = int((avail - gap * (n - 1)) / n)
    bx, bw = x + Inches(0.3), w - Inches(0.6)
    by = y + top
    for i, step in enumerate(steps):
        add_rrect(slide, bx, by, bw, step_h, LIGHT_BG, radius=0.16, line_color=RGBColor(0xB9, 0xC5, 0xD1))
        add_text(slide, bx + Inches(0.1), by, bw - Inches(0.2), step_h, step, size=9.5, color=INK,
                  align=PP_ALIGN.CENTER, anchor=MSO_ANCHOR.MIDDLE)
        if i < n - 1:
            add_connector(slide, bx + bw // 2, by + step_h, bx + bw // 2, by + step_h + gap, color=GOLD_DEEP, width=Pt(1.5))
        by += step_h + gap
    if footnote:
        add_text(slide, x + Inches(0.2), y + h - Inches(0.38), w - Inches(0.4), Inches(0.32), footnote,
                  size=8, color=MUTED)


def diagram_pipeline(slide, x, y, w, h):
    _panel_frame(slide, x, y, w, h, "QUERY PIPELINE")
    _vertical_flow(slide, x, y, w, h, ["Source (raw)", "Staging (typed)", "Transform (parameterized)", "Load (model-ready)"])


def diagram_hierarchy(slide, x, y, w, h):
    _panel_frame(slide, x, y, w, h, "NATURAL HIERARCHY")
    _vertical_flow(slide, x, y, w, h, ["Year", "Quarter", "Month", "Day"],
                    footnote="Drill down moves through levels without changing the visual.")


def diagram_diagnose(slide, x, y, w, h):
    _panel_frame(slide, x, y, w, h, "DIAGNOSE, THEN FIX")
    _vertical_flow(slide, x, y, w, h, ["Slow measure reported", "Performance Analyzer", "Query plan review", "Targeted fix"])


def diagram_deployment(slide, x, y, w, h):
    _panel_frame(slide, x, y, w, h, "DEPLOYMENT PIPELINE")
    _vertical_flow(slide, x, y, w, h, ["Dev workspace", "Test workspace (validate)", "Prod workspace (promote)"])


def diagram_usage_signal(slide, x, y, w, h):
    _panel_frame(slide, x, y, w, h, "USAGE \u2192 SIGNAL \u2192 ACTION")
    _vertical_flow(slide, x, y, w, h, ["Usage metrics (views, load time)", "App / workspace signals", "Alert or governance action"])


def diagram_cicd(slide, x, y, w, h):
    _panel_frame(slide, x, y, w, h, "CI/CD FLOW")
    _vertical_flow(slide, x, y, w, h, ["Commit (Git / PBIP)", "Build & validate (pipeline)", "Deploy: Dev \u2192 Test \u2192 Prod"])


def _two_col_compare(slide, x, y, w, h, colA_title, colA_steps, colB_title, colB_steps, arrows=True, footnote=None):
    col_w = int((w - Inches(0.3)) / 2)
    col1_x = x + Inches(0.15)
    col2_x = col1_x + col_w + Inches(0.15)

    def col(cx, title, steps):
        cy = y + Inches(0.55)
        add_text(slide, cx, cy, col_w, Inches(0.24), title, size=9.5, bold=True, color=GOLD_DEEP,
                  align=PP_ALIGN.CENTER)
        cy += Inches(0.32)
        n = len(steps)
        avail = h - Inches(0.55) - Inches(0.32) - Inches(0.2)
        gap = Inches(0.1)
        step_h = int((avail - gap * (n - 1)) / n)
        for i, step in enumerate(steps):
            add_rrect(slide, cx, cy, col_w, step_h, LIGHT_BG, radius=0.2, line_color=RGBColor(0xB9, 0xC5, 0xD1))
            add_text(slide, cx + Inches(0.08), cy, col_w - Inches(0.16), step_h, step, size=8.5, color=INK,
                      align=PP_ALIGN.CENTER, anchor=MSO_ANCHOR.MIDDLE)
            if arrows and i < n - 1:
                add_connector(slide, cx + col_w // 2, cy + step_h, cx + col_w // 2, cy + step_h + gap, color=GOLD_DEEP, width=Pt(1.25))
            cy += step_h + gap

    col(col1_x, colA_title, colA_steps)
    col(col2_x, colB_title, colB_steps)
    if footnote:
        add_text(slide, x + Inches(0.2), y + h - Inches(0.34), w - Inches(0.4), Inches(0.3), footnote, size=7.5, color=MUTED)


def diagram_context(slide, x, y, w, h):
    _panel_frame(slide, x, y, w, h, "ROW vs FILTER CONTEXT")
    _two_col_compare(slide, x, y, w, h,
                       "ROW CONTEXT", ["Iterator (SUMX, FILTER)", "Evaluates one row at a time", "Produces a per-row value"],
                       "FILTER CONTEXT", ["Visuals, slicers, CALCULATE", "Filters applied to the model", "Shapes what's visible"],
                       footnote="CALCULATE bridges row context into filter context.")


def diagram_rls(slide, x, y, w, h):
    _panel_frame(slide, x, y, w, h, "STATIC vs DYNAMIC RLS")
    _two_col_compare(slide, x, y, w, h,
                       "STATIC", ["Role: fixed region/team", "Fixed DAX filter", "Simple, auditable"],
                       "DYNAMIC", ["USERPRINCIPALNAME()", "Security mapping table", "Scales, harder to test"],
                       arrows=False)


def diagram_capacity(slide, x, y, w, h):
    _panel_frame(slide, x, y, w, h, "SHARED vs DEDICATED CAPACITY")
    _two_col_compare(slide, x, y, w, h,
                       "SHARED", ["Pay-per-user (Pro)", "Shared, variable resources", "Best for small teams"],
                       "PREMIUM / FABRIC", ["Dedicated capacity", "Predictable performance", "Enterprise scale"],
                       arrows=False)


def diagram_anomaly(slide, x, y, w, h):
    _panel_frame(slide, x, y, w, h, "ANOMALY DETECTION")
    cx, cy = x + Inches(0.45), y + Inches(0.6)
    cw, ch = w - Inches(0.85), Inches(1.65)
    add_connector(slide, cx, cy + ch, cx + cw, cy + ch, color=RGBColor(0xB9, 0xC5, 0xD1))
    add_connector(slide, cx, cy, cx, cy + ch, color=RGBColor(0xB9, 0xC5, 0xD1))
    ys = [0.55, 0.5, 0.6, 0.5, 0.45, 0.55, 0.95, 0.5]
    pts = []
    for i, yv in enumerate(ys):
        px = cx + int(cw * (i + 0.5) / len(ys))
        py = cy + int(ch * yv)
        pts.append((px, py))
    for i in range(len(pts) - 1):
        add_connector(slide, pts[i][0], pts[i][1], pts[i+1][0], pts[i+1][1], color=FAINT, width=Pt(1.75))
    for i, (px, py) in enumerate(pts):
        anom = ys[i] > 0.85
        r = Inches(0.14) if anom else Inches(0.09)
        add_oval(slide, px - r // 2, py - r // 2, r, r, GOLD if anom else GOLD_DEEP)
        if anom:
            add_text(slide, px - Inches(0.5), py - Inches(0.42), Inches(1.0), Inches(0.24), "Anomaly",
                      size=8, bold=True, color=GOLD_DEEP, align=PP_ALIGN.CENTER)
    add_text(slide, x + Inches(0.2), y + h - Inches(0.38), w - Inches(0.4), Inches(0.32),
              "Flags the point that breaks the trend, not just the highest value.", size=8, color=MUTED)


DIAGRAM_FOR_MODULE = {
    1: diagram_star_schema, 2: diagram_context, 3: diagram_pipeline, 4: diagram_hierarchy,
    5: diagram_diagnose, 6: diagram_anomaly, 7: diagram_rls, 8: diagram_deployment,
    9: diagram_usage_signal, 10: diagram_capacity, 11: diagram_cicd,
}

# ============================================================ animations / transitions
P_NS = 'http://schemas.openxmlformats.org/presentationml/2006/main'


def _pqn(tag):
    return '{%s}%s' % (P_NS, tag.split(':')[1])


def add_fade_transition(slide):
    elm = slide._element
    existing = elm.find(_pqn('p:transition'))
    if existing is not None:
        elm.remove(existing)
    trans = etree.SubElement(elm, _pqn('p:transition'))
    trans.set('spd', 'med')
    etree.SubElement(trans, _pqn('p:fade'))
    cSld = elm.find(_pqn('p:cSld'))
    elm.remove(trans)
    elm.insert(list(elm).index(cSld) + 1, trans)


def add_paragraph_build(slide, shape, n_paras):
    """Click-to-reveal build for each paragraph of `shape`."""
    shape_id = shape.shape_id
    pars_xml = []
    tn_id = 5
    for i in range(n_paras):
        pars_xml.append(f'''
        <p:par>
          <p:cTn id="{tn_id}" presetID="1" presetClass="entr" presetSubtype="0" fill="hold" grpId="0" nodeType="clickEffect">
            <p:stCondLst><p:cond delay="0"/></p:stCondLst>
            <p:childTnLst>
              <p:set>
                <p:cBhvr>
                  <p:cTn id="{tn_id+1}" dur="1" fill="hold"><p:stCondLst><p:cond delay="0"/></p:stCondLst></p:cTn>
                  <p:tgtEl><p:spTgt spid="{shape_id}"><p:txEl><p:pRg st="{i}" end="{i}"/></p:txEl></p:spTgt></p:tgtEl>
                  <p:attrNameLst><p:attrName>style.visibility</p:attrName></p:attrNameLst>
                </p:cBhvr>
                <p:to><p:strVal val="visible"/></p:to>
              </p:set>
            </p:childTnLst>
          </p:cTn>
        </p:par>''')
        tn_id += 2
    timing_xml = f'''<p:timing xmlns:p="{P_NS}">
      <p:tnLst>
        <p:par>
          <p:cTn id="1" dur="indefinite" restart="never" nodeType="tmRoot">
            <p:childTnLst>
              <p:seq concat="1" nextAc="seek">
                <p:cTn id="2" dur="indefinite" nodeType="mainSeq">
                  <p:childTnLst>{''.join(pars_xml)}</p:childTnLst>
                </p:cTn>
                <p:prevCondLst><p:cond evt="onPrev" delay="0"><p:tgtEl><p:sldTgt/></p:tgtEl></p:cond></p:prevCondLst>
                <p:nextCondLst><p:cond evt="onNext" delay="0"><p:tgtEl><p:sldTgt/></p:tgtEl></p:cond></p:nextCondLst>
              </p:seq>
            </p:childTnLst>
          </p:cTn>
        </p:par>
      </p:tnLst>
      <p:bldLst><p:bldP spid="{shape_id}" grpId="0" build="p"/></p:bldLst>
    </p:timing>'''
    elm = slide._element
    existing = elm.find(_pqn('p:timing'))
    if existing is not None:
        elm.remove(existing)
    elm.append(etree.fromstring(timing_xml.encode('utf-8')))


# ============================================================ slide builders

def build_title(s):
    set_bg(s, NAVY)
    d = MODEL['title']
    add_rect(s, 0, 0, SW, Inches(0.08), GOLD)
    add_oval(s, SW - Inches(3.2), -Inches(1.6), Inches(4.0), Inches(4.0), NAVY_PANEL)
    add_kicker(s, d['eyebrow'], Inches(0.9), Inches(1.35), Inches(10.0), color=GOLD, size=13)
    add_text(s, Inches(0.9), Inches(1.85), Inches(10.8), Inches(1.3), d['title'], size=60, bold=True,
              color=WHITE, font=HEAD_FONT)
    add_rect(s, Inches(0.92), Inches(3.05), Inches(1.4), Inches(0.045), GOLD)
    add_text(s, Inches(0.9), Inches(3.3), Inches(9.6), Inches(1.3), d['subtitle'], size=15.5,
              color=ICE, line_spacing=1.3)
    sx = Inches(0.9)
    for val, label in d['stats']:
        add_text(s, sx, Inches(5.35), Inches(2.6), Inches(0.7), val, size=34, bold=True, color=GOLD, font=HEAD_FONT)
        add_text(s, sx, Inches(6.05), Inches(2.6), Inches(0.4), label, size=11.5, color=ICE)
        sx += Inches(2.85)
    add_text(s, Inches(0.9), Inches(6.9), Inches(6), Inches(0.4), "Program Overview", size=10.5, color=MUTED)
    add_fade_transition(s)


def build_instructor(s, page_num):
    set_bg(s, LIGHT_BG)
    d = MODEL['instructor']
    add_header(s, d['eyebrow'])
    add_text(s, Inches(0.6), Inches(0.95), Inches(6.8), Inches(0.6), d['name'], size=32, bold=True,
              color=INK, font=HEAD_FONT)
    add_text(s, Inches(0.6), Inches(1.6), Inches(6.8), Inches(0.35), d['role'], size=15, color=GOLD_DEEP, bold=True)
    add_text(s, Inches(0.6), Inches(1.98), Inches(6.8), Inches(0.32), d['credential'], size=11.5, color=MUTED, italic=True)
    add_rect(s, Inches(0.6), Inches(2.42), Inches(1.0), Inches(0.035), GOLD)
    add_text(s, Inches(0.6), Inches(2.65), Inches(6.8), Inches(1.9), d['bio'], size=13, color=INK, line_spacing=1.35)

    sy = Inches(4.85)
    sx = Inches(0.6)
    for val, label in d['stats']:
        add_rrect(s, sx, sy, Inches(2.15), Inches(1.15), CARD_BG, radius=0.08, line_color=LINE_GREY)
        add_text(s, sx + Inches(0.15), sy + Inches(0.12), Inches(1.9), Inches(0.5), val, size=20, bold=True,
                  color=GOLD_DEEP, font=HEAD_FONT)
        add_text(s, sx + Inches(0.15), sy + Inches(0.68), Inches(1.9), Inches(0.4), label, size=9.5, color=MUTED)
        sx += Inches(2.3)

    px, py, pw, ph = Inches(8.1), Inches(0.95), Inches(4.6), Inches(5.05)
    add_rrect(s, px, py, pw, ph, NAVY, radius=0.05)
    img_path = os.path.join(IMG_DIR, "ProHeadShot_A.png")
    if os.path.exists(img_path):
        img_size = Inches(2.6)
        s.shapes.add_picture(img_path, px + (pw - img_size) // 2, py + Inches(0.4), height=img_size)
    add_text(s, px + Inches(0.3), py + ph - Inches(1.1), pw - Inches(0.6), Inches(0.8),
              "\u201cTaught the way I build it for clients.\u201d", size=13, italic=True, color=ICE,
              align=PP_ALIGN.CENTER, font=HEAD_FONT)

    add_footer(s, "Advanced Power BI  \u00b7  Meet Your Instructor", page_num)
    add_fade_transition(s)


def build_why(s, page_num):
    set_bg(s, LIGHT_BG)
    d = MODEL['why']
    add_header(s, d['eyebrow'])
    add_text(s, Inches(0.6), Inches(0.95), Inches(11.6), Inches(1.1), d['title'], size=27, bold=True,
              color=INK, font=HEAD_FONT, line_spacing=1.1)
    cx = Inches(0.6)
    cw = Inches(3.85)
    for i, (title, body) in enumerate(d['cards']):
        add_rrect(s, cx, Inches(2.35), cw, Inches(4.15), CARD_BG, radius=0.045, line_color=LINE_GREY)
        add_rect(s, cx, Inches(2.35), cw, Inches(0.06), GOLD)
        add_text(s, cx + Inches(0.28), Inches(2.7), cw - Inches(0.56), Inches(0.3), f"0{i+1}", size=13,
                  bold=True, color=GOLD_DEEP, font=HEAD_FONT)
        add_text(s, cx + Inches(0.28), Inches(3.1), cw - Inches(0.56), Inches(1.15), title, size=16.5,
                  bold=True, color=INK, font=HEAD_FONT, line_spacing=1.15)
        add_text(s, cx + Inches(0.28), Inches(4.3), cw - Inches(0.56), Inches(2.0), body, size=11.5,
                  color=MUTED, line_spacing=1.32)
        cx += cw + Inches(0.22)
    add_footer(s, "Advanced Power BI  \u00b7  Program Overview", page_num)
    add_fade_transition(s)


def build_curriculum_map(s, page_num, module_target, capstone_target):
    set_bg(s, LIGHT_BG)
    add_header(s, "CURRICULUM MAP")
    add_text(s, Inches(0.6), Inches(0.95), Inches(9.0), Inches(0.5), "Four arcs, eleven modules, one capstone",
              size=22, bold=True, color=INK, font=HEAD_FONT)
    add_text(s, Inches(0.6), Inches(1.42), Inches(9.5), Inches(0.35),
              "Click any module to jump straight there.", size=11, color=MUTED, italic=True)

    modules = MODEL['modules']
    col_map = {0: [], 1: [], 2: [], 3: []}
    for m in modules:
        col_map[ARC_FOR_MODULE[m['num']]].append(m)

    cx = Inches(0.6)
    cw = Inches(2.98)
    for arc_i, arc_label in enumerate(ARC_NAMES):
        cy = Inches(2.0)
        add_rrect(s, cx, cy, cw, Inches(4.55), CARD_BG, radius=0.04, line_color=LINE_GREY)
        add_rect(s, cx, cy, cw, Inches(0.5), NAVY)
        add_text(s, cx, cy, cw, Inches(0.5), arc_label, size=12.5, bold=True, color=GOLD,
                  align=PP_ALIGN.CENTER, anchor=MSO_ANCHOR.MIDDLE, letter_spacing=60)
        by = cy + Inches(0.75)
        for m in col_map[arc_i]:
            add_oval(s, cx + Inches(0.28), by + Inches(0.06), Inches(0.11), Inches(0.11), GOLD)
            tb = add_text(s, cx + Inches(0.52), by, cw - Inches(0.75), Inches(0.7),
                           f"{m['num']}. {m['title']}", size=10.5, color=INK, line_spacing=1.1)
            tb.click_action.target_slide = module_target[m['num']]
            by += Inches(0.82)
        if arc_i == 3:
            tb = add_text(s, cx + Inches(0.52), by, cw - Inches(0.75), Inches(0.7),
                           "Capstone Project", size=10.5, bold=True, color=GOLD_DEEP, line_spacing=1.1)
            tb.click_action.target_slide = capstone_target
        cx += cw + Inches(0.18)

    add_footer(s, "Advanced Power BI  \u00b7  Program Overview", page_num)
    add_fade_transition(s)


def build_arc_divider(s, arc_index, page_num, module_target, cm_target):
    arc = MODEL['arcs'][arc_index]
    modules = [m for m in MODEL['modules'] if ARC_FOR_MODULE[m['num']] == arc_index]
    set_bg(s, NAVY)
    add_rect(s, 0, 0, Inches(0.12), SH, GOLD)
    add_progress_rail(s, arc_index, ARC_NAMES, Inches(0.6), Inches(0.55), Inches(5.6), dark=True)
    add_kicker(s, arc['label'], Inches(0.9), Inches(1.5), Inches(9.0), color=GOLD, size=13)
    add_text(s, Inches(0.85), Inches(1.95), Inches(10.6), Inches(1.3), arc['headline'], size=34, bold=True,
              color=WHITE, font=HEAD_FONT, line_spacing=1.1)
    add_text(s, Inches(0.9), Inches(3.1), Inches(9.0), Inches(1.0), arc['description'], size=13.5, color=ICE,
              line_spacing=1.35)

    cx = Inches(0.9)
    cw = Inches(3.55)
    for m in modules:
        add_rrect(s, cx, Inches(4.55), cw, Inches(1.55), NAVY_PANEL, radius=0.06, line_color=NAVY_LINE)
        add_rect(s, cx, Inches(4.55), Inches(0.07), Inches(1.55), GOLD)
        tb = add_text(s, cx + Inches(0.3), Inches(4.78), cw - Inches(0.55), Inches(1.05),
                       f"{m['num']:0>2} \u00b7 {m['title']}", size=13.5, bold=True, color=WHITE,
                       font=HEAD_FONT, line_spacing=1.15)
        tb.click_action.target_slide = module_target[m['num']]
        cx += cw + Inches(0.2)

    arc_short = arc['label'].split('\u00b7')[-1].strip()
    add_footer(s, f"Advanced Power BI  \u00b7  {arc_short}", page_num, dark=True, backlink_target=cm_target)
    add_fade_transition(s)


def build_module_overview(s, m, page_num, cm_target):
    set_bg(s, LIGHT_BG)
    arc_i = ARC_FOR_MODULE[m['num']]
    add_header(s, ARC_NAMES[arc_i], breadcrumb_right=f"Module {m['num']:02d} of 11")
    add_progress_rail(s, m['num'] - 1, [f"{i:02d}" for i in range(1, 12)], Inches(0.6), Inches(0.62), Inches(11.9))

    icon = MODULE_ICON.get(m['num'], "layers")
    add_icon(s, icon, Inches(0.6), Inches(1.15), Inches(0.85), badge_color=NAVY)
    add_text(s, Inches(1.65), Inches(1.12), Inches(1.1), Inches(0.85), f"{m['num']:02d}", size=30, bold=True,
              color=GOLD_DEEP, font=HEAD_FONT, anchor=MSO_ANCHOR.MIDDLE)
    add_text(s, Inches(2.75), Inches(1.1), Inches(7.5), Inches(0.55), m['title'], size=24, bold=True,
              color=INK, font=HEAD_FONT)
    add_text(s, Inches(2.75), Inches(1.62), Inches(7.5), Inches(0.4), m['description'], size=12, color=MUTED,
              italic=True)
    add_rrect(s, Inches(2.75), Inches(2.06), Inches(2.7), Inches(0.32), NAVY, radius=0.5)
    add_text(s, Inches(2.75), Inches(2.06), Inches(2.7), Inches(0.32), m['level'], size=9.5, bold=True,
              color=GOLD, align=PP_ALIGN.CENTER, anchor=MSO_ANCHOR.MIDDLE, letter_spacing=40)

    add_text(s, Inches(0.6), Inches(2.68), Inches(6.9), Inches(0.3), "WHAT LEARNERS WALK AWAY WITH",
              size=10.5, bold=True, color=GOLD_DEEP, letter_spacing=60)

    cards = m['cards']
    card_h = Inches(3.55) / len(cards) - Inches(0.12)
    cy = Inches(3.05)
    for title, body in cards:
        add_rrect(s, Inches(0.6), cy, Inches(6.9), card_h, CARD_BG, radius=0.08, line_color=LINE_GREY)
        add_rect(s, Inches(0.6), cy, Inches(0.06), card_h, GOLD)
        add_text(s, Inches(0.85), cy + Inches(0.08), Inches(6.4), Inches(0.3), title, size=12.5, bold=True,
                  color=INK, font=HEAD_FONT)
        add_text(s, Inches(0.85), cy + Inches(0.4), Inches(6.4), card_h - Inches(0.45), body, size=10,
                  color=MUTED, line_spacing=1.2)
        cy += card_h + Inches(0.14)

    draw_fn = DIAGRAM_FOR_MODULE.get(m['num'])
    if draw_fn:
        draw_fn(s, Inches(7.75), Inches(2.68), Inches(4.95), Inches(3.95))

    add_footer(s, f"Advanced Power BI  \u00b7  {m['title']}", page_num, backlink_target=cm_target)
    add_fade_transition(s)


def build_module_deepdive(s, m, page_num, cm_target):
    set_bg(s, LIGHT_BG)
    arc_i = ARC_FOR_MODULE[m['num']]
    add_header(s, ARC_NAMES[arc_i], breadcrumb_right=f"Module {m['num']:02d} of 11")
    add_text(s, Inches(0.6), Inches(0.9), Inches(10.5), Inches(0.5),
              f"{m['title']} \u00b7 Topic Guide", size=20, bold=True, color=INK, font=HEAD_FONT)
    add_text(s, Inches(0.6), Inches(1.4), Inches(11.0), Inches(0.35),
              f"{len(m['topics'])} topics covered in this module \u2014 full detail in speaker notes.",
              size=10.5, color=MUTED, italic=True)

    topics = m['topics']
    n = len(topics)
    nrows = (n + 1) // 2
    col_w = Inches(5.85)
    row_h = Inches(4.85) / nrows
    x0 = Inches(0.6)
    y0 = Inches(1.95)
    key_shapes = []
    for i, t in enumerate(topics):
        col = i % 2
        row = i // 2
        cx = x0 + col * (col_w + Inches(0.3))
        cy = y0 + row * (row_h + Inches(0.08))
        add_rrect(s, cx, cy, col_w, row_h - Inches(0.08), CARD_BG, radius=0.06, line_color=LINE_GREY)
        add_text(s, cx + Inches(0.22), cy + Inches(0.1), Inches(0.7), Inches(0.3), t['num'], size=10,
                  bold=True, color=GOLD_DEEP, font=HEAD_FONT)
        add_text(s, cx + Inches(0.22), cy + Inches(0.38), col_w - Inches(0.44), Inches(0.32), t['title'],
                  size=12.5, bold=True, color=INK, font=HEAD_FONT)
        add_text(s, cx + Inches(0.22), cy + Inches(0.74), col_w - Inches(0.44), row_h - Inches(0.9),
                  t['why'], size=9.5, color=MUTED, line_spacing=1.22)

    notes = s.notes_slide.notes_text_frame
    lines = [f"{m['title']} -- full topic detail:\n"]
    for t in topics:
        lines.append(f"{t['num']} {t['title']}")
        lines.append(f"Why it matters: {t['why']}")
        for b in t['bullets']:
            lines.append(f"  - {b}")
        lines.append("")
    notes.text = "\n".join(lines)

    add_footer(s, f"Advanced Power BI  \u00b7  {m['title']}", page_num, backlink_target=cm_target)
    add_fade_transition(s)


def build_capstone(s, page_num, cm_target):
    set_bg(s, NAVY)
    raw = MODEL['capstone_raw']
    add_rect(s, 0, 0, Inches(0.12), SH, GOLD)
    add_kicker(s, raw[0], Inches(0.9), Inches(0.7), Inches(9), color=GOLD, size=13)
    add_text(s, Inches(0.85), Inches(1.15), Inches(10.8), Inches(1.3), raw[1], size=30, bold=True,
              color=WHITE, font=HEAD_FONT, line_spacing=1.12)
    add_text(s, Inches(0.9), Inches(2.35), Inches(10.8), Inches(0.8), raw[2], size=13, color=ICE, line_spacing=1.3)

    steps = []
    for i in range(3, len(raw) - 3, 2):
        steps.append((raw[i], raw[i + 1]))

    cx = Inches(0.9)
    cw = Inches(2.78)
    for i, (title, body) in enumerate(steps):
        add_oval(s, cx, Inches(3.55), Inches(0.55), Inches(0.55), GOLD)
        add_text(s, cx, Inches(3.55), Inches(0.55), Inches(0.55), str(i + 1), size=16, bold=True,
                  color=NAVY, align=PP_ALIGN.CENTER, anchor=MSO_ANCHOR.MIDDLE, font=HEAD_FONT)
        add_text(s, cx, Inches(4.28), cw - Inches(0.2), Inches(0.4), title, size=14, bold=True, color=WHITE,
                  font=HEAD_FONT)
        add_text(s, cx, Inches(4.72), cw - Inches(0.2), Inches(1.3), body, size=10.5, color=ICE, line_spacing=1.3)
        if i < len(steps) - 1:
            add_connector(s, cx + cw - Inches(0.2), Inches(3.82), cx + cw + Inches(0.05), Inches(3.82),
                           color=GOLD, width=Pt(1.5))
        cx += cw + Inches(0.2)

    add_footer(s, "Advanced Power BI  \u00b7  Capstone", page_num, dark=True, backlink_target=cm_target)
    add_fade_transition(s)


def build_closing(s, page_num, cm_target):
    set_bg(s, NAVY)
    d = MODEL['closing']
    add_oval(s, -Inches(1.5), SH - Inches(2.5), Inches(4.0), Inches(4.0), NAVY_PANEL)
    add_kicker(s, d['eyebrow'], Inches(0.9), Inches(2.2), Inches(9), color=GOLD, size=13)
    add_text(s, Inches(0.85), Inches(2.65), Inches(10.8), Inches(1.5), d['headline'], size=32, bold=True,
              color=WHITE, font=HEAD_FONT, line_spacing=1.15)
    add_rect(s, Inches(0.92), Inches(3.95), Inches(1.4), Inches(0.045), GOLD)
    add_text(s, Inches(0.9), Inches(4.2), Inches(9.6), Inches(1.0), d['body'], size=13.5, color=ICE, line_spacing=1.35)
    add_footer(s, "Advanced Power BI  \u00b7  Program Overview", page_num, dark=True, backlink_target=cm_target)
    add_fade_transition(s)


# ============================================================ assembly
# Pass 1: create every slide (blank) in final order so we have stable
# python-pptx Slide objects to use as click-through hyperlink targets.

order = []  # list of (key, kind, payload)
order.append(('title', 'title', None))
order.append(('instructor', 'instructor', None))
order.append(('why', 'why', None))
order.append(('curriculum_map', 'curriculum_map', None))

modules_by_arc = {0: [], 1: [], 2: [], 3: []}
for m in MODEL['modules']:
    modules_by_arc[ARC_FOR_MODULE[m['num']]].append(m)

for arc_i in range(4):
    order.append((f'arc_{arc_i}', 'arc', arc_i))
    for m in modules_by_arc[arc_i]:
        order.append((f'module_overview_{m["num"]}', 'module_overview', m))
        order.append((f'module_deepdive_{m["num"]}', 'module_deepdive', m))

order.append(('capstone', 'capstone', None))
order.append(('closing', 'closing', None))

slide_objs = {}
for key, kind, payload in order:
    slide_objs[key] = add_slide()

module_target = {m['num']: slide_objs[f'module_overview_{m["num"]}'] for m in MODEL['modules']}
cm_target = slide_objs['curriculum_map']

# Pass 2: populate content into each already-created slide, in order, so
# page numbers line up with final slide order.
for page_num, (key, kind, payload) in enumerate(order, start=1):
    s = slide_objs[key]
    if kind == 'title':
        build_title(s)
    elif kind == 'instructor':
        build_instructor(s, page_num)
    elif kind == 'why':
        build_why(s, page_num)
    elif kind == 'curriculum_map':
        build_curriculum_map(s, page_num, module_target, slide_objs['capstone'])
    elif kind == 'arc':
        build_arc_divider(s, payload, page_num, module_target, cm_target)
    elif kind == 'module_overview':
        build_module_overview(s, payload, page_num, cm_target)
    elif kind == 'module_deepdive':
        build_module_deepdive(s, payload, page_num, cm_target)
    elif kind == 'capstone':
        build_capstone(s, page_num, cm_target)
    elif kind == 'closing':
        build_closing(s, page_num, cm_target)

prs.save(OUT_PATH)
print(f"Saved {len(order)} slides to {OUT_PATH}")
