"""Advanced Power BI -- Program Overview Deck

Original content deck (not derived from the internal training-outline docs)
covering the arc of an advanced Power BI curriculum: modeling, DAX, Power
Query, report UX, performance, analytics/AI, security, Service deployment,
governance, platform architecture, DevOps, and a capstone.

Visual language: "Signal & Slate" palette -- deep slate/navy base with a
electric-green/teal accent (distinct from the navy/cyan datasheet branding),
dark title/section slides, light content slides, one repeated motif: icons in
a rounded accent-colored badge, thin left-edge accent bars on cards.
"""
from pptx import Presentation
from pptx.util import Inches, Pt, Emu
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR
from pptx.enum.shapes import MSO_SHAPE
from pptx.oxml.ns import qn
import os
import copy

HERE = os.path.dirname(__file__)
ICON_DIR = os.path.join(HERE, "assets", "icons")
OUT_PATH = os.path.join(HERE, "Advanced-PowerBI-Program-Overview.pptx")

# ---------------------------------------------------------------- palette
SLATE_DARK = RGBColor(0x0E, 0x17, 0x24)     # near-black slate (dark bg)
SLATE = RGBColor(0x1B, 0x2A, 0x3D)          # panel slate
SLATE_MED = RGBColor(0x2D, 0x41, 0x59)      # card / divider
INK = RGBColor(0x11, 0x18, 0x22)            # body text on light bg
MUTED = RGBColor(0x5B, 0x6B, 0x7E)          # secondary text
ACCENT = RGBColor(0x3D, 0xDC, 0x97)         # electric green/teal
ACCENT_DEEP = RGBColor(0x10, 0x9A, 0x6A)    # deeper green for light-bg text
GOLD = RGBColor(0xE8, 0xB4, 0x4D)           # warm secondary accent
LIGHT_BG = RGBColor(0xF3, 0xF6, 0xF9)       # content slide background
CARD_BG = RGBColor(0xFF, 0xFF, 0xFF)
WHITE = RGBColor(0xFF, 0xFF, 0xFF)
ICE = RGBColor(0xC9, 0xDA, 0xE8)

HEAD_FONT = "Cambria"
BODY_FONT = "Calibri"

prs = Presentation()
prs.slide_width = Inches(13.333)
prs.slide_height = Inches(7.5)
BLANK = prs.slide_layouts[6]
SW, SH = prs.slide_width, prs.slide_height


def add_slide():
    return prs.slides.add_slide(BLANK)


def set_bg(slide, color):
    slide.background.fill.solid()
    slide.background.fill.fore_color.rgb = color


def no_line(shape):
    shape.line.fill.background()


def add_rect(slide, x, y, w, h, color, line=False):
    shp = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, x, y, w, h)
    shp.fill.solid()
    shp.fill.fore_color.rgb = color
    if not line:
        no_line(shp)
    shp.shadow.inherit = False
    return shp


def add_rrect(slide, x, y, w, h, color, radius=0.06, line=False, shadow=False):
    shp = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, x, y, w, h)
    try:
        shp.adjustments[0] = radius
    except Exception:
        pass
    shp.fill.solid()
    shp.fill.fore_color.rgb = color
    if not line:
        no_line(shp)
    shp.shadow.inherit = False
    return shp


def add_text(slide, x, y, w, h, text, size=14, color=INK, bold=False, italic=False,
             align=PP_ALIGN.LEFT, font=BODY_FONT, anchor=MSO_ANCHOR.TOP, line_spacing=1.0,
             char_spacing=None):
    tb = slide.shapes.add_textbox(x, y, w, h)
    tf = tb.text_frame
    tf.word_wrap = True
    tf.vertical_anchor = anchor
    tf.margin_left = 0
    tf.margin_right = 0
    tf.margin_top = 0
    tf.margin_bottom = 0
    lines = text.split("\n") if isinstance(text, str) else text
    for i, line in enumerate(lines):
        p = tf.paragraphs[0] if i == 0 else tf.add_paragraph()
        p.alignment = align
        p.line_spacing = line_spacing
        r = p.add_run()
        r.text = line
        r.font.size = Pt(size)
        r.font.bold = bold
        r.font.italic = italic
        r.font.color.rgb = color
        r.font.name = font
    return tb


def add_bullets(slide, x, y, w, h, items, size=13, color=INK, font=BODY_FONT,
                 space_after=8, bullet_color=None, line_spacing=1.08):
    tb = slide.shapes.add_textbox(x, y, w, h)
    tf = tb.text_frame
    tf.word_wrap = True
    tf.margin_left = 0
    tf.margin_right = 0
    tf.margin_top = 0
    tf.margin_bottom = 0
    for i, item in enumerate(items):
        p = tf.paragraphs[0] if i == 0 else tf.add_paragraph()
        p.line_spacing = line_spacing
        p.space_after = Pt(space_after)
        pPr = p._pPr
        if pPr is None:
            pPr = p.get_or_add_pPr()
        buChar = pPr.makeelement(qn('a:buChar'), {'char': '\u2013'})
        buFont = pPr.makeelement(qn('a:buFont'), {'typeface': font})
        buClr = pPr.makeelement(qn('a:buClr'), {})
        srgb = pPr.makeelement(qn('a:srgbClr'), {'val': '%02X%02X%02X' % ((bullet_color or ACCENT_DEEP)[0], (bullet_color or ACCENT_DEEP)[1], (bullet_color or ACCENT_DEEP)[2])})
        buClr.append(srgb)
        pPr.append(buClr)
        pPr.append(buFont)
        pPr.append(buChar)
        pPr.set('marL', '182880')
        pPr.set('indent', '-182880')
        r = p.add_run()
        r.text = item
        r.font.size = Pt(size)
        r.font.color.rgb = color
        r.font.name = font
    return tb


def add_icon(slide, icon_name, x, y, size, badge_color=ACCENT, badge=True):
    if badge:
        b = slide.shapes.add_shape(MSO_SHAPE.OVAL, x, y, size, size)
        b.fill.solid()
        b.fill.fore_color.rgb = badge_color
        no_line(b)
        b.shadow.inherit = False
    pad = size * 0.22
    slide.shapes.add_picture(os.path.join(ICON_DIR, f"{icon_name}.png"), x + pad, y + pad, size - 2 * pad, size - 2 * pad)


def add_page_footer(slide, label, page_num, dark=False):
    color = ICE if dark else MUTED
    add_text(slide, Inches(0.5), Inches(7.14), Inches(8), Inches(0.3),
              f"Advanced Power BI  ·  {label}", size=9, color=color, font=BODY_FONT)
    add_text(slide, Inches(12.4), Inches(7.14), Inches(0.5), Inches(0.3),
              str(page_num), size=9, color=color, align=PP_ALIGN.RIGHT, font=BODY_FONT)


def set_notes(slide, text):
    notes_slide = slide.notes_slide
    tf = notes_slide.notes_text_frame
    tf.text = text


MODULE_COLOR = ACCENT
PAGE = [0]


def next_page():
    PAGE[0] += 1
    return PAGE[0]


# ============================================================== SLIDE 1: TITLE
def slide_title():
    s = add_slide()
    set_bg(s, SLATE_DARK)
    # background geometric accents
    add_rect(s, Emu(0), Emu(0), SW, Inches(0.12), ACCENT)
    circle = s.shapes.add_shape(MSO_SHAPE.OVAL, Inches(9.6), Inches(-2.2), Inches(6), Inches(6))
    circle.fill.solid()
    circle.fill.fore_color.rgb = SLATE
    no_line(circle)
    circle.shadow.inherit = False
    circle2 = s.shapes.add_shape(MSO_SHAPE.OVAL, Inches(10.6), Inches(-1.0), Inches(3.6), Inches(3.6))
    circle2.fill.solid()
    circle2.fill.fore_color.rgb = SLATE_MED
    no_line(circle2)
    circle2.shadow.inherit = False

    add_text(s, Inches(0.9), Inches(1.55), Inches(7.5), Inches(0.5),
              "A DATA STORYTELLING & GOVERNANCE CURRICULUM", size=13, color=ACCENT,
              bold=True, font=BODY_FONT, char_spacing=2)
    add_text(s, Inches(0.85), Inches(2.05), Inches(10.8), Inches(2.0),
              "Advanced Power BI", size=58, color=WHITE, bold=True, font=HEAD_FONT)
    add_text(s, Inches(0.9), Inches(3.35), Inches(9.5), Inches(0.9),
              "From modeling fundamentals to enterprise-grade delivery: an 11-module\narc for analysts, BI developers, and platform owners who need Power BI\nto hold up under real production weight.",
              size=15, color=ICE, font=BODY_FONT, line_spacing=1.3)

    # stat strip
    stats = [("11", "Curriculum modules"), ("40+", "Guided exercises"), ("1", "Capstone solution")]
    x = Inches(0.9)
    for label_num, label_txt in stats:
        add_text(s, x, Inches(5.15), Inches(2.6), Inches(0.7), label_num, size=34, color=ACCENT, bold=True, font=HEAD_FONT)
        add_text(s, x, Inches(5.85), Inches(2.6), Inches(0.4), label_txt, size=11.5, color=ICE, font=BODY_FONT)
        x += Inches(2.85)

    add_text(s, Inches(0.9), Inches(6.9), Inches(6), Inches(0.4), "Program Overview", size=10.5, color=MUTED, font=BODY_FONT)


# ============================================================== SLIDE 2: WHY THIS PROGRAM EXISTS
def slide_why():
    s = add_slide()
    set_bg(s, LIGHT_BG)
    add_rect(s, Emu(0), Emu(0), SW, Inches(0.12), ACCENT)
    add_text(s, Inches(0.7), Inches(0.5), Inches(10), Inches(0.4), "WHY THIS PROGRAM EXISTS", size=12, color=ACCENT_DEEP, bold=True, font=BODY_FONT, char_spacing=1.5)
    add_text(s, Inches(0.7), Inches(0.85), Inches(11.5), Inches(0.8), "Most Power BI training stops at chart-building.\nThis one doesn't.", size=28, color=INK, bold=True, font=HEAD_FONT, line_spacing=1.1)

    cards = [
        ("gauge", "Reports that survive contact with real data", "Large fact tables, messy sources, and slow refreshes expose weak modeling fast. This program builds habits that hold up at scale, not just in a demo."),
        ("shield", "Governance is part of the design, not an afterthought", "Row-level security, workspace strategy, and monitoring are taught alongside DAX and visuals -- because a report nobody trusts doesn't get used."),
        ("gear", "A path from author to platform owner", "Learners move from single-report skills toward the operational and architectural judgment that senior BI roles actually require."),
    ]
    x = Inches(0.7)
    w = Inches(3.85)
    for icon, title, body in cards:
        add_rrect(s, x, Inches(2.15), w, Inches(4.3), CARD_BG, radius=0.05, shadow=True)
        add_rect(s, x, Inches(2.15), Inches(0.09), Inches(4.3), ACCENT)
        add_icon(s, icon, x + Inches(0.35), Inches(2.5), Inches(0.9), badge_color=SLATE_DARK)
        add_text(s, x + Inches(0.35), Inches(3.65), w - Inches(0.7), Inches(1.0), title, size=15.5, color=INK, bold=True, font=HEAD_FONT, line_spacing=1.15)
        add_text(s, x + Inches(0.35), Inches(4.55), w - Inches(0.7), Inches(1.8), body, size=11.5, color=MUTED, font=BODY_FONT, line_spacing=1.35)
        x += w + Inches(0.25)

    add_page_footer(s, "Program Overview", next_page())


# ============================================================== SLIDE 3: CURRICULUM MAP (dark, timeline)
def slide_curriculum_map():
    s = add_slide()
    set_bg(s, SLATE_DARK)
    add_text(s, Inches(0.7), Inches(0.45), Inches(10), Inches(0.4), "CURRICULUM MAP", size=12, color=ACCENT, bold=True, font=BODY_FONT, char_spacing=1.5)
    add_text(s, Inches(0.7), Inches(0.8), Inches(11), Inches(0.7), "Four arcs, eleven modules, one capstone", size=26, color=WHITE, bold=True, font=HEAD_FONT)

    arcs = [
        ("BUILD", ["1. Advanced Semantic Modeling", "2. Advanced DAX", "3. Advanced Power Query"], ACCENT),
        ("PRESENT", ["4. Report Design & UX", "5. Performance Optimization", "6. Advanced Analytics & AI"], GOLD),
        ("PROTECT & OPERATE", ["7. Security Design", "8. Service Deployment", "9. Monitoring & Governance"], RGBColor(0x6E, 0xA8, 0xE0)),
        ("SCALE", ["10. Capacity & Architecture", "11. DevOps & Lifecycle", "Capstone Project"], RGBColor(0xE0, 0x7A, 0x5F)),
    ]
    x = Inches(0.7)
    w = Inches(2.95)
    for label, items, color in arcs:
        add_rrect(s, x, Inches(1.75), w, Inches(4.9), SLATE, radius=0.06)
        add_rect(s, x, Inches(1.75), w, Inches(0.55), color)
        add_text(s, x, Inches(1.75), w, Inches(0.55), label, size=13, color=SLATE_DARK, bold=True, align=PP_ALIGN.CENTER, anchor=MSO_ANCHOR.MIDDLE, font=BODY_FONT)
        yy = Inches(2.5)
        for item in items:
            dot = s.shapes.add_shape(MSO_SHAPE.OVAL, x + Inches(0.3), yy + Inches(0.08), Inches(0.12), Inches(0.12))
            dot.fill.solid(); dot.fill.fore_color.rgb = color; no_line(dot); dot.shadow.inherit = False
            add_text(s, x + Inches(0.55), yy, w - Inches(0.75), Inches(0.9), item, size=12, color=ICE, font=BODY_FONT, line_spacing=1.15)
            yy += Inches(1.35)
        x += w + Inches(0.2)

    add_page_footer(s, "Program Overview", next_page(), dark=True)


# ============================================================== MODULE SLIDE TEMPLATE
def slide_module(num, title, icon, tagline, points, skill_level, arc_label, arc_color, notes=None):
    s = add_slide()
    set_bg(s, LIGHT_BG)
    # left rail
    add_rect(s, Emu(0), Emu(0), Inches(4.3), SH, SLATE_DARK)
    add_rect(s, Inches(4.3), Emu(0), Inches(0.06), SH, arc_color)

    add_text(s, Inches(0.55), Inches(0.55), Inches(3), Inches(0.4), arc_label, size=11, color=arc_color, bold=True, font=BODY_FONT, char_spacing=1.5)
    add_text(s, Inches(0.55), Inches(0.95), Inches(1.6), Inches(0.9), f"{num:02d}", size=52, color=WHITE, bold=True, font=HEAD_FONT)
    add_icon(s, icon, Inches(0.55), Inches(2.1), Inches(1.15), badge_color=arc_color)
    add_text(s, Inches(0.55), Inches(3.5), Inches(3.3), Inches(1.6), title, size=22, color=WHITE, bold=True, font=HEAD_FONT, line_spacing=1.15)
    add_text(s, Inches(0.55), Inches(5.15), Inches(3.3), Inches(1.5), tagline, size=12.5, color=ICE, italic=True, font=BODY_FONT, line_spacing=1.35)
    add_text(s, Inches(0.55), Inches(6.55), Inches(3.3), Inches(0.4), f"LEVEL: {skill_level}", size=10, color=arc_color, bold=True, font=BODY_FONT, char_spacing=1.2)

    add_text(s, Inches(4.75), Inches(0.55), Inches(8), Inches(0.4), "WHAT LEARNERS WALK AWAY WITH", size=12, color=ACCENT_DEEP, bold=True, font=BODY_FONT, char_spacing=1.2)
    yy = Inches(1.15)
    for head, body in points:
        add_rrect(s, Inches(4.75), yy, Inches(8.0), Inches(1.35), CARD_BG, radius=0.08, shadow=True)
        add_rect(s, Inches(4.75), yy, Inches(0.08), Inches(1.35), arc_color)
        add_text(s, Inches(5.05), yy + Inches(0.14), Inches(7.4), Inches(0.4), head, size=13.5, color=INK, bold=True, font=HEAD_FONT)
        add_text(s, Inches(5.05), yy + Inches(0.55), Inches(7.4), Inches(0.75), body, size=11, color=MUTED, font=BODY_FONT, line_spacing=1.25)
        yy += Inches(1.55)

    add_page_footer(s, title, next_page())
    if notes:
        set_notes(s, notes)
    return s



# ============================================================== SUB-TOPIC SLIDE TEMPLATE
def slide_subtopic(module_num, module_title, topic_num, topic_count, topic_title, icon,
                    why_matters, key_points, arc_label, arc_color, notes):
    s = add_slide()
    set_bg(s, LIGHT_BG)
    add_rect(s, Emu(0), Emu(0), SW, Inches(0.85), SLATE_DARK)
    add_text(s, Inches(0.6), Inches(0.18), Inches(6), Inches(0.35),
              f"MODULE {module_num:02d} · {module_title.upper()}", size=11, color=arc_color, bold=True, font=BODY_FONT, char_spacing=1.2)
    add_text(s, Inches(0.6), Inches(0.46), Inches(9), Inches(0.35),
              f"Topic {topic_num} of {topic_count}", size=10, color=ICE, font=BODY_FONT)
    add_text(s, Inches(11.3), Inches(0.18), Inches(1.6), Inches(0.5),
              f"{module_num:02d}.{topic_num}", size=20, color=WHITE, bold=True, align=PP_ALIGN.RIGHT, font=HEAD_FONT)

    add_icon(s, icon, Inches(0.6), Inches(1.15), Inches(0.9), badge_color=arc_color)
    add_text(s, Inches(1.75), Inches(1.2), Inches(10.6), Inches(0.9), topic_title, size=25, color=INK, bold=True, font=HEAD_FONT, anchor=MSO_ANCHOR.MIDDLE)

    add_rrect(s, Inches(0.6), Inches(2.3), Inches(12.1), Inches(1.15), CARD_BG, radius=0.08)
    add_rect(s, Inches(0.6), Inches(2.3), Inches(0.08), Inches(1.15), arc_color)
    add_text(s, Inches(0.95), Inches(2.42), Inches(2.6), Inches(0.4), "WHY IT MATTERS", size=11, color=ACCENT_DEEP, bold=True, font=BODY_FONT, char_spacing=1.2)
    add_text(s, Inches(0.95), Inches(2.78), Inches(11.4), Inches(0.6), why_matters, size=12.5, color=INK, font=BODY_FONT, line_spacing=1.25)

    add_text(s, Inches(0.6), Inches(3.75), Inches(6), Inches(0.4), "KEY POINTS", size=11, color=ACCENT_DEEP, bold=True, font=BODY_FONT, char_spacing=1.2)
    add_bullets(s, Inches(0.6), Inches(4.2), Inches(11.9), Inches(2.6), key_points, size=13.5, color=INK,
                font=BODY_FONT, space_after=12, bullet_color=arc_color, line_spacing=1.2)

    add_page_footer(s, f"{module_title} · {topic_title}", next_page())
    set_notes(s, notes)
    return s


# ============================================================== SECTION DIVIDER
def slide_section(arc_label, arc_title, arc_desc, arc_color, module_titles):
    s = add_slide()
    set_bg(s, SLATE_DARK)
    add_rect(s, Emu(0), Emu(0), SW, Inches(0.12), arc_color)
    add_text(s, Inches(0.9), Inches(1.5), Inches(6), Inches(0.4), f"ARC · {arc_label}", size=13, color=arc_color, bold=True, font=BODY_FONT, char_spacing=2)
    add_text(s, Inches(0.85), Inches(2.0), Inches(9.5), Inches(1.3), arc_title, size=40, color=WHITE, bold=True, font=HEAD_FONT, line_spacing=1.05)
    add_text(s, Inches(0.9), Inches(3.35), Inches(8.3), Inches(1.0), arc_desc, size=14, color=ICE, font=BODY_FONT, line_spacing=1.35)

    x = Inches(0.9)
    w = Inches(3.75)
    for t in module_titles:
        add_rrect(s, x, Inches(4.7), w, Inches(1.5), SLATE, radius=0.08)
        add_rect(s, x, Inches(4.7), Inches(0.07), Inches(1.5), arc_color)
        add_text(s, x + Inches(0.3), Inches(4.95), w - Inches(0.5), Inches(1.0), t, size=13.5, color=WHITE, bold=True, font=BODY_FONT, line_spacing=1.2, anchor=MSO_ANCHOR.MIDDLE)
        x += w + Inches(0.2)

    add_page_footer(s, arc_label, next_page(), dark=True)


# ============================================================== CAPSTONE SLIDE
def slide_capstone():
    s = add_slide()
    set_bg(s, SLATE_DARK)
    add_rect(s, Emu(0), Emu(0), SW, Inches(0.12), GOLD)
    add_icon(s, "rocket", Inches(0.9), Inches(1.15), Inches(1.5), badge_color=GOLD)
    add_text(s, Inches(2.75), Inches(1.2), Inches(4), Inches(0.4), "CAPSTONE", size=13, color=GOLD, bold=True, font=BODY_FONT, char_spacing=2)
    add_text(s, Inches(2.7), Inches(1.6), Inches(9.5), Inches(0.9), "One dataset. One deadline.\nEvery skill from Modules 1-11.", size=30, color=WHITE, bold=True, font=HEAD_FONT, line_spacing=1.1)

    add_text(s, Inches(0.9), Inches(3.05), Inches(11.5), Inches(0.7),
              "Learners take a single messy source dataset from raw files to a governed, secured, production-ready Power BI solution -- end to end, without a script to follow.",
              size=14, color=ICE, font=BODY_FONT, line_spacing=1.3)

    steps = [
        ("1", "Model it", "Build a star schema, handle grain and role-playing dates, and validate relationships."),
        ("2", "Calculate it", "Write the DAX measures the business actually asked for, not the easy ones."),
        ("3", "Present it", "Design a navigable, accessible report with drillthrough and a mobile layout."),
        ("4", "Ship it", "Publish, secure with RLS, and document it like it has to survive an audit."),
    ]
    x = Inches(0.9)
    w = Inches(2.75)
    for num, head, body in steps:
        add_rrect(s, x, Inches(4.0), w, Inches(2.6), SLATE, radius=0.06)
        add_text(s, x + Inches(0.3), Inches(4.2), Inches(1), Inches(0.7), num, size=30, color=GOLD, bold=True, font=HEAD_FONT)
        add_text(s, x + Inches(0.3), Inches(4.95), w - Inches(0.6), Inches(0.4), head, size=15, color=WHITE, bold=True, font=HEAD_FONT)
        add_text(s, x + Inches(0.3), Inches(5.4), w - Inches(0.6), Inches(1.1), body, size=10.5, color=ICE, font=BODY_FONT, line_spacing=1.25)
        x += w + Inches(0.2)

    add_page_footer(s, "Capstone", next_page(), dark=True)


# ============================================================== CLOSING SLIDE
def slide_closing():
    s = add_slide()
    set_bg(s, SLATE_DARK)
    circle = s.shapes.add_shape(MSO_SHAPE.OVAL, Inches(-2.2), Inches(3.6), Inches(6), Inches(6))
    circle.fill.solid(); circle.fill.fore_color.rgb = SLATE; no_line(circle); circle.shadow.inherit = False
    add_rect(s, Emu(0), Emu(0), SW, Inches(0.12), ACCENT)

    add_text(s, Inches(0.9), Inches(2.2), Inches(10), Inches(0.4), "WHERE LEARNERS END UP", size=12, color=ACCENT, bold=True, font=BODY_FONT, char_spacing=2)
    add_text(s, Inches(0.85), Inches(2.65), Inches(11), Inches(1.6), "Confident modeling, defensible\nnumbers, and reports people trust.", size=36, color=WHITE, bold=True, font=HEAD_FONT, line_spacing=1.15)
    add_text(s, Inches(0.9), Inches(4.35), Inches(9), Inches(0.9),
              "Advanced Power BI is built for teams who are past the tutorials and need\nthe judgment, patterns, and habits that hold up in production.",
              size=14, color=ICE, font=BODY_FONT, line_spacing=1.35)

    add_page_footer(s, "Program Overview", next_page(), dark=True)


# ================================================================ BUILD DECK
slide_title()
set_notes(prs.slides[-1],
    "Welcome slide. Frame the program in one sentence before diving into detail: this is an "
    "advanced curriculum for people who already know the Power BI basics and need to make Power BI "
    "hold up in production -- at scale, under governance, and under real business pressure. Call out "
    "the three numbers (11 modules, 40+ exercises, 1 capstone) as the shape of the day(s) ahead.")

slide_why()
set_notes(prs.slides[-1],
    "This slide exists to set expectations before the curriculum map. Most Power BI training stops at "
    "'how do I build a visual' -- this program assumes that skill already exists and pushes into the "
    "three things that separate a hobby report from an enterprise one: it survives real data volume and "
    "mess, it treats governance as a design constraint from day one instead of a bolt-on, and it grows "
    "the learner from single-report author toward the judgment a platform owner needs. Use this to "
    "address any 'why do I need all this, I just want to build reports' skepticism up front.")

slide_curriculum_map()
set_notes(prs.slides[-1],
    "Orient the audience to the shape of the whole program before going deep on any one module. Four "
    "arcs: BUILD (get the data and model right), PRESENT (turn it into something people want to use), "
    "PROTECT & OPERATE (make it trustworthy and durable), SCALE (make it an organizational asset, not "
    "a single hero report). Every module slide that follows will call back to its arc color so the "
    "audience always knows where they are in the arc.")


# ---------------------------------------------------------------- ARC DATA
ARC_COLOR = {
    "BUILD": ACCENT,
    "PRESENT": GOLD,
    "PROTECT & OPERATE": RGBColor(0x6E, 0xA8, 0xE0),
    "SCALE": RGBColor(0xE0, 0x7A, 0x5F),
}


def build_module(num, title, icon, tagline, points, skill_level, arc_label, module_notes, subtopics):
    """points: list of (head, body) triples used on the module summary slide.
    subtopics: list of dicts with keys: title, icon, why, key_points, notes."""
    color = ARC_COLOR[arc_label]
    slide_module(num, title, icon, tagline, points, skill_level, arc_label, color, notes=module_notes)
    for i, st in enumerate(subtopics, start=1):
        slide_subtopic(
            num, title, i, len(subtopics), st["title"], st["icon"],
            st["why"], st["key_points"], arc_label, color, st["notes"],
        )


# ================================================================ ARC: BUILD
slide_section(
    "BUILD", "Model it right before you calculate anything",
    "The foundation arc: turning raw, flat data into a trustworthy semantic model and a clean transformation layer -- the two things every later module depends on.",
    ACCENT,
    ["01 · Advanced Semantic Modeling", "02 · Advanced DAX", "03 · Advanced Power Query"],
)
set_notes(prs.slides[-1],
    "Set up the BUILD arc: everything downstream -- every measure, every visual, every performance "
    "optimization -- inherits the quality of the model underneath it. If you get the grain, "
    "relationships, and transformation layer wrong here, every later module gets harder, not easier. "
    "This is the arc worth slowing down for.")

build_module(
    1, "Advanced Semantic Modeling", "layers",
    "Turning a flat file into a model that scales, and a schema that survives new requirements.",
    [
        ("Star schema from first principles", "Fact and dimension design, grain decisions, and why a wide flat table eventually breaks every report built on it."),
        ("Relationships that don't fight you", "Role-playing dimensions, bridge tables, and composite/DirectQuery model choices with their tradeoffs."),
        ("Models built for growth", "Handling large fact tables, incremental refresh readiness, and keeping a model maintainable as it scales."),
    ],
    "Intermediate → Advanced", "BUILD",
    "Module overview for Advanced Semantic Modeling. The core message: a flat file feels easier at "
    "first, but it silently caps what the report can ever answer well. This module is where learners "
    "make the mental shift from 'a table I can filter' to 'a model I can ask new questions of.' Three "
    "sub-topic slides follow, each expanding one of these three walk-away points.",
    [
        {
            "title": "Star Schema From First Principles",
            "icon": "layers",
            "why": "A correct grain and a clean fact/dimension split is the single decision that determines whether every later measure is easy or painful to write.",
            "key_points": [
                "Grain first: decide what one row of the fact table represents before building anything else.",
                "Facts hold numbers and foreign keys; dimensions hold the descriptive context used to slice them.",
                "A wide flat table works for a demo, but duplicates dimension data and breaks distinct-count and filter logic at scale.",
                "Refactoring a flat table into a star schema is demonstrated live, not just described.",
            ],
            "notes": (
                "This is the single most important idea in the whole program, and it's worth spending real time here. "
                "Most learners have only ever worked from a flat export, so the star schema feels like unnecessary "
                "complexity at first. Make it concrete: show what happens to a distinct customer count or an "
                "average-per-order measure when the same customer or order appears on multiple rows because the "
                "table is flat. Once they see a number silently inflate or shrink, the case for grain discipline "
                "makes itself. This is also the moment to say explicitly: every DAX problem you hit in Module 2 "
                "traces back to a modeling decision made here."
            ),
        },
        {
            "title": "Relationships That Don't Fight You",
            "icon": "layers",
            "why": "Most 'DAX bugs' reported by learners later in the program are actually relationship design problems -- ambiguous filter paths, wrong cardinality, or a missing role-playing dimension.",
            "key_points": [
                "Role-playing dimensions (e.g., order date vs. ship date) need their own dedicated date tables and active relationships.",
                "Bridge tables resolve many-to-many relationships without duplicating fact rows.",
                "Composite models mix Import and DirectQuery deliberately -- know which tables need which mode and why.",
                "Bidirectional filtering is a tool for specific problems, not a default setting.",
            ],
            "notes": (
                "Learners tend to reach for bidirectional cross-filtering to 'fix' a relationship problem instead of "
                "diagnosing the actual cause. Use this slide to build the habit of asking 'what question is this "
                "relationship supposed to answer' before touching a filter direction toggle. The role-playing date "
                "table pattern is worth a live demo -- show one fact table needing both an order date and a ship "
                "date, and why duplicating the date dimension (rather than reusing one with two active paths) is "
                "the sustainable answer."
            ),
        },
        {
            "title": "Models Built For Growth",
            "icon": "gauge",
            "why": "A model that works great with 50,000 test rows can fall over with 50 million production rows -- and by then it's expensive to fix.",
            "key_points": [
                "Large fact tables need a plan for incremental refresh before they need it operationally.",
                "Cardinality and column choices made early determine how much the model can grow before performance suffers.",
                "Composite and aggregation strategies let a model stay fast without giving up detail entirely.",
                "'Maintainable' means the next person on the team can understand the model too, not just the original author.",
            ],
            "notes": (
                "This sub-topic is the bridge into Module 5 (Performance Optimization) later in the program -- "
                "flag that connection explicitly so the audience sees the curriculum is intentionally sequenced. "
                "The key point to land: performance problems are cheap to prevent at modeling time and expensive "
                "to fix after a report is already in production with real users depending on it. If time allows, "
                "ask the group to estimate their own largest fact table's likely row count in two years -- it "
                "usually reframes 'this seems like overkill for our data' pretty quickly."
            ),
        },
    ],
)

build_module(
    2, "Advanced DAX", "formula",
    "Getting past copy-pasted measures into real command of filter and row context.",
    [
        ("Context, demystified", "Filter context vs. row context, and why the same measure can return different numbers on two visuals."),
        ("Patterns that hold up under scrutiny", "Time intelligence, semi-additive measures, ranking, and calculation groups -- built and explained, not memorized."),
        ("Measures that perform", "Diagnosing slow DAX, using variables deliberately, and validating logic before layering complexity on top."),
    ],
    "Advanced", "BUILD",
    "Module overview for Advanced DAX. The goal is not to teach more functions -- it's to teach the "
    "evaluation model underneath them, so learners can reason about any measure instead of pattern-"
    "matching from a blog post. Three sub-topics follow: context, proven patterns, and performance.",
    [
        {
            "title": "Context, Demystified",
            "icon": "formula",
            "why": "Almost every 'my measure gives a wrong number' problem is actually a filter-context or row-context misunderstanding -- once this clicks, most DAX bugs become diagnosable in seconds.",
            "key_points": [
                "Row context exists inside iterators (SUMX, FILTER); filter context comes from visuals, slicers, and CALCULATE.",
                "CALCULATE is the tool that transitions row context into filter context -- this is the concept, not the syntax, to teach.",
                "The same measure can return different totals in a card vs. a matrix because the filter context is different, not because the measure is broken.",
                "Test every new measure in a simple table visual before trusting it in a complex one.",
            ],
            "notes": (
                "This is the slide to slow down on, more than any other single sub-topic in the deck. Use a live "
                "example: put the same measure in a card, a matrix by category, and a matrix by category and year, "
                "and ask the class to predict the numbers before you show them. The goal is for learners to build "
                "an internal mental model of context transition, not to memorize that CALCULATE 'changes filters.' "
                "If you only have time to teach one concept deeply in the whole DAX module, make it this one."
            ),
        },
        {
            "title": "Patterns That Hold Up Under Scrutiny",
            "icon": "gauge",
            "why": "Time intelligence and ranking measures are the ones stakeholders actually stare at during a business review -- they need to be right and explainable, not just plausible-looking.",
            "key_points": [
                "Time intelligence functions rely on a proper, contiguous marked date table -- this is where Module 1's foundation pays off.",
                "Semi-additive measures (e.g., ending balance) need explicit handling; they don't sum correctly across time by default.",
                "Dynamic Top N/ranking measures should be built with RANKX and validated against a static, hand-checked example.",
                "Calculation groups reduce measure duplication when the same time-intelligence logic applies to many base measures.",
            ],
            "notes": (
                "Anchor this slide with a business scenario, not just a function list: a finance stakeholder asking "
                "'what were our top 5 products by margin last quarter, and how does that compare year over year' "
                "requires ranking AND time intelligence to work together correctly. Walk through why a naive SUM "
                "over a semi-additive balance measure gives a nonsensical total, since that's the mistake learners "
                "make most often once they leave this room."
            ),
        },
        {
            "title": "Measures That Perform",
            "icon": "gauge",
            "why": "A correct measure that takes eight seconds to render is still a support ticket waiting to happen -- performance is part of correctness in a production report.",
            "key_points": [
                "Use VAR to name intermediate results -- it clarifies logic and can prevent redundant evaluation.",
                "Diagnose slow measures with the query plan and Performance Analyzer before guessing at a fix.",
                "Iterators (SUMX, FILTER) over large tables are the most common source of slow DAX -- know when a simpler aggregation will do.",
                "Validate correctness first, then optimize -- a fast wrong answer is worse than a slow right one.",
            ],
            "notes": (
                "This sub-topic previews the deeper performance work in Module 5, so it's fine to keep this "
                "relatively brief and practical: a few concrete 'smells' to watch for (nested iterators, filtering "
                "a big table row by row) rather than a full optimization workshop. The sequencing point worth "
                "making out loud: we validate logic here in Module 2, and we come back with proper tooling "
                "(DAX Studio, Performance Analyzer) in Module 5 once there's a full report to profile."
            ),
        },
    ],
)

build_module(
    3, "Advanced Power Query", "funnel",
    "Building a transformation layer that's traceable, testable, and reusable -- not a black box.",
    [
        ("Staged, parameterized pipelines", "Parameters for source switching, staged queries, and functions that eliminate copy-pasted steps."),
        ("Query folding and performance", "Understanding when a query pushes work back to the source, and when it silently stops."),
        ("Data quality as a first-class step", "Error-review patterns and cleansing functions that make bad rows visible instead of silently dropped."),
    ],
    "Advanced", "BUILD",
    "Module overview for Advanced Power Query. The theme: transformation logic should be as reviewable "
    "and reusable as code, because it effectively is code. Three sub-topics: pipeline architecture, "
    "folding/performance, and data quality as a deliberate design step rather than an afterthought.",
    [
        {
            "title": "Staged, Parameterized Pipelines",
            "icon": "funnel",
            "why": "A query pipeline built from named, staged steps and parameters can move from dev to test to production data sources in minutes instead of a manual rebuild.",
            "key_points": [
                "Parameters isolate source paths/connections so the same pipeline works across environments.",
                "Staging queries (raw → staged → final) make each transformation step reviewable in isolation.",
                "Custom functions eliminate copy-pasted steps across similar queries -- one function, many callers.",
                "Descriptive step names turn the Applied Steps pane into documentation, not a mystery.",
            ],
            "notes": (
                "Make the case with a concrete failure mode: a query with 40 steps named 'Changed Type1', "
                "'Changed Type2' is unreviewable and un-debuggable six months later, even by the person who wrote "
                "it. Contrast that with a staged raw/staging/final pattern with named, purposeful steps. If there's "
                "time, live-build a small custom function that replaces three duplicated cleanup steps across two "
                "queries -- it's a concrete, memorable 'aha' moment for learners who've never used M functions."
            ),
        },
        {
            "title": "Query Folding And Performance",
            "icon": "gauge",
            "why": "When folding breaks silently, the query still returns correct results -- it just pulls every row into memory and transforms it locally, which can turn a 30-second refresh into a 30-minute one.",
            "key_points": [
                "Folding means Power Query translates steps into a query the source executes natively (e.g., SQL).",
                "Certain steps (custom columns with complex logic, some merges) can break folding -- know which ones.",
                "Use 'View Native Query' or step-by-step inspection to confirm folding is still happening after each transformation.",
                "When folding isn't possible, decide deliberately where local processing happens rather than discovering it by accident.",
            ],
            "notes": (
                "This is one of the most commonly misunderstood topics for intermediate Power Query users, because "
                "the query still 'works' when folding breaks -- there's no error, just a performance cliff. Use a "
                "before/after refresh time comparison if you have one available; it's a much stronger argument "
                "than describing folding in the abstract. The practical takeaway for learners: check folding "
                "status after any step that feels 'clever,' not just at the end of the whole query."
            ),
        },
        {
            "title": "Data Quality As A First-Class Step",
            "icon": "shield",
            "why": "Silently dropped or coerced bad rows are one of the most dangerous failure modes in a BI solution -- the report looks fine and the numbers are simply wrong.",
            "key_points": [
                "Build an explicit error-review query that captures rows failing type conversion or business rules instead of discarding them.",
                "Cleansing functions (e.g., text normalization) should be reusable and applied consistently across sources.",
                "Decide and document what happens to invalid rows: excluded, corrected, or flagged for manual review.",
                "Treat the error-review query as a report artifact stakeholders can see, not just a developer scratch pad.",
            ],
            "notes": (
                "Tell a short story here if you have one: a report that quietly excluded a currency-mismatched batch "
                "of orders for months because a type conversion failed silently, and nobody noticed until finance "
                "reconciliation caught it. That's the risk this sub-topic is designed to prevent. The instructor "
                "framing that lands well: 'invalid data doesn't go away when you ignore it, it just goes somewhere "
                "you can't see' -- an explicit error-review query makes it visible instead."
            ),
        },
    ],
)

# ================================================================ ARC: PRESENT
slide_section(
    "PRESENT", "Turn a correct model into a report people actually use",
    "The experience arc: report UX that guides rather than overwhelms, a model fast enough to stay that way under load, and analytics that go beyond a static chart.",
    GOLD,
    ["04 · Report Design & UX", "05 · Performance Optimization", "06 · Advanced Analytics & AI"],
)
set_notes(prs.slides[-1],
    "Set up the PRESENT arc: a technically correct model and measure set doesn't matter if nobody wants "
    "to use the report built on top of it. This arc is about turning correctness into something usable, "
    "fast, and genuinely insightful -- the difference between a report people tolerate and one people "
    "actually open every morning.")

build_module(
    4, "Report Design & UX", "canvas_layout",
    "Designing navigation and interaction, not just placing visuals on a canvas.",
    [
        ("Guided interaction patterns", "Drillthrough, bookmarks, custom tooltips, and navigation that leads users somewhere intentional."),
        ("Hierarchies, drill & groups", "Drill-down, drill-across, and custom groups/bins that let one visual answer questions at multiple levels of detail."),
        ("Flexible by design", "Field parameters and dynamic titles that let one report answer several questions instead of one."),
        ("Built for every screen and every user", "Mobile-optimized layouts, conditional formatting, and an accessibility pass before anything ships."),
    ],
    "Intermediate → Advanced", "PRESENT",
    "Module overview for Report Design & UX. The mindset shift: a report page is an interface, and "
    "interfaces are designed, not assembled. Four sub-topics: guided interaction, hierarchies/drill/"
    "groups, flexible/reusable report design, and designing for every device and every user.",
    [
        {
            "title": "Guided Interaction Patterns",
            "icon": "canvas_layout",
            "why": "Users don't read instructions -- the report itself has to teach them where to click and what will happen when they do.",
            "key_points": [
                "Drillthrough pages answer 'tell me more about this one thing' without cluttering the main view.",
                "Bookmarks capture a specific state (filters, visibility, selections) that can be recalled with one click.",
                "Custom tooltips replace a bare number with context -- a mini visual, a trend, a comparison.",
                "Navigation should be visibly obvious, not just technically functional -- users won't discover a hidden feature.",
            ],
            "notes": (
                "This is a good slide for a short live demo if a sample report is available: click a data point, "
                "show the drillthrough, then show a bookmark toggling a view. The instructor point worth making "
                "explicitly: a feature that works but isn't discoverable might as well not exist. Watch a first-"
                "time user try to navigate the report before you explain how -- if they get stuck, that's real "
                "signal about where the navigation design needs more visual affordance."
            ),
        },
        {
            "title": "Hierarchies, Drill & Groups",
            "icon": "layers",
            "why": "A hierarchy lets one visual answer 'show me the summary' and 'show me exactly why' without building a separate visual or page for every level of detail.",
            "key_points": [
                "Hierarchies (e.g., Category > Subcategory > Product) are built once in the model and reused across every visual that needs them.",
                "Drill-down moves a single visual deeper into its own hierarchy; drill-up reverses it -- both stay on the same visual.",
                "Drill-across sends a selection to a different visual or page (via drillthrough or cross-report drill), showing related detail in a different context.",
                "Groups and bins collapse many raw values into meaningful buckets (an 'Other' category, price bands) without changing the underlying data.",
                "Turn on 'Show next level' and 'Expand all' deliberately -- they change how a user explores versus click-drills one level at a time.",
            ],
            "notes": (
                "This sub-topic fills a real gap: drillthrough (jump to a dedicated detail page) is often taught, "
                "but in-visual drill-down/drill-up through a hierarchy and drill-across between related visuals are "
                "just as important and easy to overlook. Demo live if possible: build a Category > Subcategory > "
                "Product hierarchy on a matrix or bar chart, drill down two levels, then contrast that with a "
                "drillthrough page and a cross-report drill scenario so learners see all three interaction models "
                "side by side. For groups: show a long tail of low-volume categories collapsed into an 'Other' "
                "bucket via Power BI's built-in grouping feature -- it's a five-minute demo that immediately cleans "
                "up a busy slicer or legend. Make the connection back to Module 1 explicit: hierarchies and groups "
                "are modeling decisions with a UX payoff, not just a report-canvas trick."
            ),
        },
        {
            "title": "Flexible By Design",
            "icon": "canvas_layout",
            "why": "Field parameters turn one page into several reports' worth of value, which matters enormously when report sprawl is already a real governance problem in most organizations.",
            "key_points": [
                "Field parameters let users swap the measure or dimension a visual is built on, without duplicating pages.",
                "Dynamic titles that reflect the current selection make a flexible page still feel intentional, not generic.",
                "Fewer, more flexible pages are easier to maintain and govern than many near-duplicate ones.",
                "Flexibility should still have sensible defaults -- 'infinitely customizable' isn't the same as 'usable.'",
            ],
            "notes": (
                "Field parameters are one of the more recently introduced Power BI features and many experienced "
                "authors haven't adopted them yet -- treat this as a genuine 'new trick' moment even for learners "
                "who consider themselves advanced. The organizational argument matters as much as the technical "
                "one here: every report team has a folder full of near-duplicate pages built because nobody knew "
                "field parameters existed. That's real technical debt this feature can prevent."
            ),
        },
        {
            "title": "Built For Every Screen And Every User",
            "icon": "shield",
            "why": "A report that only works well on a widescreen monitor for a sighted user with perfect color vision has quietly excluded a meaningful share of its actual audience.",
            "key_points": [
                "Mobile layout is a separate design pass, not an automatic resize of the desktop layout.",
                "Conditional formatting should communicate meaning (status, threshold, direction), not just add color.",
                "An accessibility review checks tab order, alt text, and color contrast before a report ships.",
                "Design for the edge-case user (screen reader, phone, colorblind) and the default experience improves for everyone.",
            ],
            "notes": (
                "This sub-topic often gets treated as optional or a 'nice to have if there's time' -- push back on "
                "that framing. If the organization has any public-sector, education, or enterprise-scale customers, "
                "accessibility compliance is very likely already a requirement somewhere in their governance policy, "
                "even if the report authors don't know it yet. Consider a quick tab-order test live: tab through an "
                "existing report with the mouse untouched and see how far you get."
            ),
        },
    ],
)

build_module(
    5, "Performance Optimization", "gauge",
    "Finding out why a report is slow before guessing how to fix it.",
    [
        ("Evidence before optimization", "Performance Analyzer baselines and DAX Studio comparisons -- measure first, then change one thing."),
        ("Model size is a design decision", "Cardinality reduction, aggregation tables, and composite model choices that shrink refresh and query time."),
        ("Refresh strategy at scale", "Incremental refresh policy design and the licensing/capacity tradeoffs that come with it."),
    ],
    "Advanced", "PRESENT",
    "Module overview for Performance Optimization. The core discipline: diagnose before you fix. Three "
    "sub-topics: capturing a real baseline, model-size decisions, and refresh strategy at scale.",
    [
        {
            "title": "Evidence Before Optimization",
            "icon": "gauge",
            "why": "Optimizing based on a guess wastes time and can make performance worse -- a baseline turns 'this feels slow' into 'this specific visual takes 4.2 seconds because of this specific query.'",
            "key_points": [
                "Performance Analyzer breaks down visual load time into DAX query time, visual display time, and other time.",
                "DAX Studio shows the underlying query and server timings for a measure in isolation.",
                "Capture the baseline, change exactly one thing, and re-measure -- don't change five things and hope.",
                "Some 'performance problems' are actually visual count or interaction problems, not DAX problems at all.",
            ],
            "notes": (
                "The instructor framing that works well: performance work without a baseline is just superstition. "
                "If possible, run Performance Analyzer live on a real report page and let the class see the actual "
                "breakdown -- it's common for the surprise to be 'the DAX was fine, it's the sheer number of visuals "
                "on the page that's slow,' which is a genuinely useful and non-obvious lesson."
            ),
        },
        {
            "title": "Model Size Is A Design Decision",
            "icon": "layers",
            "why": "Reducing an unnecessary column's cardinality or removing an unused column can shrink a model more than any DAX tweak -- and it costs nothing at query time once it's done.",
            "key_points": [
                "High-cardinality columns (e.g., unique IDs, full timestamps) are expensive -- keep only what's needed at the grain required.",
                "Aggregation tables serve common high-level queries fast while detail tables remain available for drill-down.",
                "Composite models let a subset of tables use DirectQuery for freshness while the rest stay in-memory for speed.",
                "Column removal and cardinality reduction are usually the highest-leverage, lowest-risk performance wins available.",
            ],
            "notes": (
                "This connects directly back to Module 1's 'models built for growth' sub-topic -- call that out "
                "explicitly, since it reinforces that the curriculum is a connected arc and not disconnected topics. "
                "A good live example if available: show model size before and after removing a full-precision "
                "timestamp column in favor of a date key plus a separate time-of-day column, since that's one of "
                "the most common and highest-impact cardinality mistakes in real models."
            ),
        },
        {
            "title": "Refresh Strategy At Scale",
            "icon": "gear",
            "why": "A report that takes 20 minutes to refresh isn't just annoying -- it delays every downstream decision that depends on the data being current, and it can hit capacity or gateway limits outright.",
            "key_points": [
                "Incremental refresh only reprocesses new or changed partitions instead of the entire table.",
                "Refresh policy design (partition ranges, detection columns) has to match how the source data actually changes.",
                "Licensing and capacity tier determine what refresh options and frequency are even available.",
                "A refresh strategy should be documented, not just configured -- the next owner needs to understand why it's set up this way.",
            ],
            "notes": (
                "This is a good moment to be honest about a real constraint: incremental refresh setup and policy "
                "options differ by Power BI license/capacity, so what's demoed here may not be available to every "
                "learner in their own tenant. Encourage learners to confirm their organization's capacity tier "
                "before assuming a specific incremental refresh pattern will work as shown."
            ),
        },
    ],
)

build_module(
    6, "Advanced Analytics & AI", "spark",
    "Going past the default chart into scenario modeling and AI-assisted analysis -- with eyes open about where it's available.",
    [
        ("What-if, not just what-happened", "Parameter-driven scenario analysis that lets stakeholders test assumptions live in the report."),
        ("The analytics visual toolkit", "Decomposition trees, forecasting, and key influencers for driver analysis beyond a trend line."),
        ("AI-assisted authoring, used deliberately", "Python/R visuals and Azure ML integration explored with a clear-eyed view of cloud and tenant availability."),
    ],
    "Advanced", "PRESENT",
    "Module overview for Advanced Analytics & AI. The theme: move reports from describing the past to "
    "supporting a decision about the future -- while being honest about which features are available in "
    "which cloud and tenant. Three sub-topics: what-if analysis, the analytics visual toolkit, and "
    "AI-assisted features used with informed judgment.",
    [
        {
            "title": "What-If, Not Just What-Happened",
            "icon": "gauge",
            "why": "A what-if parameter turns a static report into a decision tool -- stakeholders can test 'what if margin improves 2%' live instead of asking an analyst to rebuild the report for every scenario.",
            "key_points": [
                "What-if parameters generate a slicer-driven value that flows into a measure's calculation.",
                "Pair the parameter with a clear before/after visual so the impact of the scenario is immediately obvious.",
                "Document the assumption behind the scenario -- a number without its underlying assumption is easy to misread.",
                "This is the most universally available advanced analytics feature -- it works the same way across clouds.",
            ],
            "notes": (
                "This sub-topic is deliberately positioned first in the module because it's the one advanced "
                "analytics capability that's fully available regardless of cloud (Commercial, Government, or DoD) "
                "-- make that explicit reassurance early, since the next two sub-topics carry real caveats. A live "
                "demo of a margin or pricing what-if scenario is one of the most effective 'wow' moments in the "
                "whole curriculum for a business-stakeholder audience."
            ),
        },
        {
            "title": "The Analytics Visual Toolkit",
            "icon": "spark",
            "why": "Decomposition trees, forecasting, and key influencers answer 'why' and 'what's next,' not just 'what happened' -- that's a fundamentally different, higher-value conversation with stakeholders.",
            "key_points": [
                "Decomposition trees let users explore what's driving a metric interactively, without pre-building every breakdown.",
                "Forecasting visuals project a trend forward with a confidence interval, not a false-precision single number.",
                "Key influencers surface which factors most affect an outcome, using statistical analysis under the hood.",
                "These visuals are commonly Commercial-cloud features -- confirm availability before building a delivery plan around them.",
            ],
            "notes": (
                "Flag clearly and early: per the program's cloud coverage matrix, decomposition tree, forecasting, "
                "and key influencers are treated as Not Covered in Government (GCC High) and DoD deliveries. If "
                "you're teaching a Gov/DoD cohort, this sub-topic becomes a conceptual walkthrough with screenshots "
                "rather than a hands-on lab -- say so plainly rather than letting learners discover the limitation "
                "themselves later. For a Commercial audience, this is a strong hands-on moment."
            ),
        },
        {
            "title": "AI-Assisted Authoring, Used Deliberately",
            "icon": "spark",
            "why": "Python/R visuals and Azure ML integration open real analytical power, but they also introduce dependencies (package policy, network, identity, region) that can quietly break a report outside a Commercial tenant.",
            "key_points": [
                "Python/R visuals require Desktop configuration and depend on organizational package policy in the Service.",
                "Azure Machine Learning integration depends on cloud, region, workspace, network, and identity configuration all lining up.",
                "Treat these features as 'available if validated,' not 'available by default' -- confirm in the actual delivery tenant.",
                "Where these aren't available, teach the underlying concept and workflow even if the hands-on lab isn't possible.",
            ],
            "notes": (
                "This is the most tenant-sensitive sub-topic in the entire program -- resist the temptation to "
                "promise more than can be verified. The instructor's job here is to teach good judgment about "
                "AI-assisted features, not just the mechanics: what's the validation process before you rely on "
                "one of these in a production report? Who owns confirming availability in a given cloud? That "
                "meta-skill outlasts any specific feature's availability window."
            ),
        },
    ],
)

# ================================================================ ARC: PROTECT & OPERATE
PO_COLOR = ARC_COLOR["PROTECT & OPERATE"]
slide_section(
    "PROTECT & OPERATE", "A report that can't be trusted won't be used",
    "The trust arc: locking data down to the right audience, getting a solution live in the Service, and keeping it healthy once real users depend on it.",
    PO_COLOR,
    ["07 · Security Design", "08 · Service Deployment", "09 · Monitoring & Governance"],
)
set_notes(prs.slides[-1],
    "Set up the PROTECT & OPERATE arc: a technically excellent report that exposes the wrong data to "
    "the wrong people, or that nobody can get to reliably, fails just as completely as a badly modeled "
    "one. This arc covers the trust infrastructure -- security, deployment, and ongoing operations -- "
    "that turns a good report into a solution the organization can actually depend on.")

build_module(
    7, "Security Design", "shield",
    "Designing row-level security like an architecture decision, not a checkbox.",
    [
        ("Static and dynamic RLS", "Building roles that scale from a handful of regions to a dynamic, user-driven security model."),
        ("Testing security like a user would", "Validating roles in Desktop and Service, not just trusting the DAX filter is correct."),
        ("Where Build permission fits in", "Understanding what elevated permissions expose, and designing around that risk deliberately."),
    ],
    "Advanced", "PROTECT & OPERATE",
    "Module overview for Security Design. The core message: RLS is an architecture decision made early, "
    "not a setting bolted on right before publishing. Three sub-topics: static vs. dynamic RLS, testing "
    "security properly, and understanding what Build permission actually exposes.",
    [
        {
            "title": "Static And Dynamic RLS",
            "icon": "shield",
            "why": "Choosing the wrong RLS pattern for the organization's scale means either constant manual role maintenance or an overengineered dynamic model nobody can audit.",
            "key_points": [
                "Static RLS assigns fixed filter values per role -- simple, auditable, but doesn't scale past a small number of roles.",
                "Dynamic RLS derives the filter from the logged-in user's identity against a security mapping table.",
                "Dynamic RLS scales far better but is harder to test and requires a well-maintained user-to-permission mapping table.",
                "Choose the pattern based on how many distinct access groups actually exist today and realistically will in two years.",
            ],
            "notes": (
                "A common mistake is reaching for dynamic RLS by default because it feels more 'advanced' -- "
                "push back on that instinct. If an organization genuinely only has three or four regional roles "
                "that rarely change, static RLS is simpler, more auditable, and easier for a security reviewer to "
                "reason about. Dynamic RLS earns its complexity when the number of distinct access combinations "
                "would make static roles unmanageable."
            ),
        },
        {
            "title": "Testing Security Like A User Would",
            "icon": "shield",
            "why": "An RLS role that looks correct in the DAX editor can still leak data in practice -- the only real test is viewing the report exactly as that user would see it.",
            "key_points": [
                "Use 'View As Roles' in Desktop to preview the report exactly as a given role would experience it.",
                "Test in the Service too -- Desktop and Service can behave differently depending on gateway and dataset configuration.",
                "Check every visual and every drillthrough/tooltip page under each role, not just the main landing page.",
                "Negative testing matters: confirm a role can't see what it shouldn't, not just that it can see what it should.",
            ],
            "notes": (
                "Emphasize negative testing specifically -- most authors test 'does this role see their own data "
                "correctly' and stop there, without testing 'does this role fail to see anyone else's data.' Walk "
                "through at least one live example of testing a role in both Desktop and Service, since the "
                "difference between the two is a real and common source of RLS surprises after publish."
            ),
        },
        {
            "title": "Where Build Permission Fits In",
            "icon": "shield",
            "why": "Build permission on a dataset is powerful and easy to grant too broadly -- and because it bypasses report-level RLS enforcement in certain connection scenarios, understanding it is core to real security design, not a footnote.",
            "key_points": [
                "Build permission lets a user create new reports against the shared dataset directly.",
                "Depending on connection type, a user with Build permission and a live connection can potentially see data outside their intended RLS scope.",
                "Grant Build deliberately and sparingly, understanding exactly what it unlocks for that user.",
                "Document who has Build permission on sensitive datasets as part of the ongoing governance review, not a one-time setup step.",
            ],
            "notes": (
                "This is a genuinely under-taught topic and worth calling out as such -- many experienced Power BI "
                "authors have never had Build permission's implications explained to them clearly. Keep the framing "
                "practical rather than alarmist: the point isn't 'never grant Build permission,' it's 'know exactly "
                "what you're granting when you do it.' This connects forward to Module 9's governance review "
                "practices -- Build permission audits belong on that recurring checklist."
            ),
        },
    ],
)

build_module(
    8, "Service Deployment", "gear",
    "Getting from a working PBIX/PBIP to a solution real people depend on daily.",
    [
        ("Publishing with intent", "Workspace strategy, App packaging, and deployment pipelines instead of ad hoc publish clicks."),
        ("Refresh you can trust", "Scheduled refresh, gateway considerations, and troubleshooting failures before users notice them."),
        ("A repeatable release process", "Moving a report from dev to test to production with a process, not tribal memory."),
    ],
    "Intermediate → Advanced", "PROTECT & OPERATE",
    "Module overview for Service Deployment. The core message: publishing is not the finish line -- it's "
    "the start of an operational responsibility. Three sub-topics: intentional publishing, trustworthy "
    "refresh, and a repeatable release process.",
    [
        {
            "title": "Publishing With Intent",
            "icon": "gear",
            "why": "A workspace strategy decided in advance prevents the all-too-common outcome of dozens of ungoverned personal workspaces holding business-critical reports nobody officially owns.",
            "key_points": [
                "Workspace structure (by team, by domain, by environment) should be decided before reports start accumulating in it.",
                "Apps package a curated, permissioned view of workspace content for end users, separate from author-facing workspace access.",
                "Deployment pipelines move content through dev/test/production stages without manual re-publishing at each step.",
                "Publishing should be a deliberate act with an owner, not an ad hoc click whenever a report 'feels done.'",
            ],
            "notes": (
                "If your organization has a 'personal workspace with a business-critical report in it' story, this "
                "is the slide to tell it on -- it's a scenario almost every experienced Power BI admin has seen and "
                "it lands the point better than any abstract explanation. The instructor goal here is to get "
                "learners thinking about workspace strategy before they publish their first solo report, not after "
                "workspace sprawl is already a governance headache."
            ),
        },
        {
            "title": "Refresh You Can Trust",
            "icon": "gear",
            "why": "A dashboard nobody trusts because it silently failed to refresh last Tuesday is worse than no dashboard at all -- it actively misleads decisions.",
            "key_points": [
                "Scheduled refresh needs monitoring, not just configuration -- know how failures are surfaced and to whom.",
                "Gateway configuration and credentials are common points of silent failure; document them clearly.",
                "Build a habit of checking refresh history proactively rather than waiting for a user to report stale data.",
                "Failure alerts should go to someone who can act, not just to an inbox nobody monitors.",
            ],
            "notes": (
                "This connects directly to Module 9's usage-metrics-as-a-signal sub-topic -- refresh history is "
                "exactly that kind of operational signal, and this is a good moment to preview that theme. If "
                "possible, show a real refresh history log (even a fabricated example) and walk through how to "
                "read a failure versus a warning versus a successful-but-slow refresh."
            ),
        },
        {
            "title": "A Repeatable Release Process",
            "icon": "gear",
            "why": "Moving a report from 'my test version' to 'the production version everyone relies on' without a defined process is how untested changes end up in front of executives.",
            "key_points": [
                "Define distinct dev, test, and production environments -- even a lightweight version of this separation helps.",
                "A release checklist (RLS tested, refresh validated, visuals reviewed) turns tribal memory into a repeatable process.",
                "Deployment pipelines or a documented manual process should produce the same outcome regardless of who runs it.",
                "The goal isn't heavyweight process for its own sake -- it's making sure nothing important gets skipped under deadline pressure.",
            ],
            "notes": (
                "Keep this practical and scaled to reality: not every organization has full CI/CD tooling in place "
                "for Power BI (that's covered more deeply in Module 11), but every organization can adopt a "
                "lightweight release checklist today. The instructor point worth making: the checklist doesn't "
                "need to be sophisticated, it needs to exist and be followed consistently under deadline pressure, "
                "which is exactly when steps get skipped without one."
            ),
        },
    ],
)

build_module(
    9, "Monitoring, Administration & Governance", "gauge",
    "Running Power BI like a platform, not a pile of individually-managed reports.",
    [
        ("Usage as a signal", "Usage metrics and refresh history read as operational data, not vanity numbers."),
        ("Tenant settings with intention", "Understanding what each governance lever controls and who should be allowed to change it."),
        ("An operations runbook that outlives you", "Documenting the processes so governance survives a team change, not just a tenure."),
    ],
    "Advanced", "PROTECT & OPERATE",
    "Module overview for Monitoring, Administration & Governance. The core message: Power BI at scale "
    "is a platform with real operational needs, not a collection of independently-managed reports. Three "
    "sub-topics: usage as operational signal, deliberate tenant settings, and durable documentation.",
    [
        {
            "title": "Usage As A Signal",
            "icon": "gauge",
            "why": "Usage metrics tell you which reports are actually load-bearing for the organization and which ones are quietly abandoned -- both are equally important to know.",
            "key_points": [
                "Usage metrics reports show views, unique users, and trends over time per report and workspace.",
                "A report with declining usage may be broken, outdated, or simply superseded -- investigate rather than assume.",
                "Refresh history read alongside usage tells you whether a heavily used report is also a reliably fresh one.",
                "Use usage data to justify retiring reports, not just to justify building new ones.",
            ],
            "notes": (
                "The retirement point is worth emphasizing -- most governance conversations focus on what to build "
                "next, but a mature practice also actively retires unused or duplicate reports. Ask the group: "
                "does your organization currently have a process for retiring a report, or does everything that's "
                "ever been published just accumulate indefinitely? That question alone often reveals a real gap."
            ),
        },
        {
            "title": "Tenant Settings With Intention",
            "icon": "shield",
            "why": "Tenant settings control organization-wide capabilities like external sharing and export -- getting them wrong is either a security exposure or an unnecessary blocker on legitimate work.",
            "key_points": [
                "Every tenant setting should have a known owner and a documented reason for its current configuration.",
                "Export and sharing settings are common areas where security and usability pull in opposite directions.",
                "Changes to tenant settings should go through the same review rigor as any other security-relevant change.",
                "Review tenant settings on a schedule -- they drift out of alignment with actual organizational need over time.",
            ],
            "notes": (
                "Many organizations have tenant settings configured years ago by someone no longer with the company, "
                "with no documented rationale for the current configuration -- this sub-topic is partly about "
                "surfacing that risk. If you can get a screenshot or example of the tenant admin portal, walking "
                "through a few representative settings live makes this much more concrete than describing it "
                "abstractly."
            ),
        },
        {
            "title": "An Operations Runbook That Outlives You",
            "icon": "gear",
            "why": "Governance that lives only in one administrator's head disappears the day that person changes roles -- a written runbook is what makes governance durable.",
            "key_points": [
                "Document the 'why' behind key decisions (workspace structure, tenant settings, RLS patterns), not just the 'what.'",
                "A runbook should be specific enough that a new administrator could follow it without asking the previous owner.",
                "Review and update the runbook on a cadence -- stale documentation is almost as risky as no documentation.",
                "Treat the runbook itself as a governed artifact with an owner and a review date.",
            ],
            "notes": (
                "This is a good closing sub-topic for the PROTECT & OPERATE arc because it ties everything back "
                "together: security design, deployment process, and tenant configuration all need to be written "
                "down somewhere durable, or the organization is one departure away from losing institutional "
                "knowledge about how its own BI platform actually works. Ask if anyone in the room could hand off "
                "their current Power BI environment to a new hire using only existing documentation -- the answer "
                "is usually revealing."
            ),
        },
    ],
)

# ================================================================ ARC: SCALE
SC_COLOR = ARC_COLOR["SCALE"]
slide_section(
    "SCALE", "From one great report to a platform that scales",
    "The platform arc: capacity and architecture decisions, engineering discipline around Power BI artifacts, and the capstone that proves it all fits together.",
    SC_COLOR,
    ["10 · Capacity & Architecture", "11 · DevOps & Lifecycle", "Capstone · End-to-End Build"],
)
set_notes(prs.slides[-1],
    "Set up the SCALE arc: everything so far has been about building one excellent, trustworthy report. "
    "This arc is about what changes when there are dozens or hundreds of them across an organization -- "
    "capacity architecture decisions and engineering discipline that turn Power BI from a tool into a "
    "true platform. The capstone at the end proves all eleven modules connect into one coherent skill set.")

build_module(
    10, "Premium, Fabric & Capacity Architecture", "layers",
    "Choosing the right capacity and architecture before scale forces the decision for you.",
    [
        ("Capacity models compared", "Shared, Premium, and Fabric capacity tradeoffs mapped to real workload patterns, not just price sheets."),
        ("Fabric-aware design", "OneLake and Direct Lake concepts, and what they change about how a model should be built."),
        ("Planning for growth, not just today", "Architecture decisions that anticipate the next 10x of data volume and user count."),
    ],
    "Advanced", "SCALE",
    "Module overview for Premium, Fabric & Capacity Architecture. The core message: capacity choice is an "
    "architecture decision with real cost and performance consequences, not a licensing afterthought. "
    "Three sub-topics: comparing capacity models, Fabric-aware design, and planning ahead of growth.",
    [
        {
            "title": "Capacity Models Compared",
            "icon": "layers",
            "why": "The 'right' capacity tier depends on workload shape (refresh frequency, concurrent users, dataset size), not just budget -- picking based on price alone often causes a costly re-architecture later.",
            "key_points": [
                "Shared capacity is fine for smaller, less demanding workloads but has real limits on refresh frequency and dataset size.",
                "Premium/Fabric capacity unlocks larger datasets, more frequent refresh, and dedicated performance, at a cost tradeoff.",
                "Map actual workload characteristics (concurrency, data volume, refresh needs) to capacity tiers before committing.",
                "Capacity decisions are reversible but not free -- migrating workloads after the fact has real cost and effort.",
            ],
            "notes": (
                "This is a good slide to bring in real numbers if your organization has them: current dataset "
                "sizes, refresh frequency needs, and concurrent user counts, mapped against the actual capacity "
                "tier being used today. Learners in an architect or platform-owner role will find this the most "
                "immediately useful sub-topic in the whole program -- give it room to be a real discussion, not "
                "just a slide to read through."
            ),
        },
        {
            "title": "Fabric-Aware Design",
            "icon": "layers",
            "why": "OneLake and Direct Lake change some long-held assumptions about when data needs to be imported versus queried live -- modeling decisions made without this context can be needlessly conservative.",
            "key_points": [
                "OneLake provides a single, shared storage layer that multiple Fabric workloads can reference without duplication.",
                "Direct Lake mode can deliver near-Import-mode performance while querying data directly from OneLake.",
                "Not every existing modeling pattern needs to change for Fabric -- but new projects should evaluate it deliberately.",
                "Fabric adoption is itself an architecture decision with tenant, licensing, and readiness prerequisites to confirm first.",
            ],
            "notes": (
                "Keep this conceptual and honest about maturity -- Fabric and Direct Lake are evolving quickly, so "
                "frame this as 'here's the architecture pattern and why it matters' rather than a definitive, "
                "unchanging feature list. Encourage learners to verify current Fabric capability and licensing "
                "status in their own tenant rather than assuming what's true today will still be exactly true by "
                "the time they implement it."
            ),
        },
        {
            "title": "Planning For Growth, Not Just Today",
            "icon": "gauge",
            "why": "An architecture that comfortably handles today's data volume and user count can become the org's biggest technical debt the moment either one grows tenfold.",
            "key_points": [
                "Ask 'what does this look like at 10x scale' as a standard part of any new architecture decision.",
                "Capacity, refresh strategy, and model design all have breaking points -- know roughly where they are before hitting them.",
                "Plan a capacity/architecture review cadence rather than only revisiting it after something breaks.",
                "Growth planning is cheaper before the system is load-bearing than after -- treat it as proactive work, not reactive firefighting.",
            ],
            "notes": (
                "This is the natural closing sub-topic for the SCALE arc's architecture module -- it ties together "
                "capacity choice and Fabric-aware design into a single forward-looking habit: always sanity-check "
                "an architecture decision against a materially larger future state, not just today's numbers. If "
                "time allows, ask the group to name one part of their current Power BI environment they suspect "
                "won't scale well, and use it as a live case study."
            ),
        },
    ],
)

build_module(
    11, "Automation, DevOps & Lifecycle Management", "gear",
    "Treating Power BI artifacts like source code, because eventually they have to be.",
    [
        ("PBIP and version control", "Git-friendly project files that make Power BI diffable, reviewable, and mergeable."),
        ("External tools and APIs", "Using the broader tooling ecosystem to script, validate, and automate what used to be manual."),
        ("CI/CD concepts for BI", "Deployment pipeline thinking applied to semantic models and reports, not just application code."),
    ],
    "Advanced", "SCALE",
    "Module overview for Automation, DevOps & Lifecycle Management. The core message: Power BI artifacts "
    "should get the same engineering discipline as any other production code. Three sub-topics: PBIP and "
    "version control, external tools/APIs, and CI/CD concepts applied to BI.",
    [
        {
            "title": "PBIP And Version Control",
            "icon": "gear",
            "why": "A binary PBIX file is a black box to source control -- PBIP's plain-text project structure is what makes real code review and collaboration on a Power BI file possible.",
            "key_points": [
                "PBIP splits a report and model into readable folders and files instead of one opaque binary.",
                "Plain-text definitions mean Git can show a meaningful diff -- what actually changed, not just 'the file changed.'",
                "Multiple authors can work on different parts of the same solution with real merge and review capability.",
                "Adopting PBIP is itself a process change for a team, not just a file-format setting to flip once.",
            ],
            "notes": (
                "This is a great slide to demo live if you have Git installed: open a PBIP project's folder "
                "structure and show a real diff after a small change, compared to what a PBIX-based 'diff' looks "
                "like (nothing useful at all). For teams still on PBIX, this sub-topic alone can be the single "
                "most practically useful thing they take from the whole program."
            ),
        },
        {
            "title": "External Tools And APIs",
            "icon": "gear",
            "why": "The Power BI ecosystem's external tools and REST APIs let teams script and validate things that used to require slow, error-prone manual clicking through the UI.",
            "key_points": [
                "External tools (Tabular Editor, DAX Studio, ALM Toolkit) extend what's possible beyond Desktop's built-in UI.",
                "The Power BI REST API enables scripted operations -- refresh triggers, dataset metadata, permission audits.",
                "ALM Toolkit-style comparisons make model migrations and merges safer by showing exactly what will change.",
                "Automating repetitive validation steps reduces the chance of a manual error during a release.",
            ],
            "notes": (
                "This sub-topic is a good opportunity to demo whichever external tool the delivery team is most "
                "comfortable with live -- even a five-minute walkthrough of Tabular Editor's scripting pane or "
                "the API's dataset refresh trigger makes the concept concrete instead of abstract. Frame this "
                "explicitly as 'here's the tooling ecosystem beyond Desktop' for learners who may not know it "
                "exists yet."
            ),
        },
        {
            "title": "CI/CD Concepts For BI",
            "icon": "rocket",
            "why": "Applying deployment-pipeline thinking to semantic models and reports catches breaking changes before they reach production, the same way it does for application code.",
            "key_points": [
                "A CI/CD mindset for BI means automated validation (does the model still build, do key measures still return expected results) before deployment.",
                "Deployment pipelines or scripted release processes reduce the manual steps that introduce human error.",
                "Automated testing for BI is less mature than for application code -- start with what can be automated today and grow from there.",
                "This module closes the loop with Module 8's repeatable release process -- automation is how that process scales past a handful of reports.",
            ],
            "notes": (
                "Be honest with the class: BI-specific CI/CD tooling and practices are less standardized than "
                "traditional software CI/CD, so this sub-topic is as much about the mindset and direction of "
                "travel as it is about a specific tool stack. Explicitly connect this back to Module 8's release "
                "process sub-topic -- this is the module where that process gets real automation behind it instead "
                "of remaining a manual checklist."
            ),
        },
    ],
)

slide_capstone()
set_notes(prs.slides[-1],
    "The capstone is the proof point for the entire program: one dataset, no script to follow, and every "
    "skill from Modules 1-11 has to come together into a single working, secured, production-ready "
    "solution. Use this slide to preview the four stages (Model it, Calculate it, Present it, Ship it) "
    "as a roadmap for the capstone deliverable, and frame it as the moment learners prove to themselves "
    "-- not just to an instructor -- that the whole arc actually connects.")

slide_closing()
set_notes(prs.slides[-1],
    "Closing slide. Land the outcome in one sentence: this program exists to produce confident modeling, "
    "defensible numbers, and reports people actually trust -- not just more Power BI features memorized. "
    "Use this as the natural point to transition into Q&A, next steps, or a scoping conversation about "
    "which modules and delivery format best fit this specific audience.")

prs.save(OUT_PATH)
print(f"Saved {OUT_PATH}")
print(f"Total slides: {len(prs.slides.__iter__.__self__._sldIdLst)}")
