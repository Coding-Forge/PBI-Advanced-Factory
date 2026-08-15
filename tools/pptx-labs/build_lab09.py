#!/usr/bin/env python3
"""
Builds the Lab 09 (Monitoring, Administration, and Governance) instructor deck.
Run from repo root: python tools/pptx-labs/build_lab09.py
Output: modules/09-monitoring-governance/assets/monitoring-governance.pptx
"""
import sys
from pathlib import Path

sys.path.insert(0, str(Path(__file__).resolve().parent))
from slide_kit import (
    new_presentation, blank_slide, title_slide, agenda_slide,
    content_slide, table_slide, checklist_slide, closing_slide,
    add_rect, add_text, add_page_number,
    SLIDE_W, SLIDE_H, NAVY, NAVY_DARK, ICE, WHITE, INK, SLATE, GOLD,
    LIGHT_BG, HEADER_FONT, BODY_FONT,
)
from pptx.util import Inches, Pt, Emu
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR
from pptx.enum.shapes import MSO_CONNECTOR

OUT_PATH = Path(__file__).resolve().parent.parent.parent / \
    "modules" / "09-monitoring-governance" / "assets" / "monitoring-governance.pptx"

MODULE_NO = 9
TITLE = "Monitoring, Administration, and Governance"
SUBTITLE = ("Operating deployed Power BI content: usage metrics, refresh reliability, tenant "
            "controls, audit evidence, capacity health, and Purview/DLP \u2014 built into a "
            "support-ready runbook.")

AGENDA_TOPICS = [
    "Why monitoring and governance matter",
    "Usage metrics",
    "Refresh monitoring",
    "Tenant settings",
    "Gateway monitoring",
    "Activity and audit logs",
    "Admin monitoring workspace",
    "Capacity metrics",
    "Purview and DLP",
    "Adoption tracking",
    "Operations model",
    "Module lab walkthrough",
    "Knowledge check and discussion",
]


def monitoring_flow_slide(prs, number, title, page, note=None, script=None):
    """Custom diagram: telemetry sources -> collection surfaces -> reporting/action."""
    s = blank_slide(prs)
    add_rect(s, 0, 0, SLIDE_W, SLIDE_H, WHITE)
    add_rect(s, 0, 0, SLIDE_W, Inches(1.15), NAVY)
    add_text(s, Inches(0.7), Inches(0.22), Inches(1.6), Inches(0.4), f"TOPIC {number:02d}",
             size=13, color=GOLD, bold=True, font=BODY_FONT)
    add_text(s, Inches(0.7), Inches(0.52), Inches(11.5), Inches(0.55), title, size=24,
             color=WHITE, bold=True, font=HEADER_FONT)

    # Three columns: Sources -> Surfaces -> Outcomes
    col_labels = ["Telemetry sources", "Collection surfaces", "Reporting & action"]
    col_x = [Inches(0.7), Inches(5.05), Inches(9.4)]
    col_w = Inches(3.6)
    header_y = Inches(1.55)
    header_h = Inches(0.5)
    for i, lbl in enumerate(col_labels):
        add_rect(s, col_x[i], header_y, col_w, header_h, NAVY if i == 1 else ICE)
        add_text(s, col_x[i], header_y, col_w, header_h, lbl,
                 size=15, color=(WHITE if i == 1 else NAVY_DARK), bold=True,
                 font=BODY_FONT, align=PP_ALIGN.CENTER, anchor=MSO_ANCHOR.MIDDLE)

    sources = [
        "Report & App views",
        "Semantic model refreshes",
        "Gateway operations",
        "Tenant admin actions",
        "Sensitivity label events",
    ]
    surfaces = [
        "Usage metrics reports",
        "Refresh history",
        "Activity log / audit log",
        "Admin monitoring workspace",
        "Capacity metrics app",
    ]
    outcomes = [
        "Adoption follow-up",
        "Refresh incident response",
        "Governance evidence",
        "Tenant policy tuning",
        "Capacity & DLP action",
    ]

    box_top = Inches(2.25)
    box_h = Inches(0.68)
    box_gap = Inches(0.14)
    for row, (src, surf, out) in enumerate(zip(sources, surfaces, outcomes)):
        y = box_top + row * (box_h + box_gap)
        for col, text in enumerate((src, surf, out)):
            fill = LIGHT_BG if col != 1 else ICE
            add_rect(s, col_x[col], y, col_w, box_h, fill, line_color=NAVY)
            add_text(s, col_x[col] + Inches(0.15), y, col_w - Inches(0.3), box_h, text,
                     size=13, color=INK, bold=(col == 1), font=BODY_FONT,
                     anchor=MSO_ANCHOR.MIDDLE)
        # arrows from source -> surface and surface -> outcome
        mid_y = y + box_h / 2
        for a, b in ((0, 1), (1, 2)):
            ln = s.shapes.add_connector(
                MSO_CONNECTOR.STRAIGHT,
                Emu(int(col_x[a] + col_w)), Emu(int(mid_y)),
                Emu(int(col_x[b])), Emu(int(mid_y)),
            )
            ln.line.color.rgb = SLATE
            ln.line.width = Pt(1.25)

    if note:
        add_text(s, Inches(0.7), Inches(6.7), Inches(11.9), Inches(0.7), note, size=12.5,
                 color=SLATE, italic=True, font=BODY_FONT, line_spacing=1.1)
    add_page_number(s, page)
    if script:
        from slide_kit import set_notes
        set_notes(s, script)
    return s


def build():
    prs = new_presentation()
    page = 1

    # 1. Title
    title_slide(
        prs, MODULE_NO, TITLE, SUBTITLE,
        script=[
            "Welcome learners to Module 9. Frame this module as the transition from 'building "
            "reports' to 'operating a platform.' Everything up to now has been about producing "
            "content \u2014 today is about keeping it healthy, trusted, and compliant after it's "
            "deployed.",
            "Set the stakes: unmonitored Power BI environments accumulate failing refreshes, "
            "orphaned workspaces, over-shared reports, and unlabeled sensitive data \u2014 all of "
            "which erode user trust and create governance risk. Monitoring exists specifically to "
            "catch those issues before they become incidents.",
            "Preview what students will do: they'll walk through usage metrics, refresh history, "
            "tenant settings, gateway health, activity logs, admin monitoring, capacity metrics, "
            "and Purview/DLP \u2014 then consolidate everything into a written operations runbook in "
            "Exercise 5.",
            "Flag the Azure Government angle up front: usage metrics and refresh history are "
            "Gov-ready, but gateway monitoring, activity/audit logs, admin monitoring workspace, "
            "capacity metrics, and Purview/DLP are all 'Verify for Gov.' We treat those as "
            "conceptual or documented rather than assumed to be available in every classroom "
            "tenant.",
        ]
    )
    page += 1

    # 2. Agenda
    agenda_slide(
        prs, MODULE_NO, AGENDA_TOPICS, page=page,
        script=[
            "Walk the room through the thirteen topics quickly. Group them mentally: 1 sets the "
            "'why'; 2-3 are the day-to-day operational signals (usage and refresh); 4-6 cover "
            "governance controls (tenant settings, gateway, audit logs); 7-9 are admin and "
            "compliance surfaces (admin workspace, capacity, Purview/DLP); 10-11 are the human "
            "side (adoption, operations model); 12-13 wrap up with the lab and discussion.",
            "Call out that several topics \u2014 activity logs, admin monitoring, capacity metrics, "
            "Purview/DLP \u2014 are Verify for Gov. Students may not be able to touch them in the "
            "classroom tenant, so we'll treat those as read-and-document rather than "
            "click-and-configure.",
            "Tell students the payoff of the whole module is Exercise 5's operations runbook: "
            "everything we discuss becomes a section in a real support document they could hand to "
            "an operations team.",
        ]
    )
    page += 1

    # 3. Topic 1 - Why monitoring & governance matter
    content_slide(
        prs, 1, "Why Monitoring and Governance Matter", page=page,
        lead_items=[
            "Production Power BI content needs the same discipline as any production system: "
            "known owners, monitored signals, documented incident paths, and change control.",
            "Five operational concerns drive this module: production support, adoption, risk "
            "controls, refresh reliability, and clear ownership.",
            "Governance isn't only about restrictions \u2014 it's about making trusted content easy to "
            "find (certified/promoted endorsement) and untrusted content easy to identify.",
            "Monitoring turns invisible failures visible: a silent refresh failure or a rarely-"
            "viewed 'critical' report is a governance signal, not just an IT issue.",
        ],
        why_items=[
            "Without monitoring, refresh failures are discovered by users \u2014 when a leadership "
            "dashboard is stale on Monday morning \u2014 which destroys trust in the platform.",
            "Without governance controls (tenant settings, sensitivity labels, certification), "
            "sharing and export decisions default to whatever each author chooses, which is how "
            "sensitive data leaves the tenant unintentionally.",
            "Clear ownership per workspace and per report shortens incident response time from "
            "hours (hunting for who owns this) to minutes (paging the documented owner).",
            "This is exactly the frame Exercise 5's operations runbook is built around \u2014 owners, "
            "monitoring cadence, incident paths, and change control.",
        ],
        footer="Instructor prompt: ask the room \u2014 'in your current environment, who gets paged "
               "when a semantic model refresh fails at 3am, and how do they even find out?'",
        script=[
            "Open with the shift in perspective: Modules 1-8 have been about building content. "
            "Module 9 is the first module where we treat Power BI as a *production system* that "
            "needs to be run, not just built. Every production system \u2014 database, API, website "
            "\u2014 has owners, monitoring, incident paths, and change control. Power BI is no "
            "different.",
            "Explain each of the five drivers on the slide with a concrete example: production "
            "support means 'someone is responsible when it breaks'; adoption means 'we know which "
            "reports are actually used'; risk controls means 'sharing and export aren't a "
            "free-for-all'; refresh reliability means 'we know before users do that data is "
            "stale'; ownership means 'every asset has a name attached to it.'",
            "Reframe governance so students don't hear it as 'the department of no.' Governance's "
            "dual job is to elevate trusted content (endorsement, certified badges, promoted "
            "content) so users can find it, and to identify or restrict risky content (sensitivity "
            "labels, DLP, tenant settings) so it doesn't leak. Both halves matter equally.",
            "Ask the footer prompt out loud and let two or three students answer. Most rooms will "
            "reveal that current-state monitoring is 'a user emails us.' Use that as the launching "
            "point for the rest of the module \u2014 today is about not being that team.",
        ]
    )
    page += 1

    # 4. Topic 2 - Usage metrics
    content_slide(
        prs, 2, "Usage Metrics", page=page,
        lead_items=[
            "Every report and App in the Service exposes a usage metrics report showing views, "
            "unique viewers, and view trend over the last 90 days.",
            "Open it from the report's More options menu (\u2026) \u2192 'View usage metrics report' in "
            "the workspace; App-level metrics are available from the App itself.",
            "Interpret three things together: which reports (and which pages inside them) are "
            "high-use, which are declining, and who the unique viewer population looks like versus "
            "the intended audience.",
            "If usage metrics aren't available, the blocker is tenant policy, license level, or "
            "workspace permissions \u2014 document which one, because Exercise 1 requires either the "
            "metrics or the documented blocker.",
        ],
        why_items=[
            "Usage metrics are the cheapest, fastest signal for adoption \u2014 no external telemetry "
            "system needed, and every workspace has them.",
            "High-use pages tell you where to invest in performance and UX; low-use pages tell you "
            "where to retire, redirect, or re-train users.",
            "A gap between unique viewer count and intended audience size is the earliest sign a "
            "report isn't reaching the people it was built for \u2014 a governance signal, not just "
            "an analytics one.",
            "This is Gov-ready and available in the classroom tenant, so Exercise 1 is a "
            "hands-on interpretation exercise, not just conceptual.",
        ],
        footer="Lab connection: Exercise 1 has learners open usage metrics for a report or App, "
               "identify high-use and low-use pages, and document one concrete follow-up action.",
        script=[
            "Introduce usage metrics as the workhorse signal of Power BI operations: they exist "
            "for every published report and App, they're free, and they update daily. Show "
            "students the exact click-path \u2014 the More options menu on a report in a workspace, "
            "'View usage metrics report' \u2014 so they know where to look in Exercise 1.",
            "Walk through the three fields to interpret together. Views tell you volume; unique "
            "viewers tell you reach; the trend tells you direction. All three matter \u2014 a report "
            "with high views but a shrinking unique viewer count is being re-visited by a small "
            "core, which is a very different story from broad adoption.",
            "Give a concrete adoption interpretation example: if a Sales Ops App shows 200 views "
            "per day but only 12 unique viewers, and the intended audience is 40 sales managers, "
            "you have a discoverability or enablement problem \u2014 not a performance problem. That's "
            "an adoption signal that leads to a training or communication action, not a rebuild.",
            "Close by naming the failure mode: if usage metrics are unavailable in a target "
            "tenant, don't skip the exercise \u2014 document *why* they're unavailable (tenant policy "
            "off, license too low, workspace permission missing). That documentation is itself a "
            "governance artifact and belongs in the operations runbook.",
        ]
    )
    page += 1

    # 5. Topic 3 - Refresh monitoring
    content_slide(
        prs, 3, "Refresh Monitoring", page=page,
        lead_items=[
            "Every semantic model's refresh history is one click away: workspace \u2192 semantic "
            "model \u2192 Settings or the refresh history icon.",
            "The refresh history log shows scheduled and on-demand runs with status, duration, "
            "start/end times, and \u2014 for failures \u2014 an expandable error detail per data source.",
            "Failure triage sequence: read the error text first, then check data source credentials "
            "on the semantic model settings page, then check gateway mapping if the source is "
            "on-premises or requires a private network path.",
            "Common failure patterns: expired credentials, gateway offline or version-mismatched, "
            "source system unavailable, source schema changed (renamed/removed column), and "
            "timeout on large queries.",
        ],
        why_items=[
            "A stale semantic model silently produces wrong-looking numbers to every downstream "
            "report \u2014 the report still renders, it just reflects yesterday's data.",
            "Structured refresh triage saves time: without it, teams jump straight to 'restart the "
            "gateway' every time, which fixes symptoms but hides real credential or schema drift "
            "issues.",
            "Refresh history is Gov-ready and available for every semantic model, which is why "
            "Exercise 2 uses it as the primary troubleshooting artifact.",
            "The refresh troubleshooting sequence students learn here becomes a section in the "
            "operations runbook from Exercise 5.",
        ],
        footer="Lab connection: Exercise 2 has students walk the refresh history \u2192 error detail "
               "\u2192 credentials \u2192 gateway mapping sequence and document a likely cause and next "
               "action for a real (or simulated) failure.",
        script=[
            "Start with the click-path: workspace, click into a semantic model, and either open "
            "Settings or click the refresh history icon at the top. Every student needs to know "
            "this path cold because they'll use it every time a refresh alert comes in.",
            "Explain what refresh history shows: not just success/fail, but timing (start, end, "
            "duration) and \u2014 crucially \u2014 an expandable error per data source when a run fails. "
            "That per-source error text is where triage actually starts; the top-level 'Failed' "
            "status is not enough.",
            "Teach the fixed triage sequence: read the error text first, don't jump to conclusions. "
            "Only then move to credentials (are they expired? did the account password rotate?), "
            "then gateway mapping (is the on-premises data gateway online, on a supported version, "
            "and mapped to this semantic model's data source?), then upstream (is the SQL server "
            "or file share actually reachable?). Give a concrete example \u2014 an 'invalid credential' "
            "error almost always maps to a rotated service account, not a gateway problem.",
            "Set expectations for Exercise 2: even if they don't have a failing refresh handy in "
            "the classroom, the exercise still walks them through refresh history for a real "
            "semantic model and asks them to document how they *would* triage a hypothetical "
            "failure. That documentation is what feeds into the operations runbook.",
        ]
    )
    page += 1

    # 6. Topic 4 - Tenant settings
    table_slide(
        prs, 4, "Tenant Settings That Affect Governance", page=page,
        headers=["Setting area", "What it controls", "Governance decision"],
        col_widths=[2.4, 5.2, 4.3],
        rows=[
            ["Sharing",
             "Who can share reports and Apps externally, share links, and 'Allow guests to share.'",
             "Restrict to specific security groups, not tenant-wide 'entire organization.'"],
            ["Export & download",
             "Export to Excel, CSV, PowerPoint, PDF, and 'Analyze in Excel' from the Service.",
             "Scope by group; align with sensitivity labels so labeled data can't be exported "
             "unlabeled."],
            ["Publish to web",
             "Anonymous, publicly-accessible embed of a report \u2014 no auth, no audit.",
             "Almost always disable or restrict to a single approved security group; audit "
             "existing embeds."],
            ["Certification & endorsement",
             "Who can mark content as Certified (top trust tier) or Promoted; who can request "
             "certification.",
             "Certified is reserved for reviewed content only \u2014 define the review process before "
             "enabling."],
            ["Build permission",
             "Whether users can build new reports on top of a shared semantic model.",
             "Encourage build permission on certified models so authors reuse instead of rebuild."],
            ["External users (B2B)",
             "Whether Azure AD guest users can access Power BI content and be invited to "
             "workspaces.",
             "Coordinate with identity/security team; Gov tenants often disable or tightly scope."],
        ],
        note="Tenant settings are edited by Power BI admins in the Admin portal (Settings \u2192 "
             "Admin portal \u2192 Tenant settings). Exercise 3 has students review each area and "
             "document the settings that require an explicit customer policy decision.",
        script=[
            "Frame tenant settings as the 'defaults' every user in the tenant inherits. If a "
            "setting is left at its default 'Enabled for the entire organization,' every one of "
            "the tenant's authors gets that capability \u2014 which is rarely what a governed "
            "environment actually wants.",
            "Walk the six rows one at a time. For sharing, emphasize scoping to security groups: "
            "the 'entire organization' toggle is almost never the right answer for a governed "
            "tenant. For export/download, connect it forward to sensitivity labels \u2014 you can "
            "restrict export by label so a 'Confidential' report can't be dumped to CSV.",
            "Spend extra time on 'Publish to web' because it's the highest-risk setting in the "
            "entire tenant: it publishes an anonymous, unauthenticated embed of a report to the "
            "public internet. Almost every governed tenant disables it entirely, or restricts to a "
            "single approved security group and audits any existing embeds.",
            "For certification vs promoted \u2014 clarify the hierarchy: Promoted is 'the author says "
            "this is good,' Certified is 'the organization has reviewed this as an authoritative "
            "source.' Certification should only be enabled once there's a written review process "
            "\u2014 otherwise the badge means nothing.",
            "Close by reminding students that Exercise 3 is a review-and-document exercise: they "
            "don't necessarily change anything in the classroom tenant, but they produce a list of "
            "the settings that require a customer policy decision \u2014 which is exactly what the "
            "operations runbook needs.",
        ]
    )
    page += 1

    # 7. Topic 5 - Gateway monitoring
    content_slide(
        prs, 5, "Gateway Monitoring", page=page,
        lead_items=[
            "The on-premises data gateway is the bridge between the Power BI Service and any "
            "source that isn't reachable from the cloud \u2014 SQL Servers, file shares, private "
            "APIs.",
            "Gateway clusters (a primary + one or more members) provide load distribution and "
            "high availability; monitor cluster status, each member's online state, and version.",
            "Data source mappings tie a semantic model's data source connection to a specific "
            "gateway cluster and stored credential set \u2014 mismatches here cause 'gateway not "
            "found' or credential errors at refresh time.",
            "Gateway versions are updated monthly; running an out-of-date version can break new "
            "connectors and is a common silent cause of intermittent refresh failures.",
        ],
        why_items=[
            "Most refresh failures for on-premises sources are actually gateway-layer issues \u2014 "
            "offline members, expired credentials, or version mismatch \u2014 not source-system "
            "problems.",
            "Documenting the gateway support owner and escalation path is a runbook essential: "
            "gateway machines are often owned by a different team than the Power BI content "
            "authors.",
            "Gateway monitoring is Verify for Gov because gateway topology, network path, version "
            "support, and admin permissions all vary by tenant configuration.",
            "The gateway is a shared platform component \u2014 a version upgrade or restart affects "
            "every semantic model mapped to it, so change control matters.",
        ],
        footer="Lab connection: Exercise 4 (Verify for Gov) has students review gateway cluster "
               "status, data source mappings, credentials, and version \u2014 or, if unavailable, use "
               "the operations runbook template to document gateway requirements conceptually.",
        script=[
            "Introduce the gateway plainly: it's a small Windows service running on a server (or "
            "cluster of servers) inside the customer's network that lets Power BI in the cloud "
            "reach data sources that aren't publicly reachable. Without it, on-premises SQL and "
            "file-share refreshes cannot happen at all.",
            "Explain the cluster concept: a gateway cluster is one primary and one or more member "
            "gateways, which together provide load distribution and high availability. Monitoring "
            "means checking each member is online, all are on a supported version, and the cluster "
            "as a whole is healthy \u2014 not just 'the primary is up.'",
            "Walk through data source mappings: every semantic model that refreshes through a "
            "gateway has one or more data source mappings on the semantic model's settings page, "
            "pointing to a specific gateway cluster and stored credential. When these drift \u2014 the "
            "gateway is renamed, the credential expires \u2014 the refresh fails with a 'gateway not "
            "found' or credential error, which triages very differently from a source problem.",
            "Reinforce the Verify-for-Gov note: gateway monitoring depends on having gateway admin "
            "permission plus network reachability plus the right admin surfaces enabled. In the "
            "classroom tenant, students may not have any of those, so Exercise 4 has an explicit "
            "alternate path \u2014 use the runbook template to document the gateway topology, owner, "
            "and escalation path conceptually.",
        ]
    )
    page += 1

    # 8. Topic 6 - Activity and audit logs
    table_slide(
        prs, 6, "Activity Logs and Audit Logs", page=page,
        headers=["Event type", "What it tells you", "Where it's captured"],
        col_widths=[2.6, 5.5, 3.8],
        rows=[
            ["ViewReport / ViewDashboard",
             "Who opened which report, when, and from which App or workspace.",
             "Power BI activity log; Unified audit log."],
            ["CreateApp / UpdateApp / DeleteApp",
             "App lifecycle changes \u2014 including who published a new App version.",
             "Power BI activity log."],
            ["ShareReport / AddWorkspaceUser",
             "Who shared what with whom, and who was added to a workspace and at which role.",
             "Power BI activity log \u2014 primary sharing evidence."],
            ["ExportReport / ExportTile",
             "Any export to PDF, PPTX, CSV, XLSX \u2014 who exported which item and when.",
             "Power BI activity log; correlates with DLP events."],
            ["UpdateTenantSetting",
             "Any tenant admin toggling a tenant-level setting on or off.",
             "Power BI activity log \u2014 governance change evidence."],
            ["SensitivityLabelApplied / Changed",
             "Sensitivity label activity from Purview integration.",
             "Purview audit log; correlates with Power BI activity log."],
        ],
        note="Activity logs are queried by Power BI admins via the Power BI REST API "
             "(Get-PowerBIActivityEvent) or the Fabric admin monitoring workspace; audit logs via "
             "Microsoft Purview / M365 audit search. Both are Verify for Gov \u2014 confirm admin "
             "permission, audit configuration, and cloud availability first.",
        script=[
            "Frame this topic honestly: activity logs and audit logs are the answer to almost "
            "every governance question that starts with 'who did what, when.' A tenant that "
            "doesn't collect them cannot answer basic compliance questions like 'who shared this "
            "report externally last month.'",
            "Walk each event type and connect it to a real question: 'Who's actually opening the "
            "leadership dashboard?' \u2192 ViewReport. 'Who added a guest user to that workspace?' "
            "\u2192 AddWorkspaceUser. 'Did anyone export the payroll report to Excel this quarter?' "
            "\u2192 ExportReport. 'When did the tenant admin turn Publish-to-Web back on?' \u2192 "
            "UpdateTenantSetting. These aren't hypothetical \u2014 they're the exact questions audit "
            "and compliance teams ask.",
            "Explain the two collection paths: the Power BI activity log via the REST API (or "
            "PowerShell's Get-PowerBIActivityEvent), which keeps ~30 days of history, versus the "
            "unified M365 audit log via Purview, which keeps longer but requires audit to be "
            "explicitly enabled. Most enterprises pipe both into a SIEM for long-term retention.",
            "Give a concrete investigation walkthrough: suppose an executive asks 'has our "
            "quarterly financials report been shared outside finance?' The audit path is: query "
            "activity log for ShareReport and AddWorkspaceUser events on that report and its "
            "workspace, filter by recipient domain and role, and produce a timeline. Under five "
            "minutes if the logs exist \u2014 impossible if they don't.",
            "Close with the Verify-for-Gov reminder: activity/audit logs require admin permissions "
            "and audit configuration that vary by cloud. The optional 'Activity logs' lab is "
            "conceptual for most students because they don't have tenant admin in the classroom.",
        ]
    )
    page += 1

    # 9. Topic 7 - Admin monitoring workspace
    content_slide(
        prs, 7, "Admin Monitoring Workspace", page=page,
        lead_items=[
            "The Fabric/Power BI 'Admin monitoring' workspace is an auto-provisioned workspace "
            "containing prebuilt reports on tenant-wide activity, feature usage, and inventory.",
            "Included reports typically cover Feature Usage & Adoption, Purview Hub, and Fabric "
            "Chargeback \u2014 exact contents evolve as Microsoft ships new admin surfaces.",
            "Accessed by users in the Fabric administrator or Power BI administrator role; it "
            "appears in the workspace list once the admin role is assigned and the feature is "
            "available in the target cloud.",
            "Availability varies by cloud (commercial vs. Gov vs. Gov-High) and tenant version \u2014 "
            "which is exactly why this topic is 'optional / Verify for Gov' for the lab.",
        ],
        why_items=[
            "Pre-built admin reports eliminate the need to build every governance dashboard from "
            "scratch out of raw activity log data \u2014 they're the fastest path to tenant "
            "visibility.",
            "Feature Usage & Adoption gives per-workspace, per-user activity summaries admins can "
            "use to spot dormant workspaces, over-shared content, and unused capacity.",
            "Purview Hub inside this workspace is the primary surface for sensitivity label and "
            "endorsement inventory across the tenant.",
            "Understanding what the admin workspace *can* show \u2014 even when it isn't available "
            "here \u2014 lets students specify it as a requirement in the operations runbook.",
        ],
        footer="Lab connection: The optional 'Admin monitoring workspace' lab has students open "
               "the workspace if available, note its reports, and document one tenant-level "
               "adoption or governance follow-up \u2014 or document the availability blocker.",
        script=[
            "Introduce the admin monitoring workspace as a relatively recent (and evolving) "
            "Microsoft-provided workspace: when your tenant has it enabled and you have the Fabric "
            "or Power BI administrator role, a workspace literally named 'Admin monitoring' shows "
            "up automatically in your workspace list, prepopulated with reports.",
            "Describe what students would typically see there today: a Feature Usage & Adoption "
            "report (per-workspace, per-user activity), the Purview Hub (sensitivity label and "
            "endorsement inventory), and Fabric-oriented reports like chargeback. Set expectations "
            "that the exact list evolves \u2014 don't memorize it, know the categories.",
            "Explain who can see it: only users with Fabric administrator or Power BI "
            "administrator role. A regular workspace admin or Pro user will not see this workspace "
            "at all. That role gating is intentional \u2014 the data inside is tenant-wide.",
            "Handle the availability question head-on: this workspace's contents and even its "
            "existence depend on tenant version and cloud. In classroom Gov tenants it may not be "
            "present. Even when unavailable, students should note it in the runbook as 'we would "
            "use this for tenant-level adoption tracking once available' \u2014 that's what "
            "'Verify for Gov' means in practice.",
        ]
    )
    page += 1

    # 10. Topic 8 - Capacity metrics
    content_slide(
        prs, 8, "Capacity Metrics", page=page,
        lead_items=[
            "The Microsoft Fabric Capacity Metrics app (installed by capacity admins from AppSource) "
            "reports CPU, memory, and throttling telemetry for Premium/Fabric capacities.",
            "CPU is measured in Capacity Units (CU) \u2014 track both interactive workload (report "
            "queries, dataflows) and background workload (scheduled refreshes, dataset "
            "reprocessing).",
            "Refresh impact: large scheduled refreshes are the most common cause of interactive "
            "slowdown \u2014 the app shows a stacked view of refresh CU vs. user query CU over the "
            "day.",
            "Throttling: sustained overload puts the capacity into interactive delay, then "
            "interactive rejection, then background rejection \u2014 each of which users experience "
            "differently and each requires a different remediation.",
        ],
        why_items=[
            "Slow reports on a healthy report and model design are almost always a capacity issue, "
            "not a report design issue \u2014 without capacity metrics you'll spend hours 'optimizing' "
            "the wrong layer.",
            "Understanding CU consumption per workspace and item is what enables cost allocation, "
            "capacity right-sizing, and refresh schedule staggering.",
            "Throttling states are the earliest warning a capacity is under-provisioned \u2014 acting "
            "on them before rejections start is the difference between a proactive upgrade and an "
            "outage.",
            "Capacity metrics are Verify for Gov \u2014 the app must be available in the target "
            "cloud, the user must have capacity admin, and the capacity must actually be Premium/"
            "Fabric (not shared).",
        ],
        footer="Lab connection: The optional 'Capacity metrics app' lab has students review CPU/"
               "memory, refresh workload, and throttling signals \u2014 or document capacity "
               "monitoring requirements in the runbook if the app isn't available.",
        script=[
            "Frame capacity metrics as the answer to 'is our platform sized correctly, or are we "
            "silently degrading?' The Microsoft Fabric Capacity Metrics app is the tool for it: a "
            "capacity admin installs it from AppSource and connects it to the capacity, and it "
            "then reports CPU and memory usage over time.",
            "Introduce the Capacity Unit (CU) as the metric that matters: everything in Fabric \u2014 "
            "a report query, a dataflow refresh, a semantic model reprocess \u2014 consumes CUs. The "
            "app shows CU consumption over time, split by interactive workload (user activity) and "
            "background workload (refreshes). This split is important because they compete for the "
            "same capacity.",
            "Explain the refresh impact pattern concretely: a large scheduled refresh at 8am can "
            "spike background CU right when interactive users start opening reports, which is why "
            "the app's stacked timeline view is so useful \u2014 you can literally see refresh CU "
            "displacing interactive CU. The fix is often to stagger refreshes off-peak, not to buy "
            "more capacity.",
            "Walk through the throttling states in order: sustained overload triggers interactive "
            "delay (users notice slowness), then interactive rejection (report queries fail), then "
            "background rejection (refreshes fail). Each is a distinct signal and each has "
            "distinct remediation \u2014 scaling up, scheduling changes, autoscale, or workload "
            "migration.",
            "Reinforce Verify-for-Gov: the app itself must be available in the target cloud, the "
            "user must be a capacity admin, and the workload must actually be running on a "
            "Premium/Fabric SKU. If any of those is missing the student documents capacity "
            "requirements in the runbook instead of clicking through the app.",
        ]
    )
    page += 1

    # 11. Topic 9 - Purview and DLP
    content_slide(
        prs, 9, "Purview and Data Loss Prevention (DLP)", page=page,
        lead_items=[
            "Microsoft Purview sensitivity labels (General, Confidential, Highly Confidential, "
            "etc.) can be applied to reports, semantic models, and Apps \u2014 and travel with "
            "exports to Excel, PowerPoint, and PDF.",
            "Labels are defined once in Purview and become available in Power BI once tenant "
            "settings for sensitivity labels are enabled; authors apply them from a report's "
            "Sensitivity dropdown in the Service.",
            "DLP policies for Power BI (defined in the Purview compliance portal) can detect "
            "labeled content and block or alert on risky sharing/export actions \u2014 for example, "
            "blocking export of Highly Confidential content to unmanaged devices.",
            "The Purview Hub inside the admin monitoring workspace gives a tenant-wide inventory "
            "of labeled and unlabeled content \u2014 the starting point for a compliance sweep.",
        ],
        why_items=[
            "Sensitivity labels are the enforceable link between 'this is sensitive data' and "
            "'here's what users can and can't do with it' \u2014 without them, classification is "
            "documentation-only.",
            "Labels persist through export, which closes the biggest historical gap in Power BI "
            "governance: an unlabeled CSV export used to leave the tenant with no protection at "
            "all.",
            "DLP catches the actions that policy alone can't prevent \u2014 a user attempting to share "
            "labeled content externally, or exporting from an unmanaged device \u2014 in real time.",
            "All of this is Verify for Gov: it depends on the M365/Purview cloud availability, "
            "label configuration, licensing, and tenant settings being aligned.",
        ],
        footer="Lab connection: The optional 'Purview and DLP review' lab has students review "
               "label availability, DLP policy scope, and export behavior \u2014 or document "
               "compliance requirements for the runbook.",
        script=[
            "Introduce Purview and DLP as the compliance layer sitting on top of Power BI. "
            "Sensitivity labels are the 'what' \u2014 classification of content \u2014 and DLP is the "
            "'so what' \u2014 automatic enforcement when labeled content is at risk.",
            "Walk through how labels get to Power BI: they're defined once in the Purview "
            "compliance portal (the same labels that apply to Word, Excel, SharePoint, and "
            "Outlook), and once tenant settings for sensitivity labels are enabled, they appear as "
            "a Sensitivity dropdown on every report and semantic model in the Service. Authors "
            "apply them; the label then persists into exports \u2014 an Excel export of a "
            "'Confidential' report opens as a Confidential Excel file.",
            "Explain DLP with a concrete scenario: a DLP policy can say 'if a report is labeled "
            "Highly Confidential and a user tries to share it with an external guest, block the "
            "share and notify the compliance team.' That kind of real-time enforcement is only "
            "possible because labels are attached to the content.",
            "Introduce the Purview Hub as the tenant-wide inventory view: it lives inside the "
            "admin monitoring workspace and shows how many items are labeled vs unlabeled, "
            "endorsed vs unendorsed, and where the compliance gaps are. This is the natural "
            "starting point for a governance sweep.",
            "Reinforce Verify-for-Gov: Purview cloud availability, label config, DLP licensing, "
            "and tenant settings all have to line up before any of this works. In the classroom "
            "students most likely document the compliance requirements rather than actually "
            "applying labels.",
        ]
    )
    page += 1

    # 12. Topic 10 - Adoption tracking (custom flow diagram)
    monitoring_flow_slide(
        prs, 10, "Adoption Tracking \u2014 From Telemetry to Action", page=page,
        note="Adoption tracking connects usage patterns to action. High usage \u2192 invest in "
             "training and support. Low or declining usage \u2192 enablement gap, discoverability "
             "gap, or retirement candidate. Each telemetry surface feeds a distinct governance "
             "outcome.",
        script=[
            "Use this diagram to consolidate the last several topics into one picture. Read it "
            "left to right: telemetry sources on the left are the raw events happening in the "
            "tenant \u2014 report views, refreshes, gateway ops, admin actions, label events. Those "
            "flow into the collection surfaces in the middle \u2014 usage metrics reports, refresh "
            "history, activity/audit logs, the admin monitoring workspace, the capacity metrics "
            "app. And each surface feeds a distinct governance outcome on the right.",
            "Walk each row across as a story: 'Report views' \u2192 'Usage metrics reports' \u2192 "
            "'Adoption follow-up.' 'Semantic model refreshes' \u2192 'Refresh history' \u2192 'Refresh "
            "incident response.' 'Tenant admin actions' \u2192 'Activity log' \u2192 'Governance "
            "evidence.' 'Sensitivity label events' \u2192 'Capacity metrics app / Purview Hub' \u2192 "
            "'Capacity & DLP action.' The point is that every telemetry source has an outcome; if "
            "you can't articulate the outcome, you're collecting telemetry for no reason.",
            "Give the adoption interpretation framing directly: high usage isn't automatically "
            "good \u2014 it needs training and support investment to stay reliable. Low or declining "
            "usage isn't automatically bad \u2014 you have to diagnose whether it's an enablement "
            "gap, a discoverability gap, or a genuine retirement candidate. The same telemetry "
            "leads to different actions depending on the interpretation.",
            "Close by tying to Exercise 5's operations runbook: the runbook should have a section "
            "for each collection surface on this diagram, naming the owner, the review cadence, "
            "and the action path when the signal fires.",
        ]
    )
    page += 1

    # 13. Topic 11 - Operations model
    content_slide(
        prs, 11, "Operations Model", page=page,
        lead_items=[
            "Every deployed workspace/App needs a named business owner AND a named technical owner "
            "\u2014 one for content decisions, one for platform issues.",
            "Support path: users know exactly where to report an issue (helpdesk ticket queue, "
            "Teams channel, email alias) and what SLA to expect.",
            "Incident handling: documented triage sequence for the three main incident types \u2014 "
            "refresh failure, capacity throttling, and sharing/DLP violation.",
            "Change control: report and semantic model changes go through deployment pipelines "
            "(Dev \u2192 Test \u2192 Prod), not direct edits to production workspaces.",
            "Review cadence: weekly refresh review, monthly usage review, quarterly tenant "
            "settings and access review, annual runbook refresh.",
        ],
        why_items=[
            "Named owners eliminate the 'who owns this?' delay at the start of every incident \u2014 "
            "the single biggest cause of long time-to-restore.",
            "A published support path routes issues away from the author's personal inbox and "
            "creates a ticket trail that feeds capacity and adoption planning.",
            "A repeatable review cadence turns governance from an annual audit scramble into an "
            "always-current picture \u2014 the operations runbook is never stale.",
            "Change control via deployment pipelines is the mechanism that lets you say 'yes, this "
            "report can be updated safely' instead of 'we can't touch it, it's critical.'",
        ],
        footer="Lab connection: Exercise 5 has students assemble everything from Exercises 1-4 into "
               "an operations runbook covering owners, monitoring cadence, incident response, "
               "change control, and Azure Government validation notes.",
        script=[
            "Frame the operations model as the human wrapper around all the telemetry we've been "
            "discussing. Signals without owners, cadences, and paths are just noise \u2014 the "
            "operations model turns them into a system.",
            "Insist on the dual-owner pattern: business owner (who decides what the content should "
            "show, who the audience is, when it retires) and technical owner (who fixes the "
            "refresh, updates the model, deals with gateway issues). One person can play both "
            "roles in a small team, but they're distinct responsibilities and the runbook should "
            "list them separately.",
            "Walk through the three main incident types with a triage sketch for each. Refresh "
            "failure: check refresh history \u2192 error text \u2192 credentials \u2192 gateway. Capacity "
            "throttling: check capacity metrics app for CU trend \u2192 identify the workload "
            "displacing others \u2192 reschedule or scale. Sharing/DLP violation: check activity log "
            "for the ShareReport or ExportReport event \u2192 identify actor and content \u2192 apply the "
            "documented remediation. Each is a runbook section.",
            "Explain change control briefly: production semantic models and reports should be "
            "deployed via Power BI deployment pipelines (Dev \u2192 Test \u2192 Prod stages) so changes "
            "are reviewed and reversible, not edited directly in the production workspace. That's "
            "the answer to 'we're afraid to touch that report.'",
            "Cadence: give the concrete weekly/monthly/quarterly/annual pattern on the slide and "
            "emphasize the point of Exercise 5 is to make this real for a specific piece of "
            "content \u2014 pick one App, fill in each section of the runbook template, and you'll "
            "have a working operations model artifact by the end of class.",
        ]
    )
    page += 1

    # 14. Topic 12 - Module lab walkthrough
    checklist_slide(
        prs, "Module Lab Walkthrough", kicker="Topic 12 \u2014 What you'll build", page=page,
        items=[
            "Exercise 1: Usage metrics \u2014 open metrics for a report/App, interpret views and "
            "viewers, identify one follow-up action.",
            "Exercise 2: Refresh troubleshooting \u2014 walk refresh history \u2192 error \u2192 "
            "credentials \u2192 gateway; document likely cause.",
            "Exercise 3: Tenant setting review \u2014 sharing, export, publish-to-web, "
            "certification, external users \u2014 note policy decisions.",
            "Exercise 4: Gateway monitoring (Verify for Gov) \u2014 cluster status, mappings, "
            "credentials, version, or documented alternate path.",
            "Optional labs: Activity logs, Admin monitoring workspace, Capacity metrics, "
            "Purview & DLP \u2014 all Verify for Gov.",
            "Exercise 5: Operations runbook \u2014 owners, data sources, access, monitoring "
            "cadence, incident response, Gov validation notes.",
            "Runbook template lives at modules/09-monitoring-governance/operations-runbook-"
            "template.md \u2014 use it as the scaffold.",
            "Deliverable: a completed draft runbook plus one adoption follow-up and one refresh "
            "triage note per learner.",
        ],
        script=[
            "Use this slide as the literal table of contents for the hands-on portion. Walk each "
            "row and pair it to what students will actually click.",
            "Emphasize the ordering: Exercises 1-3 are Gov-ready and hands-on; Exercise 4 is "
            "Verify-for-Gov with a documented alternate path; the four Optional labs are all "
            "Verify-for-Gov and may be conceptual/documentation-only in the classroom tenant.",
            "Point students at the operations-runbook-template.md file so they know where the "
            "scaffold lives. Exercise 5 is not writing from scratch \u2014 it's filling in a "
            "template with the observations they've collected from the earlier exercises.",
            "Set the finish-line expectation clearly: by the end of the lab, every learner should "
            "have a draft runbook, at least one adoption follow-up action, and at least one "
            "refresh triage note documented. If they don't have those artifacts, they haven't "
            "actually completed Module 9.",
        ]
    )
    page += 1

    # 15. Topic 13 - Knowledge check & discussion
    checklist_slide(
        prs, "Knowledge Check & Discussion", kicker="Topic 13 \u2014 Wrap-up", page=page,
        items=[
            "Usage metrics reviewed for a report or App (or blocker documented).",
            "One adoption follow-up, training need, support signal, or retirement candidate "
            "identified.",
            "Refresh troubleshooting sequence \u2014 history, error, credentials, gateway \u2014 "
            "documented.",
            "Tenant sharing, export, publish-to-web, certification, and external-user settings "
            "reviewed.",
            "Gateway monitoring, activity logs, admin monitoring workspace, capacity metrics, and "
            "Purview/DLP each marked Verify for Gov with a documented plan.",
            "Operations runbook draft assembled \u2014 owners, monitoring cadence, incident paths, "
            "change control.",
            "Discussion: what does 'certified' mean in your organization \u2014 and what review "
            "process would you require before enabling it?",
            "Discussion: how would you prove, to an auditor, that a specific report has not been "
            "shared externally in the last 90 days?",
        ],
        script=[
            "Run this as a discussion-driven wrap-up. Walk the checklist as a checkpoint \u2014 by a "
            "show of hands, who has each artifact completed? \u2014 and then open the two discussion "
            "questions to the room.",
            "For the certification question, push learners to describe a real review process: who "
            "reviews content, against what criteria, how often re-review happens, and how "
            "certification is revoked. The point is that 'Certified' is only meaningful if a "
            "review process is genuinely enforced \u2014 otherwise it's just a badge.",
            "For the auditor question, walk backwards from the answer: to prove 'no external "
            "sharing in the last 90 days,' you need activity log data covering that window, with "
            "ShareReport and AddWorkspaceUser events, filterable by recipient domain. If any of "
            "that is missing (log retention too short, audit not enabled, no domain filter), the "
            "organization cannot answer the question. That's how governance requirements get "
            "concrete.",
            "Close by connecting forward: Module 10 shifts from operating today's content to "
            "architecting for scale on Premium and Fabric capacity. Everything students learned "
            "today about capacity metrics, refresh impact, and throttling is the diagnostic layer "
            "underneath the architecture decisions in Module 10.",
        ]
    )
    page += 1

    # 16. Closing
    closing_slide(
        prs, MODULE_NO,
        "Module 10: Premium, Fabric, and Capacity-Aware Architecture \u2014 designing for scale "
        "on top of the monitoring and governance foundation you just built.",
        page=page,
        subtitle="Learners now have a working operations runbook and can interpret the full stack "
                 "of Power BI monitoring signals \u2014 usage, refresh, activity, capacity, and "
                 "compliance \u2014 for governed production support.",
        script=[
            "Congratulate the class on completing the operations module. This is the module where "
            "Power BI stops being 'a reporting tool' and starts being 'a governed platform.' The "
            "runbook artifact they produced today is the deliverable that platform owners actually "
            "hand to support teams in real projects.",
            "Recap the through-line: telemetry (usage, refresh, activity, capacity, labels) "
            "\u2192 collection surfaces (metrics reports, refresh history, audit logs, admin "
            "workspace, capacity app, Purview Hub) \u2192 governance outcomes (adoption action, "
            "incident response, evidence, policy tuning). If a signal doesn't have an outcome, "
            "revisit whether you're collecting it for the right reason.",
            "Take final questions, especially on the Verify-for-Gov topics \u2014 activity logs, "
            "admin monitoring, capacity metrics, Purview/DLP \u2014 since those are the areas most "
            "likely to differ between the classroom tenant and the learner's real environment.",
            "Preview Module 10 as the architectural counterpart to today: today was about "
            "operating the platform we have; Module 10 is about sizing and architecting Premium/"
            "Fabric capacity for the workload we expect. Capacity metrics from today's Exercise "
            "will come back immediately as the diagnostic evidence behind those architecture "
            "decisions.",
        ]
    )

    OUT_PATH.parent.mkdir(parents=True, exist_ok=True)
    prs.save(str(OUT_PATH))
    print(f"Saved: {OUT_PATH} ({len(prs.slides.__iter__.__self__._sldIdLst)} slides)")


if __name__ == "__main__":
    build()
