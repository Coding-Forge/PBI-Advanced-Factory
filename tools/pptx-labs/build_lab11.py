#!/usr/bin/env python3
"""
Builds the Lab 11 (Automation, DevOps & Lifecycle Management) instructor deck.
Run from repo root: python tools/pptx-labs/build_lab11.py
Output: modules/11-automation-devops/assets/automation-devops.pptx
"""
import sys
from pathlib import Path

sys.path.insert(0, str(Path(__file__).resolve().parent))
from slide_kit import (
    new_presentation, blank_slide, add_rect, add_text, add_bullets,
    add_page_number, add_kicker,
    title_slide, agenda_slide, content_slide, table_slide,
    checklist_slide, closing_slide,
    NAVY, NAVY_DARK, ICE, WHITE, INK, SLATE, GOLD, LIGHT_BG, CARD_BORDER,
    HEADER_FONT, BODY_FONT, SLIDE_W, SLIDE_H,
)
from pptx.util import Inches, Pt, Emu
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR
from pptx.enum.shapes import MSO_SHAPE, MSO_CONNECTOR

OUT_PATH = Path(__file__).resolve().parent.parent.parent / \
    "modules" / "11-automation-devops" / "assets" / "automation-devops.pptx"

MODULE_NO = 11
TITLE = "Automation, DevOps, and Lifecycle Management"
SUBTITLE = ("Making Power BI deployments repeatable, reviewable, and recoverable "
             "with PBIP, Git, external tools, REST APIs, and conceptual CI/CD.")

AGENDA_TOPICS = [
    "Lifecycle management goals",
    "PBIP as source of record",
    "Git workflow",
    "PBIP file structure",
    "External tools (Tabular Editor, ALM Toolkit)",
    "REST APIs and PowerShell",
    "Service principals",
    "Fabric workspace Git integration",
    "Azure DevOps conceptual pipeline",
    "GitHub Actions conceptual pipeline",
    "Azure Government considerations",
    "Deployment checklist",
    "Module lab walkthrough",
    "Knowledge check and discussion",
]


# ---------------------------------------------------------------------------
# Custom CI/CD pipeline flow diagram (local to lab 11 — does not touch slide_kit)
# ---------------------------------------------------------------------------
def pipeline_flow_slide(prs, number, title, stages, page, note=None):
    """Horizontal 5-stage pipeline flow: PBIP -> Git -> Build/Validate -> Deploy -> Service.

    stages: list of (label, sublabel) tuples, exactly 5.
    """
    s = blank_slide(prs)
    add_rect(s, 0, 0, SLIDE_W, SLIDE_H, WHITE)
    add_rect(s, 0, 0, SLIDE_W, Inches(1.15), NAVY)
    add_text(s, Inches(0.7), Inches(0.22), Inches(1.6), Inches(0.4), f"TOPIC {number:02d}",
              size=13, color=GOLD, bold=True, font=BODY_FONT)
    add_text(s, Inches(0.7), Inches(0.52), Inches(11.5), Inches(0.55), title, size=24,
              color=WHITE, bold=True, font=HEADER_FONT)

    n = len(stages)
    box_w = Inches(2.15)
    box_h = Inches(1.35)
    gap = Inches(0.30)
    total_w = box_w * n + gap * (n - 1)
    start_x = Emu(int((SLIDE_W - total_w) / 2))
    y = Inches(2.35)

    step_y = Inches(3.95)
    step_h = Inches(2.2)

    for i, (label, sub) in enumerate(stages):
        x = Emu(int(start_x + i * (box_w + gap)))
        # top numbered box
        add_rect(s, x, y, box_w, box_h, NAVY)
        add_text(s, x, y + Inches(0.1), box_w, Inches(0.4), f"STAGE {i+1}",
                  size=11, color=GOLD, bold=True, font=BODY_FONT,
                  align=PP_ALIGN.CENTER)
        add_text(s, x + Inches(0.1), y + Inches(0.5), box_w - Inches(0.2), Inches(0.85), label,
                  size=15, color=WHITE, bold=True, font=HEADER_FONT,
                  align=PP_ALIGN.CENTER, anchor=MSO_ANCHOR.MIDDLE)

        # detail card below
        add_rect(s, x, step_y, box_w, step_h, LIGHT_BG, line_color=CARD_BORDER)
        add_text(s, x + Inches(0.12), step_y + Inches(0.1), box_w - Inches(0.24),
                  step_h - Inches(0.2), sub, size=11.5, color=INK, font=BODY_FONT,
                  line_spacing=1.15, anchor=MSO_ANCHOR.TOP)

        # arrow connector to next box
        if i < n - 1:
            x1 = Emu(int(x + box_w))
            yc = Emu(int(y + box_h / 2))
            arr_w = Emu(int(gap))
            arr_h = Inches(0.34)
            arr = s.shapes.add_shape(
                MSO_SHAPE.RIGHT_ARROW,
                x1, Emu(int(yc - arr_h / 2)),
                arr_w, arr_h,
            )
            arr.fill.solid()
            arr.fill.fore_color.rgb = GOLD
            arr.line.fill.background()

    if note:
        add_text(s, Inches(0.7), Inches(6.75), Inches(11.9), Inches(0.6), note, size=12.5,
                  color=SLATE, italic=True, font=BODY_FONT, line_spacing=1.1)
    add_page_number(s, page)
    return s


def build():
    prs = new_presentation()
    page = 1

    # 1. Title -------------------------------------------------------------
    title_slide(
        prs, MODULE_NO, TITLE, SUBTITLE,
        script=[
            "Welcome learners to Module 11. Frame this as the module that turns everything they "
            "have built so far \u2014 the semantic model, DAX measures, Power Query steps, report "
            "pages, RLS, deployment configuration \u2014 into a repeatable, reviewable engineering "
            "asset instead of a one-off PBIX floating in someone's OneDrive.",
            "Set the stakes: a lot of Power BI work still gets 'deployed' by emailing a PBIX or "
            "clicking Publish from Desktop. That works exactly once. It doesn't survive a peer "
            "review, a rollback, or a Gov compliance audit. Today is about the engineering "
            "practices that make it survive all three.",
            "Preview the lab: students will practice PBIP source control with git, review PBIP "
            "structure, and then work through conceptual walkthroughs of external tools "
            "(Tabular Editor, ALM Toolkit), REST APIs, PowerShell, service principals, Fabric Git "
            "integration, and Azure DevOps / GitHub Actions CI/CD patterns.",
            "Call out the Gov posture up front: PBIP and local git are Gov-ready, so the required "
            "hands-on portion works everywhere. Everything else in this module \u2014 external "
            "tools, REST APIs, service principals, Fabric Git integration, CI/CD platforms \u2014 is "
            "Verify-for-Gov or Commercial-focused, and we treat it conceptually unless the "
            "instructor has confirmed the tenant supports it.",
        ]
    )
    page += 1

    # 2. Agenda ------------------------------------------------------------
    agenda_slide(
        prs, MODULE_NO, AGENDA_TOPICS, page=page,
        script=[
            "Walk the room through the fourteen items quickly. Group them mentally: 1-4 are "
            "foundations (why lifecycle matters, PBIP, git, PBIP structure), 5-8 are the tooling "
            "layer (external tools, REST APIs and PowerShell, service principals, Fabric Git), "
            "9-10 are the two conceptual CI/CD platforms (Azure DevOps and GitHub Actions), and "
            "11-12 are the governance/Gov and pre-deployment discipline.",
            "Set expectations clearly: only topics 2, 3, and 4 (PBIP, git, PBIP structure) are "
            "guaranteed hands-on for every student today. Everything else is Verify-for-Gov and "
            "will be delivered as conceptual walkthroughs unless the classroom tenant has already "
            "been validated for those features.",
            "Tell students the lab (topic 13) is where they will actually do the git workflow "
            "against a real PBIP project, and the knowledge check (topic 14) will focus on "
            "reasoning about tradeoffs \u2014 not memorizing feature lists.",
        ]
    )
    page += 1

    # 3. Topic 1 - Lifecycle management goals ------------------------------
    content_slide(
        prs, 1, "Lifecycle Management Goals", page=page,
        lead_items=[
            "Repeatability: the same PBIP source, deployed the same way, produces the same "
            "workspace state \u2014 no 'it worked when I clicked Publish from my laptop.'",
            "Reviewability: every model or report change is visible as a diff in git before it "
            "reaches an environment (semantic model TMDL and report JSON both diff cleanly).",
            "Environment promotion: Dev \u2192 Test \u2192 Prod workspaces are configured with "
            "parameters and dataset bindings, not by editing the PBIP for each environment.",
            "Governance: certified/promoted content, RLS, sensitivity labels, and workspace "
            "access are set intentionally, not left to whoever pressed Publish last.",
            "Rollback: a tagged git commit + release notes let you identify exactly what was "
            "promoted and re-deploy the previous version if a release regresses.",
        ],
        why_items=[
            "These five properties are what separates a Power BI hobby project from an "
            "engineering deliverable that can survive audit, personnel change, and outages.",
            "Every technique in the rest of this module \u2014 PBIP, git, external tools, REST APIs, "
            "service principals, CI/CD \u2014 exists to serve one or more of these five goals.",
            "In Azure Government especially, 'rollback evidence' is not optional \u2014 auditors "
            "will ask what was deployed, when, by whom, and how you would reverse it.",
        ],
        footer="Lab connection: Exercise 7 (Deployment checklist) is the concrete artifact that "
               "operationalizes these five goals \u2014 source, security, refresh, validation, "
               "rollback \u2014 for every release.",
        script=[
            "Anchor the module in the five words on the left: repeatability, reviewability, "
            "environment promotion, governance, rollback. Everything else today is in service of "
            "one of these. If a technique doesn't serve one of these goals, you probably don't "
            "need it in your release process.",
            "Give a concrete failure story to make it stick: someone republishes a PBIX from a "
            "laptop to fix a 'small' measure, the RLS role assignments get wiped, and finance "
            "sees numbers they aren't cleared to see. Every one of the five goals here would have "
            "prevented or contained that incident \u2014 reviewability catches the change, "
            "environment promotion prevents the direct-to-Prod deploy, governance preserves the "
            "role bindings, rollback lets you undo it in minutes instead of hours.",
            "Explain why the diff experience matters specifically: PBIP stores the semantic model "
            "as TMDL text and the report as structured JSON, so a git diff on a measure change "
            "actually shows the DAX before and after \u2014 a reviewer can catch a mistake without "
            "opening Power BI Desktop. A PBIX is a binary zip \u2014 a code reviewer can't see "
            "anything, so the review has to happen in a person's head, which doesn't scale.",
            "Tie forward to the lab: Exercise 7 asks students to fill in a deployment checklist "
            "covering source, security, refresh, validation, and rollback \u2014 that checklist is "
            "literally these five goals turned into a repeatable pre-release ritual.",
        ]
    )
    page += 1

    # 4. Topic 2 - PBIP as source of record --------------------------------
    content_slide(
        prs, 2, "PBIP as Source of Record", page=page,
        lead_items=[
            "PBIP is a folder-based project format: a .pbip pointer file plus sibling "
            ".Report and .SemanticModel folders that hold text-based definition files.",
            "The semantic model is stored as TMDL (Tabular Model Definition Language) \u2014 human-"
            "readable text for tables, columns, measures, relationships, and roles.",
            "The report is stored as structured JSON \u2014 pages, visuals, bookmarks, theme \u2014 "
            "which diffs cleanly, unlike the opaque binary inside a PBIX.",
            "PBIX is treated as a generated output for distribution, not the source of record. "
            "Committing PBIX to git defeats the entire point of PBIP.",
        ],
        why_items=[
            "Text-based definition files are what makes code review, pull requests, blame "
            "history, and CI validation possible for Power BI content at all.",
            "TMDL and report JSON let non-authors (peers, reviewers, security) see exactly what "
            "changed \u2014 a new measure, a modified RLS filter, a removed visual \u2014 without "
            "opening Power BI Desktop.",
            "PBIP is explicitly Gov-ready in this workshop, so this is one part of the module "
            "every student can practice hands-on regardless of tenant type.",
        ],
        footer="Lab connection: Exercise 1 has students open the actual PBIP folder in File "
               "Explorer / VS Code, locate the .Report and .SemanticModel folders, and read the "
               "TMDL and JSON to see what is source-controlled vs. generated.",
        script=[
            "Open Power BI Desktop's Save As dialog and note the format choice: .pbix (single "
            "binary file) vs .pbip (project folder). Explain that PBIP is not a different file "
            "format so much as a different unit of work \u2014 you are checking a *folder* of text "
            "into git instead of a single opaque binary.",
            "Walk the structure verbally: the .pbip file is a small pointer, the .Report folder "
            "contains report pages and visuals as JSON, and the .SemanticModel folder contains "
            "the model definition as TMDL. When a student edits a measure or renames a table in "
            "Desktop, saving the project writes textual changes into those folders \u2014 changes a "
            "reviewer can read.",
            "Make the review case concrete: 'Show me the diff for the measure change that "
            "shipped last Friday.' With PBIX that is impossible \u2014 the file is binary. With PBIP "
            "the reviewer opens the pull request and sees six lines of TMDL changed inside "
            "definition/model.tmdl, and can literally read the old DAX and the new DAX side by "
            "side.",
            "Close by warning against the common anti-pattern of committing the generated PBIX "
            "alongside the PBIP. The PBIX is a *build output*, like a compiled .exe \u2014 it "
            "belongs in a release/distribution channel, not in source control. If a PBIX has to "
            "be produced (for example for a Service that doesn't accept PBIP directly), generate "
            "it as part of the release process, not by hand-committing it.",
        ]
    )
    page += 1

    # 5. Topic 3 - Git workflow -------------------------------------------
    content_slide(
        prs, 3, "Git Workflow for Power BI Content", page=page,
        lead_items=[
            "Branch per unit of work: main is protected, feature branches (feature/measure-yoy) "
            "isolate a specific change until it's reviewed.",
            "Commits document intent: `git status` and `git diff` before every commit to avoid "
            "accidentally checking in cache files, .pbix outputs, or secrets.",
            "Pull requests are the review gate: at least one peer reviews the TMDL / report JSON "
            "diff and confirms it matches the described intent before merge to main.",
            "Tags and releases mark what was deployed: an annotated tag (v2026.11.1) plus "
            "release notes gives operations a rollback anchor.",
            "Exclude generated and local artifacts via .gitignore: .pbix, cache folders, "
            "*.tmp, and any local settings that don't belong to every environment.",
        ],
        why_items=[
            "The git workflow is what makes 'reviewable' and 'recoverable' from the previous "
            "slide real instead of aspirational \u2014 the review happens in the pull request, and "
            "the rollback anchor is the tag.",
            "Branch-and-PR discipline also acts as a lightweight change advisory board: no "
            "measure ships to Prod without a second pair of eyes on the DAX diff.",
            "This is Gov-ready and platform-agnostic \u2014 Azure DevOps, GitHub, GitHub "
            "Enterprise Server, or an internal Git host all work as long as the repository "
            "policy is customer-approved.",
        ],
        footer="Lab connection: Exercise 2 has students create a feature branch, make a small "
               "PBIP change, run git status/diff, commit, and discuss pull-request review "
               "expectations.",
        script=[
            "Start with the workflow shape verbally: main is protected (nobody pushes to it "
            "directly), you branch off main for a specific unit of work, you commit small "
            "meaningful changes on that branch, you open a pull request, a peer reviews the "
            "diff, and only then does it merge back to main. Every deployable release comes "
            "from an annotated tag on main.",
            "Stress the 'commits document intent' point. A good Power BI commit message says "
            "'Add YoY sales growth measure and update Sales Overview KPI card' \u2014 not "
            "'updates.' The commit message plus the TMDL / JSON diff is what a future reviewer "
            "or auditor uses to reconstruct why the change was made.",
            "Walk the review expectation concretely. A PBIP pull request will show diffs "
            "inside definition/model.tmdl (measures, relationships, roles) and inside the "
            "report JSON (visuals, pages, bookmarks). The reviewer's job is to confirm those "
            "diffs match the PR description \u2014 no accidental role changes, no orphan visuals, "
            "no committed PBIX or cache files. Point out .gitignore is your friend here.",
            "Wrap by connecting tags to rollback. When Prod breaks, the on-call needs to answer "
            "'what commit is currently deployed?' in seconds, not minutes. A tag plus release "
            "notes on the tagged commit is the mechanism that makes that answer instant, and "
            "that in turn is what makes rollback a five-minute operation instead of a war room.",
        ]
    )
    page += 1

    # 6. Topic 4 - PBIP file structure ------------------------------------
    table_slide(
        prs, 4, "PBIP File Structure \u2014 What Lives Where", page=page,
        headers=["Location", "Contents", "Source-control treatment"],
        col_widths=[3.2, 5.4, 3.3],
        rows=[
            ["<project>.pbip",
             "Small pointer file that ties the Report and SemanticModel folders together.",
             "Commit. Small, stable, required."],
            ["<project>.Report/",
             "Report definition: report.json, page and visual definitions, bookmarks, theme.",
             "Commit. Diffs cleanly as JSON."],
            ["<project>.SemanticModel/",
             "Semantic model as TMDL: tables, columns, measures, relationships, roles, "
             "perspectives, translations.",
             "Commit. Primary review surface."],
            ["definition/ subfolders",
             "Individual TMDL / JSON files that Power BI Desktop writes on save.",
             "Commit. These are the reviewable diff units."],
            ["Generated .pbix output",
             "Binary produced when you Build/Export PBIX from the PBIP \u2014 for distribution or "
             "Service upload only.",
             "Do NOT commit. Treat as a build artifact."],
            ["Local cache / temp / user settings",
             ".pbi/cache, *.tmp, editor-specific settings that vary per workstation.",
             "Exclude via .gitignore \u2014 never source-controlled."],
        ],
        note="Lab tie-in: Exercise 1 walks students through this table concretely against the "
             "PBIP project from earlier modules \u2014 identify each of these six categories in "
             "File Explorer and confirm which ones the .gitignore correctly excludes.",
        script=[
            "This is a reference slide \u2014 you don't have to talk through every row, but leave "
            "it on screen long enough for students to skim. The point is that PBIP is not a "
            "mystery: every file on disk falls into one of these six categories, and the source-"
            "control treatment for each category is completely mechanical.",
            "Focus attention on rows 3 and 4 (SemanticModel folder and its subfolders). This is "
            "where 95% of your reviewable diffs will live \u2014 measures, relationships, roles, "
            "translations. When a reviewer says 'show me what changed', they are almost always "
            "looking at TMDL files inside definition/ under the SemanticModel folder.",
            "Call out the anti-patterns explicitly. Two common ones: committing the generated "
            ".pbix (row 5) 'just in case,' which pollutes the diff history and can leak binary "
            "cache data \u2014 and committing local cache/user settings (row 6), which creates "
            "spurious changes on every teammate's machine. Both belong in .gitignore.",
            "Bridge to Exercise 1: send students to their PBIP folder in File Explorer or VS "
            "Code and have them physically locate each of these six categories, then open one "
            "TMDL file and one report JSON file as text so they can see that these really are "
            "just text files a human can read.",
        ]
    )
    page += 1

    # 7. Topic 5 - External tools -----------------------------------------
    content_slide(
        prs, 5, "External Tools: Tabular Editor & ALM Toolkit", page=page,
        lead_items=[
            "Tabular Editor: dedicated model editor that reads the same TMDL PBIP uses \u2014 "
            "bulk measure edits, scripting, best-practice analyzer, and calculation groups at a "
            "speed Desktop can't match.",
            "ALM Toolkit: model comparison tool that produces a schema-level diff between two "
            "models (or a model and a PBIP source) and generates a deploy script.",
            "Both tools rely on the XMLA endpoint to write back to a published dataset \u2014 "
            "that endpoint has to be enabled at tenant and capacity level.",
            "Customer policy matters: some organizations restrict which third-party tools may "
            "connect to their tenant. Confirm workstation and tenant policy before use.",
        ],
        why_items=[
            "Once a model grows past a handful of measures, doing bulk changes in Desktop is "
            "slow and error-prone \u2014 Tabular Editor scripting turns 'rename 40 measures' from "
            "a half-day of clicking into a five-line script.",
            "ALM Toolkit is what makes 'deploy only the metadata diff' possible \u2014 you can "
            "push a single measure change to Prod without republishing the entire model and "
            "resetting caches.",
            "Both are marked Verify for Gov: XMLA, tenant settings, and customer approval must "
            "line up before either tool is used against a Gov tenant.",
        ],
        footer="Lab connection: Exercises 3 and 4 are conditional \u2014 only performed when "
               "Tabular Editor / ALM Toolkit are available and customer-approved; otherwise "
               "delivered as a conceptual walkthrough.",
        script=[
            "Introduce both tools as 'community-standard, but not built in.' They aren't part "
            "of the Power BI Desktop install \u2014 they're third-party downloads that plug into "
            "the same TMDL / XMLA surfaces PBIP already exposes.",
            "For Tabular Editor, focus on the two capabilities that matter most for a lifecycle "
            "story: bulk edits (rename, reformat, refactor 40 measures at once using C# script) "
            "and the Best Practice Analyzer (a rules engine that flags common modeling mistakes "
            "\u2014 unformatted measures, missing descriptions, hidden foreign keys \u2014 as part of "
            "a review or CI step).",
            "For ALM Toolkit, focus on the deploy-the-diff case. Instead of overwriting a "
            "workspace model with a full PBIP publish, you compare source vs. deployed and "
            "generate a change script that pushes only what actually differs. That's how you "
            "ship a single measure update without disturbing production caches, refresh "
            "history, or bindings.",
            "Land the Gov point hard: both tools connect via XMLA endpoints, which are tenant- "
            "and capacity-gated. In a Gov environment, XMLA availability, workstation tooling "
            "policy, and customer approval all have to be validated before treating either "
            "tool as a real option. In this workshop, if that hasn't been validated, Exercises "
            "3 and 4 stay conceptual.",
        ]
    )
    page += 1

    # 8. Topic 6 - REST APIs and PowerShell -------------------------------
    table_slide(
        prs, 6, "Power BI REST APIs & PowerShell", page=page,
        headers=["Capability", "REST / cmdlet surface", "What it enables"],
        col_widths=[3.0, 5.2, 3.7],
        rows=[
            ["Workspace ops",
             "Groups API; New-PowerBIWorkspace / Set-PowerBIWorkspace.",
             "Provision Dev/Test/Prod workspaces programmatically."],
            ["Import / export",
             "Imports API (POST /groups/{id}/imports); "
             "New-PowerBIReport / Get-PowerBIReport -OutFile.",
             "Publish a PBIP-generated PBIX and pull artifacts for backup."],
            ["Refresh",
             "Datasets Refreshes API; Invoke-PowerBIRestMethod against "
             "/datasets/{id}/refreshes.",
             "Trigger and monitor dataset refresh from a pipeline."],
            ["Permissions",
             "Users API; Add-PowerBIWorkspaceUser; Add-PowerBIDataset..Permission.",
             "Manage workspace and content access as code, not clicks."],
            ["Dataset parameters / bindings",
             "UpdateParameters, TakeOver, UpdateDatasources APIs.",
             "Rebind a promoted dataset to the target environment's data source."],
            ["Endpoint & cloud validation",
             "api.powerbi.com (Commercial) vs api.powerbigov.us (Gov); "
             "Connect-PowerBIServiceAccount -Environment.",
             "Point every call at the right sovereign cloud; wrong endpoint = "
             "silent failure."],
        ],
        note="All rows are Verify for Gov \u2014 confirm endpoint, permissions, tenant settings, "
             "and authentication method before scripting against a Gov tenant.",
        script=[
            "Frame the REST APIs and the PowerShell cmdlets as two skins on the same surface \u2014 "
            "PowerShell is a convenient wrapper, but every cmdlet is ultimately making a REST "
            "call. Some low-level operations don't have a dedicated cmdlet and require calling "
            "the REST API directly with Invoke-PowerBIRestMethod \u2014 that's normal.",
            "Walk the capability rows as a mental checklist for 'what a deployment script "
            "actually needs to do': provision or find the workspace, upload the model, rebind "
            "data sources for the target environment, kick off a refresh, and set permissions. "
            "Every row on the slide answers one of those needs.",
            "Highlight the endpoint row specifically. This is the single most common cause of a "
            "Gov automation script silently failing \u2014 it authenticates fine, runs fine, and "
            "just talks to the Commercial cloud because nobody remembered to pass -Environment "
            "USGov to Connect-PowerBIServiceAccount, or hard-coded api.powerbi.com in the "
            "script. In a Gov engagement, treat endpoint validation as its own review step.",
            "Close with the Verify-for-Gov reminder in the note: none of these API calls is a "
            "given \u2014 tenant settings for service principal API access, permissions on the "
            "target workspace, and the cloud endpoint all have to be confirmed before you rely "
            "on them in a real release process.",
        ]
    )
    page += 1

    # 9. Topic 7 - Service principals -------------------------------------
    content_slide(
        prs, 7, "Service Principals for Automation", page=page,
        lead_items=[
            "A service principal is an Entra ID app registration that acts as a non-human "
            "identity \u2014 the correct identity for a CI/CD pipeline or scheduled job to use "
            "instead of a personal account.",
            "Tenant setting 'Allow service principals to use Power BI APIs' must be enabled and "
            "scoped to a specific security group. Without that setting, every REST call fails "
            "authorization.",
            "Grant workspace access to that security group (Member or Contributor), not to the "
            "app itself \u2014 group-based access is auditable and easier to rotate.",
            "Prefer certificate-based authentication over client secrets; if a secret is used, "
            "store it in Azure Key Vault, rotate it on a schedule, and never commit it to git.",
        ],
        why_items=[
            "Automation that runs as a person breaks the moment that person leaves, changes "
            "roles, has MFA enforced, or has their session expire \u2014 all of which happen "
            "routinely. Service principals are the only durable answer.",
            "Group-based workspace access is how a security team can rotate credentials or "
            "revoke pipeline access without ripping the pipeline apart \u2014 the app changes, "
            "the group membership doesn't.",
            "In a Gov tenant, service-principal API access is Verify for Gov \u2014 tenant setting "
            "approval, app registration policy, and workspace access all have to be confirmed "
            "before wiring one into a pipeline.",
        ],
        footer="Lab connection: this is a conceptual walkthrough in Lab 11 (Optional Lab: "
               "Service principal authentication) \u2014 students review app registration, tenant "
               "settings, workspace access, and secret handling requirements without "
               "provisioning a real principal.",
        script=[
            "Open by naming the anti-pattern: a lot of teams start automation by 'just using my "
            "account for now.' That works until the person leaves, MFA gets enforced, the "
            "password changes, or a security audit shows a personal identity running "
            "unattended jobs \u2014 all of which are inevitable. Service principals exist "
            "precisely so nothing production-critical is tied to a human's login.",
            "Walk the four requirements on the left in order, because they trip teams up in "
            "exactly this sequence. First, the Entra app registration exists. Second, the "
            "tenant setting 'Allow service principals to use Power BI APIs' is enabled and "
            "scoped to a security group. Third, that group is a member of the target "
            "workspace. Fourth, credentials (preferably a certificate) are handled outside "
            "the code repo, typically in Key Vault. Skip any one of these and the pipeline "
            "fails at a different, confusing point.",
            "Use the credential-rotation example to make group-based access click: if a "
            "secret is compromised and you need to rotate, you spin up a new app registration, "
            "add it to the security group, remove the old one, and every pipeline in the "
            "organization keeps working \u2014 no code changes, no workspace re-permissioning. If "
            "you had granted workspace access to the app directly, you'd be re-permissioning "
            "every workspace by hand.",
            "Close on the Gov point: service principals in a Gov tenant are Verify for Gov. "
            "The tenant setting approval, the app registration policy, and the customer's "
            "security-group governance all have to line up before this becomes viable. In this "
            "workshop we deliberately do not provision one \u2014 we walk the checklist so students "
            "know what to ask for later.",
        ]
    )
    page += 1

    # 10. Topic 8 - Fabric workspace Git integration ----------------------
    content_slide(
        prs, 8, "Fabric Workspace Git Integration", page=page,
        lead_items=[
            "A Fabric workspace can connect directly to a git branch and folder \u2014 workspace "
            "items sync bidirectionally with the repo instead of being published from Desktop.",
            "Sync direction matters: 'Update from Git' pulls repo changes into the workspace, "
            "'Commit to Git' pushes workspace changes back to the branch. Conflicts are "
            "resolved item by item, not file by file.",
            "Not every Fabric item type is supported \u2014 confirm which of your artifacts "
            "(semantic models, reports, notebooks, pipelines, lakehouses) round-trip cleanly "
            "before relying on it for release management.",
            "Commercial-focused / Verify for Gov: this feature is generally available in "
            "Commercial but availability, connector support, and identity requirements have to "
            "be confirmed in Gov before use.",
        ],
        why_items=[
            "Workspace Git integration collapses part of the CI/CD story \u2014 the workspace "
            "itself becomes an environment tied to a branch, so promotion can be 'merge branch, "
            "sync workspace' instead of a bespoke publish script.",
            "It also changes where the source of record lives during authoring: Fabric items "
            "authored in-workspace can be committed to git without ever passing through a "
            "PBIP on someone's desktop.",
            "Because behavior differs between Commercial and Gov, treat it as an optional "
            "commercial-enhanced path in this workshop, not a required capability.",
        ],
        footer="Lab connection: Optional commercial lab \u2014 Fabric workspace Git integration is "
               "walked through conceptually (connect workspace, select branch/folder, review "
               "sync status) unless the classroom tenant is validated for it.",
        script=[
            "Introduce this as a fundamentally different model from PBIP + Desktop + Publish. "
            "In workspace Git integration, a Fabric workspace is directly bound to a branch and "
            "folder in a git repo \u2014 there is no separate 'publish' step, because commits and "
            "syncs are how the workspace and the branch stay aligned.",
            "Explain the two sync directions clearly, because they're the source of most "
            "confusion. 'Update from Git' is repo \u2192 workspace: whoever merged the PR just "
            "changed the branch, and the workspace is now behind. 'Commit to Git' is workspace "
            "\u2192 repo: someone edited an item in-workspace, and the branch is now behind. "
            "Conflicts are resolved per item, not per file, and workspace users see them in a "
            "familiar UI, not by editing text.",
            "Warn about item-type coverage. Not everything a Fabric workspace can hold syncs "
            "cleanly, and the supported list changes over time. Before adopting workspace Git "
            "integration for release management, run through your specific artifact list "
            "(semantic models, reports, notebooks, pipelines, lakehouses) and confirm each one "
            "round-trips as expected.",
            "Land the Gov note plainly. This is Commercial-focused; Gov availability, connector "
            "and identity requirements, and network path all have to be validated separately. "
            "In this workshop we cover it because students will encounter it in commercial "
            "engagements, but we do not require it for Gov delivery.",
        ]
    )
    page += 1

    # 11. Topic 9 - Azure DevOps conceptual pipeline (custom flow diagram)
    pipeline_flow_slide(
        prs, 9, "Azure DevOps Conceptual Pipeline", page=page,
        stages=[
            ("Validate",
             "Lint PBIP folder structure. Run Tabular Editor Best "
             "Practice Analyzer against TMDL. Fail on high-severity rule violations."),
            ("Package",
             "Generate PBIX from PBIP if the target requires it. "
             "Capture build metadata (commit SHA, tag, build number)."),
            ("Deploy",
             "Authenticate as service principal. Call Power BI REST "
             "Imports API to publish to target workspace."),
            ("Configure",
             "UpdateParameters and UpdateDatasources to bind to the "
             "target environment. Set RLS role members. Enable schedule."),
            ("Smoke test",
             "Trigger refresh, poll status, run a canary DAX query. "
             "Publish evidence artifact and tag the commit on success."),
        ],
        note="Verify for Gov: network path to api.powerbigov.us, service-principal tenant "
             "settings, agent identity, and customer policy all validated before enabling.",
    )
    # Add speaker notes to this custom slide the same way slide_kit does
    from slide_kit import set_notes
    set_notes(prs.slides[-1], [
        "Walk the five stages left to right as a mental model, not a specific YAML. The point "
        "isn't which Azure DevOps task syntax to use \u2014 it's that any credible Power BI "
        "release pipeline has to answer these five questions in this order.",
        "Stage 1 (Validate) is where you catch a bad model before it reaches any environment: "
        "lint the PBIP folder shape (are the .Report and .SemanticModel folders where they "
        "should be, is the .pbix committed by mistake?) and run Tabular Editor's Best "
        "Practice Analyzer against the TMDL as a scripted check. High-severity BPA failures "
        "should fail the build.",
        "Stages 2 and 3 (Package and Deploy) are the mechanical part. If your Service target "
        "needs a PBIX, generate it from the PBIP now and stamp it with the commit SHA and "
        "build metadata so you can trace any deployed artifact back to source. Then "
        "authenticate as the service principal (never as a human), and call the Imports API "
        "to publish. Note how naturally this uses topics 6 (REST APIs) and 7 (service "
        "principals) from earlier today \u2014 that's not a coincidence.",
        "Stage 4 (Configure) is where most homegrown pipelines under-invest. Publishing the "
        "model isn't enough \u2014 you have to rebind parameters and data sources to the target "
        "environment, set RLS role memberships, and enable the refresh schedule. Do this "
        "with API calls, not by clicking around in the Service afterward.",
        "Stage 5 (Smoke test) is what earns the deployment the right to be called successful: "
        "trigger a refresh, poll until it finishes, run a canary DAX query, capture the "
        "evidence as a build artifact, and only then tag the commit. That evidence + tag is "
        "your rollback anchor from the very first topic today. Close by reminding students "
        "this is Verify for Gov \u2014 endpoint, identity, network, and customer policy all get "
        "validated before enabling this in a Gov tenant.",
    ])
    page += 1

    # 12. Topic 10 - GitHub Actions conceptual pipeline (same flow shape) --
    pipeline_flow_slide(
        prs, 10, "GitHub Actions Conceptual Pipeline", page=page,
        stages=[
            ("Trigger",
             "workflow on: push to main / pull_request / tag. Path "
             "filters scope runs to PBIP folders only."),
            ("Validate",
             "actions/checkout + PBIP structural checks + BPA rules "
             "run against TMDL. Fail fast on high-severity findings."),
            ("Authenticate",
             "OIDC federation to Entra ID (preferred) or a scoped "
             "service-principal secret from GitHub Encrypted Secrets."),
            ("Deploy",
             "Call Power BI REST Imports API against the correct "
             "cloud endpoint. Update parameters and data source bindings."),
            ("Record evidence",
             "Trigger and poll refresh. Upload logs and refresh "
             "history as workflow artifacts. Create annotated release on success."),
        ],
        note="Verify for Gov: hosted runner egress to api.powerbigov.us, OIDC federation, "
             "identity provider trust, and customer policy for GitHub availability all "
             "validated before enabling.",
    )
    set_notes(prs.slides[-1], [
        "Walk the same five-question mental model as the Azure DevOps slide, and call out "
        "explicitly that the shape is identical on purpose. Trigger \u2192 Validate \u2192 "
        "Authenticate \u2192 Deploy \u2192 Record evidence is the pattern; Azure DevOps and GitHub "
        "Actions are just two rentals of the same house.",
        "Stage 1 (Trigger) differs in flavor from Azure DevOps: GitHub Actions leans on push, "
        "pull_request, and tag events, and path filters keep the workflow from running when "
        "unrelated docs change. That path scoping matters for Power BI mono-repos where a "
        "PBIP is one folder among many.",
        "Stage 2 (Validate) is essentially the same lint + BPA gate as the Azure DevOps side. "
        "Fail fast \u2014 a broken TMDL should never reach the authentication step, both "
        "because it wastes runner minutes and because a partial deploy is far worse than a "
        "clean 'validation failed.'",
        "Stage 3 (Authenticate) is where I want students to pay attention: prefer OIDC "
        "federation to Entra ID over long-lived client secrets stored in GitHub Encrypted "
        "Secrets. OIDC means the workflow trades a short-lived token per run instead of "
        "carrying a secret at all. In a Gov context, whether OIDC federation is actually "
        "available and approved is one of the specific Verify-for-Gov items to check.",
        "Stages 4 and 5 (Deploy and Record evidence) reuse the same REST API surface \u2014 "
        "correct cloud endpoint, Imports API, parameter and datasource updates, refresh + "
        "canary check, then upload logs as workflow artifacts and cut an annotated release. "
        "Close by tying back: Gov availability of the GitHub runner network path, the OIDC "
        "trust, and the customer's overall approval for GitHub as a platform are all "
        "Verify-for-Gov gates before turning this on for real.",
    ])
    page += 1

    # 13. Topic 11 - Azure Government considerations ----------------------
    table_slide(
        prs, 11, "Azure Government Considerations", page=page,
        headers=["Area", "What to validate", "Why it matters"],
        col_widths=[2.5, 5.7, 3.7],
        rows=[
            ["Endpoints",
             "api.powerbigov.us for Power BI APIs; login.microsoftonline.us for identity; "
             "correct sovereign endpoints for every dependency.",
             "Wrong endpoint = silent auth against Commercial cloud."],
            ["Identity",
             "Entra ID Gov tenant, MFA, conditional access, service-principal tenant "
             "settings and app registration policy.",
             "Automation identity has to exist and be permitted before it can be used."],
            ["Network",
             "Egress path from build agents / runners / workstations to Gov endpoints; "
             "private endpoints or gateway where required.",
             "Runners in the wrong network can't reach the Gov endpoints at all."],
            ["Customer policy",
             "Workstation tooling policy for Tabular Editor / ALM Toolkit; repository "
             "hosting policy; approved CI/CD platforms.",
             "Tools you can't install or run don't belong in the design."],
            ["Feature availability",
             "Fabric workspace Git integration, Copilot, deployment pipelines, and other "
             "commercial-only features \u2014 confirmed available before design.",
             "Designing around a feature that isn't in Gov is a rework guarantee."],
        ],
        note="Rule of thumb: assume nothing about Gov availability. Confirm every endpoint, "
             "identity, network path, and feature before designing a release process around it.",
        script=[
            "This slide exists so no student walks out of the module thinking 'Gov is just "
            "Commercial with a different URL.' It is not. Endpoints, identity, network path, "
            "customer policy, and feature availability all have to be confirmed separately for "
            "any automation you plan to run against a Gov tenant.",
            "Walk the endpoints row first because it's the highest-impact and most common bug: "
            "hard-coded api.powerbi.com in a script, or forgetting -Environment USGov on "
            "Connect-PowerBIServiceAccount, results in the pipeline authenticating against the "
            "Commercial cloud, failing quietly, and leaving nobody sure whether the "
            "credentials, the network, or the code was the problem.",
            "For identity and network, explain that these are usually not the Power BI team's "
            "call \u2014 tenant settings, conditional access, and network egress from build agents "
            "are set by the customer's Entra and network teams. The Power BI team's job is to "
            "produce the concrete list of what needs to be enabled and negotiate for it early, "
            "not to discover it on release day.",
            "Land the feature-availability row with a specific example: designing an entire "
            "release process around Fabric workspace Git integration and then discovering it "
            "isn't approved in the target Gov tenant means starting over. The rule of thumb "
            "in the note \u2014 assume nothing, confirm everything \u2014 is not paranoia, it's the "
            "cheapest possible insurance policy.",
        ]
    )
    page += 1

    # 14. Topic 12 - Deployment checklist ---------------------------------
    checklist_slide(
        prs, "Deployment Checklist \u2014 Pre-Release Gate",
        kicker="Topic 12 \u2014 Pre-release gate", page=page,
        items=[
            "Source: PBIP committed on the correct branch; annotated tag applied; PBIX and "
            "cache excluded via .gitignore.",
            "Security: RLS roles reviewed; workspace membership matches release plan; "
            "sensitivity labels applied where required.",
            "Refresh: data source credentials configured on target dataset; schedule set; "
            "gateway (if used) validated.",
            "Validation: BPA passed on TMDL; measure totals reconciled to a trusted source; "
            "critical report pages smoke-tested.",
            "Environment: parameters and datasource bindings switched to the target "
            "environment; incremental refresh policy verified.",
            "Automation: service principal authentication tested end-to-end; cloud endpoint "
            "confirmed (Commercial vs Gov).",
            "Azure Government: endpoints, identity, network, and customer policy validated "
            "for every automation touchpoint.",
            "Rollback: previous release tag identified; rollback procedure written down and "
            "known to on-call; release notes attached to the tag.",
        ],
        script=[
            "Present this as the artifact that pulls the module together: everything we've "
            "discussed \u2014 PBIP, git, external tools, REST APIs, service principals, Fabric "
            "Git, DevOps and Actions pipelines, Gov posture \u2014 shows up here as a concrete "
            "checkbox on a pre-release checklist.",
            "Walk the eight items and explicitly map each one back to an earlier topic. "
            "Source ties to topics 2-4. Security ties to earlier modules' RLS and to topic 1's "
            "governance goal. Refresh, Validation, and Environment come out of REST APIs and "
            "the CI/CD pipeline stages. Automation and Azure Government are topics 6, 7, and "
            "11. Rollback closes the loop with topic 1.",
            "Emphasize that this checklist is not aspirational \u2014 in a production Power BI "
            "practice, no release ships without it filled out and stored as release evidence. "
            "In a Gov engagement especially, auditors expect to see this artifact per "
            "release, not summarized once a year.",
            "Bridge to Exercise 7 in the lab: students fill out this checklist against their "
            "own PBIP project as the module's culminating deliverable. Half-answers ('yes but "
            "we'll fix it in Prod') count as fails \u2014 the point of the checklist is to move "
            "the failure earlier, not to paper over it.",
        ],
    )
    page += 1

    # 15. Module lab walkthrough ------------------------------------------
    checklist_slide(
        prs, "Module Lab Walkthrough",
        kicker="Topic 13 \u2014 What you'll build", page=page,
        items=[
            "Exercise 1: PBIP file structure \u2014 identify .Report, .SemanticModel, generated "
            "and cache files in your project folder.",
            "Exercise 2: Source control workflow \u2014 feature branch, small PBIP change, "
            "git status/diff, commit, PR review discussion.",
            "Exercise 3: Tabular Editor workflow (Verify for Gov) \u2014 open model, review "
            "measures, small metadata change, review diff.",
            "Exercise 4: ALM Toolkit model comparison (Verify for Gov) \u2014 diff two model "
            "versions and discuss deployment impact.",
            "Optional labs: REST API deployment, PowerShell admin, service principal setup "
            "\u2014 all conceptual walkthroughs unless the tenant is validated.",
            "Optional commercial lab: Fabric workspace Git integration \u2014 connect workspace, "
            "select branch/folder, review sync behavior.",
            "Exercise 5 & 6: Conceptual Azure DevOps and GitHub Actions CI/CD pipelines \u2014 "
            "identify stages, auth, validation, Gov notes.",
            "Exercise 7: Deployment checklist \u2014 fill it out end to end against your PBIP.",
        ],
        script=[
            "Use this slide as the literal table of contents for the hands-on session. Only "
            "Exercises 1, 2, and 7 are guaranteed hands-on for every student today \u2014 the "
            "others depend on tenant and workstation validation.",
            "Set expectations for time: PBIP structure (Exercise 1) is quick, git workflow "
            "(Exercise 2) usually takes the most classroom time because git behavior and PR "
            "review discussion open up broader conversations, and the deployment checklist "
            "(Exercise 7) benefits from being done thoughtfully rather than rushed.",
            "For any exercise students can't do hands-on today, tell them to still read the "
            "task list and think through what they would need to validate in a real "
            "engagement \u2014 the value is in knowing the shape of the work, so they can scope "
            "it correctly when it comes up on a real project.",
            "Bridge to Exercise 7 specifically: the deployment checklist is the artifact they "
            "should walk out with. If nothing else from this module survives, that checklist "
            "should \u2014 it turns everything we discussed into a repeatable, auditable ritual.",
        ]
    )
    page += 1

    # 16. Knowledge check & discussion ------------------------------------
    checklist_slide(
        prs, "Knowledge Check & Discussion",
        kicker="Topic 14 \u2014 Wrap-up", page=page,
        items=[
            "Why is PBIP the source of record and not PBIX? Give one concrete review scenario "
            "that PBIX makes impossible.",
            "Which lifecycle goal (repeatability, reviewability, promotion, governance, "
            "rollback) is most at risk on your current projects, and what would fix it?",
            "Which parts of a REST API / PowerShell deployment change when you move from "
            "Commercial to Gov? Name three.",
            "When would you accept a client secret over certificate auth for a service "
            "principal, and how would you compensate?",
            "Fabric workspace Git integration vs. PBIP + Desktop + Publish \u2014 when would you "
            "choose each, and what changes in the review process?",
            "Walk through the five CI/CD pipeline stages from memory (Validate, Package, "
            "Deploy, Configure, Smoke test) and give one Gov note for each.",
            "Deployment checklist self-assessment: which section (source, security, refresh, "
            "validation, automation, Gov, rollback) is your team weakest on today?",
        ],
        script=[
            "Run this as an open discussion, not a quiz. The goal is to hear reasoning \u2014 "
            "especially on tradeoff questions where multiple defensible answers exist.",
            "Pick two or three questions to spend real time on. The 'lifecycle goal most at "
            "risk on your projects' question tends to surface the most honest, useful "
            "conversation because students immediately map it to real pain they're feeling.",
            "For the Commercial-vs-Gov question, expect students to name at least the "
            "endpoint change; push them to also name the identity and network-path changes "
            "\u2014 that's the completeness answer, and it reinforces topic 11.",
            "Close by tying back to the deployment checklist: the last question is deliberately "
            "self-reflective. Students should walk out with a specific weakest-link category "
            "in mind, because that's where they should invest first when they go back to "
            "their real projects. Then bridge into the closing slide's capstone framing.",
        ]
    )
    page += 1

    # 17. Closing slide (course wrap-up) ----------------------------------
    # Custom closing that reframes "Up next" as the capstone rather than another module.
    s = blank_slide(prs)
    add_rect(s, 0, 0, SLIDE_W, SLIDE_H, NAVY_DARK)
    add_text(s, Inches(0.9), Inches(1.4), Inches(11), Inches(0.7),
              "COURSE WRAP-UP", size=15, color=GOLD, bold=True, font=BODY_FONT)
    add_text(s, Inches(0.9), Inches(1.85), Inches(11), Inches(1.2),
              "Module 11 Complete \u2014 Modules 1\u201311 delivered.",
              size=34, color=WHITE, bold=True, font=HEADER_FONT)
    add_text(s, Inches(0.9), Inches(3.05), Inches(11.5), Inches(1.7),
              "You now have a complete, source-controlled Power BI lifecycle: PBIP as source of "
              "record, git-based review, external-tool and REST/PowerShell automation patterns, "
              "service-principal identity, conceptual Azure DevOps and GitHub Actions CI/CD, "
              "and an Azure Government-aware deployment checklist to gate every release.",
              size=15, color=ICE, italic=True, font=BODY_FONT, line_spacing=1.2)
    add_rect(s, Inches(0.9), Inches(4.95), Inches(11.5), Inches(1.5), NAVY)
    add_text(s, Inches(1.15), Inches(5.05), Inches(11.0), Inches(0.45),
              "Next: Capstone \u2014 Enterprise-Ready Power BI Solution", size=17, color=GOLD,
              bold=True, font=BODY_FONT)
    add_text(s, Inches(1.15), Inches(5.55), Inches(11.0), Inches(0.85),
              "Combine every module \u2014 semantic model, DAX, Power Query, report UX, "
              "performance, security, deployment, governance, monitoring, and today's "
              "automation \u2014 into one end-to-end Contoso Advanced Manufacturing delivery.",
              size=13, color=ICE, font=BODY_FONT, line_spacing=1.15)
    add_page_number(s, page, dark=True)
    set_notes(s, [
        "Congratulate the class \u2014 they've completed the full eleven-module instructor track. "
        "Everything from advanced semantic modeling in Module 1 through automation and Gov-aware "
        "deployment today is now in their toolkit.",
        "Reframe today's module specifically as the capstone-enabler: without lifecycle "
        "management, everything they built in Modules 1-10 lives as a series of one-off PBIX "
        "files. With today's material, those artifacts become a governed, reviewable, "
        "deployable engineering asset.",
        "Preview the capstone (Module 12): it is not a lecture module and has no slide deck. "
        "Instead, it is the end-to-end 'Enterprise-Ready Power BI Solution' delivery for "
        "Contoso Advanced Manufacturing, exercising every capability from the previous eleven "
        "modules \u2014 model, DAX, Power Query, report UX, RLS, deployment, governance, "
        "monitoring, and lifecycle \u2014 in a single Gov-ready solution with optional "
        "commercial-enhanced extensions.",
        "Take final questions with a bias toward capstone readiness: which module's material do "
        "students feel least confident carrying into the capstone? Point them at that module's "
        "artifacts (README, lab, deck) as their pre-capstone review and thank the class.",
    ])
    page += 1

    OUT_PATH.parent.mkdir(parents=True, exist_ok=True)
    prs.save(str(OUT_PATH))
    print(f"Saved: {OUT_PATH} ({len(prs.slides.__iter__.__self__._sldIdLst)} slides)")


if __name__ == "__main__":
    build()
