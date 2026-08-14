#!/usr/bin/env python3
"""
Builds the Lab 10 (Premium, Fabric, and Capacity-Aware Architecture)
instructor deck.
Run from repo root: python tools/pptx-labs/build_lab10.py
Output: modules/10-premium-fabric-capacity/assets/premium-fabric-capacity.pptx
"""
import sys
from pathlib import Path

from pptx.util import Inches, Pt, Emu
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR
from pptx.enum.shapes import MSO_SHAPE, MSO_CONNECTOR

sys.path.insert(0, str(Path(__file__).resolve().parent))
from slide_kit import (
    new_presentation, title_slide, agenda_slide,
    content_slide, table_slide, checklist_slide, closing_slide,
    blank_slide, add_rect, add_text, add_page_number, set_notes,
    NAVY, NAVY_DARK, ICE, WHITE, INK, SLATE, GOLD, LIGHT_BG, CARD_BORDER,
    HEADER_FONT, BODY_FONT, SLIDE_W, SLIDE_H,
)

OUT_PATH = Path(__file__).resolve().parent.parent.parent / \
    "modules" / "10-premium-fabric-capacity" / "assets" / "premium-fabric-capacity.pptx"

MODULE_NO = 10
TITLE = "Premium, Fabric, and Capacity-Aware Architecture"
SUBTITLE = ("Choosing between Pro, PPU, Premium capacity, and Fabric capacity — and "
            "designing Gov-safe fallbacks for XMLA, paginated reports, Direct Lake, "
            "OneLake, and capacity metrics.")

AGENDA_TOPICS = [
    "Capacity-aware architecture",
    "Licensing and capacity options",
    "Large semantic models",
    "XMLA endpoint",
    "Paginated reports",
    "Fabric capacity",
    "Direct Lake",
    "OneLake, Lakehouse, and Warehouse",
    "Semantic Link",
    "Capacity metrics and throttling",
    "Azure Government considerations",
    "Architecture decision review",
    "Module lab walkthrough",
    "Knowledge check and discussion",
]


# ---------------------------------------------------------------------------
# Custom diagram: throttling / capacity-smoothing states
# ---------------------------------------------------------------------------

def throttling_states_slide(prs, number, title, page, note=None, script=None):
    """Custom diagram: Fabric/Premium capacity CU consumption smoothing states —
    Healthy → Overage (smoothing) → Interactive Delay → Interactive Rejection →
    Background Rejection."""
    s = blank_slide(prs)
    add_rect(s, 0, 0, SLIDE_W, SLIDE_H, WHITE)
    add_rect(s, 0, 0, SLIDE_W, Inches(1.15), NAVY)
    add_text(s, Inches(0.7), Inches(0.22), Inches(1.6), Inches(0.4),
             f"TOPIC {number:02d}", size=13, color=GOLD, bold=True, font=BODY_FONT)
    add_text(s, Inches(0.7), Inches(0.52), Inches(11.5), Inches(0.55), title,
             size=24, color=WHITE, bold=True, font=HEADER_FONT)

    add_text(s, Inches(0.7), Inches(1.4), Inches(11.9), Inches(0.5),
             "How a Fabric/Premium capacity moves through CU-consumption states",
             size=16, color=NAVY_DARK, bold=True, font=HEADER_FONT)
    add_text(s, Inches(0.7), Inches(1.9), Inches(11.9), Inches(0.55),
             "Capacity Units (CU) are smoothed over a 24-hour window. Sustained "
             "overage escalates from delay to rejection — interactive workloads "
             "are throttled before background jobs are.",
             size=12.5, color=SLATE, italic=True, font=BODY_FONT, line_spacing=1.2)

    states = [
        ("Healthy",
         "CU usage\u00a0\u2264\u00a0100%",
         "All interactive & background jobs run normally.",
         ICE, NAVY_DARK),
        ("Overage\n(smoothed)",
         "Short bursts >100%",
         "Excess CU carried forward against the 24-h future budget.",
         ICE, NAVY_DARK),
        ("Interactive\nDelay",
         "10-min avg\u00a0>\u00a0100%",
         "Queries/report renders delayed ~20s before running.",
         GOLD, NAVY_DARK),
        ("Interactive\nRejection",
         "60-min avg\u00a0>\u00a0100%",
         "New interactive requests refused; reports fail to render.",
         NAVY, WHITE),
        ("Background\nRejection",
         "24-h avg\u00a0>\u00a0100%",
         "Scheduled refreshes and Dataflows refused until debt clears.",
         NAVY_DARK, WHITE),
    ]

    row_y = Inches(3.05)
    box_h = Inches(2.05)
    box_w = Inches(2.35)
    gap = Inches(0.16)
    start_x = Inches(0.55)

    positions = []
    for i, (label, trigger, effect, fill, text_color) in enumerate(states):
        x = start_x + i * (box_w + gap)
        add_rect(s, x, row_y, box_w, box_h, fill,
                 line_color=NAVY if fill in (ICE, GOLD) else None)
        add_text(s, x + Inches(0.08), row_y + Inches(0.12), box_w - Inches(0.16),
                 Inches(0.7), label, size=14, color=text_color, bold=True,
                 font=BODY_FONT, align=PP_ALIGN.CENTER, anchor=MSO_ANCHOR.MIDDLE,
                 line_spacing=1.05)
        add_text(s, x + Inches(0.08), row_y + Inches(0.82), box_w - Inches(0.16),
                 Inches(0.42), trigger, size=11, color=text_color, italic=True,
                 font="Consolas", align=PP_ALIGN.CENTER, anchor=MSO_ANCHOR.MIDDLE)
        add_text(s, x + Inches(0.1), row_y + Inches(1.28), box_w - Inches(0.2),
                 Inches(0.72), effect, size=10.5, color=text_color,
                 font=BODY_FONT, align=PP_ALIGN.CENTER, anchor=MSO_ANCHOR.MIDDLE,
                 line_spacing=1.15)
        positions.append((x, x + box_w))

    arrow_h = Inches(0.28)
    arrow_y = row_y + box_h + Inches(0.15)
    for i in range(len(states) - 1):
        x1 = positions[i][1]
        x2 = positions[i + 1][0]
        arrow = s.shapes.add_shape(
            MSO_SHAPE.RIGHT_ARROW,
            Emu(int(x1 - Inches(0.05))), Emu(int(arrow_y)),
            Emu(int(x2 - x1 + Inches(0.1))), arrow_h,
        )
        arrow.fill.solid()
        arrow.fill.fore_color.rgb = GOLD
        arrow.line.fill.background()

    add_text(s, Inches(0.7), arrow_y + Inches(0.4), Inches(11.9), Inches(0.45),
             "Escalation as sustained CU overage grows  \u2192  "
             "recovery flows right-to-left as usage drops back under budget.",
             size=11.5, color=SLATE, italic=True, font=BODY_FONT,
             align=PP_ALIGN.CENTER)

    if note:
        add_text(s, Inches(0.7), Inches(6.7), Inches(11.9), Inches(0.7), note,
                 size=12.5, color=SLATE, italic=True, font=BODY_FONT,
                 line_spacing=1.1)
    add_page_number(s, page)
    if script:
        set_notes(s, script)
    return s


# ---------------------------------------------------------------------------
# Custom diagram: workspace-to-capacity assignment flow
# ---------------------------------------------------------------------------

def workspace_assignment_slide(prs, number, title, page, note=None, script=None):
    """Custom diagram: workspace assignment to different license/capacity backends."""
    s = blank_slide(prs)
    add_rect(s, 0, 0, SLIDE_W, SLIDE_H, WHITE)
    add_rect(s, 0, 0, SLIDE_W, Inches(1.15), NAVY)
    add_text(s, Inches(0.7), Inches(0.22), Inches(1.6), Inches(0.4),
             f"TOPIC {number:02d}", size=13, color=GOLD, bold=True, font=BODY_FONT)
    add_text(s, Inches(0.7), Inches(0.52), Inches(11.5), Inches(0.55), title,
             size=24, color=WHITE, bold=True, font=HEADER_FONT)

    add_text(s, Inches(0.7), Inches(1.35), Inches(11.9), Inches(0.5),
             "Workspaces are assigned to a license mode; the mode gates which "
             "features light up.",
             size=15, color=NAVY_DARK, bold=True, font=HEADER_FONT)

    # Left: workspace card
    ws_x, ws_y, ws_w, ws_h = Inches(0.7), Inches(2.4), Inches(3.2), Inches(3.4)
    add_rect(s, ws_x, ws_y, ws_w, ws_h, ICE, line_color=NAVY)
    add_text(s, ws_x, ws_y + Inches(0.2), ws_w, Inches(0.45),
             "Power BI workspace", size=15, color=NAVY_DARK, bold=True,
             font=BODY_FONT, align=PP_ALIGN.CENTER)
    add_text(s, ws_x + Inches(0.2), ws_y + Inches(0.85), ws_w - Inches(0.4),
             ws_h - Inches(1.0),
             "Contains:\n  • Semantic models\n  • Reports & dashboards\n"
             "  • Dataflows / Datamarts\n  • Paginated reports\n  • Notebooks / "
             "Lakehouses (Fabric)\n\nAdmin sets one license mode per workspace.",
             size=12, color=INK, font=BODY_FONT, anchor=MSO_ANCHOR.TOP,
             line_spacing=1.3)

    # Right column: 4 target license modes
    targets = [
        ("Pro (shared capacity)",
         "Per-user Pro license. Model \u2264 1 GB, 8 refreshes/day. No XMLA "
         "write, no paginated, no large model.",
         SLATE),
        ("Premium Per User (PPU)",
         "Per-user PPU license. Premium features (XMLA, paginated, large models "
         "with tenant setting) but audience must also hold PPU.",
         NAVY),
        ("Premium capacity (P/EM SKU)",
         "Reserved capacity (P1-P5, EM1-EM3). Free-viewer sharing, XMLA r/w, "
         "paginated, large models, autoscale add-on.",
         NAVY_DARK),
        ("Fabric capacity (F SKU)",
         "F2-F2048. Fabric CU meter, OneLake, Lakehouse/Warehouse, Direct Lake, "
         "Semantic Link. F64+ unlocks free Power BI viewer sharing.",
         GOLD),
    ]
    col_x = Inches(5.35)
    col_w = Inches(7.55)
    top = Inches(2.2)
    row_h = Inches(0.95)
    gap = Inches(0.12)
    for i, (name, desc, accent) in enumerate(targets):
        y = top + i * (row_h + gap)
        add_rect(s, col_x, y, col_w, row_h, WHITE, line_color=CARD_BORDER)
        add_rect(s, col_x, y, Inches(0.14), row_h, accent)
        add_text(s, col_x + Inches(0.28), y + Inches(0.08), col_w - Inches(0.4),
                 Inches(0.35), name, size=13.5, color=NAVY_DARK, bold=True,
                 font=BODY_FONT)
        add_text(s, col_x + Inches(0.28), y + Inches(0.42), col_w - Inches(0.4),
                 row_h - Inches(0.45), desc, size=11, color=INK, font=BODY_FONT,
                 line_spacing=1.2)

        arrow_x1 = ws_x + ws_w + Inches(0.05)
        arrow_x2 = col_x - Inches(0.05)
        arrow_y = y + row_h / 2 - Inches(0.11)
        arrow = s.shapes.add_shape(
            MSO_SHAPE.RIGHT_ARROW,
            Emu(int(arrow_x1)), Emu(int(arrow_y)),
            Emu(int(arrow_x2 - arrow_x1)), Emu(int(Inches(0.22))),
        )
        arrow.fill.solid()
        arrow.fill.fore_color.rgb = accent
        arrow.line.fill.background()

    if note:
        add_text(s, Inches(0.7), Inches(6.75), Inches(11.9), Inches(0.6), note,
                 size=12.5, color=SLATE, italic=True, font=BODY_FONT,
                 line_spacing=1.1)
    add_page_number(s, page)
    if script:
        set_notes(s, script)
    return s


def build():
    prs = new_presentation()
    page = 1

    # 1. Title
    title_slide(
        prs, MODULE_NO, TITLE, SUBTITLE,
        script=[
            "Welcome learners to Module 10. Frame this as the 'architecture money "
            "module' \u2014 every decision here shows up on the bill and in the "
            "tenant's capability matrix. Pro vs PPU vs Premium capacity vs Fabric "
            "capacity is not just licensing trivia; it gates which features are "
            "even available in the workspaces they will build in later modules.",
            "Set expectations about hands-on work: most of this module is "
            "conceptual and diagram-driven because Fabric capacity, Direct Lake, "
            "OneLake, Lakehouse/Warehouse, Semantic Link, and autoscale are all "
            "'Commercial-focused / Verify for Gov'. That is a real classroom "
            "constraint, not a hedge \u2014 don't demo features the tenant may not "
            "have.",
            "Preview the lab: five short exercises \u2014 licensing/capacity "
            "comparison, XMLA endpoint (where available), paginated reports "
            "(where available), large semantic model settings (where available), "
            "and capacity metrics/throttling concepts \u2014 each with an alternate "
            "conceptual path so the class stays synchronized.",
            "Land the throughline that will repeat all module: 'What is the "
            "Gov-safe fallback?' For every Premium/Fabric feature we discuss, we "
            "will name the fallback so learners can still deliver value in an "
            "Azure Government or restricted-tenant engagement.",
        ]
    )
    page += 1

    # 2. Agenda
    agenda_slide(
        prs, MODULE_NO, AGENDA_TOPICS, page=page,
        script=[
            "Walk the room through the fourteen items quickly. Topics 1-2 frame "
            "requirements and license options; 3-5 cover Premium-era features "
            "(large models, XMLA, paginated); 6-9 cover the Fabric-era features "
            "(capacity, Direct Lake, OneLake, Semantic Link).",
            "Topic 10 (capacity metrics and throttling) is the operational "
            "reality check \u2014 it's where they'll learn what 'the capacity is "
            "slow' actually means in CU terms.",
            "Topics 11-12 pull it all together for Azure Government: how to "
            "validate, and how to run an architecture decision review with "
            "explicit fallbacks. Then lab walkthrough and knowledge check.",
        ]
    )
    page += 1

    # 3. Topic 1 — Capacity-aware architecture (content)
    content_slide(
        prs, 1, "Capacity-Aware Architecture", page=page,
        lead_items=[
            "Start every architecture conversation from six requirement axes: "
            "user scale (viewers + authors), data size (per model, per "
            "workspace), refresh needs (frequency + duration), query performance "
            "(latency SLOs), governance (certification, sharing, sovereignty), "
            "and cloud availability (Commercial vs. Government).",
            "Only after those are captured do you map to a license/capacity "
            "SKU \u2014 Pro, PPU, Premium capacity (P/EM), or Fabric capacity (F).",
            "Every 'nice-to-have' Premium/Fabric feature (XMLA, paginated, "
            "Direct Lake, OneLake) has a matching Gov-safe fallback that must "
            "be identified before the design is signed off.",
        ],
        why_items=[
            "Skipping requirements leads to buying capacity for the wrong "
            "bottleneck \u2014 e.g. paying for an F64 when the real problem was a "
            "slow refresh that would have been solved by incremental refresh on "
            "a Pro workspace.",
            "Cloud availability is not a footnote in Gov engagements; it can "
            "eliminate Fabric-only options entirely, which changes the entire "
            "recommendation.",
            "Naming the fallback up front avoids painful redesigns when a "
            "feature turns out not to be available in the target tenant.",
        ],
        footer="Lab connection: Exercise 1 walks students through exactly these "
               "six requirement axes for a defined workload scenario.",
        script=[
            "Open by pushing back on a common instinct: 'we need Premium' is "
            "not a design conclusion, it's a shortcut. Force the discussion "
            "back to the six axes: user scale, data size, refresh needs, query "
            "performance, governance, cloud availability. Every one of those is "
            "something you can quantify.",
            "Give a concrete example: a workload with 200 viewers, a 400 MB "
            "model, twice-daily refresh, and no XMLA/paginated requirement is a "
            "Pro workload \u2014 buying Premium capacity for it would waste money "
            "and add operational overhead they don't need.",
            "Now flip it: 5,000 viewers, a 6 GB model, hourly refresh, XMLA "
            "read/write for CI/CD, and free viewer sharing across the company. "
            "That's clearly Premium or Fabric F64+. The six axes make the "
            "recommendation almost mechanical.",
            "Close by tying to Gov: any 'Commercial-focused / Verify for Gov' "
            "feature (Fabric capacity, Direct Lake, OneLake, Semantic Link, "
            "autoscale) needs a documented fallback before the design goes into "
            "review. That's exactly what Exercise 1 in the lab has students "
            "practice.",
        ]
    )
    page += 1

    # 4. Topic 2 — Licensing and capacity options (table)
    table_slide(
        prs, 2, "Licensing and Capacity Options", page=page,
        headers=["Option", "Sharing model", "Key features & limits",
                 "Gov status"],
        col_widths=[2.2, 2.6, 5.4, 1.7],
        rows=[
            ["Power BI Pro",
             "Every viewer needs a Pro license",
             "Shared capacity. Semantic model \u2264 1 GB, 8 scheduled "
             "refreshes/day. No XMLA write, no paginated, no large models.",
             "Available"],
            ["Premium Per User (PPU)",
             "Every viewer needs a PPU license",
             "Premium features (XMLA r/w, paginated, deployment pipelines, "
             "large models with tenant setting) at per-user scale.",
             "Verify for Gov"],
            ["Premium capacity (P / EM SKUs)",
             "Free viewers with Pro authors (P SKUs)",
             "Reserved v-cores. P1=25 GB model, P5=400 GB. Adds XMLA r/w, "
             "paginated, large models, deployment pipelines, autoscale add-on.",
             "Verify for Gov"],
            ["Fabric capacity (F SKUs)",
             "F64+ gives free Power BI viewer sharing",
             "Pay-as-you-go or reserved CUs. F2-F2048. Unlocks OneLake, "
             "Lakehouse/Warehouse, Direct Lake, notebooks, Semantic Link, and "
             "all Premium features.",
             "Commercial-focused"],
            ["Embedded (A SKUs)",
             "For ISV / embedded scenarios",
             "Pay-as-you-go equivalent to P SKUs. Consumed by app users, not "
             "signed-in Power BI users.",
             "Verify for Gov"],
        ],
        note="Lab connection: Exercise 1 asks learners to pick an option per "
             "scenario and record required Gov validations.",
        script=[
            "Read the table top to bottom in one pass, then go back and drill "
            "in on the two rows that trip up most classes: PPU and the "
            "difference between P and F SKUs.",
            "PPU is often misunderstood as 'Premium for small orgs'. It isn't. "
            "It unlocks Premium features at per-user pricing, but every viewer "
            "must also hold a PPU license \u2014 so it does not solve 'share with "
            "the whole company for free', which is the reason most orgs move to "
            "P or F64+.",
            "For F vs P: they overlap on Premium features, but only F unlocks "
            "the Fabric-era data engineering surface (OneLake, Lakehouse, "
            "Direct Lake, notebooks, Semantic Link). If the customer wants "
            "those, Fabric is the answer \u2014 assuming it is available in their "
            "cloud.",
            "Close on the Gov column: only Pro is unambiguously available "
            "everywhere. Everything else is 'Verify for Gov' or 'Commercial-"
            "focused'. This is why every capacity decision in a Gov engagement "
            "has to be paired with a validation step and a fallback.",
        ]
    )
    page += 1

    # 5. Topic 3 — Large semantic models (content)
    content_slide(
        prs, 3, "Large Semantic Models", page=page,
        lead_items=[
            "'Large semantic model storage format' is a workspace + model "
            "toggle that lifts the 1 GB / 10 GB model-size cap for a Premium "
            "or PPU-backed workspace, subject to the SKU's memory limit (P1 = "
            "25 GB, P3 = 100 GB, P5 = 400 GB; F SKUs scale similarly).",
            "It requires: (1) Premium capacity or PPU, (2) tenant setting "
            "'Large semantic model storage format' enabled, and (3) the "
            "workspace storage format switched from small to large.",
            "Memory pressure matters more than raw file size: refresh needs "
            "\u2248 2\u00d7 model size in RAM, so a 30 GB model on P1 (25 GB) will "
            "fail to refresh even though the file 'fits'.",
        ],
        why_items=[
            "Without this setting, models silently cap at 1 GB \u2014 authors "
            "notice only when refresh starts failing with an 'exceeds the "
            "available memory' error.",
            "The setting is the prerequisite for enterprise semantic models "
            "(hundreds of GB, hundreds of measures) that many reports share.",
            "Sizing capacity by model size alone is a common mistake; the "
            "instructor should always frame it as 'peak refresh memory', not "
            "'file size'.",
        ],
        footer="Lab connection: Exercise 4 (where available) has learners "
               "inspect the workspace + semantic model settings and document "
               "model-size and refresh-memory constraints.",
        script=[
            "Anchor the topic to the concrete symptom: 'my model won't refresh "
            "even though it's only 15 GB and my capacity says 25 GB'. That is "
            "almost always a memory-during-refresh problem, not a storage-size "
            "problem.",
            "Walk through the enablement gates in order: capacity type "
            "(Premium/PPU/F), tenant admin toggle for 'Large semantic model "
            "storage format', then the per-workspace storage-format switch. "
            "All three must be on. Missing any one is the top classroom "
            "gotcha.",
            "Give the 2\u00d7 rule of thumb: peak refresh RAM is roughly twice "
            "the compressed model size, because the refresh needs to load the "
            "existing model AND materialize the new one before swapping. "
            "Sizing to file size alone gets you paged out.",
            "Close with the Gov note: large semantic models are 'Verify for "
            "Gov'. In a Government tenant, confirm the tenant setting is "
            "enabled and the target capacity supports it before promising a "
            "large-model architecture.",
        ]
    )
    page += 1

    # 6. Topic 4 — XMLA endpoint (content)
    content_slide(
        prs, 4, "XMLA Endpoint", page=page,
        lead_items=[
            "The XMLA endpoint exposes the workspace's semantic models over "
            "the Analysis Services protocol. Read-only is on by default for "
            "Premium/PPU/F; read/write is admin-toggled per capacity.",
            "Read/write unlocks external-tool workflows: Tabular Editor (bulk "
            "measure edits, calculation groups, best-practice analyzer), DAX "
            "Studio (server-side query traces), SSMS (scripted deployment), "
            "and ALM Toolkit (schema diffs).",
            "ALM patterns become possible: deploy semantic model changes from "
            "source control via TMDL/BIM scripts, run partitioned refreshes, "
            "and script role membership \u2014 all without opening Desktop.",
        ],
        why_items=[
            "It's the bridge between 'Power BI as an app' and 'Power BI as an "
            "enterprise BI platform under source control'.",
            "Without XMLA, changes are gated by opening the PBIX in Desktop "
            "\u2014 which does not scale for enterprise models with dozens of "
            "authors or CI/CD pipelines.",
            "It is also the enabler for partitioned refresh strategies that "
            "aren't expressible in the Desktop UI, e.g. refreshing a single "
            "year-partition of a 10-year fact table.",
        ],
        footer="Lab connection: Exercise 2 (where available) has learners "
               "confirm the endpoint, connect with an approved tool, and "
               "document allowed read/write operations \u2014 or take the "
               "conceptual alternate path.",
        script=[
            "Open by naming the audience: XMLA is 'enterprise BI ops' \u2014 the "
            "people building CI/CD pipelines, managing 50+ measure models, "
            "running scripted deployments. It is not a feature every author "
            "uses; it's a feature every enterprise-scale author eventually "
            "depends on.",
            "Walk the tool list: Tabular Editor is the most common (bulk "
            "measure editing, calculation groups, BPA). DAX Studio for "
            "performance traces. SSMS for scripted deployment. ALM Toolkit "
            "for diffing schema between environments. Each one is unlocked by "
            "XMLA r/w on the target capacity.",
            "Give the ALM angle explicitly: enterprise teams keep TMDL or BIM "
            "definitions in Git, and CI/CD pipelines push those definitions "
            "through XMLA to Dev / Test / Prod workspaces. This is what makes "
            "'infrastructure as code' possible for semantic models.",
            "Close with the validation checklist: capacity type, workspace "
            "settings, admin toggle for XMLA endpoint = read/write, tenant "
            "settings for external tools, and \u2014 in Gov \u2014 confirmation that "
            "the endpoint is exposed in that sovereign cloud. All 'Verify for "
            "Gov'.",
        ]
    )
    page += 1

    # 7. Topic 5 — Paginated reports (content)
    content_slide(
        prs, 5, "Paginated Reports", page=page,
        lead_items=[
            "Paginated reports are the pixel-perfect, print-oriented reporting "
            "surface built on the SSRS/RDL engine \u2014 authored in Power BI "
            "Report Builder (a separate free tool), not in Desktop.",
            "They shine for operational scenarios: invoices, purchase orders, "
            "regulatory filings, statement runs \u2014 anything that must land on "
            "an exact page with headers, footers, page numbers, and precise "
            "column alignment.",
            "They support parameterization, subscriptions, and export to PDF, "
            "Excel (native cell alignment), Word, CSV, and image formats \u2014 "
            "with much larger data volumes than interactive reports.",
        ],
        why_items=[
            "Interactive Power BI reports are terrible at print and at "
            "1,000-row exports. Paginated is the right tool for both.",
            "Regulated industries frequently have 'must produce this exact "
            "PDF layout' requirements that only paginated can meet.",
            "They can be surfaced side-by-side with interactive reports in "
            "the same workspace, so users don't need a separate portal.",
        ],
        footer="Lab connection: Exercise 3 (where available) walks through "
               "opening/creating a simple paginated report with a table and a "
               "parameter, or the conceptual alternate: when paginated beats "
               "interactive.",
        script=[
            "Start by drawing the distinction crisply: interactive reports "
            "are optimized for exploration on a screen; paginated reports are "
            "optimized for a fixed layout that lands on paper (or a PDF that "
            "acts like paper). Different job, different tool.",
            "Give the canonical examples: a monthly invoice run, a compliance "
            "filing that must match a regulator's template, a 20-page "
            "customer statement with dynamic headers and page X of Y. "
            "Interactive Power BI cannot produce those cleanly \u2014 paginated "
            "can.",
            "Mention the authoring reality: Power BI Report Builder is a "
            "separate free download, and the authoring experience feels like "
            "SSRS from a decade ago. That is by design \u2014 the engine is "
            "battle-tested; the UX just isn't modernized. Set expectations "
            "with students so they aren't surprised.",
            "Close on the Gov note: paginated is 'Verify for Gov'. Validate "
            "licensing (P/EM/F SKU or PPU, not Pro), workspace support, and "
            "cloud availability before promising it. If unavailable, the "
            "fallback is typically an SSRS-on-VM or an interactive-report "
            "'export to PDF' path with a strong 'this is not pixel-perfect' "
            "caveat.",
        ]
    )
    page += 1

    # 8. Topic 6 — Fabric capacity (table)
    table_slide(
        prs, 6, "Fabric Capacity", page=page,
        headers=["Concept", "What it is", "Instructor talking point"],
        col_widths=[2.6, 5.4, 3.9],
        rows=[
            ["Capacity Unit (CU)",
             "The single meter across every Fabric workload \u2014 Power BI "
             "queries/refreshes, Spark notebooks, Warehouse queries, Data "
             "Factory pipelines, Real-Time Analytics. Smoothed over 24 h.",
             "Everything on Fabric consumes CU; sizing means summing all "
             "workloads, not just Power BI."],
            ["F SKUs",
             "F2, F4, F8, F16, F32, F64, F128, F256, F512, F1024, F2048. "
             "F64 unlocks free Power BI viewer sharing (Pro-license "
             "equivalent for viewers).",
             "F64 is the practical minimum for enterprise Power BI on "
             "Fabric; below that, viewers still need Pro."],
            ["Workloads",
             "Power BI, Data Engineering (Spark/Lakehouse), Data Warehouse, "
             "Data Factory, Real-Time Intelligence, Data Science, Databases.",
             "All workloads share the same CU pool \u2014 a runaway Spark job "
             "can throttle Power BI reports on the same capacity."],
            ["Workspace experience",
             "One workspace can hold Power BI items and Fabric items "
             "(Lakehouse, Warehouse, notebooks, pipelines) side by side, all "
             "backed by the same F capacity.",
             "This is why 'Fabric-enabled workspace' isn't a separate thing "
             "\u2014 it's an assignment to an F capacity."],
            ["Pause / resume",
             "F SKUs can be paused (billing stops) and resumed on demand \u2014 "
             "unlike reserved P SKUs, which are always-on.",
             "Enables dev/test cost control; production capacities typically "
             "stay on 24\u00d77."],
            ["Gov validation",
             "Fabric is 'Commercial-focused / Verify for Gov' \u2014 not all "
             "workloads have Government parity yet.",
             "In a Gov engagement, confirm each Fabric workload individually "
             "before committing to it in the architecture."],
        ],
        note="Lab connection: Exercise 1's licensing/capacity comparison "
             "includes Fabric F SKUs as one of the options to evaluate.",
        script=[
            "Frame Fabric as 'one meter, many workloads' \u2014 that is the "
            "biggest mental shift from Premium. In Premium, Power BI had its "
            "own capacity. In Fabric, Power BI is one workload among many, "
            "all drawing from the same CU pool.",
            "Explain F64 specifically: it's the magic number where viewer "
            "sharing goes free for Power BI. Below F64, viewers still need "
            "Pro licenses. Above F64, sharing works like it did on P1+. "
            "This is often the actual driver of SKU selection.",
            "Warn about the cross-workload interaction: because CU is a "
            "single pool, a Spark job that misbehaves can throttle Power BI "
            "reports on the same capacity. Give the classroom takeaway: "
            "isolate production Power BI on its own F capacity when the "
            "workload mix is unpredictable.",
            "Close with the Gov reality: Fabric parity in Government clouds "
            "is a moving target, and 'available' isn't the same as 'all "
            "workloads available'. Validate each Fabric workload the design "
            "depends on before committing.",
        ]
    )
    page += 1

    # 9. Topic 7 — Direct Lake (content)
    content_slide(
        prs, 7, "Direct Lake", page=page,
        lead_items=[
            "Direct Lake is a Fabric-only semantic model storage mode that "
            "reads Delta/Parquet files directly from OneLake \u2014 no import, "
            "no scheduled refresh, no DirectQuery round-trip to a SQL engine.",
            "Compared to Import: no refresh step, always current with the "
            "Lakehouse. Compared to DirectQuery: no SQL translation, near-"
            "Import query speed. Fallback to DirectQuery kicks in when a "
            "query hits an unsupported feature or memory limit.",
            "Modeling considerations: relationships, measures, and "
            "calculation groups work as usual, but calculated columns and "
            "calculated tables are not supported \u2014 that logic must move "
            "upstream into the Lakehouse/Warehouse.",
        ],
        why_items=[
            "Removes the 'model refresh window' problem entirely for very "
            "large datasets \u2014 the Lakehouse update IS the refresh.",
            "Removes the DirectQuery latency problem for large fact tables "
            "\u2014 queries run against columnar Parquet at Import-like speed.",
            "Requires a shift in where transformations live: calculated "
            "columns move upstream, which is actually good hygiene for "
            "enterprise models but a change for existing Power BI authors.",
        ],
        footer="Lab connection: Optional commercial lab compares Import, "
               "DirectQuery, and Direct Lake; the alternate path is a "
               "concept/comparison discussion.",
        script=[
            "Frame Direct Lake as 'the third storage mode' and put it "
            "side-by-side with Import and DirectQuery. Import = load into "
            "memory, requires refresh. DirectQuery = translate every visual "
            "to SQL, no refresh but slow. Direct Lake = read columnar "
            "Parquet from OneLake on demand, no refresh AND fast.",
            "Explain the 'fallback to DirectQuery' behavior: Direct Lake "
            "will silently fall back for queries it can't handle "
            "(unsupported functions, memory pressure). Students need to "
            "know this exists so they can diagnose 'why is one visual "
            "suddenly slow' \u2014 it may be falling back.",
            "Cover the modeling delta clearly: no calculated columns, no "
            "calculated tables. Those must move into the Lakehouse or "
            "Warehouse upstream. This is actually aligned with enterprise "
            "best practice \u2014 transformation belongs in the data platform, "
            "not the semantic model \u2014 but it's a change for Import-model "
            "authors.",
            "Close on the Gov note: Direct Lake is 'Commercial-focused / "
            "Verify for Gov'. In a Gov engagement, assume it's unavailable "
            "and use the fallback \u2014 Import from a Warehouse with a "
            "scheduled refresh \u2014 unless the tenant has confirmed Direct "
            "Lake support.",
        ]
    )
    page += 1

    # 10. Topic 8 — OneLake, Lakehouse, and Warehouse (table)
    table_slide(
        prs, 8, "OneLake, Lakehouse, and Warehouse", page=page,
        headers=["Item", "What it is", "When to use it", "Power BI integration"],
        col_widths=[2.0, 4.0, 3.6, 2.3],
        rows=[
            ["OneLake",
             "The tenant-wide Delta/Parquet data lake underneath every "
             "Fabric workspace \u2014 'OneDrive for data'. One copy, many "
             "workloads.",
             "Always \u2014 it's the storage substrate; you don't choose it, "
             "you use it.",
             "Source for Direct Lake and Import."],
            ["Lakehouse",
             "OneLake area with Delta tables + Files, exposed via a Spark "
             "SQL endpoint and an auto-generated SQL analytics endpoint.",
             "Big-data / data-engineering workloads; Spark notebooks; "
             "schema-on-read files.",
             "Direct Lake semantic model auto-created; also SQL endpoint."],
            ["Warehouse",
             "T-SQL data warehouse on OneLake, with full DDL/DML and "
             "cross-database queries. Writes go through T-SQL.",
             "Classic BI ELT / T-SQL workloads; teams that want stored "
             "procs and full write support.",
             "Direct Lake or Import via SQL endpoint."],
            ["SQL analytics endpoint",
             "Read-only T-SQL endpoint auto-provisioned on top of a "
             "Lakehouse.",
             "Serving Lakehouse data to T-SQL tools without moving it.",
             "Standard SQL connection from Desktop / paginated."],
            ["Shortcut",
             "OneLake pointer to data in another workspace, ADLS Gen2, S3, "
             "or GCS \u2014 read without copying.",
             "Federating data across domains without duplication.",
             "Consumable by any Fabric workload, including Power BI."],
        ],
        note="Lab connection: Optional commercial lab has learners identify "
             "how data is stored and exposed; alternate path documents "
             "architecture and governance considerations.",
        script=[
            "Anchor the conversation on OneLake first \u2014 it is the "
            "substrate. Every Lakehouse, every Warehouse, every KQL "
            "database, every dataflow output lands in OneLake as "
            "Delta/Parquet. That's why 'one copy of data' is the Fabric "
            "pitch.",
            "Distinguish Lakehouse vs Warehouse explicitly: Lakehouse is "
            "Spark-first, files + tables, schema-on-read friendly. "
            "Warehouse is T-SQL-first, tables only, full DDL/DML including "
            "writes. Same underlying OneLake storage, different serving "
            "engines and different authoring audiences.",
            "Introduce shortcuts as the 'no-copy federation' story \u2014 you "
            "can shortcut into ADLS Gen2, S3, or GCS, and it looks like a "
            "local Fabric table to every workload. This is often the fastest "
            "path to Fabric adoption for orgs with existing lakes.",
            "Close on Gov: everything on this slide is 'Commercial-focused "
            "/ Verify for Gov'. The fallback architecture in a Gov tenant "
            "is typically ADLS Gen2 + Azure Synapse or Azure SQL, with "
            "Import-mode semantic models on top \u2014 slower to iterate but "
            "widely available.",
        ]
    )
    page += 1

    # 11. Topic 9 — Semantic Link (content)
    content_slide(
        prs, 9, "Semantic Link", page=page,
        lead_items=[
            "Semantic Link is the Fabric notebook library (SemPy) that lets "
            "Python code read Power BI semantic model metadata, evaluate "
            "DAX/measures, and pull data as pandas DataFrames \u2014 without "
            "leaving the notebook.",
            "Primary API surface: fabric.list_datasets(), "
            "fabric.read_table(), fabric.evaluate_measure(), "
            "fabric.evaluate_dax() \u2014 plus data-quality helpers for "
            "relationships and measure dependencies.",
            "Enables data-science and MLOps workflows that treat the "
            "semantic model as a governed source of truth, not just a "
            "reporting artifact.",
        ],
        why_items=[
            "Closes the loop between BI (semantic model) and data science "
            "(notebooks) using the SAME certified measures \u2014 no "
            "reimplementation in Python.",
            "Makes semantic-model quality visible to data scientists: "
            "missing relationships, ambiguous joins, and slow measures "
            "surface in the notebook they use daily.",
            "Availability caveat: Semantic Link is 'Commercial-focused / "
            "Verify for Gov' \u2014 in Gov tenants, the fallback is exporting "
            "measure results via XMLA/DAX Studio into the notebook "
            "environment.",
        ],
        footer="Lab connection: Optional commercial lab reviews Semantic "
               "Link use cases and connects to semantic model metadata "
               "where supported; otherwise conceptual only.",
        script=[
            "Frame Semantic Link as the 'data science bridge' \u2014 for years, "
            "data scientists rebuilt Power BI measures in pandas because "
            "they couldn't reach the model. Semantic Link makes that "
            "unnecessary: the notebook calls the certified measure "
            "directly.",
            "Give a concrete example: a data scientist building a "
            "forecasting model wants 'net sales excluding returns' from the "
            "certified corporate semantic model. Previously they'd "
            "reimplement the DAX in Python and risk drift. With Semantic "
            "Link, they call fabric.evaluate_measure('Net Sales ex "
            "Returns'). Same number as the report.",
            "Mention the data-quality helpers briefly: SemPy can flag "
            "missing relationships, ambiguous cross-filters, and unused "
            "columns. That is a governance win \u2014 model authors get "
            "feedback from the notebook side of the house.",
            "Close on availability: Fabric-only, Commercial-focused. In a "
            "Gov engagement, the fallback is documented DAX queries "
            "executed through the XMLA endpoint, or exports from DAX "
            "Studio, into the notebook environment. Not as elegant, but "
            "achieves the same governance outcome.",
        ]
    )
    page += 1

    # 12. Topic 10 — Capacity metrics and throttling (custom diagram)
    throttling_states_slide(
        prs, 10, "Capacity Metrics and Throttling", page=page,
        note="Lab connection: Exercise 5 opens the capacity metrics app "
             "(where available) to identify interactive vs. background "
             "pressure and throttling signals; alternate path maps "
             "symptoms to operational actions.",
        script=[
            "Introduce the Microsoft Fabric Capacity Metrics App \u2014 the "
            "canonical tool for seeing what's actually consuming CU on a "
            "capacity. Every learner responsible for a Premium or Fabric "
            "capacity should install it on day one.",
            "Walk the states left-to-right on the diagram. Healthy is "
            "steady-state. Overage is normal short bursts \u2014 the "
            "smoothing window absorbs them against the future budget. "
            "Interactive delay is the first warning sign users can feel: "
            "reports render \u224820 seconds slower.",
            "Escalate to interactive rejection and background rejection. "
            "The key point: interactive workloads are throttled BEFORE "
            "background jobs, so users see slowness before scheduled "
            "refreshes start failing. If refreshes are also failing, the "
            "capacity has been over-budget for hours, not minutes.",
            "Give the operational playbook: identify the top-consuming "
            "artifact (usually a semantic model with expensive DAX or a "
            "runaway Spark job), fix the artifact OR scale up the "
            "capacity OR add autoscale (P SKUs) OR isolate the offender "
            "to its own capacity. All 'Verify for Gov' because the metrics "
            "app itself may not be available in every sovereign cloud.",
        ]
    )
    page += 1

    # 13. Topic 11 — Azure Government considerations (content)
    content_slide(
        prs, 11, "Azure Government Considerations", page=page,
        lead_items=[
            "Validate the four gates for every Premium/Fabric feature: "
            "(1) cloud parity, (2) licensing/SKU availability, (3) tenant "
            "admin settings, (4) capacity/workspace configuration.",
            "'Commercial-focused / Verify for Gov' items in this module: "
            "Fabric capacity, Direct Lake, OneLake/Lakehouse/Warehouse, "
            "Semantic Link, autoscale. Assume unavailable until confirmed.",
            "'Verify for Gov' items: XMLA endpoint, paginated reports, "
            "large semantic models, capacity metrics app. Typically "
            "available but must be confirmed in the target tenant.",
            "Every architecture recommendation must ship with a Gov-safe "
            "fallback \u2014 usually Pro or PPU + Import from Azure "
            "SQL/Synapse/ADLS Gen2, with scheduled refresh and paginated "
            "SSRS if pixel-perfect output is required.",
        ],
        why_items=[
            "Gov customers can't consume 'coming soon' \u2014 the design must "
            "work with the features available today in their cloud.",
            "Named fallbacks avoid awkward mid-engagement redesigns when a "
            "feature turns out not to be available.",
            "The four-gate validation is repeatable across engagements and "
            "documents the risk explicitly, which is what compliance "
            "reviewers expect to see.",
        ],
        footer="Lab connection: Exercise 1 explicitly asks learners to "
               "document required Gov validations for the chosen option.",
        script=[
            "Open by naming the discipline: 'Verify for Gov' is not a "
            "phrase we say to hedge \u2014 it is a real validation step with "
            "four specific gates. Cloud parity, license/SKU availability, "
            "tenant settings, capacity/workspace settings. Miss any one, "
            "and the design fails at deployment time.",
            "Contrast the two categories: 'Commercial-focused / Verify for "
            "Gov' items are the Fabric-era ones \u2014 assume unavailable "
            "until proven otherwise. 'Verify for Gov' items are the "
            "Premium-era ones \u2014 usually available, but confirm.",
            "Push the fallback discipline: for every recommendation on the "
            "primary path, the write-up must include a Gov-safe fallback. "
            "That's usually Pro/PPU + Import from Azure "
            "SQL/Synapse/ADLS Gen2. Pixel-perfect requirements fall back "
            "to paginated on PPU or an SSRS-on-VM path.",
            "Close with the meta-point: this discipline applies far beyond "
            "Gov. Regulated commercial customers, sovereign clouds outside "
            "the US, and customers with strict tenant policies all "
            "benefit from the same 'validate + fallback' pattern.",
        ]
    )
    page += 1

    # 14. Topic 12 — Architecture decision review (custom diagram: workspace assignment)
    workspace_assignment_slide(
        prs, 12, "Architecture Decision Review", page=page,
        note="Lab connection: Exercise 1's deliverable is exactly this \u2014 "
             "a recorded recommendation with option, requirements met, "
             "risks, and Gov validation plan.",
        script=[
            "This slide is the closing framework: every architecture "
            "decision in this space ultimately assigns a workspace to a "
            "license mode. The mode gates every downstream feature, so the "
            "decision is worth reviewing formally.",
            "Walk left to right. The workspace is the unit of assignment "
            "\u2014 it holds all the artifacts (semantic models, reports, "
            "dataflows, Lakehouses if Fabric). The admin sets ONE license "
            "mode per workspace, and that choice determines which "
            "features light up.",
            "Now walk the four target modes. Pro is the safe default but "
            "caps at 1 GB models and 8 refreshes. PPU adds Premium features "
            "at per-user pricing. Premium capacity (P/EM) adds free-viewer "
            "sharing and reserved compute. Fabric (F) adds the data-"
            "engineering surface and unlocks free viewer sharing at F64+.",
            "Close by turning it into a review checklist: requirements "
            "captured against the six axes, options evaluated against "
            "features + limits + Gov status, risks named (feature "
            "unavailability, capacity throttling, cost overrun), and a "
            "validation plan documented. That is the deliverable of "
            "Exercise 1, and it's the pattern learners should apply to "
            "every real engagement.",
        ]
    )
    page += 1

    # 15. Module Lab Walkthrough (checklist)
    checklist_slide(
        prs, "Module Lab Walkthrough",
        kicker="Topic 13 \u2014 What you'll build",
        page=page,
        items=[
            "Exercise 1: Licensing and capacity comparison \u2014 map workload "
            "requirements to Pro/PPU/Premium/Fabric options.",
            "Exercise 2: XMLA endpoint (where available) \u2014 confirm setting, "
            "connect with an approved tool, inspect metadata.",
            "Exercise 3: Paginated report (where available) \u2014 build a "
            "simple report with a table and parameter in Report Builder.",
            "Exercise 4: Large semantic model settings (where available) \u2014 "
            "review workspace + model settings and document constraints.",
            "Optional commercial: Direct Lake, OneLake/Lakehouse/Warehouse, "
            "Semantic Link \u2014 inspect items where the tenant supports them.",
            "Exercise 5: Capacity metrics and throttling concepts \u2014 "
            "identify interactive vs. background pressure, symptoms, and "
            "operational follow-up.",
            "Alternate path for every exercise: conceptual walk-through + "
            "written Gov-safe fallback for classes on restricted tenants.",
            "Deliverable: a documented recommendation per scenario, "
            "including required Azure Government validations.",
        ],
        script=[
            "Bridge from lecture to hands-on. Emphasize that this lab is "
            "unusual in the workshop: most exercises have an alternate "
            "conceptual path because the features under discussion depend "
            "on licensing and tenant settings the classroom may not have.",
            "Walk the exercises in order. Exercise 1 is the anchor \u2014 "
            "every student can do it because it's paper design. Exercises "
            "2-5 branch: hands-on if the tenant supports the feature, "
            "conceptual write-up if not.",
            "For the optional commercial exercises (Direct Lake, "
            "OneLake/Lakehouse/Warehouse, Semantic Link), tell students "
            "up front whether the classroom Fabric tenant supports them, "
            "so they know which path to follow before they start.",
            "The deliverable is the same regardless of path: a documented "
            "recommendation with option, requirements, risks, validation "
            "steps, and Gov-safe fallback. That's the artifact they can "
            "reuse on real engagements.",
        ]
    )
    page += 1

    # 16. Knowledge Check & Discussion (checklist)
    checklist_slide(
        prs, "Knowledge Check & Discussion",
        kicker="Topic 14 \u2014 Wrap-up",
        page=page,
        items=[
            "Capacity comparison completed with a defensible recommendation.",
            "XMLA endpoint marked Verify for Gov with a documented fallback.",
            "Paginated reports marked Verify for Gov with a documented "
            "fallback.",
            "Large semantic models marked Verify for Gov with capacity + "
            "tenant validation steps.",
            "Direct Lake, OneLake/Lakehouse/Warehouse, Semantic Link, and "
            "autoscale marked Commercial-focused / Verify for Gov.",
            "Capacity metrics marked Verify for Gov with a symptom \u2192 "
            "action mapping.",
            "Gov-safe alternate architecture documented end-to-end.",
            "Discussion: which of today's features would you accept as "
            "hard requirements, and which would you always ship with a "
            "fallback?",
        ],
        script=[
            "Use this as a working checklist, not a quiz. Walk each item "
            "and ask a student to describe how their own scenario answers "
            "it. If they can't, the corresponding topic wasn't landed \u2014 "
            "revisit it before moving on.",
            "Push the discussion question at the bottom: for a real "
            "customer, which of today's features would you treat as hard "
            "requirements (design fails without them), and which would you "
            "always ship with a fallback? The answers differ by industry "
            "and by cloud.",
            "Tie the checklist back to the module's validation checklist "
            "verbatim \u2014 this is the same list they'll fill out for the "
            "lab deliverable, so time spent here is time saved in the "
            "lab.",
            "Close by connecting forward to Module 11: automation and "
            "DevOps assumes the capacity/licensing decisions from today "
            "are settled, because CI/CD pipelines target specific "
            "workspaces on specific capacities. Everything they decide "
            "today shapes what Module 11 automates.",
        ]
    )
    page += 1

    # 17. Closing
    closing_slide(
        prs, MODULE_NO,
        "Module 11: Automation and DevOps \u2014 building CI/CD pipelines "
        "that target the workspaces and capacities you just designed.",
        page=page,
        subtitle="Learners can now compare Pro, PPU, Premium, and Fabric "
                 "capacity options, interpret capacity metrics and "
                 "throttling, and design Gov-safe architectures with "
                 "explicit fallbacks for every Verify-for-Gov feature.",
        script=[
            "Congratulate the class on completing what is arguably the "
            "densest architectural module in the workshop. The volume of "
            "licensing, capacity, and Fabric surface area is real \u2014 no "
            "one memorizes it all in one sitting, and they aren't "
            "expected to.",
            "Reinforce the two takeaways: the six-axis requirements "
            "framework, and the four-gate Gov validation with named "
            "fallbacks. Those two habits carry across every architecture "
            "conversation, even as the specific SKUs and features "
            "evolve.",
            "Preview Module 11: Automation and DevOps builds directly on "
            "today's decisions \u2014 CI/CD pipelines target specific "
            "workspaces on specific capacities, and XMLA read/write is "
            "the enabler for the semantic-model-as-code pattern they'll "
            "practice next.",
            "Take final questions, especially on the Commercial-focused / "
            "Verify for Gov items \u2014 those are the most likely source of "
            "lingering confusion, and getting them right up front pays "
            "off through the rest of the workshop.",
        ]
    )

    OUT_PATH.parent.mkdir(parents=True, exist_ok=True)
    prs.save(str(OUT_PATH))
    print(f"Saved: {OUT_PATH} "
          f"({len(prs.slides.__iter__.__self__._sldIdLst)} slides)")


if __name__ == "__main__":
    build()
