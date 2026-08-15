#!/usr/bin/env python3
"""
Builds the Lab 05 (Performance Optimization) instructor deck.
Run from repo root: python tools/pptx-labs/build_lab05.py
Output: modules/05-performance-optimization/assets/performance-optimization.pptx
"""
import sys
from pathlib import Path

sys.path.insert(0, str(Path(__file__).resolve().parent))
from slide_kit import (
    new_presentation, title_slide, agenda_slide,
    content_slide, table_slide, checklist_slide, closing_slide,
    blank_slide, add_rect, add_text, add_page_number,
    NAVY, NAVY_DARK, ICE, WHITE, INK, SLATE, GOLD, LIGHT_BG, CARD_BORDER,
    HEADER_FONT, BODY_FONT, SLIDE_W, SLIDE_H, set_notes,
)
from pptx.util import Inches, Pt, Emu
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR
from pptx.enum.shapes import MSO_SHAPE, MSO_CONNECTOR

OUT_PATH = Path(__file__).resolve().parent.parent.parent / \
    "modules" / "05-performance-optimization" / "assets" / "performance-optimization.pptx"

MODULE_NO = 5
TITLE = "Performance Optimization"
SUBTITLE = ("A repeatable measure-diagnose-tune process across the model, DAX, Power Query, "
            "visuals, refresh, and capacity layers")

AGENDA_TOPICS = [
    "Performance optimization mindset",
    "Performance layers",
    "Performance Analyzer",
    "Model size and cardinality",
    "DAX optimization",
    "Visual optimization",
    "Power Query and refresh optimization",
    "Aggregations",
    "DirectQuery and hybrid patterns",
    "Service and capacity monitoring",
    "External tools",
    "Lab review and benchmark targets",
    "Module lab walkthrough",
    "Knowledge check and discussion",
]


def layers_pipeline_slide(prs, number, title, page, note=None, script=None):
    """Custom horizontal pipeline diagram: 7 performance layers, source -> capacity."""
    s = blank_slide(prs)
    add_rect(s, 0, 0, SLIDE_W, SLIDE_H, WHITE)
    add_rect(s, 0, 0, SLIDE_W, Inches(1.15), NAVY)
    add_text(s, Inches(0.7), Inches(0.22), Inches(1.6), Inches(0.4), f"TOPIC {number:02d}",
             size=13, color=GOLD, bold=True, font=BODY_FONT)
    add_text(s, Inches(0.7), Inches(0.52), Inches(11.5), Inches(0.55), title, size=24,
             color=WHITE, bold=True, font=HEADER_FONT)

    layers = [
        ("Source system", "SQL / API / lake\nquery cost, indexes"),
        ("Power Query", "folding, staging,\nfilter/reduce early"),
        ("Semantic model", "cardinality,\ntypes, relations"),
        ("DAX", "measures, VAR,\nfilter scope"),
        ("Visual rendering", "visual count,\ncross-highlight"),
        ("Service refresh", "incremental,\nrefresh windows"),
        ("Capacity", "SKU, memory,\nCU throttling"),
    ]

    n = len(layers)
    top = Inches(2.15)
    box_h = Inches(2.0)
    total_w = Inches(12.2)
    gap = Inches(0.12)
    box_w = Emu(int((total_w - gap * (n - 1)) / n))
    left = Inches(0.55)

    for i, (name, detail) in enumerate(layers):
        x = left + i * (box_w + gap)
        # header band
        add_rect(s, x, top, box_w, Inches(0.55), NAVY)
        add_text(s, x, top, box_w, Inches(0.55), name, size=12, color=WHITE, bold=True,
                 font=BODY_FONT, align=PP_ALIGN.CENTER, anchor=MSO_ANCHOR.MIDDLE)
        # body
        add_rect(s, x, top + Inches(0.55), box_w, box_h - Inches(0.55), ICE,
                 line_color=CARD_BORDER)
        add_text(s, x + Inches(0.05), top + Inches(0.6), box_w - Inches(0.1),
                 box_h - Inches(0.6), detail, size=10.5, color=NAVY_DARK,
                 font=BODY_FONT, align=PP_ALIGN.CENTER, anchor=MSO_ANCHOR.MIDDLE,
                 line_spacing=1.1)
        # arrow between boxes
        if i < n - 1:
            ax1 = x + box_w
            ax2 = ax1 + gap
            ay = top + box_h / 2
            line = s.shapes.add_connector(
                MSO_CONNECTOR.STRAIGHT,
                Emu(int(ax1)), Emu(int(ay)), Emu(int(ax2)), Emu(int(ay)))
            line.line.color.rgb = SLATE
            line.line.width = Pt(1.25)

    # bottom band: measure-first arrow
    band_y = Inches(4.55)
    add_rect(s, Inches(0.55), band_y, Inches(12.2), Inches(0.55), LIGHT_BG,
             line_color=CARD_BORDER)
    add_text(s, Inches(0.55), band_y, Inches(12.2), Inches(0.55),
             "\u2190  Measure first: the bottleneck is usually one layer, not everywhere at once  \u2192",
             size=13, color=NAVY_DARK, bold=True, font=BODY_FONT, align=PP_ALIGN.CENTER,
             anchor=MSO_ANCHOR.MIDDLE)

    # takeaway strip
    add_rect(s, Inches(0.55), Inches(5.35), Inches(12.2), Inches(1.25), WHITE,
             line_color=CARD_BORDER)
    add_text(s, Inches(0.8), Inches(5.45), Inches(11.8), Inches(1.1),
             "Performance Analyzer separates DAX query time (semantic model + DAX layers) from "
             "visual display time (rendering layer) from other/render time. That split tells you "
             "which layer to open next \u2014 don't rewrite DAX if the real cost is a 40-visual page "
             "or a non-folding Power Query step.",
             size=13, color=INK, font=BODY_FONT, line_spacing=1.2, anchor=MSO_ANCHOR.TOP)

    if note:
        add_text(s, Inches(0.7), Inches(6.8), Inches(11.9), Inches(0.6), note, size=12.5,
                 color=SLATE, italic=True, font=BODY_FONT, line_spacing=1.1)
    add_page_number(s, page)
    set_notes(s, script)
    return s


def build():
    prs = new_presentation()
    page = 1

    # 1. Title
    title_slide(
        prs, MODULE_NO, TITLE, SUBTITLE,
        script=[
            "Welcome learners to Module 5. Frame this as the module where the skills from Modules "
            "1-4 get stress-tested: the star schema, DAX measures, Power Query steps, and report "
            "pages they've built now have to stay fast at real data volume and real user load.",
            "Set the core mindset up front: performance work is a diagnostic discipline, not a "
            "guessing game. Measure first, identify the actual bottleneck, change one thing, "
            "measure again. The whole module reinforces that loop \u2014 Performance Analyzer, "
            "VertiPaq Analyzer, DAX Studio, and Service capacity metrics are all just tools for "
            "steps 1 and 4 of that loop.",
            "Preview the lab: students will run Performance Analyzer against their existing report, "
            "identify a slow visual, then apply the appropriate fix \u2014 model cardinality, DAX "
            "rewrite with VAR, visual reduction, or aggregation/incremental refresh design.",
            "Flag the Azure Government angle now: Performance Analyzer and core modeling/DAX/visual "
            "practices are Gov-ready. DAX Studio, VertiPaq Analyzer, incremental refresh in the "
            "Service, and capacity metrics are all Verify-for-Gov \u2014 they depend on workstation "
            "policy, license, tenant, and cloud, so we teach them as conditional required work.",
        ]
    )
    page += 1

    # 2. Agenda
    agenda_slide(
        prs, MODULE_NO, AGENDA_TOPICS, page=page,
        script=[
            "Walk the class quickly through the fourteen topics \u2014 don't over-explain any one, "
            "just orient them to the shape: mindset and the layered mental model up front (1-2), "
            "the primary diagnostic tool (3), then the four main optimization surfaces \u2014 model, "
            "DAX, visuals, Power Query/refresh (4-7).",
            "Topics 8-10 are the architectural and scale conversations \u2014 aggregations, "
            "DirectQuery/hybrid, and Service capacity monitoring. Call out that these are the "
            "topics most likely to be Verify-for-Gov and most likely to require documentation "
            "rather than hands-on work in this classroom.",
            "Topic 11 covers external tools (DAX Studio, VertiPaq Analyzer, Tabular Editor). "
            "Topic 12 is the benchmark discipline. Topics 13-14 are the hands-on lab walkthrough "
            "and knowledge-check discussion.",
        ]
    )
    page += 1

    # 3. Topic 1 - Performance optimization mindset
    content_slide(
        prs, 1, "Performance Optimization Mindset", page=page,
        lead_items=[
            "Follow a four-step loop every time: measure first, identify the bottleneck, change "
            "one thing, then measure again with the same interaction.",
            "Capture a baseline timing before touching anything \u2014 without a baseline, you can't "
            "prove any fix actually improved the report.",
            "Change exactly one variable per iteration: one measure rewrite, one column removal, "
            "one visual reduction. Bundling changes hides which one helped.",
            "Every change carries a tradeoff \u2014 less detail, less interactivity, or a "
            "tenant-dependent feature \u2014 and the tradeoff must be documented alongside the win.",
        ],
        why_items=[
            "Guessing at fixes wastes hours and often makes reports slower (e.g., adding a "
            "CALCULATE around already-optimized logic, or splitting a fast visual into three).",
            "A documented baseline is what turns 'I think it feels faster' into 'DAX query time "
            "dropped from 4.2s to 0.6s on the same slicer interaction'.",
            "Changing one thing at a time is how you build a repeatable playbook the rest of the "
            "team can follow \u2014 not a hero-fix nobody can reproduce.",
            "Documenting tradeoffs is what makes a fix acceptable to governance: reviewers need to "
            "see that lower detail or a Verify-for-Gov feature was a conscious decision.",
        ],
        footer="Lab tie-in: Exercise 1 forces the measure-first discipline \u2014 you can't propose a "
               "fix until you've captured Performance Analyzer numbers for the slow visual.",
        script=[
            "Open by naming the anti-pattern most students have lived through: a report is 'slow', "
            "someone rewrites five measures, adds an aggregation table, and swaps three visuals \u2014 "
            "and nobody can say which change (if any) actually helped. This module is the antidote.",
            "Walk through the four-step loop explicitly and slowly: measure, identify, change one "
            "thing, measure again. Emphasize that 'measure' means Performance Analyzer numbers \u2014 "
            "DAX query time, visual display time, other time \u2014 not gut feel. A baseline you can "
            "quote in seconds is the price of admission for any optimization work.",
            "Give a concrete example: suppose the baseline is DAX query time 3.8s on a Sales by "
            "Territory slicer click. You suspect a measure. You rewrite it with VAR, re-record the "
            "same slicer click, and see 0.5s. That's a defensible win. If instead you'd also "
            "removed two visuals and changed a relationship in the same pass, you'd have no idea "
            "which change did the work.",
            "Close on tradeoff documentation: every real optimization costs something \u2014 lower "
            "granularity in an aggregation table, disabled cross-highlight, or a Verify-for-Gov "
            "feature like incremental refresh. Tell students the lab expects them to write the "
            "tradeoff next to every before/after number they capture. Then bridge into Topic 2: "
            "before you can identify the bottleneck, you need a mental map of the layers where "
            "bottlenecks live.",
        ]
    )
    page += 1

    # 4. Topic 2 - Performance layers (custom pipeline diagram)
    layers_pipeline_slide(
        prs, 2, "Performance Layers", page=page,
        note="Lab tie-in: Performance Analyzer's DAX query time vs. visual display time vs. other "
             "time maps directly onto these layers \u2014 use the split to decide which layer to open.",
        script=[
            "Introduce this as the mental map of where performance work actually happens. There "
            "are seven layers a query passes through, and any given slow report has a bottleneck "
            "in one or two of them \u2014 rarely all seven. Naming the layers explicitly is how you "
            "avoid rewriting DAX when the real problem is Power Query, or tuning the model when "
            "the real problem is a 40-visual page.",
            "Walk the pipeline left to right. Source system is your SQL/API/lake \u2014 slow there "
            "and everything downstream inherits it. Power Query is where folding, staging, and "
            "filter/column reduction happen (Topic 7 goes deep here). Semantic model is cardinality, "
            "data types, and relationship design (Topic 4). DAX is measure logic \u2014 VAR, filter "
            "scope, iterator caution (Topic 5). Visual rendering is visual count, cross-highlight, "
            "high-cardinality visuals (Topic 6). Service refresh is incremental refresh and "
            "refresh windows. Capacity is SKU memory and CU throttling (Topic 10).",
            "Use an analogy: it's like a highway with seven tollbooths. If one tollbooth has a "
            "long line, the whole trip is slow \u2014 but widening the other six lanes does nothing. "
            "Performance Analyzer's job is to tell you which tollbooth has the line, and the "
            "DAX-query-time vs. visual-display-time vs. other-time split is exactly that signal.",
            "Transition: 'So how do we actually see which layer is slow? That's Performance "
            "Analyzer \u2014 the next topic \u2014 and it's the single most important tool in this whole "
            "module.'",
        ]
    )
    page += 1

    # 5. Topic 3 - Performance Analyzer
    content_slide(
        prs, 3, "Performance Analyzer", page=page,
        lead_items=[
            "Open View > Performance Analyzer, then Start recording, then Refresh visuals (or "
            "click slicers / drill) to capture a real interaction \u2014 not just page load.",
            "Expand each visual result to see the three timings: DAX query time (engine work), "
            "visual display time (rendering), and other/render time (data transfer + wait).",
            "Sort visuals by total time and pick the slowest as the first optimization candidate.",
            "Right-click a visual's DAX query and choose Copy query \u2014 paste into DAX Studio "
            "(Verify for Gov) for Server Timings and query plan analysis when available.",
            "Export the log or screenshot the panel so the before/after numbers are documented "
            "against the exact interaction that produced them.",
        ],
        why_items=[
            "Recording a real interaction (slicer click, drill, cross-filter) is the only way to "
            "measure what a user actually feels \u2014 page load alone can hide slicer-triggered "
            "recalculations.",
            "The three-timing split is the layer diagnosis from Topic 2 in numeric form: high DAX "
            "query time \u2192 model or measure work; high visual display time \u2192 rendering or visual "
            "count; high other time \u2192 data transfer or non-folding Power Query.",
            "Copying the DAX query out means the fix can be validated in isolation, not against "
            "the whole cluttered page.",
            "Documented before/after numbers on the same interaction are what convert 'I made "
            "changes' into a defensible optimization record for governance review.",
        ],
        footer="Lab tie-in: Exercise 1 requires learners to capture a Performance Analyzer "
               "baseline and identify at least one candidate visual before proposing any change.",
        script=[
            "Introduce Performance Analyzer as the built-in, Gov-ready tool that turns "
            "'this report feels slow' into concrete numbers. Every learner has this available in "
            "Desktop today \u2014 no external tool, no tenant setting, no license.",
            "Walk the exact steps: View > Performance Analyzer, Start recording, then perform the "
            "interaction you care about. Emphasize that the interaction matters \u2014 refreshing "
            "visuals measures cold load, but clicking a slicer measures the slicer-triggered "
            "recalculation, and those can be very different. Encourage students to record the "
            "actual user workflow, not just page open.",
            "Explain the three timings using the layer map from Topic 2. If DAX query time is "
            "80% of the total, the fix lives in the semantic model or a measure. If visual "
            "display time dominates, the fix is on the report canvas (Topic 6). If other/render "
            "time is high, look at data volume, Power Query folding, or connection type. Give a "
            "concrete example: a 4.5s total that splits 4.1s DAX + 0.3s visual + 0.1s other says "
            "'go rewrite the measure', not 'reduce the visual count'.",
            "Introduce the Copy query trick and set expectations honestly: DAX Studio is "
            "Verify-for-Gov \u2014 great when workstation policy allows it, but the Alternate path in "
            "the lab (simpler visuals + Performance Analyzer alone) is what students in "
            "policy-restricted environments will use. Transition into Topic 4 by pointing out "
            "that when DAX query time is the villain, the fix often starts one layer down, in the "
            "model itself.",
        ]
    )
    page += 1

    # 6. Topic 4 - Model size and cardinality (table_slide - reference data)
    table_slide(
        prs, 4, "Model Size and Cardinality", page=page,
        headers=["Model characteristic", "Why it inflates size / slows queries",
                 "Concrete fix in this lab"],
        col_widths=[3.0, 4.9, 4.0],
        rows=[
            ["High-cardinality text column",
             "VertiPaq stores each distinct value; free-text fields (Notes, GUIDs, full "
             "timestamps) compress poorly.",
             "Remove the column, hash it, or split into lower-cardinality parts."],
            ["Full DateTime column",
             "Second/millisecond precision creates millions of distinct values in one column.",
             "Split into a Date column (joins DimDate) and a separate Time column."],
            ["Excess numeric precision",
             "Decimal Number (double) is larger and less compressible than Fixed Decimal or "
             "Whole Number.",
             "Downgrade to Fixed Decimal or Whole Number where business precision allows."],
            ["Unused columns kept in the model",
             "Every loaded column costs memory and refresh time even if no visual uses it.",
             "Remove in Power Query (Remove Columns step) \u2014 not just hide in report view."],
            ["Wrong data type",
             "A number stored as Text disables aggregation compression and forces string joins.",
             "Fix the type in Power Query, not with a DAX conversion at query time."],
            ["Bidirectional or many-to-many relationships",
             "Force the engine to evaluate filter propagation in both directions on every query.",
             "Use single-direction one-to-many; document any exception."],
            ["Duplicate business logic across tables",
             "Same measure re-implemented per table forces redundant scans.",
             "Consolidate into one measure on the fact table; reuse via measure branching."],
        ],
        note="Lab tie-in: Exercise 2 asks learners to produce a model reduction plan documenting "
             "exactly these categories \u2014 unused columns, high-cardinality fields, precision, and "
             "date/time splits \u2014 before making any change.",
        script=[
            "Frame this as the highest-leverage optimization surface: model-level fixes usually "
            "give the biggest win per hour of effort, because they benefit every measure and "
            "every visual across every report that uses this model.",
            "Walk down the table row by row, tying each to what VertiPaq (the storage engine) "
            "actually does. VertiPaq is a columnar, dictionary-encoded store \u2014 compression is a "
            "function of how many distinct values a column has. A DateTime with 20 million "
            "distinct values compresses badly; the same data split into Date (a few thousand "
            "distinct values) and Time (86,400 distinct seconds max) compresses beautifully.",
            "Use the unused-columns row as the teaching moment: it is the single most common "
            "quick win in the field. Students often 'hide' columns in report view thinking that "
            "saves memory \u2014 it does not. The column is still loaded and still costs refresh "
            "time. The fix is Remove Columns in Power Query, upstream of the model.",
            "Bridge to Topic 5: once the model is lean, the next place bottlenecks hide is in "
            "measure logic \u2014 which is where DAX optimization comes in. Prompt: 'Which of these "
            "seven categories do you think is most common in the reports your team ships? Why?'",
        ]
    )
    page += 1

    # 7. Topic 5 - DAX optimization
    content_slide(
        prs, 5, "DAX Optimization", page=page,
        lead_items=[
            "Wrap repeated expressions in VAR / RETURN \u2014 the engine evaluates the VAR once, "
            "instead of every time the expression appears in the measure.",
            "Use measure branching: build small measures that other measures reference (e.g., "
            "[Total Sales] \u2192 [Sales YTD] \u2192 [Sales YTD vs LY]) instead of copying filter logic.",
            "Narrow filter scope with KEEPFILTERS or explicit column filters rather than "
            "over-broad table filters that force full scans.",
            "Be cautious with iterators (SUMX, FILTER, ADDCOLUMNS) over large tables \u2014 they run "
            "row by row and can multiply cost dramatically when nested.",
            "Copy the original measure text before editing, and test the rewrite side by side "
            "with the original to confirm the result is identical.",
        ],
        why_items=[
            "VAR eliminates duplicated engine work: a common expression evaluated three times in "
            "one measure now runs once, and the intent of the measure becomes readable.",
            "Measure branching means the filter logic lives in exactly one place. Fix a bug in "
            "[Total Sales] and every YTD, MoM, and vs-LY measure downstream is fixed too.",
            "Narrow filter scope keeps the storage engine in fast, cache-friendly paths and "
            "avoids expensive filter-context expansions.",
            "Iterator caution prevents the classic 'measure worked in dev, timed out in prod' "
            "surprise when the fact table grew from 100k to 100M rows.",
            "Side-by-side validation is what catches subtle semantic drift \u2014 an 'optimized' "
            "measure that returns a slightly different number is not an optimization, it's a bug.",
        ],
        footer="Lab tie-in: Exercise 2 (measure rewrite with VAR/RETURN) explicitly requires "
               "keeping the original alongside the rewrite until the numbers reconcile.",
        script=[
            "Position DAX optimization as the layer you attack after the model is lean \u2014 not "
            "before. A slow measure on a bloated model is often the model's fault, not the "
            "measure's, and fixing the measure first hides the real bottleneck.",
            "Walk through the four techniques in order of frequency. VAR/RETURN is by far the "
            "most common win: it deduplicates work and doubles as documentation. Show a concrete "
            "sketch on the board \u2014 an inline expression repeated three times in one measure vs. "
            "the same expression stored in VAR SalesBase and referenced by name.",
            "Measure branching is the governance win: instead of five measures each redefining "
            "'sales in current filter context', you have one [Total Sales] and every other "
            "measure references it. When the business changes what 'sales' means (e.g., excludes "
            "returns), you fix it in one place.",
            "Warn about the iterator trap with a real-sounding scenario: SUMX over FILTER over "
            "a 200-million-row fact table, nested inside CALCULATE, will crush the engine. "
            "Encourage students to ask 'is there a set-based equivalent?' before reaching for an "
            "iterator. Transition to Topic 6: even the best DAX can't save a page that's asking "
            "the engine to render forty visuals at once.",
        ]
    )
    page += 1

    # 8. Topic 6 - Visual optimization
    content_slide(
        prs, 6, "Visual Optimization", page=page,
        lead_items=[
            "Count visuals per page \u2014 every visual issues its own DAX query, and page load is "
            "gated by the slowest one on the page.",
            "Replace dense table/matrix visuals with summary visuals (cards, bars) where "
            "learners don't actually need row-level detail on the landing page.",
            "Be cautious with high-cardinality visuals (tables with 50k rows, slicers over 10k "
            "distinct values) \u2014 they force large result sets across the wire.",
            "Disable cross-highlight interactions between visuals that don't need to filter each "
            "other (Format > Edit interactions > None).",
            "Rethink page complexity: a 40-visual 'dashboard of dashboards' is almost always "
            "slower and less useful than three focused pages linked by drillthrough.",
        ],
        why_items=[
            "Fewer visuals means fewer parallel DAX queries competing for the engine, and a "
            "faster feel for the same underlying data.",
            "Summary visuals also cost less DAX \u2014 SUM over a whole column is cheaper than "
            "returning 5,000 detail rows the user will never scroll through.",
            "High-cardinality visuals are the most common cause of 'other/render time' spikes in "
            "Performance Analyzer \u2014 the engine is fast, but the payload is huge.",
            "Turning off unnecessary cross-highlight is a free win: no data change, no measure "
            "change, just fewer queries triggered on every click.",
            "Focused pages linked by drillthrough (Module 4) both perform better and match how "
            "users actually consume reports \u2014 headline first, detail on demand.",
        ],
        footer="Lab tie-in: Exercise 4 walks learners through counting visuals, consolidating, "
               "disabling unnecessary interactions, and re-measuring page performance.",
        script=[
            "Set expectations up front: visual optimization is often the fastest win in the whole "
            "module, because it needs no DAX rewrite and no model change \u2014 just editorial "
            "discipline about what belongs on a page.",
            "Start with the visual-count point. Power BI issues one DAX query per visual, and the "
            "page waits for the slowest. Cutting a page from 20 visuals to 8 can halve the "
            "perceived load time even if no individual visual got faster. Ask the room: 'What's "
            "the visual count on your busiest production page?' \u2014 the number is usually higher "
            "than people realize.",
            "Use the drillthrough pattern from Module 4 as the design bridge: instead of putting "
            "the detail table on the landing page, hide it behind a drillthrough or tooltip page. "
            "The DAX only runs when the user asks for it. This is where Modules 4 and 5 "
            "reinforce each other.",
            "Close on cross-highlight interactions \u2014 the free-win optimization most authors "
            "forget. Every visual on a page cross-filters every other visual by default, which "
            "means every click triggers N-1 requeries. Turn off the ones that don't add analytic "
            "value. Transition to Topic 7: even a lean model and lean page can be slow if data "
            "is refreshed the wrong way.",
        ]
    )
    page += 1

    # 9. Topic 7 - Power Query and refresh optimization
    content_slide(
        prs, 7, "Power Query and Refresh Optimization", page=page,
        lead_items=[
            "Preserve query folding: keep source-native steps (Filter Rows, Remove Columns, "
            "Group By) at the top so they push down to the source; check 'View Native Query' to "
            "confirm folding is still intact.",
            "Filter early and remove columns early \u2014 every row and column dropped upstream is "
            "one less thing the model has to load, compress, and refresh.",
            "Use staging queries (referenced, not duplicated) for shared cleanup logic so it "
            "runs once per refresh, not once per downstream query.",
            "Prepare for incremental refresh: fact query must filter on a Date/DateTime column "
            "using the RangeStart and RangeEnd DateTime parameters exactly as named.",
            "Define archive and refresh windows (e.g., archive 5 years, incrementally refresh "
            "10 days) and discuss Detect data changes when the source supports it.",
        ],
        why_items=[
            "Query folding is the difference between the source doing the work (fast, indexed) "
            "and Power Query pulling everything down and doing it locally (slow, memory-heavy).",
            "Filtering and column removal upstream compound: every downstream step, every model "
            "refresh, and every VertiPaq compression pass benefits.",
            "Staging prevents the same expensive transformation from running three times because "
            "three tables reference the same source.",
            "The RangeStart/RangeEnd naming is not optional \u2014 the Service uses those exact "
            "parameter names to inject the incremental filter, so a typo silently disables "
            "incremental refresh.",
            "Well-defined archive/refresh windows are what convert a 45-minute nightly full "
            "refresh into a 90-second incremental one \u2014 and are required for Service publishing "
            "of large models.",
        ],
        footer="Lab tie-in: Exercise 6 walks learners through validating RangeStart/RangeEnd, the "
               "fact date filter, and archive/refresh windows before enabling the policy. "
               "Incremental refresh in the Service is Verify-for-Gov.",
        script=[
            "Introduce this topic as the 'upstream' side of performance: the model can be "
            "perfect, but if Power Query is pulling 400 million rows through non-folding steps "
            "every night, refresh will still be slow and fragile.",
            "Explain folding in plain terms: when a step folds, it becomes part of the SQL sent "
            "to the source, and the source (which is indexed and tuned) does the work. When "
            "folding breaks \u2014 usually because of a custom function, a Merge in the wrong place, "
            "or an added column that references the whole table \u2014 Power Query has to pull the "
            "whole table down and do the step locally. The View Native Query indicator is how "
            "you check.",
            "Use a concrete story: a fact query with a Filter Rows step at the top folds and "
            "sends 'WHERE OrderDate >= '2024-01-01'' to SQL; the same Filter Rows moved after a "
            "custom column breaks folding and now the whole 200M-row table comes down first. Same "
            "logical result, radically different cost.",
            "For incremental refresh, be precise: the parameter names must be exactly RangeStart "
            "and RangeEnd, both DateTime, and the fact query must filter on a Date/DateTime "
            "column using them. Remind students this is Verify-for-Gov \u2014 policy setup is a "
            "conditional required task depending on license and Service availability. Transition "
            "into Topic 8: sometimes even a well-folded, incrementally-refreshed model is too big "
            "for interactive query \u2014 that's where aggregations come in.",
        ]
    )
    page += 1

    # 10. Topic 8 - Aggregations (table_slide reference)
    table_slide(
        prs, 8, "Aggregations", page=page,
        headers=["Piece", "What it is", "Design rule / gotcha"],
        col_widths=[2.6, 4.6, 4.7],
        rows=[
            ["Import aggregation table",
             "Pre-summarized copy of the fact at a coarser grain (e.g., Month + Category + "
             "Territory).",
             "Store as Import mode so summary queries hit fast in-memory data."],
            ["DirectQuery detail table",
             "The original fact table stays in DirectQuery mode for row-level detail and "
             "drill-through.",
             "Detail queries fall through to the source when the agg doesn't match."],
            ["Group-by columns",
             "The dimension columns that define the aggregation grain (Month, Category, "
             "Territory).",
             "Every group-by column in the agg must map to the corresponding detail column."],
            ["Manage aggregations dialog",
             "Where you declare each agg column's Summarization (Sum, Count, GroupBy) and its "
             "detail-column mapping.",
             "Precedence controls which agg is tried first when multiple aggs exist."],
            ["Agg hit vs. agg miss",
             "Hit: query is answered from the agg (fast). Miss: query falls through to the "
             "detail table (slow).",
             "Test with Performance Analyzer + DAX Studio (Verify for Gov) to confirm hits."],
            ["Source validation",
             "DirectQuery source support, gateway readiness, connector, and tenant settings.",
             "Aggregations are Gov-ready as a pattern, but the DirectQuery half is "
             "Verify-for-source."],
        ],
        note="Lab tie-in: Exercise 5 asks learners to design (and where available, build) an "
             "aggregation table \u2014 identifying the fact, choosing a summary grain, and mapping "
             "columns \u2014 before discussing when it will and won't hit.",
        script=[
            "Introduce aggregations as the standard pattern for making very large fact tables "
            "feel interactive: keep the detail available for drill-through, but pre-summarize the "
            "common queries so 95% of clicks hit a small, fast in-memory copy.",
            "Walk the table row by row. Emphasize the agg-hit / agg-miss mechanic: if the "
            "summary grain is Month + Category + Territory, then any query grouping by Month or "
            "Category or Territory (or a subset) hits the agg. Any query at Day grain, or by "
            "SKU, misses \u2014 and falls through to the DirectQuery detail table. So grain choice is "
            "everything: too fine and the agg is huge; too coarse and it misses most queries.",
            "Give a concrete example: an 800M-row order-line fact in DirectQuery, plus a 4M-row "
            "Month/Category/Territory Import agg. Dashboard tiles that show 'sales by month by "
            "category' answer in under a second from the agg. A drill to 'sales by day by SKU' "
            "falls through to DirectQuery and takes 8 seconds \u2014 acceptable because it's "
            "occasional.",
            "Close on the governance angle: aggregations as a pattern are Gov-ready, but the "
            "DirectQuery source, gateway, connector, and tenant settings are Verify-for-source "
            "and must be documented before production. Transition to Topic 9, which is the "
            "broader Import/DirectQuery/Dual/Hybrid conversation.",
        ]
    )
    page += 1

    # 11. Topic 9 - DirectQuery and hybrid patterns (table_slide)
    table_slide(
        prs, 9, "DirectQuery and Hybrid Patterns", page=page,
        headers=["Storage pattern", "When to use it", "Tradeoff to document"],
        col_widths=[2.8, 4.8, 4.3],
        rows=[
            ["Import (default)",
             "Fits in memory; refresh cadence matches business needs (nightly, hourly).",
             "Fastest queries, but data is only as fresh as the last refresh."],
            ["DirectQuery",
             "Data too large to import, or near-real-time freshness is required.",
             "Every visual triggers a source query; only source-supported DAX runs; slower."],
            ["Dual mode (per table)",
             "Small dimensions used with both Import fact and DirectQuery fact tables.",
             "Engine chooses per query \u2014 requires disciplined relationship design."],
            ["Hybrid table (partitions)",
             "Recent partitions in DirectQuery for freshness; historical in Import for speed.",
             "Verify-for-Gov: license, workspace, and Service capabilities gate this."],
            ["Import agg + DQ detail",
             "Standard scale pattern: summary in-memory, detail on demand (Topic 8).",
             "Agg-miss falls through to DQ; grain choice determines hit rate."],
            ["Large semantic model",
             "Model exceeds standard memory limits; requires Premium/Fabric capacity.",
             "Verify-for-Gov: SKU, tenant settings, and format (Large model storage) required."],
        ],
        note="Lab tie-in: Exercise 5 (Concept note) frames these tradeoffs explicitly \u2014 source "
             "support, relationship behavior, cache-hit rules, gateway, and Gov validation must "
             "be documented before production use.",
        script=[
            "Frame this as an architecture-decision table, not a hands-on build \u2014 in the "
            "classroom environment most students won't have a DirectQuery source or Premium "
            "capacity available, so the goal is that they can reason about which pattern fits a "
            "future scenario.",
            "Anchor on Import as the default: unless there is a specific reason to leave Import, "
            "Import wins on query speed, DAX capability, and simplicity. Every other pattern in "
            "this table is a considered exception with documented tradeoffs.",
            "Walk through the exceptions in order. DirectQuery for volume or freshness \u2014 but "
            "you pay per-visual query cost and lose some DAX functions. Dual mode for shared "
            "dimensions across Import and DirectQuery fact tables. Hybrid tables for recent-hot / "
            "historical-cold partitioning. Large semantic model for datasets that simply don't "
            "fit standard memory.",
            "Emphasize the recurring Verify-for-Gov flag: hybrid tables, large models, and much "
            "of DirectQuery gateway behavior depend on tenant, cloud, and SKU. The teaching goal "
            "is that students leave able to write the tradeoff decision down for a reviewer, not "
            "that they've built it end to end today. Transition to Topic 10: once the model is "
            "deployed, monitoring is how you know it stays healthy.",
        ]
    )
    page += 1

    # 12. Topic 10 - Service and capacity monitoring
    content_slide(
        prs, 10, "Service and Capacity Monitoring", page=page,
        lead_items=[
            "Refresh history (dataset settings): confirms refresh success, duration, and error "
            "detail per run \u2014 the first place to look when a scheduled refresh fails.",
            "Dataset settings: gateway binding, scheduled refresh cadence, incremental refresh "
            "policy status, and query caching options.",
            "Capacity metrics app (Premium/Fabric): CPU (CU) utilization, memory pressure, and "
            "throttling / autoscale events per workload.",
            "Admin monitoring: tenant-level usage of datasets, refresh volume, and largest "
            "datasets by memory \u2014 signals which models to optimize next.",
            "Gov validation: capacity metrics app, incremental refresh in Service, and some "
            "admin views are Verify-for-Gov and depend on cloud/tenant availability.",
        ],
        why_items=[
            "Refresh history is what turns 'the report is stale' from a user complaint into a "
            "specific error message and timestamp you can act on.",
            "Dataset settings surface silent misconfigurations \u2014 e.g., an incremental refresh "
            "policy that reports 'success' but is actually doing a full refresh because a "
            "parameter is wrong.",
            "Capacity CU throttling is the production explanation for 'reports slow at 9 AM but "
            "fast at 3 PM' \u2014 without capacity metrics you're guessing.",
            "Admin monitoring is how a governance team decides which of 300 workspace models is "
            "worth spending an optimization sprint on \u2014 usually the top 5 by memory + refresh "
            "cost.",
            "Explicit Gov labeling prevents accidentally promising a monitoring capability the "
            "target tenant doesn't yet expose.",
        ],
        footer="Lab tie-in: the Validation checklist requires capacity metrics and incremental "
               "refresh both to be marked Verify-for-Gov in learner documentation.",
        script=[
            "Position this topic as 'performance work after the report ships'. Everything before "
            "this was authoring-time optimization; this is operational monitoring so problems get "
            "caught before the business notices them.",
            "Walk through the four surfaces in order of accessibility. Refresh history and "
            "dataset settings are available on any workspace \u2014 Gov-ready, always the first place "
            "to look. Capacity metrics and admin monitoring require Premium/Fabric and admin role "
            "respectively, and both are Verify-for-Gov depending on tenant.",
            "Use a concrete production story: 'reports slow every morning at 9 AM' turned out to "
            "be capacity CU throttling during a large scheduled refresh window overlapping with "
            "user login peak. Without the capacity metrics app you would have chased the report "
            "authors for months when the real fix was moving the refresh schedule.",
            "Close on the governance framing: this is where 'is this dataset actually being "
            "used?' gets answered. Admin monitoring surfaces datasets with heavy refresh cost and "
            "zero read activity \u2014 candidates to retire, not optimize. Transition into Topic 11: "
            "when the built-in tools aren't enough, external tools give deeper diagnostic power \u2014 "
            "with Gov caveats.",
        ]
    )
    page += 1

    # 13. Topic 11 - External tools (table_slide)
    table_slide(
        prs, 11, "External Tools", page=page,
        headers=["Tool", "What it adds beyond Performance Analyzer", "Gov / policy note"],
        col_widths=[2.4, 5.6, 3.9],
        rows=[
            ["DAX Studio",
             "Runs copied Performance Analyzer DAX queries with Server Timings, storage engine "
             "vs. formula engine split, and query plan indicators.",
             "Verify-for-Gov: workstation policy for external tools; XMLA endpoint access."],
            ["VertiPaq Analyzer",
             "Reports column-level cardinality, table size, dictionary size, and encoding \u2014 "
             "shows exactly which columns bloat the model.",
             "Verify-for-Gov: usually loaded inside DAX Studio or Tabular Editor."],
            ["Tabular Editor",
             "Bulk metadata edits, best-practice analyzer rules, and scripted model changes "
             "beyond what Desktop UI supports.",
             "Verify-for-Gov: workstation policy; version 2 free, version 3 licensed."],
            ["Customer workstation policy",
             "Enterprise IT often restricts unsigned or external tools on managed devices \u2014 "
             "must be validated per environment.",
             "Verify before making any of the above tools required in a curriculum."],
            ["Tenant / XMLA policy",
             "Read-write XMLA endpoints, external tools setting, and workspace roles gate what "
             "these tools can actually do against the Service.",
             "Tenant admin scope \u2014 confirm before promising Service-side diagnostics."],
        ],
        note="Lab tie-in: Exercise 3 offers a DAX Studio path AND an alternate path (Performance "
             "Analyzer + simpler visuals) so learners in policy-restricted environments still "
             "complete the exercise.",
        script=[
            "Frame this deliberately: Performance Analyzer answers 'which visual is slow?'. "
            "These external tools answer 'why is it slow?' at a deeper level \u2014 storage engine "
            "vs. formula engine, exact column bloat, query plan choice. Powerful, but every one "
            "of them is Verify-for-Gov because they depend on workstation policy, tenant "
            "settings, and cloud availability.",
            "DAX Studio is the deep-diagnostics companion. Server Timings tells you whether a "
            "measure is bottlenecked in the storage engine (usually a model/cardinality problem) "
            "or the formula engine (usually a measure logic problem). That split is often the "
            "difference between rewriting a measure and reshaping the model.",
            "VertiPaq Analyzer answers the model-bloat question with numbers: 'This one Notes "
            "column is 38% of your entire model's memory footprint.' That's the level of "
            "evidence that unblocks a hard 'we don't need that column' conversation with a "
            "business owner.",
            "Emphasize the alternate path: in a policy-restricted classroom or a customer where "
            "external tools aren't approved, students still get real optimization work done with "
            "just Performance Analyzer and disciplined single-variable changes. Transition to "
            "Topic 12: however you measured, the discipline is documenting before/after benchmark "
            "numbers.",
        ]
    )
    page += 1

    # 14. Topic 12 - Lab review and benchmark targets
    content_slide(
        prs, 12, "Lab Review and Benchmark Targets", page=page,
        lead_items=[
            "Every optimization gets a before number, an after number, and the exact interaction "
            "used to produce both \u2014 captured from Performance Analyzer.",
            "Document what was changed (one thing per iteration) and what the tradeoff was: "
            "lower detail, disabled interaction, or a Verify-for-Gov feature.",
            "Production readiness signals: baseline captured, at least one visual optimization "
            "applied, model reduction plan documented, DAX rewrite uses VAR / measure branching, "
            "aggregation grain defined, incremental refresh parameters documented.",
            "Documentation expectations: reviewers see the before/after numbers, the change, the "
            "tradeoff, and the Gov status of any Verify-for-Gov feature used.",
        ],
        why_items=[
            "Undocumented optimization is indistinguishable from luck \u2014 the next author can't "
            "reproduce it and the next report can't benefit from the lesson.",
            "Explicit tradeoff notes are what makes an optimization acceptable to a governance "
            "reviewer instead of 'why did you disable cross-highlighting on that visual?'.",
            "The production-readiness list doubles as the definition of 'done' for this module \u2014 "
            "and it maps directly to the lab's validation checklist.",
            "Gov labeling on Verify-for-Gov features prevents promising something in a customer "
            "environment where the tenant doesn't yet support it.",
        ],
        footer="Lab tie-in: the lab's own Validation checklist enumerates exactly these before/"
               "after, tradeoff, and Gov-status expectations \u2014 revisit it as the module wraps.",
        script=[
            "Use this topic to consolidate the discipline before the hands-on walkthrough. "
            "Performance optimization is not just a set of techniques \u2014 it is a documentation "
            "practice. If the before/after numbers, the change, and the tradeoff aren't written "
            "down, the work doesn't compound.",
            "Walk through what a good benchmark entry looks like out loud: 'Baseline: Sales by "
            "Territory slicer click, DAX query time 4.2s, visual display 0.3s. Change: rewrote "
            "[Sales YTD] with VAR. After: DAX query time 0.6s, visual display 0.3s. Tradeoff: "
            "none.' That format \u2014 interaction, before, change, after, tradeoff \u2014 is the "
            "instructor's target.",
            "Tie the production-readiness list back to real deployments. In a governance review, "
            "the reviewer is not going to re-run Performance Analyzer for you \u2014 they will look "
            "at your documented numbers, your one-thing-per-change discipline, and your "
            "Verify-for-Gov labeling, and decide whether to approve.",
            "Transition into the lab walkthrough: 'You now have the mindset (Topic 1), the layer "
            "map (Topic 2), the tools (Topics 3, 11), the surface techniques (Topics 4-7), the "
            "scale patterns (Topics 8-9), the monitoring loop (Topic 10), and the documentation "
            "discipline (this one). Time to run the loop on the actual report.'",
        ]
    )
    page += 1

    # 15. Module lab walkthrough
    checklist_slide(
        prs, "Module Lab Walkthrough", kicker="Topic 13 \u2014 What you'll build", page=page,
        items=[
            "Exercise 1: Performance Analyzer \u2014 record baseline, capture DAX / visual / other "
            "time per visual, identify slowest candidate.",
            "Exercise 2: Model size and cardinality \u2014 flag unused columns, high-cardinality "
            "text, precision, date/time splits; produce a written reduction plan.",
            "Exercise 3: DAX Studio query timings (Verify for Gov) \u2014 Server Timings on a copied "
            "PA query. Alternate: simpler visuals + PA comparison.",
            "Exercise 4: Visual optimization \u2014 count visuals, consolidate, disable unneeded "
            "cross-highlight, re-measure page performance.",
            "Exercise 5: Aggregation table (Gov-ready / Verify for source) \u2014 pick fact, define "
            "Month/Category/Territory grain, map columns, discuss hit/miss.",
            "Exercise 6: Incremental refresh policy (Verify for Gov) \u2014 confirm RangeStart/"
            "RangeEnd DateTime params, archive/refresh windows, Service prerequisites.",
            "Deliverable: before/after Performance Analyzer numbers plus tradeoff notes for "
            "every change made.",
            "Optional: DAX Studio + VertiPaq Analyzer path for deeper column bloat and query-"
            "plan analysis where workstation policy allows.",
        ],
        script=[
            "Use this slide as the literal running order for the hands-on portion of the module. "
            "Read it top to bottom and confirm students know which exercise maps to which topic "
            "you just taught \u2014 Exercise 1 to Topic 3, Exercise 2 to Topic 4, and so on.",
            "Emphasize the sequencing: Exercise 1 must come first because every subsequent "
            "exercise needs a baseline to compare against. Skipping the baseline turns the rest "
            "of the lab into guesswork.",
            "Call out the Verify-for-Gov exercises (3, 5 partial, 6) and set expectations for "
            "the two paths: if the environment supports it, do the full exercise; if not, do the "
            "alternate path (documentation and reasoning) and mark the artifact Verify-for-Gov. "
            "Both paths count as complete.",
            "Remind them of the deliverable format from Topic 12 \u2014 interaction, before number, "
            "change, after number, tradeoff \u2014 and that this is what will be reviewed, not the "
            "raw PBIP alone. Then release them to the lab.",
        ]
    )
    page += 1

    # 16. Knowledge check and discussion
    checklist_slide(
        prs, "Knowledge Check & Discussion", kicker="Topic 14 \u2014 Wrap-up", page=page,
        items=[
            "Performance Analyzer baseline was captured with a specific interaction, not just "
            "page load.",
            "At least one visual optimization is documented with before/after numbers.",
            "Model size / cardinality recommendations are written down (unused columns, high-"
            "cardinality fields, precision, date/time splits).",
            "DAX optimization uses VAR / RETURN or measure branching \u2014 and the rewrite matches "
            "the original result.",
            "Aggregation table grain is defined and mapped to detail columns; hit/miss behavior "
            "is understood.",
            "Incremental refresh parameters (RangeStart, RangeEnd) and archive/refresh windows "
            "are documented.",
            "DAX Studio, VertiPaq Analyzer, capacity metrics, and Service-side incremental "
            "refresh are all labeled Verify-for-Gov.",
            "Before/after observations follow the interaction / before / change / after / "
            "tradeoff format and are ready for governance review.",
        ],
        script=[
            "Use this as discussion-driven wrap-up, not a quiz. Pick two or three items and ask "
            "specific students to walk the room through their own answer \u2014 especially the "
            "before/after documentation item, since that is where undisciplined teams fail "
            "governance review.",
            "For the DAX rewrite item, push on the 'matches the original result' clause. Ask: "
            "'How did you prove the rewrite is semantically identical, not just faster?' The "
            "expected answer is side-by-side comparison against the original measure on the same "
            "visual and slicer state.",
            "For the aggregation item, ask a student to describe when their aggregation table "
            "would miss and what happens then. If they can explain the fall-through to the "
            "detail source and the resulting cost, they've internalized the pattern; if they "
            "can't, review Topic 8 briefly.",
            "Close by connecting forward: the reports that leave this module are measured, "
            "documented, and defensible. Module 6 (Advanced Analytics and AI-Assisted Insights) "
            "adds AI visuals and analytical features on top \u2014 many of which are Verify-for-Gov "
            "and require the same measure-and-document discipline you just practiced today.",
        ]
    )
    page += 1

    # 17. Closing
    closing_slide(
        prs, MODULE_NO,
        "Module 06: Advanced Analytics and AI-Assisted Insights \u2014 applying analytical visuals "
        "and AI features on the fast, well-tuned models you just built.",
        page=page,
        subtitle="Learners now have a repeatable measure-diagnose-tune loop and documented "
                 "before/after evidence across the model, DAX, visual, refresh, and capacity layers.",
        script=[
            "Congratulate the class on completing the performance module \u2014 this is the module "
            "that turns 'it worked in dev' into 'it stays fast in production', and the "
            "measure-first discipline they practiced today is the single most transferable skill "
            "in the whole workshop.",
            "Remind students to keep their before/after benchmark notes with the PBIP artifact \u2014 "
            "Module 6 will add analytical and AI visuals on top of the same model, and any "
            "regression will be easier to spot because they now have baseline numbers.",
            "Take final questions, especially on the Verify-for-Gov topics \u2014 DAX Studio, "
            "VertiPaq Analyzer, incremental refresh in the Service, capacity metrics \u2014 since "
            "those are where students are most likely to hit environment-specific blockers back "
            "at their own desks.",
        ]
    )

    OUT_PATH.parent.mkdir(parents=True, exist_ok=True)
    prs.save(str(OUT_PATH))
    print(f"Saved: {OUT_PATH} ({len(prs.slides.__iter__.__self__._sldIdLst)} slides)")


if __name__ == "__main__":
    build()
