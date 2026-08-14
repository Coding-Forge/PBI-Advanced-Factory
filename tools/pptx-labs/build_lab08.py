#!/usr/bin/env python3
"""
Builds the Lab 08 (Power BI Service Enterprise Deployment) instructor deck.
Run from repo root: python tools/pptx-labs/build_lab08.py
Output: modules/08-service-enterprise-deployment/assets/service-enterprise-deployment.pptx
"""
import sys
from pathlib import Path

sys.path.insert(0, str(Path(__file__).resolve().parent))
from pptx.util import Inches, Pt, Emu
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR
from pptx.enum.shapes import MSO_SHAPE, MSO_CONNECTOR

from slide_kit import (
    new_presentation, title_slide, agenda_slide,
    content_slide, table_slide, checklist_slide, closing_slide,
    blank_slide, add_rect, add_text, add_page_number, set_notes,
    NAVY, NAVY_DARK, ICE, WHITE, INK, SLATE, GOLD, LIGHT_BG, CARD_BORDER,
    HEADER_FONT, BODY_FONT, SLIDE_W, SLIDE_H,
)

OUT_PATH = Path(__file__).resolve().parent.parent.parent / \
    "modules" / "08-service-enterprise-deployment" / "assets" / "service-enterprise-deployment.pptx"

MODULE_NO = 8
TITLE = "Power BI Service Enterprise Deployment"
SUBTITLE = ("Publishing, governing, refreshing, and distributing PBIP-authored "
            "content in the Power BI Service")

AGENDA_TOPICS = [
    "From authoring to enterprise deployment",
    "Workspace design (dev / test / prod)",
    "Workspace roles and security",
    "Publishing reports and semantic models",
    "Refresh configuration and credentials",
    "Gateways and cloud connections",
    "Shared semantic models and thin reports",
    "Power BI Apps",
    "Deployment pipelines",
    "Endorsement: Promoted and Certified",
    "Azure Government considerations",
    "Module lab walkthrough",
    "Knowledge check and discussion",
]


def deployment_pipeline_slide(prs, number, title, page, note=None, script=None):
    """Custom diagram: Dev -> Test -> Prod deployment pipeline stages with
    deployment rules and capacity/licensing callouts."""
    s = blank_slide(prs)
    add_rect(s, 0, 0, SLIDE_W, SLIDE_H, WHITE)
    add_rect(s, 0, 0, SLIDE_W, Inches(1.15), NAVY)
    add_text(s, Inches(0.7), Inches(0.22), Inches(1.6), Inches(0.4),
             f"TOPIC {number:02d}", size=13, color=GOLD, bold=True, font=BODY_FONT)
    add_text(s, Inches(0.7), Inches(0.52), Inches(11.5), Inches(0.55), title,
             size=24, color=WHITE, bold=True, font=HEADER_FONT)

    add_text(s, Inches(0.7), Inches(1.4), Inches(11.9), Inches(0.5),
             "Promote the same PBIP-originated content across three "
             "capacity-backed workspaces",
             size=16, color=NAVY_DARK, bold=True, font=HEADER_FONT)

    stages = [
        ("DEV", "Development workspace",
         "PBIP author publishes here first",
         "Owner: report/model author",
         ICE, NAVY_DARK),
        ("TEST", "Test / UAT workspace",
         "Business reviewers validate data & UX",
         "Deployment rules swap data source",
         NAVY, WHITE),
        ("PROD", "Production workspace",
         "App is published from this workspace",
         "Change-controlled, restricted access",
         GOLD, NAVY_DARK),
    ]

    row_y = Inches(2.2)
    box_h = Inches(2.2)
    box_w = Inches(3.55)
    gap = Inches(0.55)
    start_x = Inches(0.55)
    positions = []
    for i, (stage, name, line1, line2, fill, text_color) in enumerate(stages):
        x = start_x + i * (box_w + gap)
        add_rect(s, x, row_y, box_w, box_h, fill,
                 line_color=NAVY if fill in (ICE, GOLD) else None)
        add_text(s, x + Inches(0.15), row_y + Inches(0.15),
                 box_w - Inches(0.3), Inches(0.55), stage,
                 size=22, color=text_color, bold=True, font=HEADER_FONT,
                 align=PP_ALIGN.CENTER)
        add_text(s, x + Inches(0.15), row_y + Inches(0.75),
                 box_w - Inches(0.3), Inches(0.45), name,
                 size=14, color=text_color, bold=True, font=BODY_FONT,
                 align=PP_ALIGN.CENTER)
        sub_color = SLATE if fill in (ICE, GOLD) else ICE
        add_text(s, x + Inches(0.18), row_y + Inches(1.25),
                 box_w - Inches(0.36), Inches(0.4), line1,
                 size=11.5, color=sub_color, italic=True, font=BODY_FONT,
                 align=PP_ALIGN.CENTER, line_spacing=1.15)
        add_text(s, x + Inches(0.18), row_y + Inches(1.65),
                 box_w - Inches(0.36), Inches(0.45), line2,
                 size=11.5, color=sub_color, italic=True, font=BODY_FONT,
                 align=PP_ALIGN.CENTER, line_spacing=1.15)
        positions.append((x, x + box_w))

    arrow_h = Inches(0.42)
    arrow_y = row_y + box_h / 2 - arrow_h / 2
    for i in range(len(stages) - 1):
        x1 = positions[i][1]
        x2 = positions[i + 1][0]
        arrow = s.shapes.add_shape(
            MSO_SHAPE.RIGHT_ARROW,
            Emu(int(x1 + Inches(0.03))), Emu(int(arrow_y)),
            Emu(int(x2 - x1 - Inches(0.06))), arrow_h,
        )
        arrow.fill.solid()
        arrow.fill.fore_color.rgb = NAVY
        arrow.line.fill.background()
        add_text(s, Emu(int(x1)), Emu(int(arrow_y - Inches(0.35))),
                 Emu(int(x2 - x1)), Inches(0.3),
                 "Deploy", size=11, color=NAVY_DARK, bold=True,
                 font=BODY_FONT, align=PP_ALIGN.CENTER)

    # Requirements strip
    req_y = Inches(4.9)
    add_rect(s, Inches(0.55), req_y, Inches(12.25), Inches(1.55), LIGHT_BG,
             line_color=NAVY)
    add_text(s, Inches(0.8), req_y + Inches(0.12), Inches(11.75), Inches(0.4),
             "What the pipeline needs to work",
             size=13, color=NAVY_DARK, bold=True, font=BODY_FONT)
    reqs = [
        "\u2022  Premium / Fabric capacity (or PPU) assigned to each stage workspace",
        "\u2022  Deployment rules to swap data sources, parameters, and connection strings per stage",
        "\u2022  Least-privileged workspace roles per stage \u2014 fewer Admins in PROD than in DEV",
        "\u2022  \u26a0  Verify for Gov: confirm licensing, capacity, and Service availability in Azure Government",
    ]
    for i, r in enumerate(reqs):
        add_text(s, Inches(0.9), req_y + Inches(0.5) + Inches(0.25) * i,
                 Inches(12.0), Inches(0.3), r,
                 size=11.5, color=INK, font=BODY_FONT)

    if note:
        add_text(s, Inches(0.7), Inches(6.65), Inches(11.9), Inches(0.7),
                 note, size=12.5, color=SLATE, italic=True, font=BODY_FONT,
                 line_spacing=1.1)
    add_page_number(s, page)
    if script:
        set_notes(s, script)
    return s


def build():
    prs = new_presentation()
    page = 1

    # 1. Title
    title_slide(
        prs, MODULE_NO, TITLE, SUBTITLE,
        script=[
            "Welcome learners to Module 8. Frame this module as the moment where the artifacts they "
            "built in Modules 1 through 7 \u2014 the semantic model, the DAX measures, the report "
            "pages, the security roles \u2014 stop living on a laptop and start living in the Power BI "
            "Service, where real users consume them.",
            "Set expectations: this is not primarily a Desktop authoring module. Most of what we "
            "cover here happens in the browser at app.powerbi.com \u2014 workspaces, roles, semantic "
            "model settings, gateway configuration, Apps, deployment pipelines, and endorsement. "
            "The lab is deliberately tenant-shaped: students publish PBIP content, review refresh "
            "and gateway settings, package an App, and walk through the endorsement checklist.",
            "Preview the lab: students take the PBIP report and model from earlier modules, publish "
            "them to a training workspace, review scheduled refresh and credentials, discuss "
            "gateway architecture, build a thin report against the shared semantic model, distribute "
            "the content as a Power BI App, and complete the endorsement governance checklist.",
            "Call out the Azure Government angle up front. Core deployment \u2014 workspaces, roles, "
            "publishing, Apps, endorsement \u2014 is Gov-ready. Gateways, cloud connections, App "
            "audiences, and deployment pipelines are Verify for Gov: we teach them, but students "
            "may not be able to run them hands-on depending on tenant policy and licensing.",
        ],
    )
    page += 1

    # 2. Agenda
    agenda_slide(
        prs, MODULE_NO, AGENDA_TOPICS, page=page,
        script=[
            "Walk through the thirteen agenda items at a brisk pace. Group them mentally: topics 1-3 "
            "are the deployment target \u2014 workspaces, roles, ownership. Topics 4-6 are the "
            "operational plumbing \u2014 publishing, refresh, gateways. Topics 7-10 are the "
            "distribution and governance layer \u2014 shared models, Apps, pipelines, endorsement.",
            "Flag that topics 6, 8-audience, 9, and 11 (gateways, App audiences, deployment "
            "pipelines, Azure Government) contain material that is Verify-for-Gov \u2014 we will "
            "teach the pattern and architecture even where the classroom tenant cannot support a "
            "live demo.",
            "Tell students that topic 12 is where all of this becomes real \u2014 they publish "
            "actual PBIP-authored content, wire up refresh, build a thin report on a shared model, "
            "and distribute an App. Topic 13 is the knowledge check and open discussion.",
        ],
    )
    page += 1

    # Topic 1 — From authoring to enterprise deployment
    content_slide(
        prs, 1, "From Authoring to Enterprise Deployment", page=page,
        lead_items=[
            "PBIP is the source of record: report.json, model definitions, and TMDL under source "
            "control \u2014 not the .pbix binary.",
            "Publishing targets in the Service are workspaces; consumers get content through Apps, "
            "shared semantic models, or direct workspace access.",
            "Ownership splits into two roles: content owner (author of the report/model) and "
            "support owner (who fields the ticket when refresh fails at 2 a.m.).",
            "Supportability means every deployed artifact has a documented refresh cadence, data "
            "source list, gateway or cloud connection, and known consumer audience.",
        ],
        why_items=[
            "Without a source-of-record convention, teams end up with three different .pbix files "
            "all named 'Sales_Final_v2' and no way to answer 'what actually shipped'.",
            "Naming the two ownership roles explicitly prevents the classic failure where the "
            "author leaves the company and refresh silently breaks for months.",
            "Supportability is what turns a report into a product: the difference between 'nice "
            "dashboard' and 'certified enterprise asset' is almost entirely documentation and "
            "process, not visual polish.",
        ],
        footer="Lab tie-in: Exercise 1 asks students to record content owner, support owner, and "
               "workspace name \u2014 that documentation is the deployment artifact.",
        script=[
            "Open by drawing the boundary line for the module: everything to the left of publish is "
            "authoring \u2014 Desktop, PBIP, source control. Everything to the right is deployment \u2014 "
            "the Service, workspaces, refresh, Apps, endorsement. Module 8 is the right-hand side.",
            "Emphasize the PBIP shift. In the old world .pbix was the deliverable; in the current "
            "world PBIP is the source of record and .pbix is a generated artifact. That matters "
            "because Git-based review, code-review-style approval, and reproducible builds are only "
            "possible against PBIP. If your team is still emailing .pbix files around, deployment "
            "governance is impossible.",
            "Give the concrete owner example: a report author moves teams and the semantic model "
            "credential expires six months later. Without a named support owner, the refresh "
            "failure email goes to a dead mailbox and consumers see stale data for a week before "
            "anyone notices. This is why we insist on two names, not one.",
            "Transition: the next several topics unpack each piece of that supportability contract \u2014 "
            "starting with the container everything lives in, the workspace.",
        ],
    )
    page += 1

    # Topic 2 — Workspace design
    content_slide(
        prs, 2, "Workspace Design: Dev, Test, Prod, and Subject Areas", page=page,
        lead_items=[
            "Separate workspaces for Development, Test, and Production \u2014 not folders, not tabs, "
            "actual distinct workspaces with distinct role assignments.",
            "Organize by domain or subject area for large tenants: 'Finance \u2014 Revenue \u2014 PROD' is "
            "clearer than 'BI Team Workspace 4'.",
            "Naming convention should encode {Domain} \u2014 {Subject} \u2014 {Stage} so a workspace "
            "picker at 200 items is still navigable.",
            "Record content owner, support owner, and business sponsor on every workspace \u2014 in "
            "the description field at minimum.",
        ],
        why_items=[
            "Dev/Test/Prod separation is what lets you change a measure or a data source without "
            "breaking the report that the CFO opens every Monday morning.",
            "Domain naming lets tenant admins see the tenant as a business map, not an "
            "alphabet soup \u2014 and lets adoption metrics roll up meaningfully.",
            "Explicit ownership is a Certified-endorsement prerequisite \u2014 you cannot certify "
            "content whose owner is 'the team that used to own this'.",
        ],
        footer="Lab tie-in: Exercise 1 documents the intended DEV/TEST/PROD pattern even when the "
               "classroom tenant provides a single workshop workspace.",
        script=[
            "Open by acknowledging that workspaces are unglamorous \u2014 they look like folders \u2014 "
            "but they are the fundamental security and lifecycle boundary in the Service. Every "
            "role assignment, every App, every deployment pipeline stage attaches to a workspace.",
            "Explain the three-stage pattern concretely. DEV is where the author publishes freely "
            "and breaks things. TEST is where a business SME validates the numbers against a known "
            "source of truth. PROD is where the App audience lives and where change is controlled. "
            "Give the analogy: if you would not run SQL updates directly against a production "
            "database, you should not republish reports directly to a production workspace.",
            "Walk through the naming convention with an example: 'Finance \u2014 Revenue Reporting \u2014 "
            "PROD'. Contrast with the anti-pattern of 'John's Workspace' or 'New Workspace (2)'. "
            "In a tenant with 500 workspaces, naming is not aesthetic \u2014 it is search.",
            "Transition: once you have workspaces, the next question is who can do what inside "
            "them, which is roles \u2014 the next topic.",
        ],
    )
    page += 1

    # Topic 3 — Workspace roles (TABLE)
    table_slide(
        prs, 3, "Workspace Roles \u2014 What Each One Can Do", page=page,
        headers=["Role", "Can do", "Cannot do", "Typical assignment"],
        rows=[
            ["Admin",
             "Everything: add/remove users, delete workspace, change all settings, publish App.",
             "\u2014 (no restrictions inside the workspace)",
             "1\u20132 named platform owners per workspace."],
            ["Member",
             "Publish/edit content, share reports, publish and update the App, add Contributors/Viewers.",
             "Cannot add other Members/Admins or delete the workspace.",
             "Senior content authors and BI leads."],
            ["Contributor",
             "Create, edit, and publish reports & semantic models; schedule refresh.",
             "Cannot manage users, share reports externally, or publish the App.",
             "Day-to-day report authors."],
            ["Viewer",
             "View reports, interact with visuals, export data if allowed.",
             "Cannot edit, publish, or see the workspace list of items.",
             "Rarely used in PROD \u2014 consumers get the App instead."],
        ],
        col_widths=[1.1, 3.5, 2.6, 2.6],
        note="Security implication: viewers should almost never be added directly to a production "
             "workspace \u2014 distribute through an App so audience membership is manageable "
             "separately from workspace membership.",
        script=[
            "Walk down the table left to right. Emphasize that these four roles are the entire "
            "security surface of a workspace \u2014 there are no custom roles, no fine-grained "
            "per-report permissions inside a workspace. If you need more granularity than this, "
            "you use a shared semantic model with Build permission, or App audiences, or RLS.",
            "Stress the Admin/Member distinction: it is easy to over-grant Admin because 'they "
            "need to publish' \u2014 but Member can publish too. Reserve Admin for the one or two "
            "people who own the workspace lifecycle. In a Certified content workspace, having six "
            "Admins is a governance red flag.",
            "Explain the Viewer anti-pattern in the footer: adding fifty Viewers directly to a "
            "PROD workspace works, but it means the audience list is now tangled up with the "
            "workspace membership list, and you cannot revoke consumer access without touching the "
            "workspace itself. Apps decouple these two concerns.",
            "Transition: now that we know who can act inside a workspace, let's actually publish "
            "something into one.",
        ],
    )
    page += 1

    # Topic 4 — Publishing
    content_slide(
        prs, 4, "Publishing Reports and Semantic Models", page=page,
        lead_items=[
            "From Power BI Desktop, Home \u2192 Publish \u2192 choose target workspace. Publishing pushes "
            "the report AND its semantic model as two separate Service artifacts.",
            "PBIP-originated content publishes the same way \u2014 open the .pbip in Desktop, then "
            "publish; the .pbix is generated in-flight and not the source of record.",
            "Semantic model ownership defaults to the publishing user's identity \u2014 use "
            "'Take over' in dataset settings to transfer to a service account for stable refresh.",
            "Overwriting an existing report keeps its Service ID (bookmarks, comments, subscriptions "
            "survive); deleting and republishing loses them.",
        ],
        why_items=[
            "Understanding that publish creates two artifacts \u2014 report AND model \u2014 is critical "
            "for the next topic: shared semantic models let many reports point at one model.",
            "Ownership tied to a user account is the single most common cause of silent refresh "
            "failure \u2014 when that user's password rotates or account is disabled, refresh dies.",
            "Preserving Service IDs on overwrite is why we publish updates rather than "
            "delete-and-republish \u2014 subscribers don't get orphaned.",
        ],
        footer="Lab tie-in: Exercise 1 publishes the PBIP report to the training workspace and "
               "verifies both the report and semantic model appear in the workspace list.",
        script=[
            "Demo the mechanics quickly and then focus on the two things students always get wrong. "
            "First: publish creates two artifacts, not one \u2014 the report file and the semantic "
            "model. They appear as separate rows in the workspace, they can be managed "
            "independently, and the report row is really just a pointer to the model plus a page "
            "layout definition.",
            "Second: whoever hits Publish becomes the semantic model owner by default. That is fine "
            "for DEV. In PROD it is a time bomb \u2014 an individual account is a single point of "
            "failure. Show the 'Take over' button in the semantic model Settings page and explain "
            "that Certified content typically has its ownership taken over by a service principal "
            "or a shared team account.",
            "Give the overwrite-vs-delete example: a subscriber has a weekly email subscription to "
            "the report. If you delete the report and republish, their subscription is silently "
            "gone. If you overwrite (publish with the same name), the subscription survives. "
            "'Publish' is an upsert, not an insert.",
            "Transition: once content is in the Service, the interesting question is how it stays "
            "current \u2014 which is refresh.",
        ],
    )
    page += 1

    # Topic 5 — Refresh
    content_slide(
        prs, 5, "Refresh Configuration and Credentials", page=page,
        lead_items=[
            "Scheduled refresh lives on the semantic model, not the report. Configure it in the "
            "model's Settings page: up to 8 refreshes/day on Pro, 48/day on Premium/Fabric.",
            "Every data source in the model needs credentials configured in the Service \u2014 the "
            "Desktop credentials do not travel with publish.",
            "Privacy levels (Public, Organizational, Private) prevent Power Query from folding "
            "queries across sources with different trust levels \u2014 misconfiguration causes "
            "'formula.firewall' errors on refresh.",
            "Refresh history shows the last ~60 runs with duration, status, and error \u2014 first "
            "stop when a user says 'the numbers are stale'.",
        ],
        why_items=[
            "Refresh on the semantic model, not the report, is what enables one model to feed many "
            "reports \u2014 you don't want to refresh the same data ten times because ten reports use it.",
            "The credential handoff at publish time is the #1 support ticket in early "
            "deployments \u2014 users don't realize Desktop credentials are local-only.",
            "Refresh history is your audit log: a Certified model must be able to answer 'when did "
            "this last refresh successfully and how long did it take?' at any moment.",
        ],
        footer="Lab tie-in: Exercise 2 opens the semantic model settings, reviews credentials, sets "
               "a schedule, runs refresh, and captures refresh history.",
        script=[
            "Start with the mental model split: a semantic model has data, a report does not. "
            "Therefore refresh belongs to the model. This is one of those things that is obvious "
            "once you know it and confusing until then \u2014 students will look for the refresh "
            "button on the report and not find it.",
            "Walk through the credential story slowly. When you develop in Desktop, credentials are "
            "cached on your machine \u2014 refresh 'just works'. When you publish, none of that "
            "travels with the model. The Service pops up asking for credentials again, and if you "
            "skip that dialog, scheduled refresh will fail with 'credentials not configured'. Make "
            "this the first thing students check in Exercise 2.",
            "Explain privacy levels with a concrete example: a query that joins a public CSV to an "
            "internal SQL Server will refuse to fold if their privacy levels are incompatible, and "
            "Power BI shows a 'Formula.Firewall' error. The fix is to declare both sources' privacy "
            "levels consistently in the data source settings. Students who see this error in the "
            "wild almost always assume it is a code bug \u2014 it isn't.",
            "Close with refresh history: this is the operational dashboard for the model. When a "
            "consumer says 'the data looks wrong', the first check is not the DAX \u2014 it is when "
            "the model last refreshed successfully.",
        ],
    )
    page += 1

    # Topic 6 — Gateways
    content_slide(
        prs, 6, "Gateways and Cloud Connections", page=page,
        lead_items=[
            "On-premises data gateway is a Windows service installed inside the corporate network; "
            "it bridges Service refresh requests to on-prem sources (SQL Server, file shares, SAP).",
            "Gateway clusters (2+ gateway nodes with the same registration) provide high "
            "availability and load balancing \u2014 a single-node gateway is a single point of failure.",
            "Data source mappings link the semantic model's source (server + database) to a "
            "specific gateway data source with stored credentials \u2014 no mapping, no refresh.",
            "Cloud connections replace legacy 'Manage connections' for cloud sources and are the "
            "modern surface for OAuth-based data sources.",
            "\u26a0  Verify for Gov: gateway version, network path, connector support, and tenant "
            "policy all need explicit validation in Azure Government.",
        ],
        why_items=[
            "The gateway is where on-prem meets cloud. Every 'refresh works in Desktop but fails in "
            "Service' ticket for an on-prem source ends up here.",
            "Clustering is not optional in real production \u2014 the gateway host rebooting for "
            "Windows Update should not take down refresh for the whole tenant.",
            "Cloud connections centralize connection strings and credentials so that ten semantic "
            "models pointing at the same Snowflake instance don't each store their own copy.",
            "The Gov note matters because gateway installers and connector matrices differ between "
            "commercial and government clouds \u2014 don't assume parity.",
        ],
        footer="Lab tie-in: Exercise 3 (gateway-backed refresh) is hands-on when a gateway is "
               "available, otherwise students document the required setup as an alternate path.",
        script=[
            "Frame the gateway as a translator that lives behind the corporate firewall. The Power "
            "BI Service cannot reach an internal SQL Server directly \u2014 the network doesn't "
            "allow it. So a small Windows service runs on-prem, holds an outbound connection to "
            "the Service, and forwards refresh queries into the internal source. That's it. It's "
            "not magic and it's not a virtual machine \u2014 it's a Windows service on a domain-joined "
            "box.",
            "Explain clustering with the operational example: a single-node gateway means every "
            "Patch Tuesday reboot silently breaks refresh for hours. Two nodes registered as a "
            "cluster give you high availability and load balancing at no extra license cost \u2014 "
            "just a second VM. Any real production deployment uses at least two.",
            "Cover data source mappings and cloud connections at a lighter touch: the important "
            "conceptual point is that the semantic model does not store credentials to the source "
            "directly \u2014 it references a named gateway data source or cloud connection, and the "
            "credentials live there. This is what lets a Power BI admin rotate a credential once "
            "and affect every model that uses it.",
            "Close with the Gov-verification note. In Azure Government the gateway installer is "
            "different, the connector matrix is smaller, and some cloud connections may not be "
            "generally available yet. This is why Exercise 3 has an explicit 'alternate path' \u2014 "
            "we teach the architecture even where we cannot demo it.",
        ],
    )
    page += 1

    # Topic 7 — Shared semantic models & thin reports
    content_slide(
        prs, 7, "Shared Semantic Models and Thin Reports", page=page,
        lead_items=[
            "One published semantic model can feed many 'thin' reports \u2014 reports that contain "
            "only visuals and page layout, no data or model definition.",
            "A thin report author needs Build permission on the shared semantic model (granted "
            "from the model's Manage permissions page) \u2014 workspace membership alone is not enough.",
            "Certified semantic models are the strongest signal a tenant offers that a model is "
            "the authoritative source \u2014 use them as the target for thin reports.",
            "Impact analysis (from the semantic model page) shows every dependent report and app "
            "before you make a breaking change \u2014 always run it before renaming a column.",
        ],
        why_items=[
            "This is the reuse pattern that makes governance economical: one governed model, ten "
            "teams building their own reports off it, one refresh cycle, one set of measures.",
            "Build permission decouples 'who can consume the model' (Read) from 'who can build "
            "new reports on it' (Build) \u2014 exactly the split enterprises need.",
            "Certification directs report authors to the right model instead of everyone forking "
            "their own copy \u2014 it is discoverability, not just a badge.",
            "Impact analysis is what turns a scary rename into a routine change \u2014 you know "
            "exactly which reports need to be retested.",
        ],
        footer="Lab tie-in: Exercise 4 grants Build permission and creates a thin report against "
               "the shared semantic model published in Exercise 1.",
        script=[
            "Open with the shift in mental model: for years, every .pbix was a self-contained unit "
            "\u2014 model plus report bundled together. Shared semantic models break that: the model "
            "is a first-class Service artifact that other reports connect to, the way many "
            "applications connect to one database.",
            "Explain Build permission concretely. Read permission means 'you can open reports "
            "built on this model'. Build permission means 'you can create new reports that connect "
            "to this model'. The distinction matters because most consumers should have Read only \u2014 "
            "we do not want ten thousand shadow reports proliferating. Build is granted "
            "deliberately, usually to a small group of trained authors.",
            "Certification is the discoverability half of the story. When a report author picks "
            "'Connect to Power BI semantic model' in Desktop, Certified models are surfaced first "
            "with a visible badge. That is how you steer authors toward the sanctioned model "
            "instead of the seventeen unofficial copies that inevitably appear.",
            "Emphasize impact analysis: before you rename [Revenue] to [Net Revenue] on the "
            "shared model, you click Impact Analysis and see 'this affects 14 reports across 3 "
            "workspaces'. That's the difference between a controlled change and a Monday-morning "
            "outage. Have students run impact analysis in the lab even if they don't ship the "
            "change.",
        ],
    )
    page += 1

    # Topic 8 — Apps
    content_slide(
        prs, 8, "Power BI Apps \u2014 Packaging Content for Consumers", page=page,
        lead_items=[
            "An App is a curated, versioned publication of selected workspace content: reports, "
            "dashboards, and a defined navigation, distributed to a consumer audience.",
            "App consumers are assigned separately from workspace membership \u2014 users, security "
            "groups, or Microsoft 365 groups. Consumers do not see the workspace itself.",
            "The App has a staged update model: changes in the workspace are not visible to "
            "consumers until an author clicks 'Update app'. This is the release valve.",
            "Audience targeting (optional) shows different pages/reports to different audience "
            "groups within one App \u2014 e.g. Regional Managers see two extra pages the wider "
            "audience does not. \u26a0  Verify for Gov.",
        ],
        why_items=[
            "Apps are the answer to 'how do 500 consumers get the report without being added to "
            "the workspace one by one'. They are the enterprise distribution channel.",
            "The consumer/workspace decoupling means you can change workspace membership (add a "
            "new author, remove a departed one) without affecting who sees the App.",
            "The staged update is what lets you fix a report in DEV, promote through TEST, land "
            "it in PROD, and still choose when consumers see the change \u2014 no surprise updates.",
            "Audiences replace the old pattern of publishing three near-duplicate Apps for three "
            "audience segments \u2014 one App, three views.",
        ],
        footer="Lab tie-in: Exercise 5 packages workspace content into an App, configures name and "
               "navigation, assigns consumers, and publishes.",
        script=[
            "Frame the App as the consumer-facing product. The workspace is the workshop where "
            "authors build; the App is the shop window where consumers browse. Consumers do not "
            "want to see all fourteen work-in-progress reports \u2014 they want a curated set with "
            "clear navigation, published deliberately.",
            "Walk through the update model, because this is where teams get confused. When an "
            "author republishes a report to the workspace, App consumers do NOT see that change "
            "immediately. The workspace and the App are decoupled by design \u2014 the author must "
            "click 'Update app' to push the workspace state to consumers. This is a feature: it "
            "means work-in-progress does not leak to consumers.",
            "Give the audience example concretely. Suppose Finance publishes one App with three "
            "audiences \u2014 'Executives' see the summary page, 'Regional Managers' see the summary "
            "plus regional detail, and 'Analysts' see everything including the reconciliation "
            "tabs. Each audience is defined once, and the App renders the appropriate pages per "
            "user. Historically you would have shipped three separate Apps; audiences collapse "
            "that into one.",
            "Close with the Gov note on audiences. Audiences are Verify for Gov \u2014 confirm "
            "availability before promising them to a Gov customer. If unavailable, fall back to "
            "the older pattern of a single audience per App and rely on RLS for row-level "
            "differentiation.",
        ],
    )
    page += 1

    # Topic 9 — Deployment pipelines (CUSTOM DIAGRAM)
    deployment_pipeline_slide(
        prs, 9, "Deployment Pipelines \u2014 Dev / Test / Prod Promotion",
        page=page,
        note="Lab tie-in: the optional deployment pipelines lab has students walk through the "
             "three-stage promotion pattern; hands-on availability depends on capacity and tenant "
             "policy.",
        script=[
            "Open by connecting this back to Topic 2. We designed three workspaces \u2014 DEV, TEST, "
            "PROD \u2014 for a reason. Deployment pipelines are the Service feature that automates "
            "the promotion of content between them. Without a pipeline you can still promote "
            "manually by republishing to each workspace; with a pipeline you promote with a single "
            "click and get a comparison view of what changed.",
            "Walk the diagram left to right. DEV is where the author publishes freely \u2014 the model "
            "points at a dev database. TEST is where a business SME validates \u2014 the pipeline's "
            "'deployment rule' automatically swaps the connection string to the test database when "
            "content is promoted. PROD is where the App audience lives \u2014 the deployment rule "
            "swaps to the prod database and access is tightly controlled.",
            "Explain deployment rules as the key mechanism. Without them, promoting DEV to PROD "
            "would carry the dev database connection string into PROD \u2014 which is wrong and "
            "sometimes dangerous. Rules let you say 'when this dataset lands in PROD, replace "
            "server=dev-sql with server=prod-sql'. Same for parameter values, credentials, and so "
            "on.",
            "Close with the requirements strip. Pipelines need Premium or Fabric capacity (or PPU) "
            "on each stage \u2014 this is the licensing constraint that most often blocks classroom "
            "hands-on. And they are Verify for Gov: confirm capacity availability and Service "
            "parity in Azure Government before promising a pipeline demo to a Gov customer. That "
            "is why the deployment pipelines lab is optional.",
        ],
    )
    page += 1

    # Topic 10 — Endorsement (TABLE)
    table_slide(
        prs, 10, "Endorsement \u2014 Promoted vs Certified",
        headers=["Aspect", "Promoted", "Certified"],
        rows=[
            ["Who grants it",
             "Any workspace Contributor+ can self-promote their own content.",
             "Only tenant-designated reviewers (defined in Admin Portal) can certify."],
            ["Signal to consumers",
             "'This is content the author considers ready for a wider audience.'",
             "'This is the authoritative, tenant-endorsed source for this topic.'"],
            ["Governance evidence required",
             "Basic \u2014 self-attestation of quality and support.",
             "Ownership, refresh reliability, security review, and support model must all be documented."],
            ["Typical use",
             "Team-level or department-level reference content.",
             "Enterprise-wide 'source of truth' semantic models and reports."],
            ["Discoverability effect",
             "Ribbon badge in workspace and shared model picker.",
             "Ribbon badge plus priority sort in the shared model picker; steers new authors here."],
        ],
        col_widths=[1.6, 3.4, 3.4],
        script=[
            "Reinforce that endorsement is a governance process, not just a visual badge. The badge "
            "is the last five percent \u2014 the ninety-five percent is the review that earned it. "
            "Consumers see the badge; authors and reviewers experience the process.",
            "Walk the two columns side by side. Promoted is lightweight: the content author asserts "
            "'this is ready to share'. It's a good discoverability signal for team-scale content and "
            "requires no tenant approval. Anyone with edit rights can promote their own work.",
            "Certified is heavyweight and deliberately so. Only a small set of tenant-designated "
            "reviewers can grant it. Before they grant it, they check the boxes on the "
            "endorsement checklist \u2014 documented ownership, reliable refresh history, security "
            "review, defined support model. These are the same items on the Certified section of "
            "the lab's endorsement checklist.",
            "Give the discoverability point at the end: Certified content is not just badged, it "
            "is prioritized in the shared model picker in Desktop. That is how the tenant steers "
            "the next report author toward the sanctioned semantic model instead of them forking "
            "yet another copy. Endorsement is how governance scales.",
        ],
        page=page,
        note="Lab tie-in: Exercise 6 walks through the endorsement checklist and decides whether "
             "the lab content is eligible for Promoted or Certified.",
    )
    page += 1

    # Topic 11 — Azure Government considerations (TABLE)
    table_slide(
        prs, 11, "Azure Government \u2014 What to Validate",
        headers=["Feature area", "Status", "What to verify before promising it"],
        rows=[
            ["Workspaces, roles, publishing",
             "Gov-ready",
             "Confirm tenant settings, licensing (Pro/PPU/Premium/Fabric), and org naming policy."],
            ["Scheduled refresh",
             "Gov-ready / verify per source",
             "Data source connector must be supported in Gov cloud; validate credentials work."],
            ["Gateways",
             "Verify for Gov",
             "Gateway installer version, network path, connector matrix, tenant policy."],
            ["Cloud connections",
             "Verify for Gov",
             "Confirm connector GA status in Gov; some cloud connections lag commercial."],
            ["Power BI Apps",
             "Gov-ready",
             "Confirm tenant external-sharing restrictions and audience licensing."],
            ["App audiences",
             "Verify for Gov",
             "Confirm Service parity; fall back to single-audience App + RLS if unavailable."],
            ["Deployment pipelines",
             "Verify for Gov",
             "Requires Premium/Fabric/PPU capacity per stage; confirm Service availability."],
            ["Promoted / Certified",
             "Gov-ready (process)",
             "Confirm tenant has designated Certified reviewers configured in Admin Portal."],
        ],
        col_widths=[2.4, 1.8, 4.2],
        script=[
            "Set the frame: Azure Government is not 'Power BI minus features' \u2014 it is Power BI "
            "with a delivery lag and a different validation matrix. Core features (workspaces, "
            "roles, publishing, refresh, Apps, endorsement) are Gov-ready. Advanced integration "
            "features (gateways, cloud connections, App audiences, deployment pipelines) are "
            "Verify for Gov \u2014 they may work, but confirm before you promise.",
            "Walk down the table. Point out that 'Verify for Gov' does not mean 'unavailable' \u2014 "
            "it means 'do not assume commercial-cloud parity, check the current state'. Feature "
            "GA dates shift; a feature that was Verify last quarter may be Gov-ready this quarter.",
            "Give the operational implication: when a customer asks 'can we do deployment "
            "pipelines in Gov?', the correct answer is not yes or no \u2014 it is 'let me confirm "
            "the current status in your specific tenant with your current licensing'. That is what "
            "Verify for Gov means in practice.",
            "Transition: this table is the reference students should keep after the course. When "
            "they are planning a Gov deployment, this is the first-pass triage list.",
        ],
        page=page,
        note="Keep this table as a reference sheet: any Gov customer conversation should start by "
             "triaging features against these status categories.",
    )
    page += 1

    # Topic 12 — Module lab walkthrough (CHECKLIST)
    checklist_slide(
        prs, "Module Lab Walkthrough",
        items=[
            "Exercise 1: Publish PBIP report and semantic model to the training workspace; record owner and role.",
            "Exercise 2: Open semantic model settings, review credentials, set a schedule, run refresh, capture history.",
            "Exercise 3: Gateway-backed refresh \u2014 map data source to gateway cluster and validate, or document as an alternate path.",
            "Exercise 4: Grant Build permission on the shared semantic model and create a thin report against it.",
            "Exercise 5: Configure and publish a Power BI App with navigation and consumer assignments.",
            "Optional: Explore App audiences \u2014 create multiple audiences and validate audience-specific views (Verify for Gov).",
            "Optional: Walk through a deployment pipeline \u2014 assign dev/test/prod workspaces and review deployment rules (Verify for Gov).",
            "Exercise 6: Complete the endorsement governance checklist and decide Promoted vs Certified eligibility.",
        ],
        script=[
            "Frame the lab as the tenant-side counterpart to everything students have built in "
            "Modules 1-7. The PBIP report and semantic model they authored now become published "
            "Service artifacts with refresh schedules, credentials, gateways, Apps, and "
            "endorsement.",
            "Walk each exercise briefly: publish (Ex 1) proves the deployment path works; refresh "
            "(Ex 2) proves the model can stay current; gateway (Ex 3) covers the on-prem case "
            "even when we can only document it; thin report (Ex 4) proves the shared-model reuse "
            "pattern; App (Ex 5) proves consumer distribution; endorsement (Ex 6) proves the "
            "governance layer.",
            "Call out the two optional labs explicitly \u2014 App audiences and deployment "
            "pipelines. Both are Verify for Gov. If the classroom tenant supports them, do them "
            "hands-on. If not, walk the architecture and have students document what setup would "
            "be required in their target tenant.",
            "Finish with a time hint: this lab tends to run long because the Service UI is slower "
            "than Desktop and refresh operations have wait time. Encourage students to move on to "
            "the next exercise while a refresh is running rather than staring at the progress bar.",
        ],
        page=page,
        kicker="Lab exercises \u2014 in order",
    )
    page += 1

    # Topic 13 — Knowledge check & discussion (CHECKLIST)
    checklist_slide(
        prs, "Knowledge Check and Discussion",
        items=[
            "Workspace role and ownership documented \u2014 who is Admin, who is support owner?",
            "Report and semantic model published (or demonstrated) to the training workspace.",
            "Refresh settings reviewed \u2014 credentials configured, schedule set, history captured.",
            "Gateway requirements documented \u2014 cluster, data source mapping, network path.",
            "Shared semantic model + Build permission pattern reviewed \u2014 who can build vs. who can read?",
            "App distribution reviewed \u2014 workspace membership vs. App consumer audience.",
            "App audiences marked Verify for Gov \u2014 what would validation entail in your tenant?",
            "Deployment pipelines marked Verify for Gov \u2014 what capacity/licensing gates apply?",
            "Endorsement checklist completed \u2014 Promoted, Certified, or neither, and why?",
            "Discussion: what is the single biggest deployment risk in your current environment?",
        ],
        script=[
            "Use this slide as both a knowledge check and a discussion opener. Walk the checklist "
            "items in order \u2014 for each one ask a student to answer with the specifics from "
            "their lab attempt, not a generic definition.",
            "Look for these red-flag answers: 'I don't remember who the owner was', 'I skipped the "
            "credential dialog', 'I added all consumers as workspace Viewers'. Each of these is a "
            "real production anti-pattern the module was designed to prevent. Use the answer as a "
            "teaching moment, not a gotcha.",
            "For the two Verify-for-Gov items, don't ask 'did it work' \u2014 ask 'what would you "
            "need to confirm before recommending this to a Gov customer'. The correct answer "
            "mentions licensing/capacity, Service availability, and tenant policy \u2014 the same "
            "three questions we've been repeating all module.",
            "Close with the final discussion prompt about the single biggest deployment risk in "
            "the student's environment. This surfaces the real-world constraints they carry back "
            "to work \u2014 shadow reports, unowned models, expired credentials, ungoverned "
            "gateways \u2014 and often becomes the most valuable part of the whole module for them.",
        ],
        page=page,
        kicker="Validation checklist and discussion",
    )
    page += 1

    # Closing
    closing_slide(
        prs, MODULE_NO,
        next_module="Module 09 \u2014 Monitoring, Administration, and Governance",
        subtitle="Learners can now publish PBIP-authored content, govern workspaces and roles, "
                 "configure refresh and gateways, distribute Apps, and evaluate content for "
                 "Promoted or Certified endorsement \u2014 in commercial and Gov clouds.",
        page=page,
        script=[
            "Close by naming what students can now do: publish PBIP content into a governed "
            "workspace, configure and troubleshoot refresh, distribute an App to consumers, and "
            "walk a piece of content through the endorsement process. That is the operational "
            "half of the workshop \u2014 the authoring half is now paired with the deployment half.",
            "Reinforce the governance thread that ran through the whole module: owner, support "
            "owner, workspace role, refresh cadence, data source, gateway, App audience, "
            "endorsement. Every one of those is a documentation artifact, and together they are "
            "what turns a report into a supportable enterprise asset.",
            "Preview Module 9 \u2014 Monitoring, Administration, and Governance. Where Module 8 "
            "asked 'how do we deploy this?', Module 9 asks 'how do we run this at tenant scale?': "
            "usage metrics, activity logs, admin monitoring, capacity metrics, Purview, DLP, and "
            "the support operating model. It is the natural continuation for platform owners and "
            "support teams.",
            "End with a practical prompt: encourage students to pick one workspace they own back "
            "at work and, before Module 9, spend fifteen minutes documenting its owner, support "
            "owner, refresh history, and endorsement status. That small act closes the loop "
            "between the classroom and their real environment.",
        ],
    )

    OUT_PATH.parent.mkdir(parents=True, exist_ok=True)
    prs.save(str(OUT_PATH))
    total = len(prs.slides)
    print(f"Saved: {OUT_PATH} ({total} slides)")


if __name__ == "__main__":
    build()
