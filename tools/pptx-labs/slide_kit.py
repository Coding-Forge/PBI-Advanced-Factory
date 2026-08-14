#!/usr/bin/env python3
"""
Builds instructor-delivery PowerPoint decks for the Advanced Power BI labs.

Each deck is generated directly from:
  - modules/<module>/slide-outline.md      (topic list / sequence)
  - Student/Labs/Source/<module>/README.md (technical substance)
  - Student/Labs/Web/<module>.html         (student-facing framing/wording)

Design goals:
  - Every slide has real, specific technical content (no generic "Key ideas /
    Why it helps" filler repeated across slides).
  - Speaker notes are a full delivery script an instructor can read/paraphrase
    from, not just a bullet re-statement.
  - Visual variety: title/section/content/table/diagram/checklist layouts.

Usage:
    python tools/pptx-labs/build_lab01.py
"""
from pptx import Presentation
from pptx.util import Inches, Pt, Emu
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR
from pptx.enum.shapes import MSO_SHAPE, MSO_CONNECTOR
from pptx.oxml.ns import qn
import copy

# ---------------------------------------------------------------------------
# Palette: "Midnight Executive" (navy/ice-blue) — matches the analytical /
# enterprise-modeling subject matter of Lab 01.
# ---------------------------------------------------------------------------
NAVY = RGBColor(0x1E, 0x27, 0x61)
NAVY_DARK = RGBColor(0x14, 0x1B, 0x45)
ICE = RGBColor(0xCA, 0xDC, 0xFC)
WHITE = RGBColor(0xFF, 0xFF, 0xFF)
INK = RGBColor(0x1F, 0x23, 0x28)
SLATE = RGBColor(0x4B, 0x55, 0x63)
GOLD = RGBColor(0xE8, 0xB4, 0x4A)
LIGHT_BG = RGBColor(0xF6, 0xF8, 0xFC)
CARD_BORDER = RGBColor(0xD8, 0xDE, 0xE6)

HEADER_FONT = "Cambria"
BODY_FONT = "Calibri"

SLIDE_W = Inches(13.333)
SLIDE_H = Inches(7.5)


def new_presentation():
    prs = Presentation()
    prs.slide_width = SLIDE_W
    prs.slide_height = SLIDE_H
    return prs


def blank_slide(prs):
    return prs.slides.add_slide(prs.slide_layouts[6])


def set_notes(slide, script_lines):
    """script_lines: list of paragraph strings forming the delivery script."""
    notes = slide.notes_slide
    tf = notes.notes_text_frame
    tf.clear()
    tf.paragraphs[0].text = script_lines[0]
    for line in script_lines[1:]:
        p = tf.add_paragraph()
        p.text = line


def add_rect(slide, x, y, w, h, color, line_color=None, shadow=False):
    shape = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, x, y, w, h)
    shape.fill.solid()
    shape.fill.fore_color.rgb = color
    if line_color is None:
        shape.line.fill.background()
    else:
        shape.line.color.rgb = line_color
        shape.line.width = Pt(0.75)
    shape.shadow.inherit = False
    return shape


def add_text(slide, x, y, w, h, text, size=16, color=INK, bold=False, italic=False,
             font=BODY_FONT, align=PP_ALIGN.LEFT, anchor=MSO_ANCHOR.TOP, line_spacing=1.0,
             space_after=0):
    box = slide.shapes.add_textbox(x, y, w, h)
    tf = box.text_frame
    tf.word_wrap = True
    tf.vertical_anchor = anchor
    tf.margin_left = 0
    tf.margin_right = 0
    tf.margin_top = 0
    tf.margin_bottom = 0
    p = tf.paragraphs[0]
    p.text = text
    p.alignment = align
    p.line_spacing = line_spacing
    p.space_after = Pt(space_after)
    for run in p.runs:
        run.font.size = Pt(size)
        run.font.bold = bold
        run.font.italic = italic
        run.font.name = font
        run.font.color.rgb = color
    return box


def add_bullets(slide, x, y, w, h, items, size=14, color=INK, font=BODY_FONT,
                 space_after=8, line_spacing=1.05, bullet_color=None, bold_lead=False):
    """items: list of str, or list of (str, level) tuples for nesting."""
    box = slide.shapes.add_textbox(x, y, w, h)
    tf = box.text_frame
    tf.word_wrap = True
    tf.margin_left = 0
    tf.margin_right = 0
    tf.margin_top = 0
    tf.margin_bottom = 0
    first = True
    for item in items:
        if isinstance(item, tuple):
            text, level = item
        else:
            text, level = item, 0
        p = tf.paragraphs[0] if first else tf.add_paragraph()
        first = False
        p.level = level
        p.space_after = Pt(space_after)
        p.line_spacing = line_spacing
        run = p.add_run()
        run.text = text
        run.font.size = Pt(size - (2 if level else 0))
        run.font.name = font
        run.font.color.rgb = color
        _set_bullet_char(p, "\u25AA" if level else "\u2022",
                          bullet_color or color)
    return box


def _set_bullet_char(paragraph, char, color):
    pPr = paragraph._pPr
    if pPr is None:
        pPr = paragraph._p.get_or_add_pPr()
    # Remove existing bullet defs
    for tag in ("a:buChar", "a:buAutoNum", "a:buNone"):
        el = pPr.find(qn(tag))
        if el is not None:
            pPr.remove(el)
    buFont = pPr.makeelement(qn("a:buFont"), {"typeface": "Arial"})
    buClr = pPr.makeelement(qn("a:buClr"), {})
    srgb = pPr.makeelement(qn("a:srgbClr"), {"val": str(color)})
    buClr.append(srgb)
    buChar = pPr.makeelement(qn("a:buChar"), {"char": char})
    pPr.append(buClr)
    pPr.append(buFont)
    pPr.append(buChar)


def add_page_number(slide, number, dark=False):
    color = ICE if dark else SLATE
    add_text(slide, Inches(12.6), Inches(7.05), Inches(0.6), Inches(0.35),
              f"{number:02d}", size=11, color=color, font=BODY_FONT, align=PP_ALIGN.RIGHT)


def add_kicker(slide, text, dark=False):
    color = GOLD if dark else NAVY
    add_text(slide, Inches(0.7), Inches(0.55), Inches(8), Inches(0.35), text.upper(),
              size=12, color=color, bold=True, font=BODY_FONT)


# ---------------------------------------------------------------------------
# Slide builders
# ---------------------------------------------------------------------------

def title_slide(prs, module_no, title, subtitle, script):
    s = blank_slide(prs)
    add_rect(s, 0, 0, SLIDE_W, SLIDE_H, NAVY_DARK)
    add_rect(s, 0, Inches(6.35), SLIDE_W, Inches(1.15), NAVY)
    add_text(s, Inches(0.9), Inches(1.55), Inches(6), Inches(0.5),
              f"MODULE {module_no:02d}", size=16, color=GOLD, bold=True, font=BODY_FONT)
    add_text(s, Inches(0.9), Inches(2.05), Inches(10.5), Inches(2.1), title,
              size=44, color=WHITE, bold=True, font=HEADER_FONT, line_spacing=1.05)
    add_text(s, Inches(0.9), Inches(4.05), Inches(10.5), Inches(1.0), subtitle,
              size=20, color=ICE, italic=True, font=BODY_FONT)
    add_text(s, Inches(0.9), Inches(6.6), Inches(10), Inches(0.6),
              "Advanced Power BI Training \u2014 Instructor Deck", size=13, color=ICE, font=BODY_FONT)
    set_notes(s, script)
    return s


def agenda_slide(prs, module_no, topics, script, page):
    s = blank_slide(prs)
    add_rect(s, 0, 0, SLIDE_W, SLIDE_H, LIGHT_BG)
    add_rect(s, 0, 0, Inches(0.18), SLIDE_H, NAVY)
    add_kicker(s, f"Module {module_no:02d} Agenda")
    add_text(s, Inches(0.7), Inches(0.95), Inches(9), Inches(0.7), "What We'll Cover",
              size=32, color=NAVY_DARK, bold=True, font=HEADER_FONT)
    col_w = Inches(5.6)
    left_x = Inches(0.7)
    right_x = Inches(6.9)
    half = (len(topics) + 1) // 2
    top = Inches(1.95)
    available_h = Inches(7.5) - top - Inches(0.5)
    row_step = Emu(int(min(Inches(0.85), available_h / max(half, 1))))
    for i, topic in enumerate(topics):
        col = left_x if i < half else right_x
        row = i if i < half else i - half
        y = top + row * row_step
        add_rect(s, col, y, Inches(0.42), Inches(0.42), NAVY)
        add_text(s, col, y, Inches(0.42), Inches(0.42), str(i + 1), size=16, color=WHITE,
                  bold=True, font=HEADER_FONT, align=PP_ALIGN.CENTER, anchor=MSO_ANCHOR.MIDDLE)
        add_text(s, col + Inches(0.58), y + Inches(0.02), col_w - Inches(0.6), Inches(0.7), topic,
                  size=14.5, color=INK, bold=False, font=BODY_FONT, anchor=MSO_ANCHOR.MIDDLE)
    add_page_number(s, page)
    set_notes(s, script)
    return s


def section_break_slide(prs, number, title, tagline, script, page):
    s = blank_slide(prs)
    add_rect(s, 0, 0, SLIDE_W, SLIDE_H, NAVY_DARK)
    add_rect(s, Inches(0.9), Inches(2.55), Inches(1.7), Inches(1.7), NAVY)
    add_text(s, Inches(0.9), Inches(2.55), Inches(1.7), Inches(1.7), f"{number:02d}",
              size=48, color=GOLD, bold=True, font=HEADER_FONT, align=PP_ALIGN.CENTER,
              anchor=MSO_ANCHOR.MIDDLE)
    add_text(s, Inches(2.9), Inches(2.55), Inches(9.5), Inches(1.0), title, size=34,
              color=WHITE, bold=True, font=HEADER_FONT, anchor=MSO_ANCHOR.MIDDLE)
    add_text(s, Inches(2.9), Inches(3.55), Inches(9.5), Inches(0.7), tagline, size=17,
              color=ICE, italic=True, font=BODY_FONT)
    add_page_number(s, page, dark=True)
    set_notes(s, script)
    return s


def content_slide(prs, number, title, lead_items, why_items, script, page,
                   footer=None):
    """Two-column concept slide: 'What & how' vs 'Why it matters'."""
    s = blank_slide(prs)
    add_rect(s, 0, 0, SLIDE_W, SLIDE_H, WHITE)
    add_rect(s, 0, 0, SLIDE_W, Inches(1.35), NAVY)
    add_text(s, Inches(0.7), Inches(0.28), Inches(1.6), Inches(0.4), f"TOPIC {number:02d}",
              size=13, color=GOLD, bold=True, font=BODY_FONT)
    add_text(s, Inches(0.7), Inches(0.62), Inches(11.5), Inches(0.65), title, size=27,
              color=WHITE, bold=True, font=HEADER_FONT)

    col_w = Inches(5.55)
    left_x = Inches(0.7)
    right_x = Inches(6.7)
    top_y = Inches(1.75)

    add_rect(s, left_x, top_y, col_w, Inches(0.5), ICE)
    add_text(s, left_x + Inches(0.2), top_y, col_w - Inches(0.4), Inches(0.5), "What & how",
              size=15, color=NAVY_DARK, bold=True, font=BODY_FONT, anchor=MSO_ANCHOR.MIDDLE)
    add_bullets(s, left_x + Inches(0.05), top_y + Inches(0.65), col_w - Inches(0.1), Inches(4.6),
                lead_items, size=14.5, color=INK, bullet_color=NAVY, space_after=10)

    add_rect(s, right_x, top_y, col_w, Inches(0.5), NAVY)
    add_text(s, right_x + Inches(0.2), top_y, col_w - Inches(0.4), Inches(0.5), "Why it matters",
              size=15, color=WHITE, bold=True, font=BODY_FONT, anchor=MSO_ANCHOR.MIDDLE)
    add_bullets(s, right_x + Inches(0.05), top_y + Inches(0.65), col_w - Inches(0.1), Inches(4.6),
                why_items, size=14.5, color=INK, bullet_color=GOLD, space_after=10)

    if footer:
        add_rect(s, 0, Inches(6.75), SLIDE_W, Inches(0.75), LIGHT_BG)
        add_text(s, Inches(0.7), Inches(6.87), Inches(11.9), Inches(0.55), footer, size=12.5,
                  color=SLATE, italic=True, font=BODY_FONT, line_spacing=1.1)
    add_page_number(s, page)
    set_notes(s, script)
    return s


def table_slide(prs, number, title, headers, rows, script, page, col_widths=None, note=None):
    s = blank_slide(prs)
    add_rect(s, 0, 0, SLIDE_W, SLIDE_H, WHITE)
    add_rect(s, 0, 0, SLIDE_W, Inches(1.15), NAVY)
    add_text(s, Inches(0.7), Inches(0.22), Inches(1.6), Inches(0.4), f"TOPIC {number:02d}",
              size=13, color=GOLD, bold=True, font=BODY_FONT)
    add_text(s, Inches(0.7), Inches(0.52), Inches(11.5), Inches(0.55), title, size=24,
              color=WHITE, bold=True, font=HEADER_FONT)

    n_rows = len(rows) + 1
    n_cols = len(headers)
    table_top = Inches(1.45)
    table_h = Inches(5.1)
    gfx = s.shapes.add_table(n_rows, n_cols, Inches(0.7), table_top, Inches(11.9), table_h)
    table = gfx.table
    if col_widths:
        total = sum(col_widths)
        for i, cw in enumerate(col_widths):
            table.columns[i].width = Inches(11.9 * cw / total)
    for c, h in enumerate(headers):
        cell = table.cell(0, c)
        cell.text = h
        cell.fill.solid()
        cell.fill.fore_color.rgb = NAVY
        for p in cell.text_frame.paragraphs:
            p.alignment = PP_ALIGN.LEFT
            for r in p.runs:
                r.font.bold = True
                r.font.size = Pt(12.5)
                r.font.color.rgb = WHITE
                r.font.name = BODY_FONT
    for r_idx, row in enumerate(rows, start=1):
        for c_idx, val in enumerate(row):
            cell = table.cell(r_idx, c_idx)
            cell.text = val
            cell.fill.solid()
            cell.fill.fore_color.rgb = LIGHT_BG if r_idx % 2 == 0 else WHITE
            for p in cell.text_frame.paragraphs:
                for rn in p.runs:
                    rn.font.size = Pt(11.5)
                    rn.font.name = BODY_FONT
                    rn.font.color.rgb = INK
    if note:
        add_text(s, Inches(0.7), Inches(6.7), Inches(11.9), Inches(0.7), note, size=12.5,
                  color=SLATE, italic=True, font=BODY_FONT, line_spacing=1.1)
    add_page_number(s, page)
    set_notes(s, script)
    return s


def diagram_slide(prs, number, title, script, page, note=None):
    """Star-schema diagram: dims around a central fact box."""
    s = blank_slide(prs)
    add_rect(s, 0, 0, SLIDE_W, SLIDE_H, WHITE)
    add_rect(s, 0, 0, SLIDE_W, Inches(1.15), NAVY)
    add_text(s, Inches(0.7), Inches(0.22), Inches(1.6), Inches(0.4), f"TOPIC {number:02d}",
              size=13, color=GOLD, bold=True, font=BODY_FONT)
    add_text(s, Inches(0.7), Inches(0.52), Inches(11.5), Inches(0.55), title, size=24,
              color=WHITE, bold=True, font=HEADER_FONT)

    cx, cy = Inches(5.55), Inches(3.4)
    fw, fh = Inches(2.4), Inches(1.1)
    add_rect(s, cx, cy, fw, fh, NAVY)
    add_text(s, cx, cy, fw, fh, "FactSales\n(transaction grain)", size=13, color=WHITE,
              bold=True, font=BODY_FONT, align=PP_ALIGN.CENTER, anchor=MSO_ANCHOR.MIDDLE)

    dims = [
        ("DimCustomer", Inches(0.7), Inches(1.55)),
        ("DimProduct", Inches(0.7), Inches(5.1)),
        ("DimTerritory", Inches(9.9), Inches(1.55)),
        ("DimOrderDate /\nDimShipDate", Inches(9.9), Inches(5.1)),
    ]
    dw, dh = Inches(2.1), Inches(1.0)
    for label, x, y in dims:
        add_rect(s, x, y, dw, dh, ICE, line_color=NAVY)
        add_text(s, x, y, dw, dh, label, size=12.5, color=NAVY_DARK, bold=True,
                  font=BODY_FONT, align=PP_ALIGN.CENTER, anchor=MSO_ANCHOR.MIDDLE)
        # connector line (coerce to int EMU to avoid float attribute corruption)
        line = s.shapes.add_connector(
            MSO_CONNECTOR.STRAIGHT,
            Emu(int(x + (dw if x < cx else 0))),
            Emu(int(y + dh / 2)),
            Emu(int(cx + (0 if x < cx else fw))),
            Emu(int(cy + fh / 2)),
        )
        line.line.color.rgb = SLATE
        line.line.width = Pt(1.5)

    if note:
        add_text(s, Inches(0.7), Inches(6.65), Inches(11.9), Inches(0.7), note, size=13,
                  color=SLATE, italic=True, font=BODY_FONT, line_spacing=1.1)
    add_page_number(s, page)
    set_notes(s, script)
    return s


def bridge_diagram_slide(prs, number, title, script, page, note=None):
    """Bridge-table diagram: DimCustomer <-> BridgeCustomerSegment <-> DimSegment,
    explicitly showing the bridge does NOT relate directly to the fact table."""
    s = blank_slide(prs)
    add_rect(s, 0, 0, SLIDE_W, SLIDE_H, WHITE)
    add_rect(s, 0, 0, SLIDE_W, Inches(1.15), NAVY)
    add_text(s, Inches(0.7), Inches(0.22), Inches(1.6), Inches(0.4), f"TOPIC {number:02d}",
              size=13, color=GOLD, bold=True, font=BODY_FONT)
    add_text(s, Inches(0.7), Inches(0.52), Inches(11.5), Inches(0.55), title, size=24,
              color=WHITE, bold=True, font=HEADER_FONT)

    box_h = Inches(1.15)
    top_y = Inches(2.15)
    bottom_y = Inches(4.55)

    dim_w = Inches(3.0)
    bridge_w = Inches(3.4)

    cust_x = Inches(0.9)
    bridge_x = Inches(4.95)
    seg_x = Inches(9.4)

    add_rect(s, cust_x, top_y, dim_w, box_h, ICE, line_color=NAVY)
    add_text(s, cust_x, top_y, dim_w, box_h, "DimCustomer\n(one row per customer)", size=13,
              color=NAVY_DARK, bold=True, font=BODY_FONT, align=PP_ALIGN.CENTER,
              anchor=MSO_ANCHOR.MIDDLE)

    add_rect(s, seg_x, top_y, dim_w, box_h, ICE, line_color=NAVY)
    add_text(s, seg_x, top_y, dim_w, box_h, "DimSegment\n(one row per segment)", size=13,
              color=NAVY_DARK, bold=True, font=BODY_FONT, align=PP_ALIGN.CENTER,
              anchor=MSO_ANCHOR.MIDDLE)

    add_rect(s, bridge_x, top_y, bridge_w, box_h, NAVY)
    add_text(s, bridge_x, top_y, bridge_w, box_h,
              "BridgeCustomerSegment\n(one row per customer-segment pair)", size=12.5,
              color=WHITE, bold=True, font=BODY_FONT, align=PP_ALIGN.CENTER,
              anchor=MSO_ANCHOR.MIDDLE)

    def straight(x1, y1, x2, y2, color=SLATE, width=1.5):
        ln = s.shapes.add_connector(MSO_CONNECTOR.STRAIGHT, Emu(int(x1)), Emu(int(y1)),
                                     Emu(int(x2)), Emu(int(y2)))
        ln.line.color.rgb = color
        ln.line.width = Pt(width)
        return ln

    mid_y = top_y + box_h / 2
    straight(cust_x + dim_w, mid_y, bridge_x, mid_y)
    straight(bridge_x + bridge_w, mid_y, seg_x, mid_y)

    add_rect(s, cust_x, bottom_y, dim_w, box_h, NAVY)
    add_text(s, cust_x, bottom_y, dim_w, box_h, "FactSales\n(transaction grain)", size=13,
              color=WHITE, bold=True, font=BODY_FONT, align=PP_ALIGN.CENTER,
              anchor=MSO_ANCHOR.MIDDLE)

    straight(cust_x + dim_w / 2, top_y + box_h, cust_x + dim_w / 2, bottom_y,
              color=NAVY, width=2.25)

    add_rect(s, seg_x, bottom_y, dim_w, box_h, RGBColor(0xE8, 0xEC, 0xF3), line_color=SLATE)
    add_text(s, seg_x, bottom_y, dim_w, box_h,
              "No direct relationship\n(bridge \u2192 dim \u2192 fact only)", size=11.5,
              color=SLATE, italic=True, font=BODY_FONT, align=PP_ALIGN.CENTER,
              anchor=MSO_ANCHOR.MIDDLE)

    mid_between_y = top_y + box_h + (bottom_y - top_y - box_h) / 2
    add_text(s, cust_x + dim_w / 2 + Inches(0.15), mid_between_y - Inches(0.15), Inches(1.4),
              Inches(0.3), "filters \u2193", size=11, color=NAVY, bold=True, font=BODY_FONT)

    if note:
        add_text(s, Inches(0.7), Inches(6.65), Inches(11.9), Inches(0.7), note, size=13,
                  color=SLATE, italic=True, font=BODY_FONT, line_spacing=1.1)
    add_page_number(s, page)
    set_notes(s, script)
    return s


def checklist_slide(prs, title, items, script, page, kicker="Validation checklist"):
    s = blank_slide(prs)
    add_rect(s, 0, 0, SLIDE_W, SLIDE_H, LIGHT_BG)
    add_rect(s, 0, 0, Inches(0.18), SLIDE_H, GOLD)
    add_kicker(s, kicker)
    add_text(s, Inches(0.7), Inches(0.95), Inches(10.5), Inches(0.7), title, size=30,
              color=NAVY_DARK, bold=True, font=HEADER_FONT)
    col_w = Inches(5.6)
    half = (len(items) + 1) // 2
    # Distribute rows as card blocks across the full available vertical space
    # so short lists don't leave a large empty band at the bottom of the slide.
    top = Inches(1.95)
    bottom_margin = Inches(0.5)
    available_h = Inches(7.5) - top - bottom_margin
    row_h = Emu(int(available_h / max(half, 1)))
    card_gap = Inches(0.14)
    card_h = Emu(int(row_h - card_gap))
    for i, item in enumerate(items):
        col = Inches(0.7) if i < half else Inches(6.9)
        row = i if i < half else i - half
        y = top + row * row_h
        add_rect(s, col, y, col_w, card_h, WHITE, line_color=CARD_BORDER)
        box_size = Inches(0.34)
        box_y = y + (card_h - box_size) / 2
        add_rect(s, col + Inches(0.2), box_y, box_size, box_size, LIGHT_BG, line_color=NAVY)
        add_text(s, col + Inches(0.72), y, col_w - Inches(0.92), card_h, item, size=13.5,
                  color=INK, font=BODY_FONT, anchor=MSO_ANCHOR.MIDDLE, line_spacing=1.05)
    add_page_number(s, page)
    set_notes(s, script)
    return s


def closing_slide(prs, module_no, next_module, script, page, subtitle=None):
    s = blank_slide(prs)
    add_rect(s, 0, 0, SLIDE_W, SLIDE_H, NAVY_DARK)
    add_text(s, Inches(0.9), Inches(2.3), Inches(11), Inches(1.0),
              f"Module {module_no:02d} Complete", size=36, color=WHITE, bold=True, font=HEADER_FONT)
    if subtitle is None:
        subtitle = "Learners now have a governed, reusable semantic model foundation to build on."
    add_text(s, Inches(0.9), Inches(3.3), Inches(11), Inches(1.4),
              subtitle,
              size=17, color=ICE, italic=True, font=BODY_FONT, line_spacing=1.15)
    add_rect(s, Inches(0.9), Inches(4.4), Inches(11.3), Inches(1.1), NAVY)
    add_text(s, Inches(1.15), Inches(4.55), Inches(10.8), Inches(0.85),
              f"Up next: {next_module}", size=16, color=GOLD, bold=True, font=BODY_FONT,
              anchor=MSO_ANCHOR.MIDDLE)
    add_page_number(s, page, dark=True)
    set_notes(s, script)
    return s
