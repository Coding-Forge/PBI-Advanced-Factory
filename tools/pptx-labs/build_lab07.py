#!/usr/bin/env python3
"""
Builds the Lab 07 (Security Design - RLS/OLS) instructor deck.
Run from repo root: python tools/pptx-labs/build_lab07.py
Output: modules/07-security-design/assets/security-design.pptx
"""
import sys
from pathlib import Path

sys.path.insert(0, str(Path(__file__).resolve().parent))
from pptx.util import Inches, Pt, Emu
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR
from pptx.enum.shapes import MSO_SHAPE, MSO_CONNECTOR

from slide_kit import (
    new_presentation, title_slide, agenda_slide,
    content_slide, table_slide, checklist_slide, closing_slide,
    blank_slide, add_rect, add_text, add_page_number, set_notes,
    SLIDE_W, SLIDE_H, NAVY, NAVY_DARK, ICE, WHITE, INK, SLATE, GOLD,
    LIGHT_BG, HEADER_FONT, BODY_FONT,
)

OUT_PATH = Path(__file__).resolve().parent.parent.parent / \
    "modules" / "07-security-design" / "assets" / "security-design.pptx"

MODULE_NO = 7
TITLE = "Security Design"
SUBTITLE = ("Layered Power BI security: RLS, dynamic role filters, OLS, Build permission, "
            "sensitivity labels, and Gov-aware sharing patterns.")

AGENDA_TOPICS = [
    "Power BI security layers",
    "Content access vs. data access",
    "Static RLS",
    "Dynamic RLS",
    "Service role assignment",
    "Build permission",
    "Object-level security (OLS)",
    "Sensitivity labels and Purview",
    "Sharing and external users",
    "Security testing",
    "Azure Government considerations",
    "Security review checklist",
    "Module lab walkthrough",
    "Knowledge check and discussion",
]


def dynamic_rls_flow_slide(prs, number, title, page, note=None, script=None):
    """Custom diagram: dynamic RLS filter propagation from the signed-in user's
    UPN through the security mapping table into DimTerritory and FactSales."""
    s = blank_slide(prs)
    add_rect(s, 0, 0, SLIDE_W, SLIDE_H, WHITE)
    add_rect(s, 0, 0, SLIDE_W, Inches(1.15), NAVY)
    add_text(s, Inches(0.7), Inches(0.22), Inches(1.6), Inches(0.4), f"TOPIC {number:02d}",
             size=13, color=GOLD, bold=True, font=BODY_FONT)
    add_text(s, Inches(0.7), Inches(0.52), Inches(11.5), Inches(0.55), title, size=24,
             color=WHITE, bold=True, font=HEADER_FONT)

    # Row layout: User identity  ->  SecurityUserTerritory  ->  DimTerritory  ->  FactSales
    row_y = Inches(2.9)
    box_h = Inches(1.4)
    box_w = Inches(2.55)
    gap = Inches(0.45)
    start_x = Inches(0.55)

    boxes = [
        ("Signed-in user\nUSERPRINCIPALNAME()",
         "e.g. jsmith@contoso.gov", ICE, NAVY_DARK),
        ("SecurityUserTerritory\n(mapping table)",
         "[UserPrincipalName] = USERPRINCIPALNAME()", NAVY, WHITE),
        ("DimTerritory\n(filtered rows)",
         "Only TerritoryKeys the user maps to", ICE, NAVY_DARK),
        ("FactSales\n(filtered rows)",
         "Only rows for allowed territories", NAVY, WHITE),
    ]

    positions = []
    for i, (label, sub, fill, text_color) in enumerate(boxes):
        x = start_x + i * (box_w + gap)
        add_rect(s, x, row_y, box_w, box_h, fill,
                 line_color=NAVY if fill == ICE else None)
        add_text(s, x + Inches(0.1), row_y + Inches(0.15), box_w - Inches(0.2),
                 Inches(0.65), label, size=13, color=text_color, bold=True,
                 font=BODY_FONT, align=PP_ALIGN.CENTER, anchor=MSO_ANCHOR.MIDDLE)
        add_text(s, x + Inches(0.1), row_y + Inches(0.78), box_w - Inches(0.2),
                 Inches(0.55), sub, size=10.5,
                 color=SLATE if fill == ICE else ICE,
                 italic=True, font=BODY_FONT, align=PP_ALIGN.CENTER,
                 anchor=MSO_ANCHOR.MIDDLE, line_spacing=1.1)
        positions.append((x, x + box_w))

    # Right-pointing block arrows between boxes
    arrow_h = Inches(0.32)
    arrow_y = row_y + box_h / 2 - arrow_h / 2
    for i in range(len(boxes) - 1):
        x1 = positions[i][1]
        x2 = positions[i + 1][0]
        arrow = s.shapes.add_shape(
            MSO_SHAPE.RIGHT_ARROW,
            Emu(int(x1 + Inches(0.02))), Emu(int(arrow_y)),
            Emu(int(x2 - x1 - Inches(0.04))), arrow_h,
        )
        arrow.fill.solid()
        arrow.fill.fore_color.rgb = GOLD
        arrow.line.fill.background()

    # Caption band above the row
    add_text(s, Inches(0.7), Inches(1.55), Inches(11.9), Inches(0.55),
             "Filter propagation for a single dynamic-RLS query",
             size=16, color=NAVY_DARK, bold=True, font=HEADER_FONT)
    add_text(s, Inches(0.7), Inches(2.05), Inches(11.9), Inches(0.85),
             "The role filter fires on the security mapping table using the "
             "signed-in user's UPN, then propagates one-to-many into "
             "DimTerritory, and finally into FactSales. The role's filter "
             "expression lives on the security table \u2014 never on FactSales.",
             size=12.5, color=SLATE, italic=True, font=BODY_FONT, line_spacing=1.2)

    # DAX callout block below the flow
    dax_y = Inches(4.85)
    add_rect(s, Inches(0.55), dax_y, Inches(12.25), Inches(1.55), LIGHT_BG,
             line_color=NAVY)
    add_text(s, Inches(0.8), dax_y + Inches(0.1), Inches(11.8), Inches(0.4),
             "Role: 'Dynamic Territory Security'  \u2014  Filter on "
             "SecurityUserTerritory:",
             size=13, color=NAVY_DARK, bold=True, font=BODY_FONT)
    add_text(s, Inches(0.8), dax_y + Inches(0.55), Inches(11.8), Inches(0.5),
             "[UserPrincipalName] = USERPRINCIPALNAME()",
             size=17, color=NAVY_DARK, bold=True, font="Consolas",
             align=PP_ALIGN.CENTER)
    add_text(s, Inches(0.8), dax_y + Inches(1.05), Inches(11.8), Inches(0.4),
             "Relationship path: SecurityUserTerritory[TerritoryKey] 1 \u2192 * "
             "DimTerritory[TerritoryKey] 1 \u2192 * FactSales[TerritoryKey]",
             size=11.5, color=SLATE, italic=True, font=BODY_FONT,
             align=PP_ALIGN.CENTER)

    if note:
        add_text(s, Inches(0.7), Inches(6.65), Inches(11.9), Inches(0.7), note,
                 size=12.5, color=SLATE, italic=True, font=BODY_FONT,
                 line_spacing=1.1)
    add_page_number(s, page)
    if script:
        set_notes(s, script)
    return s


def build():
    prs = new_presentation()
    page = 1

    # 1. Title
    title_slide(
        prs, MODULE_NO, TITLE, SUBTITLE,
        script=[
            "Welcome learners to Module 7. Frame security as the module that decides who can see "
            "what \u2014 not just at the report level, but all the way down to which rows and columns "
            "of the semantic model a signed-in user is even allowed to read.",
            "Set expectations up front: Power BI security is layered. Workspace roles, App "
            "audiences, sharing links, Build permission, RLS, OLS, and sensitivity labels are all "
            "separate decisions that combine to produce the effective access a user gets. Today's "
            "goal is to make each layer explicit instead of leaving it as an accident.",
            "Preview the lab: students will build a static RLS role that filters DimTerritory to a "
            "single region, then a dynamic RLS role that uses USERPRINCIPALNAME() against a "
            "SecurityUserTerritory mapping table imported from the workshop's synthetic CSV, and "
            "will test both roles using View as in Desktop.",
            "Call the Azure Government angle out immediately: static and dynamic RLS are Gov-ready, "
            "but OLS, sensitivity labels, Purview integration, external sharing and B2B are all "
            "'Verify for Gov' \u2014 they depend on tenant, capacity, and Purview configuration that "
            "the classroom tenant may or may not enable.",
        ]
    )
    page += 1

    # 2. Agenda
    agenda_slide(
        prs, MODULE_NO, AGENDA_TOPICS, page=page,
        script=[
            "Walk the fourteen agenda items quickly. Group them mentally for the room: topics 1-2 "
            "frame the security-layers model; 3-5 are the RLS core (static, dynamic, Service "
            "assignment); 6-9 cover the surrounding governance surface (Build permission, OLS, "
            "labels, sharing); 10-12 are validation, Gov posture, and the review checklist; "
            "13-14 are the lab and wrap-up.",
            "Flag that topics 7, 8, and parts of 9 are Verify-for-Gov \u2014 they are still important "
            "to teach, but the hands-on lab treats OLS and sensitivity labels as conceptual only.",
            "Tell students the lab (topic 13) is where dynamic RLS with a real mapping table and "
            "USERPRINCIPALNAME() becomes concrete: they will build the role, wire the relationship, "
            "and test it as sample UPN values from the mapping CSV.",
        ]
    )
    page += 1

    # 3. Topic 1 - Power BI security layers
    table_slide(
        prs, 1, "Power BI Security Layers", page=page,
        headers=["Layer", "Controls", "Where it is configured"],
        col_widths=[2.6, 5.6, 3.7],
        rows=[
            ["Tenant settings",
             "Global enablement of features like external sharing, export, Analyze in Excel, and "
             "custom visuals.",
             "Power BI Admin portal (tenant admin)."],
            ["Workspace roles",
             "Admin / Member / Contributor / Viewer inside a workspace \u2014 broad content and "
             "publishing rights.",
             "Workspace \u2192 Access."],
            ["App audiences",
             "Which users see which pages and which artifacts through a published App.",
             "Workspace \u2192 Create/Update App \u2192 Audiences."],
            ["Sharing links & Build permission",
             "Direct per-item share; Build lets a user reuse the semantic model in a new report or "
             "Analyze in Excel.",
             "Item \u2192 Share / Manage permissions."],
            ["RLS and OLS",
             "RLS filters rows a user can read; OLS hides entire tables or columns from a user.",
             "Desktop \u2192 Modeling \u2192 Manage roles; OLS via XMLA/tools."],
            ["Sensitivity labels",
             "Classify items (e.g. Confidential) and enforce export / DLP behavior downstream.",
             "Requires Microsoft Purview Information Protection."],
        ],
        note="Lab tie-in: the lab focuses on the RLS row of this table (Exercises 1-3). All other "
             "layers are set up outside Desktop \u2014 do not confuse a well-built RLS role with a "
             "complete security design.",
        script=[
            "Open with the mental model: Power BI security is not one setting \u2014 it is a stack. If "
            "any single layer is misconfigured, the effective access a user gets is wrong, even if "
            "every other layer looks correct.",
            "Walk the rows top-down. Tenant settings are the outermost gate \u2014 if the tenant "
            "admin has disabled external sharing at the tenant level, no amount of workspace or "
            "item sharing overrides that. Workspace roles are broad and blunt: a Contributor can "
            "publish content into the workspace, which means Contributor is effectively an authoring "
            "role, not a consumer role.",
            "For App audiences, use a concrete example: the same workspace can publish one App with "
            "an 'Executives' audience that sees a summary page and a 'Sales analyst' audience that "
            "sees the full detail page \u2014 App audiences do NOT filter data, only which pages and "
            "items are visible; RLS is what filters data.",
            "Land the RLS/OLS row hard because that's the module core: RLS filters rows a user is "
            "allowed to see inside a table; OLS hides an entire table or column from the metadata "
            "so the user doesn't even know it exists. They solve different problems and are "
            "configured very differently \u2014 RLS in Manage roles, OLS via XMLA-compatible tooling.",
            "Close with the footer note: passing the lab's static and dynamic RLS exercises is not "
            "the same as shipping a secure report. All the other layers still need to be reviewed "
            "before a real deployment.",
        ]
    )
    page += 1

    # 4. Topic 2 - Content access vs. data access
    content_slide(
        prs, 2, "Content Access vs. Data Access", page=page,
        lead_items=[
            "Content access answers 'can this user open the report or see the semantic model in a "
            "workspace or App?' \u2014 workspace roles, App audiences, and direct sharing decide this.",
            "Build permission answers 'can this user connect to the underlying semantic model from "
            "a new report, Excel, or XMLA endpoint?' \u2014 it is separate from view access.",
            "Data access answers 'once inside, which rows and columns does this user actually "
            "read?' \u2014 RLS filters rows, OLS hides tables and columns.",
            "These three questions are answered by three different layers. Granting one does not "
            "grant the others, and denying one does not automatically deny the others.",
        ],
        why_items=[
            "Treating 'they can open the report' as equivalent to 'they can only see their own "
            "data' is the single most common security-design mistake in Power BI.",
            "A user with only Read access to a report can still be given Build permission on the "
            "underlying semantic model \u2014 at which point they can reuse the whole model and see "
            "every row RLS would have filtered in the report.",
            "Separating the three questions makes review checklists tractable: for every "
            "user/group, you check content access, Build permission, and RLS role assignment as "
            "three distinct rows \u2014 not one blurred column.",
        ],
        footer="Lab connection: Exercise 4 is a discussion exercise that walks students through "
               "exactly this distinction using the security-role-matrix.csv persona-to-role file.",
        script=[
            "Introduce this as the mental frame that keeps the rest of the module clear. Three "
            "questions, three layers, three separate decisions.",
            "Give the concrete failure story: a well-meaning author grants a colleague Viewer on "
            "the workspace and Build on the semantic model 'so they can play with it in Excel.' The "
            "colleague opens Analyze in Excel and pulls every row of FactSales \u2014 the report's RLS "
            "role never applied to their Excel connection because Build permission bypasses the "
            "report layer entirely. That is not a bug; that is how Build permission is defined.",
            "Reinforce that RLS is a data-access mechanism, not a content-access one. RLS only "
            "kicks in for users who query the semantic model as members of an assigned RLS role. "
            "If they connect a different way \u2014 or if they are a workspace Admin/Member \u2014 the "
            "role may not apply.",
            "Tell students Exercise 4 has them walk the security-role-matrix.csv persona rows and "
            "identify which layer each row's expectation lives in \u2014 view, Build, RLS, or a "
            "combination. It's a paper exercise, but a critical one for review discipline.",
        ]
    )
    page += 1

    # 5. Topic 3 - Static RLS
    table_slide(
        prs, 3, "Static RLS \u2014 Role-to-Filter Mapping", page=page,
        headers=["Role name", "Table filtered", "DAX filter expression"],
        col_widths=[3.0, 3.2, 6.7],
        rows=[
            ["East Region", "DimTerritory", "[TerritoryRegion] = \"East\""],
            ["West Region", "DimTerritory", "[TerritoryRegion] = \"West\""],
            ["North Region", "DimTerritory", "[TerritoryRegion] = \"North\""],
            ["US Only", "DimTerritory", "[TerritoryCountry] = \"United States\""],
            ["Read-only Sample", "DimCustomer", "[CustomerType] = \"Reseller\""],
        ],
        note="Lab tie-in: Exercise 1 has students create the East Region role exactly as shown "
             "above via Modeling \u2192 Manage roles, then verify it using Modeling \u2192 View as. "
             "Filters are placed on the dimension \u2014 not on FactSales \u2014 so they propagate "
             "through the star schema.",
        script=[
            "Static RLS is the simplest pattern: one role per hard-coded scope, with a DAX filter "
            "expression baked directly into the role definition. Walk the table so students see "
            "that every row of it uses the same recipe \u2014 role name, target table, filter expression.",
            "Anchor to the lab: the East Region role in row 1 is the exact role Exercise 1 has "
            "them build. Show the click path out loud \u2014 Modeling \u2192 Manage roles \u2192 New \u2192 name "
            "the role 'East Region' \u2192 select DimTerritory \u2192 enter the filter "
            "[TerritoryRegion] = \"East\" \u2192 Save \u2014 so they know exactly what to expect when "
            "they get to the lab.",
            "Stress the placement rule: the filter goes on the dimension table (DimTerritory), not "
            "on FactSales. Because DimTerritory relates one-to-many into FactSales, filtering "
            "DimTerritory automatically filters FactSales through relationship propagation. Putting "
            "the filter on FactSales directly is a common mistake and is much slower.",
            "Point out the maintenance tradeoff explicitly: you need one role per region, and if a "
            "new region is added or an existing one is renamed, someone has to edit the role "
            "definition. That is the exact pain point dynamic RLS \u2014 the next topic \u2014 is designed "
            "to solve.",
            "Wrap by mentioning testing: Modeling \u2192 View as lets an author simulate any single "
            "role or combination of roles in Desktop before publishing. Students should always View "
            "as before assigning users in the Service.",
        ]
    )
    page += 1

    # 6. Topic 4 - Dynamic RLS (custom diagram slide)
    dynamic_rls_flow_slide(
        prs, 4, "Dynamic RLS \u2014 USERPRINCIPALNAME() Filter Flow",
        page=page,
        note="Lab tie-in: Exercise 2 has students import security-user-territory.csv, relate it to "
             "DimTerritory on TerritoryKey, and add the filter shown above to the security table "
             "inside the 'Dynamic Territory Security' role.",
        script=[
            "Dynamic RLS replaces N static roles with a single role that reads the current user's "
            "identity at query time and looks up their allowed rows from a mapping table. Introduce "
            "the pieces in the order they appear on the slide, left to right.",
            "Piece 1: USERPRINCIPALNAME() is a DAX function that returns the signed-in user's UPN "
            "\u2014 the same string as their Entra ID sign-in, e.g. jsmith@contoso.gov. In Desktop, "
            "View as lets you spoof any UPN string for testing.",
            "Piece 2: SecurityUserTerritory is a mapping table imported from "
            "security-user-territory.csv via the Web connector. Each row is one "
            "(UserPrincipalName, TerritoryKey) pair. A user with two territories has two rows; a "
            "user with one territory has one row. The role's DAX filter \u2014 "
            "[UserPrincipalName] = USERPRINCIPALNAME() \u2014 is placed on THIS table, not on "
            "DimTerritory or FactSales.",
            "Piece 3 and 4: Once the security table is filtered to just the current user's rows, "
            "the one-to-many relationship into DimTerritory propagates that filter forward, and "
            "DimTerritory's one-to-many relationship into FactSales completes the chain. The "
            "signed-in user only sees FactSales rows for territories they are mapped to.",
            "Close with the operational payoff and the identity gotcha: one role handles every "
            "user, and adding or removing access is a data update to the mapping table \u2014 no "
            "role edit. But identity must match exactly: the UPN in the mapping table has to be "
            "identical to what USERPRINCIPALNAME() returns for that user in the target tenant, "
            "including domain and case. Guest users and B2B accounts especially need validation.",
        ]
    )
    page += 1

    # 7. Topic 5 - Service role assignment
    content_slide(
        prs, 5, "Service Role Assignment", page=page,
        lead_items=[
            "In the Service, open the semantic model \u2192 Security to assign users or Entra ID "
            "groups (preferred) to each RLS role defined in Desktop.",
            "Use 'Test as role' in the Service to preview data for a specific user or role \u2014 "
            "the Service equivalent of Desktop's View as.",
            "Workspace roles override RLS in one direction: Admin, Member, and Contributor on the "
            "workspace bypass RLS on models in that workspace. Only Viewer respects RLS.",
            "Distribute content via an App with audiences instead of granting broad workspace "
            "access \u2014 App consumers get Viewer semantics, so RLS applies.",
        ],
        why_items=[
            "Assigning Entra ID groups instead of individual users means role membership is managed "
            "in the identity system, not inside Power BI \u2014 a huge governance win.",
            "The workspace-role bypass is the most surprising failure mode: if you 'quickly add' a "
            "tester as Member so they can help you QA, RLS silently stops applying to them and any "
            "test result is invalid.",
            "App distribution is the pattern that scales: authors and stewards stay in the "
            "workspace; consumers only touch the App, where RLS reliably applies.",
        ],
        footer="Lab connection: Exercise 3 has students publish to an approved training workspace, "
               "assign a user or group to a role, and test as role in the Service where available.",
        script=[
            "Bridge from Desktop to Service: everything in the previous two topics defined roles; "
            "this topic is about who is IN each role once the model is published.",
            "Walk the Service click path \u2014 workspace \u2192 semantic model \u2192 Security \u2014 and "
            "emphasize the assignment box takes users OR groups. Push groups hard: assigning "
            "'Sales-East-Region' as an Entra group means adding or removing salespeople is an "
            "identity-team operation, not a Power BI operation, which is exactly the separation of "
            "duties an enterprise wants.",
            "Now the failure mode: workspace roles Admin/Member/Contributor bypass RLS. If your "
            "tester is a workspace Member, they see all the data regardless of role assignment "
            "and the test result is meaningless. Only Viewer respects RLS. Tell students they must "
            "have a Viewer test account (or use Test as role) \u2014 they cannot QA RLS from their "
            "author account.",
            "Close on distribution shape: publish through Apps with audiences instead of granting "
            "the whole business Viewer on the workspace. The workspace stays a controlled authoring "
            "environment; the App is the consumption surface where RLS reliably applies for "
            "everyone.",
        ]
    )
    page += 1

    # 8. Topic 6 - Build permission
    content_slide(
        prs, 6, "Build Permission", page=page,
        lead_items=[
            "Build is a per-semantic-model permission that lets a user connect to the model from a "
            "new report, Analyze in Excel, or the XMLA endpoint \u2014 not just view an existing "
            "report.",
            "Build is what makes 'thin reports' possible: one shared, certified semantic model, "
            "many small reports built by different authors on top of it.",
            "RLS still applies to Build users, but only if they are assigned to an RLS role and "
            "are not a workspace Admin/Member/Contributor.",
            "Analyze in Excel connects through the XMLA endpoint, respects RLS the same way a "
            "Power BI report does, and is subject to tenant setting enablement.",
        ],
        why_items=[
            "Build permission is the primary path an author uses to reuse a certified semantic "
            "model without republishing it \u2014 core to a shared-model governance strategy.",
            "Granting Read on a report but Build on the model is a real risk: the user can pull "
            "data out of the model in ways the report itself never exposed, so RLS coverage of the "
            "underlying tables becomes the last line of defense.",
            "This is why Build permission belongs on a security review checklist as its own row, "
            "not folded under 'sharing' \u2014 it grants an entirely different capability.",
        ],
        footer="Lab connection: Exercise 4 is the Build permission discussion \u2014 which personas "
               "in security-role-matrix.csv should have Build, and what governance controls should "
               "wrap it.",
        script=[
            "Frame Build as the 'surprise' permission: most authors don't realize it's a separate "
            "grant until a stakeholder asks 'can I pull this into Excel?' and it turns out sharing "
            "the report isn't enough.",
            "Walk the mechanics: Build is granted on the semantic model, not on the report. Once "
            "granted, the user can point a new report, Excel, or an XMLA client at the model as if "
            "they had built it themselves. RLS still applies to their queries as long as they are "
            "an assigned role member and not a workspace Admin/Member/Contributor.",
            "Give the thin-report vision explicitly: the goal in an enterprise deployment is often "
            "one certified 'Sales' semantic model with Build permission granted to a controlled "
            "set of analyst authors, who then each ship small purpose-built reports \u2014 not five "
            "copies of the same model.",
            "Set up Exercise 4 as a persona walk-through: for each persona in "
            "security-role-matrix.csv, decide whether they should have Build, and if yes, which "
            "RLS role they belong to. There is no single right answer \u2014 the exercise is about "
            "the reasoning being explicit and documented.",
        ]
    )
    page += 1

    # 9. Topic 7 - Object-level security
    content_slide(
        prs, 7, "Object-Level Security (OLS)", page=page,
        lead_items=[
            "OLS hides entire tables or columns from users assigned to a role \u2014 the objects "
            "disappear from field lists, from DAX intellisense, and from downstream tools.",
            "OLS is defined in the Tabular Model Definition Language (TMDL) via XMLA-compatible "
            "tooling \u2014 Tabular Editor is the standard authoring surface; there is no built-in "
            "OLS editor in Desktop.",
            "Common use cases: hide a Salary or SSN column from a role that should see the rest of "
            "an HR table; hide an entire Compensation table from a role that only needs the "
            "Employee dimension.",
            "OLS layers with RLS \u2014 a role can filter rows AND hide columns/tables. Test both "
            "layers together, not separately.",
        ],
        why_items=[
            "RLS alone cannot solve the 'they can see the customer table but not the credit-card "
            "column' problem \u2014 that is exactly what OLS is designed for.",
            "Hiding a column, not just visually removing it from a report, is what prevents an "
            "author with Build permission from selecting it in Excel or a thin report.",
            "OLS is 'Verify for Gov' because it depends on XMLA-compatible tooling, capacity SKU, "
            "and tenant settings that are not universally available in government tenants.",
        ],
        footer="Lab connection: OLS is a conceptual/optional lab \u2014 identify sensitive fields "
               "and document whether they should be removed, hidden, or protected with OLS. Do not "
               "attempt an OLS build without validating tenant, capacity, and tooling first.",
        script=[
            "Open with the boundary that separates OLS from RLS: RLS filters rows within a table; "
            "OLS hides the whole table or specific columns so the user cannot even see the "
            "structure exists. Different problem, different tool.",
            "Explain the tooling realistically: OLS is authored in TMDL via an XMLA-compatible "
            "editor like Tabular Editor \u2014 Desktop's Manage roles surface does not expose OLS. "
            "That tooling requirement is one of the reasons OLS gets labeled Verify for Gov: not "
            "every tenant permits or supports the external tooling path.",
            "Give the classic use case pair: an HR semantic model where the 'Compensation' table is "
            "hidden entirely from most analysts (whole-table OLS), and a 'Customer' table where a "
            "sensitive PII column is hidden from the marketing role even though the rest of the "
            "table is visible (column-level OLS).",
            "Land the Build-permission connection: hiding a column visually in a report is not "
            "the same thing as OLS \u2014 an author with Build permission and no OLS restriction can "
            "still see the hidden column in Excel or a new report. Real protection requires OLS on "
            "the model, not visual hiding on the report.",
            "Wrap by explicitly telling students the lab treats OLS as conceptual: identify what "
            "you'd protect and how, but don't attempt the actual OLS build in this environment "
            "without confirming tenant, capacity, and Tabular Editor access first.",
        ]
    )
    page += 1

    # 10. Topic 8 - Sensitivity labels and Purview
    table_slide(
        prs, 8, "Sensitivity Labels and Purview", page=page,
        headers=["Aspect", "Behavior in Power BI", "Dependency / Gov note"],
        col_widths=[2.6, 6.0, 4.3],
        rows=[
            ["Label inheritance",
             "A sensitivity label applied to a semantic model propagates to downstream reports, "
             "dashboards, and exports built from that model.",
             "Requires Microsoft Purview Information Protection tenant configuration."],
            ["Export controls",
             "Labels can enforce DLP behavior on Export to Excel, PDF, PowerPoint, and Analyze in "
             "Excel downloads.",
             "DLP policy authoring happens in Purview, not in Power BI."],
            ["Analyze in Excel",
             "Downloaded workbook inherits the model's label; Excel enforces it via MIP.",
             "Requires MIP client on the user's machine; behavior varies by Excel version."],
            ["Label vs. RLS",
             "Labels classify and control export/DLP; they do NOT filter rows. RLS still does the "
             "row filtering.",
             "Do not treat a label as a substitute for RLS."],
            ["Tenant policy dependency",
             "Available labels, mandatory labeling, and default labels are all set at the tenant "
             "level, not per workspace.",
             "Verify for Gov \u2014 confirm Purview and MIP availability in target Gov tenant."],
        ],
        note="Lab tie-in: Sensitivity labels are a conceptual/optional lab \u2014 review available "
             "labels, apply one where the tenant permits, and document Purview and DLP "
             "dependencies. Do not gate the module on hands-on label configuration.",
        script=[
            "Introduce sensitivity labels as the classification-and-DLP layer: they mark an item "
            "with a policy (Public, Internal, Confidential, Highly Confidential, etc.) and let "
            "downstream systems \u2014 Excel, Outlook, DLP \u2014 enforce that policy at export time.",
            "Walk the inheritance row: applying a label to a semantic model is the highest-leverage "
            "move, because every report and every export built off that model inherits it. If you "
            "only label individual reports, you'll miss exports through Build permission.",
            "Contrast with RLS very deliberately: sensitivity labels do NOT filter data. A "
            "'Confidential' label does not hide any rows \u2014 it just classifies the artifact and "
            "controls export behavior. RLS still has to filter what a user is allowed to read. "
            "Students routinely conflate these; correct it early.",
            "Cover the dependency chain: labels come from Purview Information Protection, DLP "
            "policies are authored in Purview, and Excel enforcement uses the MIP client on the "
            "user's machine. All three have to be present and configured for the story to work end "
            "to end \u2014 which is exactly why this is Verify for Gov.",
            "Close by noting the tenant-level nature: available labels, mandatory labeling, and "
            "default labels are tenant-level policy decisions, not per-workspace, so an author "
            "cannot 'just add a new label' \u2014 it has to be provisioned through the Purview and "
            "tenant admin path.",
        ]
    )
    page += 1

    # 11. Topic 9 - Sharing and external users
    content_slide(
        prs, 9, "Sharing and External Users", page=page,
        lead_items=[
            "Direct share (per-item Share button) grants view access to one item to a specific "
            "user \u2014 fast, but hard to audit at scale.",
            "App access with audiences is the recommended distribution pattern: consumers see only "
            "the App, RLS applies cleanly, and the workspace stays a controlled authoring space.",
            "Workspace access (Admin/Member/Contributor/Viewer) is broad; only Viewer respects RLS "
            "\u2014 the other three bypass it.",
            "B2B guest users are external identities added to your tenant via Entra ID B2B; their "
            "UPN in Power BI is typically the guest #EXT# form, which must match what your dynamic "
            "RLS mapping expects.",
            "Government and regulated tenants often restrict or disable external sharing and B2B "
            "entirely at the tenant policy level.",
        ],
        why_items=[
            "The four sharing paths above answer different questions and combine unpredictably \u2014 "
            "review each user's effective access as the SUM of all paths, not just the most "
            "obvious one.",
            "Because B2B guests carry a non-obvious UPN form (with #EXT#), dynamic RLS mapping "
            "tables that were built assuming internal UPNs will silently exclude every guest "
            "user \u2014 a common production bug.",
            "Assuming external sharing 'just works' in a Gov tenant is one of the most common "
            "planning mistakes \u2014 always confirm tenant policy first.",
        ],
        footer="Concept note from README: Direct sharing, App access, workspace access, and B2B "
               "guest access are separate decisions \u2014 external and guest behavior are "
               "Verify for Gov.",
        script=[
            "Sharing is deceptively simple looking \u2014 a Share button on every item \u2014 but there "
            "are actually four distinct sharing paths, and a real user's effective access is the "
            "sum of whichever paths they have been granted.",
            "Walk each path with one clarifying detail. Direct share: per-item, fast, audit-nightmare "
            "at scale. App with audiences: the recommended pattern \u2014 the audience decides what "
            "the consumer sees, and RLS applies. Workspace roles: broad and blunt, and only Viewer "
            "respects RLS \u2014 which we covered in the Service assignment topic.",
            "Now B2B guests. This is the highest-value gotcha in the whole module. External guest "
            "users get an Entra ID B2B account whose UPN in Power BI is often of the form "
            "guest_externaldomain.com#EXT#@yourtenant.onmicrosoft.com. If your dynamic RLS mapping "
            "table has their original external UPN, the filter [UserPrincipalName] = "
            "USERPRINCIPALNAME() will never match \u2014 the guest sees zero rows. Fix it by "
            "populating the guest UPN form your tenant actually issues.",
            "Land the Gov posture explicitly: external sharing, B2B, and guest-user behavior are "
            "commonly restricted in GCC, GCC High, and DoD tenants. Do not design a workflow that "
            "depends on external sharing until tenant policy has been confirmed in writing.",
        ]
    )
    page += 1

    # 12. Topic 10 - Security testing
    content_slide(
        prs, 10, "Security Testing", page=page,
        lead_items=[
            "Desktop: Modeling \u2192 View as \u2192 Other user + one or more roles \u2014 lets you "
            "simulate any UPN and any role combination locally before publishing.",
            "Service: semantic model \u2192 Security \u2192 Test as role \u2014 the Service equivalent, "
            "and the one that will actually pick up any workspace-role bypass behavior.",
            "Positive test: a user assigned to their intended role sees exactly the rows they "
            "should. Negative test: the same user does NOT see rows from any other role's scope.",
            "Test with a real Viewer account, not just the author account \u2014 workspace "
            "Admin/Member/Contributor bypass RLS and will give false-clean test results.",
            "Document each test result: the UPN tested, the role, and the visible customer/sales "
            "counts \u2014 the README's Exercise 1-3 all ask for this documentation.",
        ],
        why_items=[
            "Untested roles are worse than no roles \u2014 they create a false sense of safety while "
            "silently letting users through or blocking legitimate access.",
            "Negative testing (confirming a user CANNOT see what they shouldn't) is what catches "
            "misplaced filter direction, wrong relationship keys, or a typo in the DAX filter "
            "expression.",
            "The workspace-Admin false-clean result is the single biggest source of 'RLS looked "
            "fine in test, failed in production' incidents.",
        ],
        footer="Lab connection: Exercise 3 has students test static and dynamic roles in Desktop "
               "and, where available, publish and test in the Service with sample UPNs from "
               "security-user-territory.csv.",
        script=[
            "Set the tone: security testing is not optional. An RLS role that has not been tested "
            "positively AND negatively should never be assigned to real users.",
            "Walk the Desktop path: Modeling \u2192 View as, tick 'Other user' and enter a UPN, then "
            "tick one or more roles. Every visual in Desktop now re-renders as if that user were "
            "signed in. Point out that Other user is what makes dynamic RLS testable in Desktop \u2014 "
            "otherwise USERPRINCIPALNAME() would return the author's own UPN and always return "
            "author-visible data.",
            "Now the Service path and its gotcha: Test as role is the Service equivalent, but it "
            "only reveals what a Viewer-level assignment would see. If your test account is a "
            "workspace Member, the actual sign-in test will still bypass RLS \u2014 which is why you "
            "need a dedicated Viewer test account for real end-to-end validation.",
            "Positive vs. negative test: positive is 'jsmith@contoso.gov sees the East region's "
            "12,437 customer rows and $8.4M in sales, as expected.' Negative is 'jsmith@contoso.gov "
            "sees ZERO rows from the West region, confirming their filter is working the other "
            "direction too.' Both must pass. Documenting both is what the lab's expected-result "
            "sections are asking for.",
        ]
    )
    page += 1

    # 13. Topic 11 - Azure Government considerations
    table_slide(
        prs, 11, "Azure Government Considerations", page=page,
        headers=["Feature", "Gov status", "Validation focus"],
        col_widths=[3.4, 2.6, 6.9],
        rows=[
            ["Static RLS",
             "Gov-ready",
             "Core Power BI capability \u2014 no extra tenant validation required."],
            ["Dynamic RLS",
             "Gov-ready",
             "Validate UPN format, Entra ID sync, and guest-user UPN mapping in target tenant."],
            ["Role testing in Desktop",
             "Gov-ready",
             "Core Desktop capability \u2014 always available."],
            ["Role assignment in Service",
             "Gov-ready",
             "Requires workspace/model permission model to be in place."],
            ["Build permission",
             "Gov-ready",
             "Validate tenant governance and downstream (Excel, XMLA) access policy."],
            ["OLS",
             "Verify for Gov",
             "Requires XMLA-compatible tooling and compatible capacity SKU."],
            ["Sensitivity labels / Purview",
             "Verify for Gov",
             "Requires Purview Information Protection configuration and cloud support."],
            ["External sharing / B2B",
             "Verify for Gov",
             "Often restricted by GCC / GCC High / DoD policy and data-handling rules."],
        ],
        note="Rule of thumb: customer policy may be stricter than what the platform technically "
             "supports. Always validate customer-specific policy alongside platform capability "
             "before treating a Verify-for-Gov feature as available.",
        script=[
            "This slide is the module's Gov-readiness truth table. Walk it top to bottom; it is "
            "the row-by-row justification for why the lab treats some topics as hands-on and "
            "others as conceptual.",
            "Anchor the Gov-ready rows first: static RLS, dynamic RLS, Desktop testing, Service "
            "role assignment, and Build permission are all Gov-ready. Anything students build in "
            "Exercises 1-4 is on that Gov-ready foundation.",
            "Now the Verify-for-Gov rows. OLS depends on XMLA-compatible tooling and capacity SKU "
            "that may or may not be enabled. Sensitivity labels require Purview Information "
            "Protection configuration end to end. External sharing and B2B are frequently "
            "restricted by GCC / GCC High / DoD policy at the tenant level.",
            "Land the closing rule of thumb: customer policy can be stricter than platform "
            "capability. A feature being technically supported in a Gov tenant does not mean the "
            "customer's own policy permits it. Validate both. That single habit prevents the "
            "vast majority of 'why can't we ship this?' surprises late in a Gov deployment.",
        ]
    )
    page += 1

    # 14. Topic 12 - Security review checklist
    content_slide(
        prs, 12, "Security Review Checklist", page=page,
        lead_items=[
            "Access path: for each persona, list every path they have to the content \u2014 direct "
            "share, App audience, workspace role, external/B2B \u2014 not just the intended path.",
            "Data filter path: confirm each RLS role's filter is on the correct table, uses the "
            "expected DAX expression, and that relationship direction propagates the filter into "
            "FactSales.",
            "Export behavior: identify every export surface (Export to Excel, Analyze in Excel, "
            "PDF, PowerPoint, XMLA) and confirm sensitivity labels / DLP behave as expected.",
            "Build permission: list every user or group with Build on the semantic model and "
            "confirm each one is in an appropriate RLS role \u2014 or has a documented reason not "
            "to be.",
            "Validation evidence: for each role, keep a record of the positive and negative test "
            "results \u2014 UPN tested, role assigned, expected visible rows, actual visible rows.",
        ],
        why_items=[
            "A checklist turns 'is this secure?' from a subjective judgement into a repeatable, "
            "audit-friendly artifact.",
            "Every item on the list maps directly to a failure mode this module covered \u2014 "
            "workspace-role bypass, Build-permission surprise, guest UPN mismatch, label vs. RLS "
            "confusion, untested roles.",
            "Documented evidence is what makes 'we tested this' defensible in a security or "
            "compliance review \u2014 without evidence, testing may as well not have happened.",
        ],
        footer="Lab connection: the README's validation checklist is a subset of this review "
               "\u2014 the review here is what you'd add on top for a production deployment.",
        script=[
            "Frame this as the operating checklist for any real deployment review, not just a "
            "training artifact. Every bullet on this slide is derived from a specific failure mode "
            "we covered in the last hour.",
            "Access path: this is the sum-of-paths habit from the Sharing topic. Do not judge a "
            "user's access by their most obvious entitlement \u2014 total up every share, App, "
            "workspace, and B2B grant they have and evaluate the union.",
            "Data filter path and Build permission bullets tie back to the topics with the same "
            "names \u2014 walk each RLS role and each Build grant explicitly, once per row, so nothing "
            "is judged by memory or vibe.",
            "Export behavior is the label / Purview / DLP integration layer \u2014 confirm every "
            "export surface behaves the way policy expects, including Excel via MIP.",
            "Validation evidence is the discipline that makes the whole thing defensible: for "
            "every role you assigned, keep the UPN-tested / role-assigned / expected-vs-actual "
            "record from the Security Testing topic. That evidence is what you show in a "
            "compliance review, not a screenshot of the Manage roles dialog.",
        ]
    )
    page += 1

    # 15. Module lab walkthrough
    checklist_slide(
        prs, "Module Lab Walkthrough",
        kicker="Topic 13 \u2014 What you'll build", page=page,
        items=[
            "Data overview: security-user-territory.csv and security-role-matrix.csv via raw "
            "GitHub URL (Get data \u2192 Web).",
            "Exercise 1: Static RLS \u2014 create 'East Region' role filtering "
            "DimTerritory[TerritoryRegion] = \"East\".",
            "Exercise 1 test: Modeling \u2192 View as \u2192 East Region; document expected visible "
            "customer and sales totals.",
            "Exercise 2: Dynamic RLS \u2014 import security-user-territory.csv and relate to "
            "DimTerritory on TerritoryKey.",
            "Exercise 2 role: create 'Dynamic Territory Security' with filter "
            "[UserPrincipalName] = USERPRINCIPALNAME() on the security table.",
            "Exercise 2 test: View as with sample UPNs from the mapping table (single-territory "
            "and multi-territory users).",
            "Exercise 3: Publish to approved training workspace, assign users/groups, Test as role "
            "in Service where available.",
            "Exercise 4: Build permission discussion \u2014 walk security-role-matrix.csv personas "
            "and decide Build + RLS for each.",
            "Optional: OLS \u2014 identify sensitive fields; document remove/hide/OLS decision "
            "(Verify for Gov, conceptual only).",
            "Optional: Sensitivity labels \u2014 review labels, apply where tenant permits, document "
            "Purview and DLP dependencies.",
        ],
        script=[
            "Use this as the literal table of contents for the hands-on portion \u2014 walk it top to "
            "bottom so students know exactly what is coming and how the exercises stack.",
            "Highlight the two CSVs and how each maps to an exercise: security-user-territory.csv "
            "is what makes Exercise 2's dynamic RLS real \u2014 it's the mapping table the role reads "
            "against USERPRINCIPALNAME(). security-role-matrix.csv is the persona-expectation "
            "reference for Exercises 3 and 4.",
            "Set the expectation for the two optional labs: OLS and sensitivity labels are marked "
            "conceptual on purpose. In this classroom environment we don't want students to try "
            "the hands-on path unless they have already confirmed tenant, capacity, and Purview "
            "prerequisites \u2014 which is unlikely inside a workshop tenant.",
            "Close by reminding them documentation is part of every exercise \u2014 the expected "
            "customer / sales counts for the East region role, the results per UPN for the dynamic "
            "role, and the persona-by-persona Build decision. If they don't write it down, they "
            "haven't finished the exercise.",
        ]
    )
    page += 1

    # 16. Knowledge check & discussion
    checklist_slide(
        prs, "Knowledge Check & Discussion",
        kicker="Topic 14 \u2014 Wrap-up", page=page,
        items=[
            "Static RLS role filters the expected data on DimTerritory.",
            "Dynamic RLS role uses USERPRINCIPALNAME() on the security mapping table.",
            "SecurityUserTerritory relates correctly to DimTerritory via TerritoryKey.",
            "Desktop role testing (View as) is documented with expected counts.",
            "Service role testing (Test as role) is documented where available.",
            "Build permission behavior is documented per persona in the role matrix.",
            "OLS is marked Verify for Gov and documented conceptually only.",
            "Sensitivity labels and Purview are marked Verify for Gov.",
            "External sharing and B2B limitations are documented for the target tenant.",
            "Discussion: when would you accept a workspace Member bypass, and how would you audit "
            "it?",
            "Discussion: how would you validate guest / B2B UPN mapping before assigning dynamic "
            "RLS to external users?",
            "Discussion: which items on the Security Review Checklist would you add to a "
            "production release gate?",
        ],
        script=[
            "Run this as a mixed check-and-discussion segment: the first nine items are validation "
            "confirmations pulled straight from the README's checklist \u2014 walk them out loud and "
            "ask 'did everyone confirm this?' for each row.",
            "The last three items are the discussion questions \u2014 pick two, call on specific "
            "students, and let the answer take two or three minutes each. The workspace-Member "
            "bypass question and the guest-UPN question especially reward being answered out loud "
            "because both are places students often think they understand until they try to "
            "articulate the reasoning.",
            "For the production release gate question, connect it back to the Security Review "
            "Checklist topic \u2014 tell students the answer should include access path, data filter "
            "path, export behavior, Build permission review, and validation evidence, at minimum.",
            "Close by connecting forward: Module 8 (Service enterprise deployment) picks up "
            "exactly where this module ended \u2014 workspaces, Apps, refresh, and distribution \u2014 "
            "so anything they leave unresolved here about workspace roles and App audiences will "
            "resurface immediately in Module 8.",
        ]
    )
    page += 1

    # 17. Closing
    closing_slide(
        prs, MODULE_NO,
        "Module 08: Power BI Service Enterprise Deployment \u2014 publishing, workspaces, Apps, "
        "gateways, and refresh on top of today's secured semantic model.",
        page=page,
        subtitle=("Learners can now design layered security: static and dynamic RLS, Service "
                  "role assignment, Build permission, OLS, sensitivity labels, and a Gov-aware "
                  "review checklist."),
        script=[
            "Congratulate the class on completing the security module \u2014 this is one of the "
            "highest-stakes modules in the workshop because the failure modes it covers are the "
            "ones that show up as data-exposure incidents in production, not just as bugs.",
            "Remind them to keep the RLS roles they built \u2014 static East Region and dynamic "
            "Territory Security \u2014 in the PBIP model, because Module 8 will publish and "
            "distribute that same model and the roles need to still be in place when it lands in "
            "the Service.",
            "Take final questions before moving on, especially anything about workspace-role "
            "bypass, Build permission, and B2B UPN mapping \u2014 those are the three areas most "
            "likely to cause quiet production incidents if left ambiguous.",
        ]
    )

    OUT_PATH.parent.mkdir(parents=True, exist_ok=True)
    prs.save(str(OUT_PATH))
    print(f"Saved: {OUT_PATH} ({len(prs.slides.__iter__.__self__._sldIdLst)} slides)")


if __name__ == "__main__":
    build()
